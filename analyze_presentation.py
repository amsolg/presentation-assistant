#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Analyse rapide d'une présentation PowerPoint"""

from pptx import Presentation

def analyze_presentation(pptx_path):
    """Analyse le contenu d'une présentation"""
    prs = Presentation(pptx_path)

    print(f"Nombre total de slides: {len(prs.slides)}")
    print("-" * 60)

    for idx, slide in enumerate(prs.slides):
        print(f"\nSlide {idx + 1}:")

        # Titre
        if slide.shapes.title:
            print(f"  Titre: {slide.shapes.title.text}")

        # Compter les shapes avec du texte
        text_shapes = 0
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text and shape != slide.shapes.title:
                text_shapes += 1

        print(f"  Shapes avec texte: {text_shapes}")

        # Afficher les premiers mots du contenu
        for shape in slide.shapes[:5]:  # Limiter aux 5 premiers shapes
            if hasattr(shape, 'text') and shape.text and shape != slide.shapes.title:
                preview = shape.text[:100] + "..." if len(shape.text) > 100 else shape.text
                print(f"    - {preview}")

if __name__ == "__main__":
    analyze_presentation("presentations/presentation-assistant/equipe-architecture/output/presentation_assistant_architecture.pptx")