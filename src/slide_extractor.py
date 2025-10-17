#!/usr/bin/env python3
"""
Slide Extractor for Presentation Assistant
Extracteur détaillé de contenu PowerPoint pour génération de narration avec Sam AI

Ce module extrait de façon granulaire le contenu de chaque slide d'une présentation
PowerPoint et génère des fichiers JSON structurés par scènes pour la narration.
"""

import json
import os
import sys
from pathlib import Path
from typing import Dict, List, Any, Optional
from xml.etree import ElementTree as ET

from pptx import Presentation
from pptx.shapes.base import BaseShape
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches, Pt


class SlideExtractor:
    """
    Extracteur détaillé de contenu PowerPoint avec support des animations
    et création de fichiers JSON structurés par scènes pour Sam AI.
    """
    
    def __init__(self, presentation_path: str, output_dir: str):
        """
        Initialise l'extracteur avec le chemin de la présentation et le dossier de sortie.
        
        Args:
            presentation_path: Chemin vers le fichier .pptx
            output_dir: Dossier de sortie pour les fichiers JSON générés
        """
        self.presentation_path = Path(presentation_path)
        self.output_dir = Path(output_dir)
        self.presentation = None
        self.metadata = {
            "presentation_name": self.presentation_path.stem,
            "total_slides": 0,
            "extraction_version": "1.0.0",
            "supported_features": [
                "text_extraction",
                "shape_positions",
                "speaker_notes",
                "basic_formatting",
                "shape_properties"
            ],
            "limitations": [
                "complex_animations_via_xml_only",
                "transitions_limited",
                "embedded_media_basic"
            ]
        }
    
    def load_presentation(self) -> bool:
        """
        Charge la présentation PowerPoint et initialise les métadonnées.
        
        Returns:
            bool: True si le chargement est réussi
        """
        try:
            self.presentation = Presentation(self.presentation_path)
            self.metadata["total_slides"] = len(self.presentation.slides)
            print(f"[OK] Présentation chargée : {self.metadata['total_slides']} slides")
            return True
        except Exception as e:
            print(f"[ERROR] Erreur lors du chargement de la présentation : {e}")
            return False
    
    def extract_all_slides(self) -> bool:
        """
        Extrait le contenu de toutes les slides et génère les fichiers JSON.
        
        Returns:
            bool: True si l'extraction est réussie
        """
        if not self.presentation:
            print("[ERROR] Présentation non chargée")
            return False
        
        success_count = 0
        
        for slide_num, slide in enumerate(self.presentation.slides, 1):
            try:
                slide_data = self.extract_slide_content(slide, slide_num)
                self.save_slide_json(slide_data, slide_num)
                success_count += 1
                print(f"[OK] Slide {slide_num} extraite et sauvegardée")
            except Exception as e:
                print(f"[ERROR] Erreur lors de l'extraction de la slide {slide_num} : {e}")
        
        # Sauvegarder les métadonnées globales
        self.save_presentation_metadata()
        
        print(f"[OK] Extraction terminée : {success_count}/{self.metadata['total_slides']} slides")
        return success_count > 0
    
    def extract_slide_content(self, slide, slide_number: int) -> Dict[str, Any]:
        """
        Extrait le contenu détaillé d'une slide individuelle.
        
        Args:
            slide: Objet slide python-pptx
            slide_number: Numéro de la slide (1-indexé)
            
        Returns:
            Dict contenant toutes les données structurées de la slide
        """
        slide_data = {
            "slide_number": slide_number,
            "slide_title": self.extract_slide_title(slide),
            "layout_name": slide.slide_layout.name if hasattr(slide.slide_layout, 'name') else "Unknown",
            "total_scenes": 1,  # Minimum une scène, sera ajusté si animations
            "scenes": []
        }
        
        # Scène principale avec contenu statique
        main_scene = {
            "scene_id": 1,
            "scene_type": "static_content",
            "context": "",  # À remplir manuellement
            "technical_description": {
                "visual_elements": self.extract_visual_elements(slide),
                "layout": slide_data["layout_name"],
                "background": self.extract_background_info(slide),
                "automatic_animations": [],  # À compléter avec XML si nécessaire
                "slide_dimensions": self.get_slide_dimensions()
            },
            "speaker_notes": self.extract_speaker_notes(slide)
        }
        
        slide_data["scenes"].append(main_scene)
        
        # TODO: Analyser les animations manuelles et créer des scènes supplémentaires
        # Cette fonctionnalité nécessitera l'analyse XML pour les animations complexes
        
        return slide_data
    
    def extract_slide_title(self, slide) -> str:
        """
        Extrait le titre de la slide si disponible.
        
        Args:
            slide: Objet slide python-pptx
            
        Returns:
            str: Titre de la slide ou titre généré
        """
        try:
            for shape in slide.shapes:
                if hasattr(shape, "placeholder_format") and shape.placeholder_format:
                    try:
                        if shape.placeholder_format.type.name in ['TITLE', 'SUBTITLE']:
                            if hasattr(shape, 'text') and shape.text.strip():
                                return shape.text.strip()
                    except:
                        continue
                elif hasattr(shape, 'text') and shape.text.strip():
                    # Premier text box non-vide comme titre de fallback
                    text = shape.text.strip()
                    if len(text) < 100:  # Probablement un titre
                        return text
        except Exception:
            pass
        
        return f"Slide {getattr(slide, 'slide_id', 'Unknown')}"
    
    def extract_visual_elements(self, slide) -> List[Dict[str, Any]]:
        """
        Extrait tous les éléments visuels de la slide avec leurs propriétés détaillées.
        
        Args:
            slide: Objet slide python-pptx
            
        Returns:
            List des éléments visuels avec leurs propriétés
        """
        elements = []
        
        for shape_idx, shape in enumerate(slide.shapes):
            try:
                element = {
                    "element_id": f"shape_{shape_idx + 1}",
                    "shape_id": getattr(shape, 'shape_id', shape_idx + 1),
                    "type": self.get_shape_type(shape),
                    "position": self.get_shape_position(shape),
                    "content": self.extract_shape_content(shape),
                    "formatting": self.extract_shape_formatting(shape),
                    "properties": self.extract_shape_properties(shape)
                }
                elements.append(element)
            except Exception as e:
                # Ajouter un élément avec erreur plutôt que d'échouer complètement
                elements.append({
                    "element_id": f"shape_{shape_idx + 1}",
                    "shape_id": shape_idx + 1,
                    "type": "error",
                    "error": str(e),
                    "position": {"left": 0, "top": 0, "width": 0, "height": 0},
                    "content": {"type": "none", "data": None},
                    "formatting": {},
                    "properties": {}
                })
        
        return elements
    
    def get_shape_type(self, shape: BaseShape) -> str:
        """
        Détermine le type de forme PowerPoint.
        
        Args:
            shape: Forme PowerPoint
            
        Returns:
            str: Type de forme lisible
        """
        try:
            if hasattr(shape, 'shape_type'):
                shape_type = shape.shape_type
                type_mapping = {
                    MSO_SHAPE_TYPE.AUTO_SHAPE: "auto_shape",
                    MSO_SHAPE_TYPE.CALLOUT: "callout",
                    MSO_SHAPE_TYPE.CHART: "chart",
                    MSO_SHAPE_TYPE.COMMENT: "comment",
                    MSO_SHAPE_TYPE.CONNECTOR: "connector",
                    MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT: "embedded_ole",
                    MSO_SHAPE_TYPE.FORM_CONTROL: "form_control",
                    MSO_SHAPE_TYPE.FREEFORM: "freeform",
                    MSO_SHAPE_TYPE.GROUP: "group",
                    MSO_SHAPE_TYPE.LINE: "line",
                    MSO_SHAPE_TYPE.LINKED_OLE_OBJECT: "linked_ole",
                    MSO_SHAPE_TYPE.LINKED_PICTURE: "linked_picture",
                    MSO_SHAPE_TYPE.MEDIA: "media",
                    MSO_SHAPE_TYPE.OLE_CONTROL_OBJECT: "ole_control",
                    MSO_SHAPE_TYPE.PICTURE: "picture",
                    MSO_SHAPE_TYPE.PLACEHOLDER: "placeholder",
                    MSO_SHAPE_TYPE.SCRIPT_ANCHOR: "script_anchor",
                    MSO_SHAPE_TYPE.TABLE: "table",
                    MSO_SHAPE_TYPE.TEXT_EFFECT: "text_effect",
                    MSO_SHAPE_TYPE.TEXT_BOX: "text_box"
                }
                return type_mapping.get(shape_type, f"unknown_{shape_type}")
        except Exception:
            # Fallback : identifier le type par les propriétés disponibles
            return self.identify_shape_by_properties(shape)
        
        return "unknown"
    
    def identify_shape_by_properties(self, shape: BaseShape) -> str:
        """
        Identifie le type de forme par ses propriétés quand shape_type échoue.
        
        Args:
            shape: Forme PowerPoint
            
        Returns:
            str: Type de forme identifié
        """
        try:
            # Identifier par les propriétés spécifiques
            if hasattr(shape, 'text') and shape.text:
                return "text_box"
            elif hasattr(shape, 'table'):
                return "table"
            elif hasattr(shape, 'chart'):
                return "chart"
            elif hasattr(shape, 'image'):
                return "picture"
            elif hasattr(shape, 'placeholder_format') and shape.placeholder_format:
                return "placeholder"
            elif hasattr(shape, 'connection_site_count'):
                return "connector"  # Propriété spécifique aux connecteurs
            elif hasattr(shape, 'auto_shape_type'):
                return "auto_shape"
            else:
                # Analyser le nom de classe pour plus d'informations
                class_name = shape.__class__.__name__.lower()
                if 'connector' in class_name:
                    return "connector"
                elif 'group' in class_name:
                    return "group"
                elif 'picture' in class_name:
                    return "picture"
                elif 'textbox' in class_name or 'textframe' in class_name:
                    return "text_box"
                else:
                    return f"identified_as_{class_name}"
        except Exception:
            return "unknown_shape"
    
    def get_shape_position(self, shape: BaseShape) -> Dict[str, float]:
        """
        Extrait la position et les dimensions d'une forme.
        
        Args:
            shape: Forme PowerPoint
            
        Returns:
            Dict avec position et dimensions en pixels
        """
        try:
            # Conversion EMU vers pixels (approximation)
            emu_to_px = 1 / 9525  # 1 pixel ≈ 9525 EMU
            
            return {
                "left": round(shape.left * emu_to_px, 2),
                "top": round(shape.top * emu_to_px, 2),
                "width": round(shape.width * emu_to_px, 2),
                "height": round(shape.height * emu_to_px, 2),
                "rotation": getattr(shape, 'rotation', 0)
            }
        except Exception:
            return {"left": 0, "top": 0, "width": 0, "height": 0, "rotation": 0}
    
    def extract_shape_content(self, shape: BaseShape) -> Dict[str, Any]:
        """
        Extrait le contenu textuel ou autres données d'une forme.
        
        Args:
            shape: Forme PowerPoint
            
        Returns:
            Dict avec le contenu de la forme
        """
        content = {"type": "none", "data": None}
        shape_type = self.get_shape_type(shape)
        
        try:
            # Texte - vérifier en premier
            if hasattr(shape, 'text'):
                try:
                    text_content = shape.text
                    if text_content and text_content.strip():
                        content = {
                            "type": "text",
                            "data": text_content.strip(),
                            "has_text_frame": hasattr(shape, 'text_frame'),
                            "paragraphs_count": len(shape.text_frame.paragraphs) if hasattr(shape, 'text_frame') else 0
                        }
                        return content
                except:
                    pass
            
            # Table
            if shape_type == "table" and hasattr(shape, 'table'):
                try:
                    content = {
                        "type": "table",
                        "data": self.extract_table_content(shape.table),
                        "rows": len(shape.table.rows),
                        "columns": len(shape.table.columns)
                    }
                    return content
                except:
                    pass
            
            # Image/Picture
            if shape_type == "picture" or hasattr(shape, 'image'):
                try:
                    content = {
                        "type": "image",
                        "data": {
                            "filename": getattr(shape.image, 'filename', 'unknown') if hasattr(shape, 'image') else 'unknown',
                            "content_type": getattr(shape.image, 'content_type', 'unknown') if hasattr(shape, 'image') else 'unknown'
                        }
                    }
                    return content
                except:
                    pass
            
            # Chart
            if shape_type == "chart" or hasattr(shape, 'chart'):
                try:
                    content = {
                        "type": "chart",
                        "data": {
                            "chart_type": str(shape.chart.chart_type) if hasattr(shape, 'chart') and shape.chart else "unknown",
                            "has_title": hasattr(shape.chart, 'chart_title') if hasattr(shape, 'chart') and shape.chart else False
                        }
                    }
                    return content
                except:
                    pass
            
            # Connector ou autres formes géométriques
            if shape_type == "connector":
                content = {
                    "type": "connector",
                    "data": {
                        "description": "Ligne de connexion ou forme géométrique",
                        "shape_category": "connector"
                    }
                }
                return content
            
            # Auto Shape
            if shape_type == "auto_shape":
                try:
                    content = {
                        "type": "auto_shape",
                        "data": {
                            "auto_shape_type": str(shape.auto_shape_type) if hasattr(shape, 'auto_shape_type') else "unknown",
                            "description": "Forme automatique PowerPoint"
                        }
                    }
                    return content
                except:
                    pass
            
            # Placeholder
            if shape_type == "placeholder":
                try:
                    content = {
                        "type": "placeholder",
                        "data": {
                            "placeholder_type": str(shape.placeholder_format.type) if hasattr(shape, 'placeholder_format') and shape.placeholder_format else "unknown",
                            "description": "Zone de texte de mise en page"
                        }
                    }
                    return content
                except:
                    pass
            
            # Type générique basé sur le shape_type détecté
            content = {
                "type": shape_type,
                "data": {
                    "description": f"Élément de type {shape_type}",
                    "shape_category": shape_type
                }
            }
        
        except Exception as e:
            content = {
                "type": "error",
                "data": None,
                "error": str(e)
            }
        
        return content
    
    def extract_table_content(self, table) -> List[List[str]]:
        """
        Extrait le contenu d'une table PowerPoint.
        
        Args:
            table: Objet table python-pptx
            
        Returns:
            List 2D représentant le contenu de la table
        """
        table_data = []
        try:
            for row in table.rows:
                row_data = []
                for cell in row.cells:
                    row_data.append(cell.text if hasattr(cell, 'text') else "")
                table_data.append(row_data)
        except Exception:
            pass
        return table_data
    
    def extract_shape_formatting(self, shape: BaseShape) -> Dict[str, Any]:
        """
        Extrait les informations de formatage d'une forme.

        Args:
            shape: Forme PowerPoint

        Returns:
            Dict avec les propriétés de formatage
        """
        formatting = {}

        try:
            # Formatage du texte si disponible
            if hasattr(shape, 'text_frame') and shape.text_frame:
                # Pour les placeholders, essayer d'obtenir le formatage depuis le placeholder_format
                if hasattr(shape, 'placeholder_format') and shape.placeholder_format:
                    formatting["text"] = self.extract_placeholder_text_formatting(shape)
                else:
                    formatting["text"] = self.extract_text_formatting(shape.text_frame)

            # Couleur de remplissage
            if hasattr(shape, 'fill'):
                formatting["fill"] = self.extract_fill_properties(shape.fill)

            # Contour
            if hasattr(shape, 'line'):
                formatting["line"] = self.extract_line_properties(shape.line)

        except Exception as e:
            formatting["error"] = str(e)

        return formatting
    
    def extract_placeholder_text_formatting(self, shape) -> Dict[str, Any]:
        """
        Extrait le formatage du texte d'un placeholder.
        IMPORTANT: Cette méthode indique maintenant quand les valeurs ne peuvent pas être extraites.

        Args:
            shape: Shape avec placeholder_format

        Returns:
            Dict avec propriétés de formatage du texte
        """
        formatting = {
            "font_name": None,
            "font_size": None,
            "bold": None,
            "italic": None,
            "underline": None,
            "color": None,
            "alignment": None,
            "_extraction_note": "Values cannot be extracted from placeholder without slide context"
        }

        try:
            # D'abord, essayer d'extraire les valeurs directes du text_frame
            if hasattr(shape, 'text_frame') and shape.text_frame:
                direct_formatting = self.extract_text_formatting(shape.text_frame)

                # Mettre à jour seulement les valeurs non-null
                for key, value in direct_formatting.items():
                    if value is not None:
                        formatting[key] = value
                        if "_extraction_note" in formatting:
                            formatting["_extraction_note"] = "Some values extracted directly from text_frame"

            # Si toujours des valeurs manquantes, l'indiquer clairement
            missing_values = [k for k, v in formatting.items() if v is None and k != "_extraction_note"]
            if missing_values:
                formatting["_missing_values"] = missing_values
                formatting["_extraction_limitation"] = (
                    "python-pptx cannot access inherited values from layout/master. "
                    "Values would need to be extracted via direct XML parsing with layout context."
                )

        except Exception as e:
            formatting["_extraction_error"] = str(e)

        return formatting

    def extract_text_formatting(self, text_frame) -> Dict[str, Any]:
        """
        Extrait le formatage du texte d'un text_frame.

        Args:
            text_frame: Objet text_frame python-pptx

        Returns:
            Dict avec propriétés de formatage du texte
        """
        formatting = {
            "font_name": None,
            "font_size": None,
            "bold": None,
            "italic": None,
            "underline": None,
            "color": None,
            "alignment": None
        }

        try:
            # Prendre le formatage du premier paragraphe comme référence
            if text_frame and text_frame.paragraphs:
                para = text_frame.paragraphs[0]

                # Extraire l'alignement du paragraphe
                if hasattr(para, 'alignment') and para.alignment is not None:
                    formatting["alignment"] = str(para.alignment)

                # Extraire le formatage du texte depuis les runs
                if para.runs:
                    run = para.runs[0]
                    if hasattr(run, 'font'):
                        font = run.font

                        # Extraire chaque propriété séparément avec gestion d'erreur
                        try:
                            if hasattr(font, 'name') and font.name is not None:
                                formatting["font_name"] = font.name
                        except:
                            pass

                        try:
                            if hasattr(font, 'size') and font.size is not None:
                                formatting["font_size"] = font.size.pt
                        except:
                            pass

                        try:
                            if hasattr(font, 'bold') and font.bold is not None:
                                formatting["bold"] = font.bold
                        except:
                            pass

                        try:
                            if hasattr(font, 'italic') and font.italic is not None:
                                formatting["italic"] = font.italic
                        except:
                            pass

                        try:
                            if hasattr(font, 'underline') and font.underline is not None:
                                formatting["underline"] = bool(font.underline)
                        except:
                            pass

                        try:
                            if hasattr(font, 'color') and font.color:
                                color_hex = self.get_color_hex(font.color)
                                if color_hex:
                                    formatting["color"] = color_hex
                        except:
                            pass

                # Si pas de runs, essayer d'obtenir le formatage directement du paragraphe
                elif hasattr(para, 'font'):
                    font = para.font
                    try:
                        if font.name:
                            formatting["font_name"] = font.name
                        if font.size:
                            formatting["font_size"] = font.size.pt
                    except:
                        pass

        except Exception as e:
            # Log l'erreur mais retourne quand même le dictionnaire avec les valeurs par défaut
            print(f"[WARNING] Erreur lors de l'extraction du formatage: {e}")

        return formatting
    
    def get_color_hex(self, color_obj) -> Optional[str]:
        """
        Convertit un objet couleur PowerPoint en hex.

        Args:
            color_obj: Objet couleur python-pptx

        Returns:
            String hex de la couleur ou None
        """
        try:
            if color_obj is None:
                return None

            # Essayer d'accéder à la propriété rgb
            if hasattr(color_obj, 'rgb') and color_obj.rgb:
                rgb = color_obj.rgb
                if hasattr(rgb, 'r') and hasattr(rgb, 'g') and hasattr(rgb, 'b'):
                    return f"#{rgb.r:02x}{rgb.g:02x}{rgb.b:02x}"
                elif isinstance(rgb, (list, tuple)) and len(rgb) >= 3:
                    return f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}"

            # Essayer theme_color si rgb n'est pas disponible
            if hasattr(color_obj, 'theme_color') and color_obj.theme_color:
                # Mapper les couleurs de thème aux couleurs Premier Tech approximatives
                theme_colors = {
                    0: "#000000",  # Background1
                    1: "#FFFFFF",  # Text1
                    2: "#0066CC",  # Accent1 (Bleu Premier Tech)
                    3: "#ED7D31",  # Accent2
                    4: "#A5A5A5",  # Accent3
                    5: "#FFC000",  # Accent4
                    6: "#5B9BD5",  # Accent5
                    7: "#70AD47",  # Accent6
                }
                theme_idx = color_obj.theme_color
                if isinstance(theme_idx, int) and theme_idx in theme_colors:
                    return theme_colors[theme_idx]

        except Exception as e:
            print(f"[DEBUG] Erreur lors de la conversion de couleur: {e}")

        return None
    
    def extract_fill_properties(self, fill) -> Dict[str, Any]:
        """Extrait les propriétés de remplissage d'une forme."""
        properties = {
            "type": "unknown",
            "color": None,
            "transparency": None
        }

        try:
            if fill is None:
                return properties

            # Type de remplissage
            if hasattr(fill, 'type'):
                fill_type = fill.type
                if fill_type is not None:
                    properties["type"] = str(fill_type)

            # Couleur de remplissage
            if hasattr(fill, 'fore_color'):
                color = self.get_color_hex(fill.fore_color)
                if color:
                    properties["color"] = color

            # Transparence
            if hasattr(fill, 'transparency') and fill.transparency is not None:
                properties["transparency"] = fill.transparency

        except Exception as e:
            print(f"[DEBUG] Erreur extraction fill: {e}")

        return properties
    
    def extract_line_properties(self, line) -> Dict[str, Any]:
        """Extrait les propriétés de contour d'une forme."""
        properties = {
            "color": None,
            "width": None,
            "style": None
        }

        try:
            if line is None:
                return properties

            # Couleur de ligne
            if hasattr(line, 'color'):
                color = self.get_color_hex(line.color)
                if color:
                    properties["color"] = color
                elif hasattr(line, 'fill') and hasattr(line.fill, 'fore_color'):
                    # Essayer via fill pour certains types de lignes
                    color = self.get_color_hex(line.fill.fore_color)
                    if color:
                        properties["color"] = color

            # Largeur de ligne
            if hasattr(line, 'width') and line.width is not None:
                try:
                    properties["width"] = line.width.pt
                except:
                    # Essayer de convertir directement si pt échoue
                    properties["width"] = float(line.width) / 12700  # EMU to pt approximation

            # Style de ligne
            if hasattr(line, 'dash_style') and line.dash_style is not None:
                properties["style"] = str(line.dash_style)

        except Exception as e:
            print(f"[DEBUG] Erreur extraction line: {e}")

        return properties
    
    def extract_shape_properties(self, shape: BaseShape) -> Dict[str, Any]:
        """
        Extrait les propriétés additionnelles d'une forme.

        Args:
            shape: Forme PowerPoint

        Returns:
            Dict avec propriétés additionnelles
        """
        properties = {}

        try:
            # Propriétés de base
            properties.update({
                "name": shape.name if hasattr(shape, 'name') else None,
                "shape_id": shape.shape_id if hasattr(shape, 'shape_id') else None,
                "is_placeholder": hasattr(shape, 'placeholder_format') and shape.placeholder_format is not None,
                "has_text_frame": hasattr(shape, 'text_frame') and shape.text_frame is not None,
                "visible": True  # PowerPoint n'expose pas directement la visibilité
            })

            # Type d'autoshape si applicable
            if hasattr(shape, 'auto_shape_type'):
                try:
                    properties["auto_shape_type"] = str(shape.auto_shape_type)
                    # Pour les cercles/rectangles utilisés comme numéros dans la TOC
                    if shape.auto_shape_type in [1, 5]:  # 1=Rectangle, 5=Oval
                        properties["possible_use"] = "number_container"
                except:
                    pass

            # Propriétés spécifiques selon le type
            if hasattr(shape, 'placeholder_format') and shape.placeholder_format:
                try:
                    properties["placeholder_type"] = str(shape.placeholder_format.type)
                    properties["placeholder_idx"] = shape.placeholder_format.idx
                except Exception:
                    properties["placeholder_type"] = "unknown_placeholder"

            # Essayer d'extraire plus d'infos si c'est un autoshape avec du texte
            if properties.get("has_text_frame") and hasattr(shape, 'text_frame'):
                try:
                    tf = shape.text_frame
                    if hasattr(tf, 'vertical_anchor'):
                        properties["text_vertical_anchor"] = str(tf.vertical_anchor)
                    if hasattr(tf, 'word_wrap'):
                        properties["text_word_wrap"] = tf.word_wrap
                    if hasattr(tf, 'auto_size'):
                        properties["text_auto_size"] = str(tf.auto_size)
                except:
                    pass

        except Exception as e:
            properties["error"] = str(e)

        return properties
    
    def extract_background_info(self, slide) -> Dict[str, Any]:
        """
        Extrait les informations de background de la slide.
        
        Args:
            slide: Objet slide python-pptx
            
        Returns:
            Dict avec informations du background
        """
        # python-pptx a des limitations pour l'accès au background
        # Cette fonction pourrait être étendue avec l'analyse XML
        return {
            "type": "unknown",
            "note": "Background extraction requires XML analysis"
        }
    
    def get_slide_dimensions(self) -> Dict[str, float]:
        """
        Obtient les dimensions de la slide.
        
        Returns:
            Dict avec largeur et hauteur de la slide
        """
        try:
            if self.presentation:
                # Conversion EMU vers pixels
                emu_to_px = 1 / 9525
                return {
                    "width": round(self.presentation.slide_width * emu_to_px, 2),
                    "height": round(self.presentation.slide_height * emu_to_px, 2)
                }
        except Exception:
            pass
        
        # Dimensions standard PowerPoint 16:9
        return {"width": 1280, "height": 720}
    
    def extract_speaker_notes(self, slide) -> str:
        """
        Extrait les notes de présentation de la slide.
        
        Args:
            slide: Objet slide python-pptx
            
        Returns:
            String avec les notes ou chaîne vide
        """
        try:
            if hasattr(slide, 'notes_slide') and slide.notes_slide:
                notes_text_frame = slide.notes_slide.notes_text_frame
                if notes_text_frame and hasattr(notes_text_frame, 'text'):
                    return notes_text_frame.text.strip()
        except Exception:
            pass
        
        return ""
    
    def save_slide_json(self, slide_data: Dict[str, Any], slide_number: int) -> bool:
        """
        Sauvegarde les données d'une slide en fichier JSON.
        
        Args:
            slide_data: Données structurées de la slide
            slide_number: Numéro de la slide
            
        Returns:
            bool: True si la sauvegarde est réussie
        """
        try:
            output_file = self.output_dir / f"slide_{slide_number:02d}.json"
            with open(output_file, 'w', encoding='utf-8') as f:
                json.dump(slide_data, f, indent=2, ensure_ascii=False)
            return True
        except Exception as e:
            print(f"[ERROR] Erreur lors de la sauvegarde de la slide {slide_number} : {e}")
            return False
    
    def save_presentation_metadata(self) -> bool:
        """
        Sauvegarde les métadonnées globales de la présentation.
        
        Returns:
            bool: True si la sauvegarde est réussie
        """
        try:
            metadata_file = self.output_dir.parent / "metadata" / "presentation_info.json"
            metadata_file.parent.mkdir(exist_ok=True)
            
            with open(metadata_file, 'w', encoding='utf-8') as f:
                json.dump(self.metadata, f, indent=2, ensure_ascii=False)
            return True
        except Exception as e:
            print(f"[ERROR] Erreur lors de la sauvegarde des métadonnées : {e}")
            return False


def main():
    """Fonction principale pour l'extraction en ligne de commande."""
    if len(sys.argv) != 3:
        print("Usage: python slide_extractor.py <presentation.pptx> <output_directory>")
        sys.exit(1)
    
    presentation_path = sys.argv[1]
    output_dir = sys.argv[2]
    
    # Vérifier que le fichier existe
    if not os.path.exists(presentation_path):
        print(f"[ERROR] Fichier non trouvé : {presentation_path}")
        sys.exit(1)
    
    # Créer le dossier de sortie
    os.makedirs(output_dir, exist_ok=True)
    
    # Initialiser et exécuter l'extraction
    extractor = SlideExtractor(presentation_path, output_dir)
    
    print(f"[INFO] Chargement de la présentation : {presentation_path}")
    if not extractor.load_presentation():
        sys.exit(1)
    
    print(f"[INFO] Extraction vers : {output_dir}")
    if extractor.extract_all_slides():
        print("[SUCCESS] Extraction terminée avec succès !")
    else:
        print("[ERROR] Extraction échouée")
        sys.exit(1)


if __name__ == "__main__":
    main()