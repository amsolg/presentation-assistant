#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Detailed Explanation Builder - Construction d'explications détaillées Premier Tech
Version JSON-native pour l'architecture 2025 du presentation_builder.
Utilise les slides 35, 39-44 du template Premier Tech pour créer des explications approfondies.
"""

import os
import json
import shutil
from datetime import datetime
from pathlib import Path
from typing import List, Dict, Optional, Any
from pptx import Presentation
from pptx.util import Inches


class DetailedExplanationBuilder:
    """
    Classe pour construire des explications détaillées Premier Tech.
    Version modernisée pour l'architecture JSON 2025.
    """

    def __init__(self, template_path: str = "templates/Template_PT.pptx"):
        """
        Initialise le constructeur avec le template Premier Tech.

        Args:
            template_path: Chemin vers le template Premier Tech
        """
        self.template_path = template_path

        # Mapping des slides disponibles pour les explications détaillées
        self.explanation_slides = {
            34: {  # Slide 35 (index 34) - 4 énoncés & Mots-clés
                "name": "4 énoncés & Mots-clés",
                "usage": "Explication détaillée avec 4 points clés",
                "audience": "Managers, Executives",
                "style": "four_points",
                "elements": 4,
                "explanation_type": "structured_points"
            },
            38: {  # Slide 39 (index 38) - 2 énoncés avec sous-titres et ligne bleue
                "name": "2 énoncés avec sous-titres et ligne bleue",
                "usage": "Comparaison détaillée avec ligne bleue",
                "audience": "Managers, Analystes",
                "style": "dual_detailed_blue",
                "elements": 2,
                "explanation_type": "dual_comparison"
            },
            39: {  # Slide 40 (index 39) - 2 énoncés avec sous-titres et ligne grise
                "name": "2 énoncés avec sous-titres et ligne grise",
                "usage": "Comparaison détaillée avec ligne grise",
                "audience": "Managers, Analystes",
                "style": "dual_detailed_grey",
                "elements": 2,
                "explanation_type": "dual_comparison"
            },
            40: {  # Slide 41 (index 40) - 2 énoncés avec titre et ligne bleue
                "name": "2 énoncés avec titre et ligne bleue",
                "usage": "Comparaison avec titre et ligne bleue",
                "audience": "Toutes audiences",
                "style": "dual_titled_blue",
                "elements": 2,
                "explanation_type": "titled_comparison"
            },
            41: {  # Slide 42 (index 41) - 2 énoncés avec titre et ligne grise
                "name": "2 énoncés avec titre et ligne grise",
                "usage": "Comparaison avec titre et ligne grise",
                "audience": "Toutes audiences",
                "style": "dual_titled_grey",
                "elements": 2,
                "explanation_type": "titled_comparison"
            },
            42: {  # Slide 43 (index 42) - 2 listes avec ligne bleue
                "name": "2 listes avec ligne bleue",
                "usage": "Deux listes détaillées avec séparation bleue",
                "audience": "Formations, Audiences techniques",
                "style": "dual_lists_blue",
                "elements": "multiple",
                "explanation_type": "dual_lists"
            },
            43: {  # Slide 44 (index 43) - 2 listes avec ligne grise
                "name": "2 listes avec ligne grise",
                "usage": "Deux listes détaillées avec séparation grise",
                "audience": "Formations, Audiences techniques",
                "style": "dual_lists_grey",
                "elements": "multiple",
                "explanation_type": "dual_lists"
            }
        }

        if not os.path.exists(template_path):
            raise FileNotFoundError(f"Template Premier Tech non trouvé: {template_path}")

        # Analyser la structure des slides d'explication de référence
        self._analyze_explanation_structure()

    def _analyze_explanation_structure(self):
        """Analyse la structure des slides d'explication de référence"""
        try:
            pres = Presentation(self.template_path)
            self.explanation_info = {}

            for slide_index, slide_data in self.explanation_slides.items():
                if len(pres.slides) <= slide_index:
                    raise ValueError(f"Template ne contient pas de slide {slide_index + 1}")

                explanation_slide = pres.slides[slide_index]

                self.explanation_info[slide_index] = {
                    'layout_name': explanation_slide.slide_layout.name,
                    'shape_count': len(explanation_slide.shapes),
                    'slide_index': slide_index,
                    'slide_number': slide_index + 1,
                    'style': slide_data['style'],
                    'usage': slide_data['usage'],
                    'elements': slide_data['elements'],
                    'explanation_type': slide_data['explanation_type']
                }

            print(f"[INFO] {len(self.explanation_info)} slides d'explications détaillées analysées")

        except Exception as e:
            raise Exception(f"Erreur analyse templates explications détaillées: {e}")

    def process_detailed_explanation_config(self, config: Dict[str, Any], presentation_path: str) -> Dict[str, Any]:
        """
        Traite la configuration d'explication détaillée et insère la slide dans la présentation.

        Args:
            config: Configuration JSON pour l'explication détaillée
            presentation_path: Chemin vers la présentation cible

        Returns:
            Dict contenant le résultat de l'opération
        """
        try:
            print(f"[INFO] Traitement configuration detailed_explanation_builder")

            # Validation de la configuration
            validation_result = self._validate_detailed_explanation_config(config)
            if not validation_result["valid"]:
                return {
                    "success": False,
                    "error": f"Configuration invalide: {validation_result['errors']}",
                    "slide_added": False
                }

            # Extraction des paramètres
            content = config.get("content", "")
            title = config.get("title")
            subtitle = config.get("subtitle")
            additional_content = config.get("additional_content", [])
            explanation_style = config.get("explanation_style", "four_points")

            # Trouver le slide template correspondant
            slide_index = None
            for idx, slide_data in self.explanation_slides.items():
                if slide_data["style"] == explanation_style:
                    slide_index = idx
                    break

            if slide_index is None:
                available_styles = [s["style"] for s in self.explanation_slides.values()]
                return {
                    "success": False,
                    "error": f"Style '{explanation_style}' non trouvé. Styles disponibles: {', '.join(available_styles)}",
                    "slide_added": False
                }

            # Insertion de la slide d'explication
            return self._insert_detailed_explanation_slide(
                presentation_path,
                slide_index,
                content,
                title,
                subtitle,
                additional_content,
                explanation_style
            )

        except Exception as e:
            return {
                "success": False,
                "error": f"Erreur traitement detailed_explanation_builder: {str(e)}",
                "slide_added": False
            }

    def _validate_detailed_explanation_config(self, config: Dict[str, Any]) -> Dict[str, Any]:
        """
        Valide la configuration d'explication détaillée.

        Args:
            config: Configuration à valider

        Returns:
            Dict avec résultat de validation
        """
        errors = []

        # Vérification des champs requis
        required_fields = ["content", "explanation_style"]
        for field in required_fields:
            if field not in config or not config[field]:
                errors.append(f"Champ requis manquant: {field}")

        # Validation du style
        if "explanation_style" in config:
            available_styles = [slide_data["style"] for slide_data in self.explanation_slides.values()]
            if config["explanation_style"] not in available_styles:
                errors.append(f"Style invalide: {config['explanation_style']}. Styles disponibles: {', '.join(available_styles)}")

        # Validation du contenu
        if "content" in config:
            content_length = len(config["content"])
            if content_length < 10:
                errors.append("Le contenu doit contenir au moins 10 caractères")
            elif content_length > 500:
                errors.append("Le contenu ne peut dépasser 500 caractères")

        # Validation spécifique selon le style
        if "explanation_style" in config and "additional_content" in config:
            style = config["explanation_style"]
            additional_content = config["additional_content"]

            # Validation spécialisée pour dual_detailed_explanation
            # Note: Ces slides ont des shapes pour contenu principal (plus long) et sous-titres (max 32 chars)
            # La limite de 32 caractères ne s'applique PAS au contenu principal
            if style in ["dual_detailed_blue", "dual_detailed_grey"]:
                # Pour dual_detailed, les additional_content sont le contenu principal des sections
                # Pas de limitation stricte de 32 caractères pour le contenu principal
                for i, item in enumerate(additional_content):
                    if len(str(item)) > 150:  # Limite raisonnable pour le contenu des sections
                        errors.append(f"L'élément {i+1} du contenu additionnel dépasse 150 caractères (limite recommandée pour lisibilité)")

            # Validation du nombre d'éléments selon le style
            if style == "four_points" and len(additional_content) != 4:
                errors.append("Le style 'four_points' requiert exactement 4 éléments dans additional_content")
            elif style in ["dual_detailed_blue", "dual_detailed_grey", "dual_titled_blue", "dual_titled_grey"] and len(additional_content) != 2:
                errors.append(f"Le style '{style}' requiert exactement 2 éléments dans additional_content")
            elif style in ["dual_lists_blue", "dual_lists_grey"] and len(additional_content) != 2:
                errors.append(f"Le style '{style}' requiert exactement 2 listes dans additional_content")

        # Validation du titre (optionnel)
        if "title" in config and config["title"]:
            if len(config["title"]) > 100:
                errors.append("Le titre ne peut dépasser 100 caractères")

        # Validation du sous-titre (optionnel)
        if "subtitle" in config and config["subtitle"]:
            if len(config["subtitle"]) > 150:
                errors.append("Le sous-titre ne peut dépasser 150 caractères")

        # Validation du contenu additionnel (optionnel)
        if "additional_content" in config and config["additional_content"]:
            if len(config["additional_content"]) > 6:
                errors.append("Le contenu additionnel ne peut contenir plus de 6 éléments")
            for i, item in enumerate(config["additional_content"]):
                if len(str(item)) > 200:
                    errors.append(f"L'élément {i+1} du contenu additionnel dépasse 200 caractères")

        return {
            "valid": len(errors) == 0,
            "errors": errors
        }

    def _insert_detailed_explanation_slide(self,
                                         presentation_path: str,
                                         slide_index: int,
                                         content: str,
                                         title: Optional[str],
                                         subtitle: Optional[str],
                                         additional_content: List[str],
                                         explanation_style: str) -> Dict[str, Any]:
        """
        Insère une slide d'explication détaillée dans la présentation.

        Args:
            presentation_path: Chemin vers la présentation
            slide_index: Index du template de slide
            content: Contenu principal
            title: Titre optionnel
            subtitle: Sous-titre optionnel
            additional_content: Contenu additionnel
            explanation_style: Style de l'explication

        Returns:
            Dict avec résultat de l'insertion
        """
        try:
            # Charger la présentation cible
            if not os.path.exists(presentation_path):
                return {
                    "success": False,
                    "error": f"Présentation non trouvée: {presentation_path}",
                    "slide_added": False
                }

            target_pres = Presentation(presentation_path)
            template_pres = Presentation(self.template_path)

            # Cloner la slide template
            if len(template_pres.slides) <= slide_index:
                return {
                    "success": False,
                    "error": f"Template ne contient pas de slide {slide_index + 1}",
                    "slide_added": False
                }

            template_slide = template_pres.slides[slide_index]

            # Ajouter la slide clonée
            slide_layout = template_slide.slide_layout
            new_slide = target_pres.slides.add_slide(slide_layout)

            # Copier tous les éléments de la slide template
            for shape in template_slide.shapes:
                if hasattr(shape, 'element'):
                    # Cloner les shapes de la slide template
                    try:
                        # Les shapes sont automatiquement copiés avec le layout
                        pass
                    except Exception as e:
                        print(f"[WARNING] Impossible de cloner shape: {e}")

            # Personnaliser le contenu
            self._customize_detailed_explanation_content(
                new_slide, content, title, subtitle, additional_content, explanation_style
            )

            # Amélioration du text wrapping et centrage
            self._improve_text_wrapping(new_slide)

            # Sauvegarder
            target_pres.save(presentation_path)

            slide_number = len(target_pres.slides)
            print(f"[SUCCESS] Slide d'explication détaillée ajoutée (slide {slide_number}, style: {explanation_style})")

            return {
                "success": True,
                "message": f"Slide d'explication détaillée ajoutée avec succès",
                "slide_added": True,
                "slide_number": slide_number,
                "explanation_style": explanation_style,
                "content_preview": content[:50] + "..." if len(content) > 50 else content
            }

        except Exception as e:
            return {
                "success": False,
                "error": f"Erreur insertion slide d'explication: {str(e)}",
                "slide_added": False
            }

    def _customize_detailed_explanation_content(self,
                                              slide,
                                              content: str,
                                              title: Optional[str],
                                              subtitle: Optional[str],
                                              additional_content: List[str],
                                              explanation_style: str):
        """
        Personnalise le contenu de la slide d'explication détaillée.

        Args:
            slide: Slide à personnaliser
            content: Contenu principal
            title: Titre optionnel
            subtitle: Sous-titre optionnel
            additional_content: Contenu additionnel
            explanation_style: Style de l'explication
        """
        try:
            # Identifier et personnaliser les shapes selon le style
            if explanation_style == "four_points":
                self._customize_four_points_slide(slide, content, title, additional_content)
            elif explanation_style in ["dual_detailed_blue", "dual_detailed_grey"]:
                self._customize_dual_detailed_slide(slide, content, title, subtitle, additional_content)
            elif explanation_style in ["dual_titled_blue", "dual_titled_grey"]:
                self._customize_dual_titled_slide(slide, content, title, additional_content)
            elif explanation_style in ["dual_lists_blue", "dual_lists_grey"]:
                self._customize_dual_lists_slide(slide, content, title, additional_content)
            else:
                # Fallback générique
                self._customize_generic_explanation_slide(slide, content, title)

        except Exception as e:
            print(f"[WARNING] Erreur personnalisation contenu: {e}")
            # Fallback vers personnalisation générique
            self._customize_generic_explanation_slide(slide, content, title)

    def _customize_four_points_slide(self, slide, content: str, title: Optional[str], additional_content: List[str]):
        """Personnalise une slide de style four_points (Slide 35)"""
        # Titre principal
        if slide.shapes.title and title:
            slide.shapes.title.text = title

        # Contenu principal et points avec mappage spécifique
        text_shapes = [shape for shape in slide.shapes if hasattr(shape, 'text_frame') and shape != slide.shapes.title]

        print(f"[DEBUG] Four points slide - Text shapes trouvés: {len(text_shapes)}")

        if text_shapes and len(text_shapes) >= 5:  # S'assurer qu'on a assez de shapes
            # Premier text shape = contenu principal
            text_shapes[0].text = content
            print(f"[DEBUG] Contenu principal assigné au shape 0")

            # Assigner les 4 points aux shapes spécifiques
            for i, point in enumerate(additional_content[:4]):
                if i + 1 < len(text_shapes):
                    text_shapes[i + 1].text = point
                    print(f"[DEBUG] Point {i+1} assigné au shape {i+1}: {point[:30]}...")
        else:
            print(f"[WARNING] Pas assez de text shapes ({len(text_shapes)}) pour four_points")
            # Fallback vers assignation générique
            if text_shapes:
                text_shapes[0].text = content
                for i, point in enumerate(additional_content[:min(4, len(text_shapes)-1)]):
                    if i + 1 < len(text_shapes):
                        text_shapes[i + 1].text = point

    def _customize_dual_detailed_slide(self, slide, content: str, title: Optional[str],
                                     subtitle: Optional[str], additional_content: List[str]):
        """Personnalise une slide de style dual_detailed (Slides 39/40)"""
        # Titre
        if slide.shapes.title and title:
            slide.shapes.title.text = title

        # Contenu principal
        text_shapes = [shape for shape in slide.shapes if hasattr(shape, 'text_frame') and shape != slide.shapes.title]

        if text_shapes and len(text_shapes) > 0:
            # Contenu principal
            text_shapes[0].text = content

            # Sous-titre si disponible
            if subtitle and len(text_shapes) > 1:
                text_shapes[1].text = subtitle

            # Contenu additionnel (2 éléments max pour dual)
            start_idx = 2 if subtitle else 1
            for i, item in enumerate(additional_content[:2]):
                if start_idx + i < len(text_shapes):
                    text_shapes[start_idx + i].text = item

    def _customize_dual_titled_slide(self, slide, content: str, title: Optional[str], additional_content: List[str]):
        """Personnalise une slide de style dual_titled (Slides 41/42)"""
        # Titre principal
        if slide.shapes.title and title:
            slide.shapes.title.text = title

        # Contenu et sections
        text_shapes = [shape for shape in slide.shapes if hasattr(shape, 'text_frame') and shape != slide.shapes.title]

        if text_shapes and len(text_shapes) > 0:
            # Contenu principal
            text_shapes[0].text = content

            # Deux sections titrées
            for i, item in enumerate(additional_content[:2]):
                if i + 1 < len(text_shapes):
                    text_shapes[i + 1].text = item

    def _customize_dual_lists_slide(self, slide, content: str, title: Optional[str], additional_content: List[str]):
        """Personnalise une slide de style dual_lists (Slides 43/44)"""
        # Titre
        if slide.shapes.title and title:
            slide.shapes.title.text = title

        # Contenu principal et listes avec mappage correct
        text_shapes = [shape for shape in slide.shapes if hasattr(shape, 'text_frame') and shape != slide.shapes.title]

        print(f"[DEBUG] Dual lists slide - Text shapes trouvés: {len(text_shapes)}")

        if text_shapes and len(text_shapes) >= 3:  # S'assurer qu'on a assez de shapes
            # Shape 0 = contenu principal (description générale)
            text_shapes[0].text = content
            print(f"[DEBUG] Contenu principal assigné au shape 0")

            # Traitement des listes duales avec positions correctes
            # Pour les slides dual lists, les positions sont spécifiques:
            # Shape 1 = Titre liste 1, Shape 2 = Contenu liste 1
            # Shape 3 = Titre liste 2, Shape 4 = Contenu liste 2
            if additional_content and len(additional_content) >= 2:
                for i, list_content in enumerate(additional_content[:2]):
                    if '|' in list_content:
                        parts = list_content.split('|')
                        list_title = parts[0]
                        list_items = parts[1:]

                        # Position correcte : i*2+1 pour titre, i*2+2 pour contenu
                        title_shape_idx = i * 2 + 1
                        content_shape_idx = i * 2 + 2

                        if title_shape_idx < len(text_shapes):
                            text_shapes[title_shape_idx].text = list_title
                            print(f"[DEBUG] Titre liste {i+1} assigné au shape {title_shape_idx}: {list_title}")

                        if content_shape_idx < len(text_shapes):
                            list_text = "\n".join([f"• {item}" for item in list_items])
                            text_shapes[content_shape_idx].text = list_text
                            print(f"[DEBUG] Contenu liste {i+1} assigné au shape {content_shape_idx}")
                    else:
                        # Format simple sans séparateur
                        if i + 1 < len(text_shapes):
                            text_shapes[i + 1].text = list_content
                            print(f"[DEBUG] Liste simple {i+1} assignée au shape {i+1}")
        else:
            print(f"[WARNING] Pas assez de text shapes ({len(text_shapes)}) pour dual_lists")
            # Fallback vers assignation simple
            if text_shapes:
                text_shapes[0].text = content
                for i, item in enumerate(additional_content[:min(2, len(text_shapes)-1)]):
                    if i + 1 < len(text_shapes):
                        text_shapes[i + 1].text = str(item)

    def _customize_generic_explanation_slide(self, slide, content: str, title: Optional[str]):
        """Personnalisation générique pour tous les styles"""
        # Titre
        if slide.shapes.title and title:
            slide.shapes.title.text = title

        # Contenu principal dans le premier text shape disponible
        text_shapes = [shape for shape in slide.shapes if hasattr(shape, 'text_frame') and shape != slide.shapes.title]
        if text_shapes:
            text_shapes[0].text = content

    def _improve_text_wrapping(self, slide):
        """
        Améliore le text wrapping et le centrage du contenu.
        Inspiré de section_header_builder.py
        """
        try:
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame') and shape.text_frame.text.strip():
                    text_frame = shape.text_frame

                    # Configuration du text wrapping
                    text_frame.word_wrap = True
                    text_frame.auto_size = None

                    # Amélioration de l'espacement
                    for paragraph in text_frame.paragraphs:
                        paragraph.space_after = Inches(0.1)

                        for run in paragraph.runs:
                            # Assurer une taille de police lisible
                            if hasattr(run.font, 'size') and run.font.size and run.font.size.pt < 14:
                                run.font.size = Inches(0.02)  # ~14pt

        except Exception as e:
            print(f"[WARNING] Erreur amélioration text wrapping: {e}")


def load_detailed_explanation_payload(payload_file_path: str) -> Dict[str, Any]:
    """
    Charge un payload d'explication détaillée depuis un fichier JSON.

    Args:
        payload_file_path: Chemin vers le fichier JSON de payload

    Returns:
        Dict contenant la configuration de l'explication détaillée

    Raises:
        FileNotFoundError: Si le fichier n'existe pas
        ValueError: Si le JSON est invalide
    """
    try:
        if not os.path.exists(payload_file_path):
            raise FileNotFoundError(f"Fichier payload non trouvé: {payload_file_path}")

        with open(payload_file_path, 'r', encoding='utf-8') as f:
            payload = json.load(f)

        print(f"[INFO] Payload detailed_explanation_builder chargé: {payload_file_path}")
        return payload

    except json.JSONDecodeError as e:
        raise ValueError(f"JSON invalide dans {payload_file_path}: {e}")
    except Exception as e:
        raise Exception(f"Erreur chargement payload detailed_explanation_builder: {e}")


def process_detailed_explanation_from_payload_file(payload_path: str, presentation_path: str,
                                                 template_path: str = "templates/Template_PT.pptx") -> Dict[str, Any]:
    """
    Traite une explication détaillée en chargeant le payload depuis un fichier JSON.
    Point d'entrée pour l'architecture nouvelle avec fichiers payload séparés.

    Args:
        payload_path: Chemin vers le fichier JSON contenant le payload
        presentation_path: Chemin vers la présentation cible
        template_path: Chemin vers le template Premier Tech

    Returns:
        Dict contenant les résultats du traitement
    """
    try:
        # Charger le payload
        payload = load_detailed_explanation_payload(payload_path)

        # Traiter avec le payload chargé
        builder = DetailedExplanationBuilder(template_path)
        result = builder.process_detailed_explanation_config(payload, presentation_path)

        # Ajouter les informations de payload
        result["payload_source"] = payload_path
        result["payload_loaded"] = True

        return result

    except Exception as e:
        return {
            "success": False,
            "error": f"Erreur traitement depuis fichier payload: {e}",
            "payload_source": payload_path,
            "payload_loaded": False,
            "slide_added": False
        }


def main():
    """Point d'entrée principal pour tests en mode standalone"""
    import argparse

    parser = argparse.ArgumentParser(description='Test Detailed Explanation Builder')
    parser.add_argument('--config', required=True, help='Fichier JSON de configuration')
    parser.add_argument('--presentation', required=True, help='Chemin vers la présentation cible')

    args = parser.parse_args()

    try:
        # Charger la configuration
        config = load_detailed_explanation_payload(args.config)

        # Créer le builder
        builder = DetailedExplanationBuilder()

        # Traiter la configuration
        result = builder.process_detailed_explanation_config(config, args.presentation)

        if result["success"]:
            print(f"[SUCCESS] {result['message']}")
        else:
            print(f"[ERROR] {result['error']}")

    except Exception as e:
        print(f"[ERROR] Erreur main: {e}")


if __name__ == "__main__":
    main()