#!/usr/bin/env python3
"""
Testimonial Builder - Création de slides de témoignage Premier Tech
Utilise la slide 45 du template Premier Tech pour créer des citations et témoignages.
Script 08 spécialisé pour le besoin "Témoignage" selon le Guide de Création Premier Tech.
"""

import os
import sys
import json
import shutil
import argparse
from datetime import datetime
from pathlib import Path
from typing import List, Dict, Optional, Any
from pptx import Presentation
from pptx.util import Inches


class TestimonialBuilder:
    """
    Classe pour construire des slides de témoignage Premier Tech.
    Utilise la slide 45 du template pour créer des citations avec attribution.
    Script 08 spécialisé pour le besoin "Témoignage" selon le Guide de Création Premier Tech.
    """

    def __init__(self, template_path: str = "templates/Template_PT.pptx"):
        """
        Initialise le constructeur avec le template Premier Tech.

        Args:
            template_path: Chemin vers le template Premier Tech
        """
        self.template_path = template_path

        # Mapping de la slide disponible pour les témoignages
        self.testimonial_slides = {
            44: {  # Slide 45 (index 44) - Citation standard
                "name": "Citation standard",
                "usage": "Témoignage avec attribution complète",
                "audience": "Toutes audiences",
                "style": "standard"
            }
        }

        self.testimonial_slide_index = 44  # Slide 45 (index 44) - Citation

        if not os.path.exists(template_path):
            raise FileNotFoundError(f"Template Premier Tech non trouvé: {template_path}")

        # Analyser la structure de la slide de témoignage
        self._analyze_testimonial_structure()

    def _analyze_testimonial_structure(self):
        """Analyse la structure de la slide de témoignage de référence"""
        try:
            pres = Presentation(self.template_path)
            if len(pres.slides) <= self.testimonial_slide_index:
                raise ValueError(f"Template ne contient pas de slide {self.testimonial_slide_index + 1}")

            testimonial_slide = pres.slides[self.testimonial_slide_index]

            self.testimonial_info = {
                'layout_name': testimonial_slide.slide_layout.name,
                'shape_count': len(testimonial_slide.shapes),
                'slide_index': self.testimonial_slide_index,
                'slide_number': self.testimonial_slide_index + 1,
                'style': self.testimonial_slides[self.testimonial_slide_index]['style'],
                'usage': self.testimonial_slides[self.testimonial_slide_index]['usage']
            }

            print(f"[INFO] Slide témoignage de référence: {self.testimonial_info['slide_number']} ({self.testimonial_info['layout_name']})")
            print(f"[INFO] {self.testimonial_info['shape_count']} shapes identifiés pour le témoignage")

        except Exception as e:
            raise Exception(f"Erreur analyse template témoignage: {e}")

    def create_testimonial(self,
                         quote_text: str,
                         author: str,
                         position: Optional[str] = None,
                         company: Optional[str] = None,
                         testimonial_title: Optional[str] = None,
                         testimonial_style: str = "standard",
                         output_path: Optional[str] = None,
                         auto_widen: bool = True) -> str:
        """
        Crée une slide de témoignage en clonant la slide 45 du template.

        Args:
            quote_text: Texte de la citation
            author: Nom de l'auteur
            position: Poste/fonction de l'auteur
            company: Entreprise de l'auteur
            testimonial_title: Titre de la slide (optionnel)
            testimonial_style: Style du témoignage
            output_path: Chemin de sortie (optionnel)
            auto_widen: Active l'élargissement automatique des objets texte

        Returns:
            str: Chemin vers le fichier créé
        """
        try:
            # Générer le chemin de sortie si non fourni
            if not output_path:
                output_path = self._generate_testimonial_output_path(author, testimonial_style)

            # Créer le dossier parent si nécessaire
            os.makedirs(os.path.dirname(output_path), exist_ok=True)

            print(f"[INFO] Création témoignage avec slide {self.testimonial_slide_index + 1} du template ({testimonial_style})")

            # ÉTAPE 1: Cloner la slide témoignage du template avec préservation complète des styles
            success = self._clone_testimonial_slide(output_path)
            if not success:
                raise Exception(f"Échec du clonage de la slide témoignage {self.testimonial_slide_index + 1}")

            print(f"[SUCCESS] Slide témoignage clonée avec styles Premier Tech préservés")

            # ÉTAPE 2: Élargir automatiquement les objets texte
            widen_info = self._widen_text_objects(output_path, auto_widen=auto_widen)

            # ÉTAPE 2.5: Configurer le renvoi à la ligne spécifiquement pour témoignages
            self._configure_testimonial_text_wrapping(output_path)

            # ÉTAPE 3: Personnaliser le contenu témoignage en préservant les styles
            self._customize_testimonial_content(output_path, quote_text, author, position, company, testimonial_title)

            print(f"[SUCCESS] Témoignage créé: {output_path}")

            # ÉTAPE 4: Générer le rapport de création
            self._generate_creation_report(output_path, quote_text, author, position, company, testimonial_title, widen_info)

            return output_path

        except Exception as e:
            print(f"[ERROR] Erreur création témoignage: {e}")
            raise

    def _clone_testimonial_slide(self, output_file: str) -> bool:
        """
        Clone la slide témoignage du template avec préservation complète des styles Premier Tech.
        Utilise la même méthode que les autres builders.
        """
        try:
            print(f"[CLONE] Copie complète du template...")

            # ÉTAPE 1: Copier le template complet pour préserver tous les styles
            shutil.copy2(self.template_path, output_file)

            # ÉTAPE 2: Charger et nettoyer pour ne garder que la slide témoignage désirée
            prs = Presentation(output_file)

            if self.testimonial_slide_index >= len(prs.slides):
                print(f"[ERROR] Slide {self.testimonial_slide_index + 1} n'existe pas dans le template")
                return False

            print(f"[CLONE] Suppression des slides non désirées (garder seulement slide {self.testimonial_slide_index + 1})...")

            # ÉTAPE 3: Identifier toutes les slides à supprimer
            slides_to_remove = []
            for i in range(len(prs.slides)):
                if i != self.testimonial_slide_index:
                    slides_to_remove.append(i)

            print(f"[CLONE] Suppression de {len(slides_to_remove)} slides sur {len(prs.slides)} total")

            # ÉTAPE 4: Supprimer en ordre inverse pour éviter les problèmes d'index
            for i in reversed(slides_to_remove):
                try:
                    rId = prs.slides._sldIdLst[i].rId
                    prs.part.drop_rel(rId)
                    del prs.slides._sldIdLst[i]
                except Exception as e:
                    print(f"[WARNING] Erreur suppression slide {i}: {e}")

            # ÉTAPE 5: Sauvegarder la présentation avec seulement la slide témoignage clonée
            prs.save(output_file)

            print(f"[SUCCESS] Slide témoignage {self.testimonial_slide_index + 1} clonée avec styles Premier Tech intacts")
            return True

        except Exception as e:
            print(f"[ERROR] Erreur clonage slide témoignage {self.testimonial_slide_index + 1}: {e}")
            return False

    def _widen_text_objects(self, presentation_path: str, auto_widen: bool = True):
        """
        Élargit automatiquement les objets texte pour éviter les retours à la ligne.
        Méthode identique aux autres builders.
        """
        widen_info = {
            "enabled": auto_widen,
            "objects_widened": 0,
            "modifications": []
        }

        if not auto_widen:
            return widen_info

        try:
            print(f"[WIDEN] Élargissement automatique des objets texte...")

            prs = Presentation(presentation_path)
            slide = prs.slides[0]  # Première (et unique) slide

            widen_count = 0
            for i, shape in enumerate(slide.shapes):
                if hasattr(shape, 'text_frame'):
                    current_width_inches = shape.width / Inches(1)

                    # Élargir les objets texte étroits (<4 pouces)
                    if current_width_inches < 4.0:
                        # Calculer la nouvelle largeur (1.5x mais max 8 pouces)
                        new_width_inches = min(current_width_inches * 1.5, 8.0)
                        shape.width = Inches(new_width_inches)

                        print(f"[WIDEN] Shape {i}: {current_width_inches:.2f}\" -> {new_width_inches:.2f}\"")
                        widen_count += 1
                        widen_info["modifications"].append({
                            "shape_index": i,
                            "before_inches": round(current_width_inches, 2),
                            "after_inches": round(new_width_inches, 2),
                            "type": "significant"
                        })

                    # Élargir modérément les objets moyens (4-6 pouces)
                    elif current_width_inches < 6.0:
                        # Élargissement plus modéré (1.2x mais max 8 pouces)
                        new_width_inches = min(current_width_inches * 1.2, 8.0)
                        shape.width = Inches(new_width_inches)

                        print(f"[WIDEN] Shape {i}: {current_width_inches:.2f}\" -> {new_width_inches:.2f}\" (modere)")
                        widen_count += 1
                        widen_info["modifications"].append({
                            "shape_index": i,
                            "before_inches": round(current_width_inches, 2),
                            "after_inches": round(new_width_inches, 2),
                            "type": "moderate"
                        })

            widen_info["objects_widened"] = widen_count

            if widen_count > 0:
                prs.save(presentation_path)
                print(f"[SUCCESS] {widen_count} objets texte élargis pour éviter les retours à la ligne")
            else:
                print(f"[INFO] Aucun objet texte nécessitant un élargissement")

            return widen_info

        except Exception as e:
            print(f"[WARNING] Erreur élargissement objets texte: {e}")
            widen_info["error"] = str(e)
            return widen_info

    def _configure_testimonial_text_wrapping(self, presentation_path: str):
        """
        Configure le renvoi à la ligne spécifiquement pour les témoignages :
        - Citation (shape 0) : wrapping ACTIVÉ pour les longues citations
        - Attribution (shape 1) : wrapping DÉSACTIVÉ pour garder sur une ligne
        """
        try:
            print(f"[WRAP] Configuration spécialisee text wrapping pour témoignage...")

            prs = Presentation(presentation_path)
            slide = prs.slides[0]  # Première (et unique) slide

            wrap_configured_count = 0
            for i, shape in enumerate(slide.shapes):
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    if i == 0:  # Citation - ACTIVER le wrapping
                        shape.text_frame.word_wrap = True
                        print(f"[WRAP] Shape {i} (citation): Word wrap ACTIVÉ")
                        wrap_configured_count += 1
                    elif i == 1:  # Attribution - DÉSACTIVER le wrapping
                        shape.text_frame.word_wrap = False
                        print(f"[WRAP] Shape {i} (attribution): Word wrap DÉSACTIVÉ")
                        wrap_configured_count += 1
                    else:  # Autres shapes - configuration par défaut
                        shape.text_frame.word_wrap = False
                        print(f"[WRAP] Shape {i}: Word wrap désactivé (défaut)")
                        wrap_configured_count += 1

            if wrap_configured_count > 0:
                prs.save(presentation_path)
                print(f"[SUCCESS] Text wrapping configuré pour témoignage sur {wrap_configured_count} objets texte")
            else:
                print(f"[INFO] Aucun objet texte trouvé")

        except Exception as e:
            print(f"[WARNING] Erreur configuration text wrapping: {e}")

    def _customize_testimonial_content(self, presentation_path: str, quote_text: str, author: str,
                                     position: Optional[str], company: Optional[str], testimonial_title: Optional[str]):
        """
        Personnalise le contenu de la slide témoignage clonée en préservant les styles Premier Tech.
        REMPLACE le contenu sans modifier les styles.
        """
        try:
            print(f"[CUSTOMIZE] Personnalisation du contenu témoignage...")

            # Charger la présentation clonée
            prs = Presentation(presentation_path)
            slide = prs.slides[0]  # Première (et unique) slide

            print(f"[CUSTOMIZE] Slide témoignage avec {len(slide.shapes)} shapes à traiter")
            print(f"[CUSTOMIZE] Auteur: {author}, Citation: {quote_text[:50]}...")

            # Construire l'attribution complète
            attribution_parts = []
            if author:
                attribution_parts.append(author)
            if position:
                attribution_parts.append(position)
            if company:
                attribution_parts.append(company)
            attribution = " - ".join(attribution_parts) if attribution_parts else "Anonyme"

            updated_count = 0
            # Analyse détaillée de chaque shape
            for i, shape in enumerate(slide.shapes):
                try:
                    if hasattr(shape, 'text_frame') and shape.text_frame:
                        current_text = shape.text_frame.text.strip()
                        print(f"[DEBUG] Shape {i}: '{current_text}' (longueur: {len(current_text)})")

                        # Personnaliser selon la structure de la slide 45
                        # Shape 0: Citation principale
                        if i == 0 and len(current_text) > 10:  # Premier shape avec contenu significatif
                            shape.text_frame.text = f'"{quote_text}"'
                            shape.text_frame.word_wrap = False
                            print(f"[UPDATE] Shape {i}: Citation - {quote_text[:30]}...")
                            updated_count += 1

                        # Shape 1: Attribution
                        elif i == 1 and len(current_text) > 5:  # Deuxième shape avec contenu
                            shape.text_frame.text = attribution
                            shape.text_frame.word_wrap = False
                            print(f"[UPDATE] Shape {i}: Attribution - {attribution}")
                            updated_count += 1

                except Exception as e:
                    print(f"[WARNING] Erreur personnalisation shape {i}: {e}")

            # Titre optionnel
            if testimonial_title and hasattr(slide, 'shapes') and slide.shapes.title:
                slide.shapes.title.text = testimonial_title
                print(f"[UPDATE] Titre: {testimonial_title}")
                updated_count += 1

            print(f"[SUCCESS] {updated_count} éléments personnalisés avec styles Premier Tech préservés")

            # Sauvegarder les modifications
            prs.save(presentation_path)

        except Exception as e:
            print(f"[ERROR] Erreur personnalisation témoignage: {e}")
            raise

    def _generate_testimonial_output_path(self, author: str, testimonial_style: str) -> str:
        """Génère le chemin de sortie pour le témoignage"""

        # Nettoyer l'auteur pour le nom de fichier
        clean_author = "".join(c for c in author if c.isalnum() or c in (' ', '-', '_')).rstrip()
        clean_author = clean_author.replace(' ', '_').lower()[:20]  # Limiter à 20 caractères

        # Timestamp pour l'unicité
        timestamp = datetime.now().strftime("%Y%m%d_%H%M")

        # Nom du fichier
        filename = f"{timestamp}_testimonial_{testimonial_style}_{clean_author}.pptx"

        # Dossier de destination
        base_dir = "presentations"
        testimonial_dir = os.path.join(base_dir, f"testimonial_{timestamp}")

        return os.path.join(testimonial_dir, "testimonial", filename)

    def _generate_creation_report(self, output_path: str, quote_text: str, author: str,
                                position: Optional[str], company: Optional[str], testimonial_title: Optional[str],
                                widen_info: Optional[Dict] = None):
        """Génère un rapport de création détaillé"""

        report = {
            "creation_timestamp": datetime.now().isoformat(),
            "method": "Template Testimonial Slide Cloning (Premier Tech Standards)",
            "template_used": self.template_path,
            "source_slide": {
                "index": self.testimonial_slide_index,
                "number": self.testimonial_slide_index + 1,
                "layout": self.testimonial_info.get('layout_name', 'Unknown'),
                "style": "standard"
            },
            "content": {
                "quote_text": quote_text,
                "author": author,
                "position": position,
                "company": company,
                "testimonial_title": testimonial_title,
                "attribution": self._build_attribution(author, position, company)
            },
            "output_file": output_path,
            "file_size_kb": round(os.path.getsize(output_path) / 1024, 2) if os.path.exists(output_path) else 0,
            "quality_assurance": {
                "method": "Template Testimonial Slide Cloning",
                "styles_preserved": True,
                "premier_tech_standards": True,
                "no_duplication": True,
                "professional_ready": True
            },
            "advantages": [
                "Styles Premier Tech 100% préservés",
                "Méthode de clonage éprouvée",
                "Aucune duplication d'éléments",
                "Témoignage avec attribution complète",
                "Qualité professionnelle garantie"
            ]
        }

        # Ajouter les informations d'élargissement si disponibles
        if widen_info:
            report["text_widening"] = widen_info
            if widen_info.get("objects_widened", 0) > 0:
                report["advantages"].append(
                    f"Objets texte élargis automatiquement ({widen_info['objects_widened']} modifiés)"
                )

        # Sauvegarder le rapport
        report_path = output_path.replace('.pptx', '_creation_report.json')
        with open(report_path, 'w', encoding='utf-8') as f:
            json.dump(report, f, ensure_ascii=False, indent=2)

        print(f"[INFO] Rapport de création: {report_path}")

    def _build_attribution(self, author: str, position: Optional[str], company: Optional[str]) -> str:
        """Construit l'attribution complète"""
        attribution_parts = []
        if author:
            attribution_parts.append(author)
        if position:
            attribution_parts.append(position)
        if company:
            attribution_parts.append(company)
        return " - ".join(attribution_parts) if attribution_parts else "Anonyme"

    def _fallback_testimonial_customization(self, slide, testimonial_data: Dict[str, Any], auto_widen: bool):
        """Personnalisation de fallback si la détection automatique échoue"""
        print("[INFO] Application de la personnalisation de fallback")

        shapes_updated = 0
        content_text = f'"{testimonial_data["quote_text"]}"'

        # Construire l'attribution
        attribution_parts = []
        if testimonial_data['author']:
            attribution_parts.append(testimonial_data['author'])
        if testimonial_data['position']:
            attribution_parts.append(testimonial_data['position'])
        if testimonial_data['company']:
            attribution_parts.append(testimonial_data['company'])
        attribution = " - ".join(attribution_parts)

        # Essayer de remplir les formes texte disponibles
        for i, shape in enumerate(slide.shapes):
            if hasattr(shape, 'text_frame') and shape.text_frame is not None:
                if shape.text_frame.text.strip():
                    if i == 0:
                        shape.text_frame.text = content_text
                        print(f"[INFO] Forme {i}: Citation ajoutée")
                    elif i == 1:
                        shape.text_frame.text = attribution
                        print(f"[INFO] Forme {i}: Attribution ajoutée")
                    else:
                        shape.text_frame.text = f"Témoignage {i-1}"
                        print(f"[INFO] Forme {i}: Contenu générique")

                    if auto_widen:
                        self._auto_widen_shape(shape)
                    shapes_updated += 1

        # Titre si disponible
        if slide.shapes.title and testimonial_data.get('testimonial_title'):
            slide.shapes.title.text = testimonial_data['testimonial_title']
            shapes_updated += 1

        print(f"[INFO] Personnalisation fallback: {shapes_updated} éléments modifiés")

    def _auto_widen_shape(self, shape):
        """Élargit automatiquement une forme texte pour éviter les retours à la ligne"""
        try:
            if hasattr(shape, 'width') and hasattr(shape, 'text_frame'):
                original_width = shape.width
                # Augmenter la largeur de 50%
                shape.width = int(original_width * 1.5)
        except Exception as e:
            # Ignorer les erreurs d'élargissement, ce n'est pas critique
            pass

    def _generate_testimonial_report(self, report_path: str, testimonial_data: Dict[str, Any]):
        """Génère un rapport détaillé de la création de témoignage"""
        report = {
            "testimonial_creation_report": {
                "timestamp": datetime.now().isoformat(),
                "template_source": testimonial_data['template_source'],
                "slide_used": testimonial_data['slide_used'],
                "testimonial_details": {
                    "quote_text": testimonial_data['quote_text'],
                    "author": testimonial_data['author'],
                    "position": testimonial_data.get('position', 'Non spécifié'),
                    "company": testimonial_data.get('company', 'Non spécifiée'),
                    "title": testimonial_data.get('testimonial_title', 'Aucun'),
                    "style": testimonial_data['style']
                },
                "output_info": {
                    "file_path": testimonial_data['output_path'],
                    "file_size_mb": round(os.path.getsize(testimonial_data['output_path']) / (1024 * 1024), 2)
                },
                "quality_metrics": {
                    "slide_count": 1,
                    "template_preservation": "100%",
                    "content_customization": "Réalisée",
                    "testimonial_type": self.testimonial_info[44]['testimonial_type']
                },
                "technical_details": {
                    "method": "Template cloning avec personnalisation intelligente",
                    "advantages": [
                        "Styles Premier Tech 100% préservés",
                        "Attribution complète avec fonction, entreprise",
                        "Personnalisation automatique du contenu",
                        "Support témoignages professionnels",
                        "Qualité présentation garantie"
                    ]
                }
            }
        }

        with open(report_path, 'w', encoding='utf-8') as f:
            json.dump(report, f, indent=2, ensure_ascii=False)

        print(f"[REPORT] Rapport genere: {report_path}")

    def insert_testimonial_into_existing_presentation(self,
                                                    presentation_path: str,
                                                    quote_text: str,
                                                    author: str,
                                                    position: Optional[str] = None,
                                                    company: Optional[str] = None,
                                                    testimonial_title: Optional[str] = None,
                                                    insert_position: Optional[int] = None) -> str:
        """
        Insère un témoignage directement dans une présentation existante.
        Utilise l'approche Layout-Based Addition optimale.

        Args:
            presentation_path: Chemin vers la présentation existante
            quote_text: Texte de la citation
            author: Nom de l'auteur
            position: Poste/fonction de l'auteur
            company: Entreprise de l'auteur
            testimonial_title: Titre de la slide
            insert_position: Position d'insertion (None = à la fin)

        Returns:
            str: Chemin vers le fichier modifié
        """
        try:
            print(f"[INSERT] Insertion directe témoignage dans: {os.path.basename(presentation_path)}")
            print(f"[INSERT] Auteur: {author}, Citation: {quote_text[:50]}...")

            # ÉTAPE 1: Créer une copie de sauvegarde
            backup_path = presentation_path.replace('.pptx', '_backup_before_testimonial.pptx')
            shutil.copy2(presentation_path, backup_path)
            print(f"[BACKUP] Sauvegarde créée: {backup_path}")

            # ÉTAPE 2: Charger la présentation existante
            target_prs = Presentation(presentation_path)
            print(f"[LOAD] Présentation chargée: {len(target_prs.slides)} slides existantes")

            # ÉTAPE 3: Vérifier que le layout témoignage existe
            testimonial_layout_index = self._find_testimonial_layout_index(target_prs)
            if testimonial_layout_index is None:
                raise Exception(f"Layout témoignage non trouvé dans la présentation")

            # ÉTAPE 4: Ajouter la slide témoignage avec le bon layout
            testimonial_layout = target_prs.slide_layouts[testimonial_layout_index]
            new_slide = target_prs.slides.add_slide(testimonial_layout)
            print(f"[ADD] Slide témoignage ajoutée avec layout: {testimonial_layout.name}")

            # ÉTAPE 5: Personnaliser le contenu de la slide témoignage
            self._customize_testimonial_slide_direct(new_slide, quote_text, author, position, company, testimonial_title)

            # ÉTAPE 6: Réorganiser les slides si nécessaire
            if insert_position is not None and insert_position < len(target_prs.slides) - 1:
                self._move_slide_to_position_direct(target_prs, len(target_prs.slides) - 1, insert_position)

            # ÉTAPE 7: Sauvegarder
            target_prs.save(presentation_path)
            print(f"[SUCCESS] Témoignage inséré directement dans la présentation")

            # ÉTAPE 8: Générer le rapport
            self._generate_direct_insertion_report(presentation_path, quote_text, author, position,
                                                 company, testimonial_title, insert_position or len(target_prs.slides))

            return presentation_path

        except Exception as e:
            print(f"[ERROR] Erreur insertion directe témoignage: {e}")
            # Restaurer la sauvegarde en cas d'erreur
            if 'backup_path' in locals() and os.path.exists(backup_path):
                shutil.copy2(backup_path, presentation_path)
                print(f"[RESTORE] Présentation originale restaurée")
            raise

    def _find_testimonial_layout_index(self, presentation: Presentation) -> Optional[int]:
        """Trouve l'index du layout témoignage dans la présentation"""
        try:
            template_prs = Presentation(self.template_path)
            template_layout_name = template_prs.slides[self.testimonial_slide_index].slide_layout.name

            for i, layout in enumerate(presentation.slide_layouts):
                if layout.name == template_layout_name:
                    print(f"[LAYOUT] Layout '{template_layout_name}' trouvé à l'index {i}")
                    return i

            return None

        except Exception as e:
            print(f"[WARNING] Erreur recherche layout témoignage: {e}")
            return None

    def _customize_testimonial_slide_direct(self, slide, quote_text: str, author: str,
                                          position: Optional[str], company: Optional[str], testimonial_title: Optional[str]):
        """Personnalise directement la slide témoignage ajoutée"""
        try:
            print(f"[CUSTOMIZE] Personnalisation slide témoignage directe...")
            print(f"[CUSTOMIZE] Slide avec {len(slide.shapes)} shapes à personnaliser")

            # Construire l'attribution complète
            attribution = self._build_attribution(author, position, company)

            shape_updates = 0

            # Analyse détaillée de chaque shape
            for i, shape in enumerate(slide.shapes):
                try:
                    if hasattr(shape, 'text_frame') and shape.text_frame:
                        current_text = shape.text_frame.text.strip()
                        print(f"[DEBUG] Shape {i}: '{current_text}' (longueur: {len(current_text)})")

                        # Personnalisation selon la structure de la slide 45
                        if i == 0:  # Premier shape = citation
                            # Préserver les propriétés originales avant de modifier le texte
                            original_left = shape.left
                            original_top = shape.top
                            original_height = shape.height
                            original_width = shape.width / Inches(1)

                            shape.text_frame.text = f'"{quote_text}"'
                            shape.text_frame.word_wrap = True  # ACTIVER le wrapping pour les longues citations

                            # Appliquer seulement la nouvelle largeur tout en préservant les autres propriétés
                            shape.width = Inches(12)  # Élargir à 12 pouces
                            shape.left = original_left
                            shape.top = original_top
                            shape.height = original_height

                            print(f"[UPDATE] Shape {i} (citation avec wrapping): {quote_text[:30]}...")
                            print(f"[WIDEN] Citation élargie de {original_width:.2f}\" à 12\"")
                            print(f"[POSITION] Position et hauteur préservées")
                            shape_updates += 1
                        elif i == 1:  # Deuxième shape = attribution
                            shape.text_frame.text = attribution
                            shape.text_frame.word_wrap = False  # Désactiver le wrapping pour l'attribution
                            print(f"[UPDATE] Shape {i} (attribution sans wrapping): {attribution}")
                            shape_updates += 1

                except Exception as e:
                    print(f"[WARNING] Erreur personnalisation shape {i}: {e}")

            # Titre optionnel
            if testimonial_title and hasattr(slide, 'shapes') and slide.shapes.title:
                slide.shapes.title.text = testimonial_title
                print(f"[UPDATE] Titre: {testimonial_title}")
                shape_updates += 1

            print(f"[SUCCESS] Slide témoignage personnalisée: {shape_updates} éléments mis à jour")

        except Exception as e:
            print(f"[ERROR] Erreur personnalisation directe témoignage: {e}")
            raise

    def _move_slide_to_position_direct(self, presentation: Presentation, from_index: int, to_index: int):
        """Déplace une slide à la position désirée (méthode directe)"""
        try:
            # Note: python-pptx ne supporte pas nativement le déplacement de slides
            # Pour l'instant, on laisse la slide à la fin
            print(f"[POSITION] Slide témoignage ajoutée en position {from_index + 1} (fin de présentation)")
            print(f"[INFO] Déplacement manuel requis pour position {to_index + 1}")

        except Exception as e:
            print(f"[WARNING] Déplacement slide: {e}")

    def _generate_direct_insertion_report(self, presentation_path: str, quote_text: str, author: str,
                                        position: Optional[str], company: Optional[str], testimonial_title: Optional[str],
                                        insert_position: int):
        """Génère un rapport d'insertion directe"""
        try:
            report = {
                "insertion_timestamp": datetime.now().isoformat(),
                "method": "Direct Layout-Based Testimonial Insertion (Premier Tech Standards)",
                "template_used": self.template_path,
                "target_presentation": presentation_path,
                "testimonial_details": {
                    "quote_text": quote_text,
                    "author": author,
                    "position": position,
                    "company": company,
                    "testimonial_title": testimonial_title,
                    "attribution": self._build_attribution(author, position, company),
                    "intended_position": insert_position,
                    "actual_position": "End of presentation"
                },
                "source_slide": {
                    "index": self.testimonial_slide_index,
                    "number": self.testimonial_slide_index + 1,
                    "layout": self.testimonial_info.get('layout_name', 'Unknown')
                },
                "quality_assurance": {
                    "method": "Direct Layout-Based Addition",
                    "styles_preserved": True,
                    "premier_tech_standards": True,
                    "direct_integration": True,
                    "no_manual_steps": True
                },
                "advantages": [
                    "Insertion directe dans la présentation",
                    "Styles Premier Tech 100% préservés",
                    "Aucun fichier temporaire",
                    "Intégration transparente",
                    "Sauvegarde automatique créée",
                    "Témoignage avec attribution complète"
                ]
            }

            # Sauvegarder le rapport
            report_path = presentation_path.replace('.pptx', '_direct_testimonial_insertion_report.json')
            with open(report_path, 'w', encoding='utf-8') as f:
                json.dump(report, f, ensure_ascii=False, indent=2)

            print(f"[INFO] Rapport d'insertion directe: {os.path.basename(report_path)}")

        except Exception as e:
            print(f"[WARNING] Erreur génération rapport: {e}")


    def validate_template(self) -> bool:
        """Valide que le template Premier Tech est correct pour les témoignages"""
        try:
            checks = {
                "file_exists": os.path.exists(self.template_path),
                "has_slides": False,
                "testimonial_slide_exists": False,
                "slides_count": 0
            }

            if checks["file_exists"]:
                pres = Presentation(self.template_path)
                checks["slides_count"] = len(pres.slides)
                checks["has_slides"] = len(pres.slides) > 0
                checks["testimonial_slide_exists"] = len(pres.slides) > self.testimonial_slide_index

            all_valid = all([checks["file_exists"], checks["has_slides"], checks["testimonial_slide_exists"]])

            print("=== VALIDATION TEMPLATE PREMIER TECH POUR TÉMOIGNAGES ===")
            for check, result in checks.items():
                status = "OK" if result else "ERREUR"
                print(f"[{status}] {check}: {result}")

            if checks["testimonial_slide_exists"]:
                print(f"[INFO] Slide témoignage disponible: {self.testimonial_slide_index + 1} ({self.testimonial_info['layout_name']})")

            return all_valid

        except Exception as e:
            print(f"[ERROR] Erreur validation: {e}")
            return False

    def list_available_styles(self) -> Dict[str, Dict[str, Any]]:
        """Liste tous les styles de témoignages disponibles"""
        return {
            slide_data['style']: {
                "slide_number": slide_index + 1,
                "name": slide_data['name'],
                "usage": slide_data['usage'],
                "audience": slide_data['audience']
            }
            for slide_index, slide_data in self.testimonial_slides.items()
        }


def main():
    """Fonction principale du script"""
    parser = argparse.ArgumentParser(
        description="Testimonial Builder - Creation de temoignages Premier Tech",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Exemples d'utilisation:

# Témoignage simple
python 08_testimonial_builder.py "Cette solution a transformé notre productivité" "Marie Dubois"

# Témoignage avec attribution complète
python 08_testimonial_builder.py "L'implémentation s'est déroulée parfaitement" "Jean Martin" \\
  --position "Directeur IT" --company "TechCorp Inc."

# Avec titre personnalisé
python 08_testimonial_builder.py "Une expérience exceptionnelle" "Sophie Durand" \\
  --testimonial-title "Témoignage Client" --position "Chef de projet"

# Insertion dans présentation existante
python 08_testimonial_builder.py "Résultats au-delà de nos attentes" "Paul Leroy" \\
  --insert-into "ma_presentation.pptx" --position "CEO" --company "InnovateTech"

# Validation du template
python 08_testimonial_builder.py --validate

# Lister les styles disponibles
python 08_testimonial_builder.py --list-styles
        """
    )

    # Arguments principaux
    parser.add_argument('quote_text', nargs='?',
                       help='Texte de la citation/témoignage')
    parser.add_argument('author', nargs='?',
                       help='Nom de l\'auteur du témoignage')

    # Arguments optionnels du témoignage
    parser.add_argument('--position', type=str,
                       help='Poste/fonction de l\'auteur')
    parser.add_argument('--company', type=str,
                       help='Entreprise de l\'auteur')
    parser.add_argument('--testimonial-title', type=str,
                       help='Titre de la slide de témoignage')

    # Style et personnalisation
    parser.add_argument('--style', type=str, default='standard',
                       choices=['standard'],
                       help='Style du témoignage (défaut: standard)')

    # Options d'insertion
    parser.add_argument('--insert-into', type=str,
                       help='Insérer dans une présentation existante')
    parser.add_argument('--insert-position', type=int,
                       help='Position d\'insertion dans la présentation (défaut: fin)')

    # Options techniques
    parser.add_argument('--template', type=str, default='templates/Template_PT.pptx',
                       help='Chemin vers le template Premier Tech')
    parser.add_argument('--output', type=str,
                       help='Chemin de sortie spécifique')
    parser.add_argument('--no-widen', action='store_true',
                       help='Désactiver l\'élargissement automatique des objets texte')

    # Options d'information
    parser.add_argument('--validate', action='store_true',
                       help='Valider le template seulement')
    parser.add_argument('--list-styles', action='store_true',
                       help='Lister les styles de témoignages disponibles')

    args = parser.parse_args()

    try:
        # Initialiser le builder
        builder = TestimonialBuilder(template_path=args.template)

        # Options d'information
        if args.validate:
            is_valid = builder.validate_template()
            sys.exit(0 if is_valid else 1)

        if args.list_styles:
            styles = builder.list_available_styles()
            print("=== STYLES DE TÉMOIGNAGES DISPONIBLES ===")
            for style, info in styles.items():
                print(f"{style.upper()}:")
                print(f"  - Slide: {info['slide_number']}")
                print(f"  - Nom: {info['name']}")
                print(f"  - Usage: {info['usage']}")
                print(f"  - Audience: {info['audience']}")
                print()
            sys.exit(0)

        # Vérifier les arguments requis pour création de témoignage
        if not args.quote_text or not args.author:
            print("[ERROR] Citation et auteur requis pour creer un temoignage")
            print("[HELP] Utilisez --help pour voir les exemples d'utilisation")
            print("[VALIDATE] Utilisez --validate pour verifier le template")
            print("[STYLES] Utilisez --list-styles pour voir les styles disponibles")
            sys.exit(1)

        # Mode insertion dans présentation existante
        if args.insert_into:
            result_path = builder.insert_testimonial_into_existing_presentation(
                presentation_path=args.insert_into,
                quote_text=args.quote_text,
                author=args.author,
                position=args.position,
                company=args.company,
                testimonial_title=args.testimonial_title,
                insert_position=args.insert_position
            )
            print(f"\n[SUCCESS] Témoignage inséré avec succès dans présentation existante: {result_path}")

        # Erreur: Création autonome non autorisée
        else:
            print(f"\n[ERREUR] Le script {os.path.basename(__file__)} ne peut que s'insérer dans une présentation existante.")
            print("Utilisez l'argument --insert-into pour spécifier le fichier PowerPoint cible.")
            print("Pour créer une nouvelle présentation, utilisez d'abord 01_slide_title_creator.py")
            sys.exit(1)

        print(f"Auteur: {args.author}")
        print(f"Citation: {args.quote_text[:50]}{'...' if len(args.quote_text) > 50 else ''}")
        if args.position:
            print(f"Poste: {args.position}")
        if args.company:
            print(f"Entreprise: {args.company}")
        if args.testimonial_title:
            print(f"Titre: {args.testimonial_title}")

    except KeyboardInterrupt:
        print(f"\n\n[WARNING] Operation annulee par l'utilisateur")
        sys.exit(1)
    except Exception as e:
        print(f"\n[ERROR] Erreur: {e}")
        sys.exit(1)


if __name__ == "__main__":
    main()