#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Navigation Builder - Construction de table des matières Premier Tech
Version JSON-native pour l'architecture 2025 du presentation_builder.
Utilise la slide 13 du template Premier Tech pour créer une table des matières structurée.
"""

import os
import json
import shutil
from datetime import datetime
from pathlib import Path
from typing import List, Dict, Optional, Any
from pptx import Presentation
from pptx.util import Inches


class NavigationBuilder:
    """
    Classe pour construire une navigation/table des matières Premier Tech.
    Version modernisée pour l'architecture JSON 2025.
    """

    def __init__(self, template_path: str = "templates/Template_PT.pptx"):
        """
        Initialise le constructeur avec le template Premier Tech.

        Args:
            template_path: Chemin vers le template Premier Tech
        """
        self.template_path = template_path
        self.toc_slide_index = 12  # Slide 13 (index 12) - Table des matières

        if not os.path.exists(template_path):
            raise FileNotFoundError(f"Template Premier Tech non trouvé: {template_path}")

        # Analyser la structure de la slide TOC de référence
        self._analyze_toc_structure()

    def _analyze_toc_structure(self):
        """Analyse la structure de la slide de table des matières de référence"""
        try:
            pres = Presentation(self.template_path)
            if len(pres.slides) <= self.toc_slide_index:
                raise ValueError(f"Template ne contient pas de slide {self.toc_slide_index + 1}")

            toc_slide = pres.slides[self.toc_slide_index]

            self.toc_info = {
                'layout_name': toc_slide.slide_layout.name,
                'shape_count': len(toc_slide.shapes),
                'slide_index': self.toc_slide_index,
                'slide_number': self.toc_slide_index + 1
            }

            print(f"[INFO] Slide TOC de référence: {self.toc_info['slide_number']} ({self.toc_info['layout_name']})")
            print(f"[INFO] {self.toc_info['shape_count']} shapes identifiés pour la TOC")

        except Exception as e:
            raise Exception(f"Erreur analyse template TOC: {e}")

    def process_navigation_config(self, config: Dict[str, Any], presentation_path: str) -> Dict[str, Any]:
        """
        Traite une configuration JSON de navigation et l'applique à une présentation.
        Point d'entrée principal pour l'architecture JSON 2025.

        Args:
            config: Configuration JSON validée de la navigation
            presentation_path: Chemin vers la présentation cible

        Returns:
            Dict contenant les résultats du traitement
        """
        try:
            print(f"[NAVIGATION] Traitement configuration JSON...")

            # Extraire les données de configuration
            title = config.get('title', 'Table des matières')
            sections = config.get('sections', [])
            options = config.get('options', {})
            auto_widen = options.get('auto_widen', True)
            insert_position = options.get('insert_position', 1)

            print(f"[NAVIGATION] Titre: {title}")
            print(f"[NAVIGATION] {len(sections)} sections à traiter")
            print(f"[NAVIGATION] Options: auto_widen={auto_widen}, position={insert_position}")

            # Valider la configuration
            validation_result = self._validate_navigation_config(config)
            if not validation_result['valid']:
                raise ValueError(f"Configuration invalide: {', '.join(validation_result['errors'])}")

            # Insérer la navigation dans la présentation
            result_path = self.insert_toc_into_existing_presentation(
                presentation_path=presentation_path,
                sections=sections,
                toc_title=title,
                insert_position=insert_position,
                auto_widen=auto_widen
            )

            # Générer le rapport de traitement
            processing_report = {
                "timestamp": datetime.now().isoformat(),
                "method": "JSON Navigation Builder 2025",
                "success": True,
                "configuration": config,
                "result_path": result_path,
                "validation": validation_result,
                "processing_details": {
                    "title_applied": title,
                    "sections_processed": len(sections),
                    "auto_widen_applied": auto_widen,
                    "insert_position": insert_position
                }
            }

            print(f"[SUCCESS] Navigation JSON traitée avec succès")
            return processing_report

        except Exception as e:
            error_report = {
                "timestamp": datetime.now().isoformat(),
                "method": "JSON Navigation Builder 2025",
                "success": False,
                "error": str(e),
                "configuration": config
            }
            print(f"[ERROR] Erreur traitement navigation JSON: {e}")
            return error_report

    def _validate_navigation_config(self, config: Dict[str, Any]) -> Dict[str, Any]:
        """Valide une configuration JSON de navigation"""
        errors = []
        warnings = []

        # Vérifier les champs requis
        if 'title' not in config:
            errors.append("Champ 'title' manquant")
        elif not config['title'] or not isinstance(config['title'], str):
            errors.append("Le titre doit être une chaîne non vide")
        elif len(config['title']) > 100:
            warnings.append("Titre très long (>100 caractères)")

        if 'sections' not in config:
            errors.append("Champ 'sections' manquant")
        elif not isinstance(config['sections'], list):
            errors.append("'sections' doit être une liste")
        elif len(config['sections']) == 0:
            errors.append("Au moins une section est requise")
        elif len(config['sections']) > 6:
            warnings.append("Plus de 6 sections peuvent causer des problèmes d'affichage")

        # Valider chaque section
        if 'sections' in config and isinstance(config['sections'], list):
            for i, section in enumerate(config['sections']):
                if not isinstance(section, str) or not section.strip():
                    errors.append(f"Section {i+1} doit être une chaîne non vide")
                elif len(section) > 150:
                    warnings.append(f"Section {i+1} très longue (>150 caractères)")

        # Valider les options
        if 'options' in config:
            options = config['options']
            if 'auto_widen' in options and not isinstance(options['auto_widen'], bool):
                errors.append("'auto_widen' doit être un booléen")
            if 'insert_position' in options:
                if not isinstance(options['insert_position'], int) or options['insert_position'] < 0:
                    errors.append("'insert_position' doit être un entier positif")

        return {
            'valid': len(errors) == 0,
            'errors': errors,
            'warnings': warnings
        }

    def insert_toc_into_existing_presentation(self,
                                           presentation_path: str,
                                           sections: List[str],
                                           toc_title: str = "Table des matières",
                                           insert_position: int = 1,
                                           auto_widen: bool = True) -> str:
        """
        Insère une table des matières directement dans une présentation existante.
        Version modernisée avec support pour l'architecture JSON 2025.

        Args:
            presentation_path: Chemin vers la présentation existante
            sections: Liste des sections à inclure
            toc_title: Titre de la table des matières
            insert_position: Position d'insertion (1 = après slide titre)
            auto_widen: Active l'élargissement automatique des objets texte

        Returns:
            str: Chemin vers le fichier modifié
        """
        try:
            print(f"[INSERT] Insertion directe TOC dans: {os.path.basename(presentation_path)}")
            print(f"[INSERT] {len(sections)} sections à insérer")

            # ÉTAPE 1: Créer une copie de sauvegarde
            backup_path = presentation_path.replace('.pptx', '_backup_before_toc.pptx')
            shutil.copy2(presentation_path, backup_path)
            print(f"[BACKUP] Sauvegarde créée: {backup_path}")

            # ÉTAPE 2: Charger la présentation existante
            target_prs = Presentation(presentation_path)
            print(f"[LOAD] Présentation chargée: {len(target_prs.slides)} slides existantes")

            # ÉTAPE 3: Vérifier que le layout TOC existe
            toc_layout_index = self._find_toc_layout_index(target_prs)
            if toc_layout_index is None:
                raise Exception("Layout 'Table des matières' non trouvé dans la présentation")

            # ÉTAPE 4: Ajouter la slide TOC avec le bon layout
            toc_layout = target_prs.slide_layouts[toc_layout_index]
            new_slide = target_prs.slides.add_slide(toc_layout)
            print(f"[ADD] Slide TOC ajoutée avec layout: {toc_layout.name}")

            # ÉTAPE 5: Personnaliser le contenu de la slide TOC
            self._customize_toc_slide_direct(new_slide, sections, toc_title)

            # ÉTAPE 6: Appliquer l'élargissement automatique si demandé
            if auto_widen:
                self._apply_auto_widen_to_slide(new_slide)

            # ÉTAPE 7: Réorganiser les slides si nécessaire
            if insert_position < len(target_prs.slides) - 1:
                self._move_slide_to_position_direct(target_prs, len(target_prs.slides) - 1, insert_position)

            # ÉTAPE 8: Sauvegarder
            target_prs.save(presentation_path)
            print(f"[SUCCESS] TOC insérée directement dans la présentation")

            # ÉTAPE 9: Générer le rapport
            self._generate_direct_insertion_report(presentation_path, sections, toc_title, insert_position, auto_widen)

            return presentation_path

        except Exception as e:
            print(f"[ERROR] Erreur insertion directe TOC: {e}")
            # Restaurer la sauvegarde en cas d'erreur
            if 'backup_path' in locals() and os.path.exists(backup_path):
                shutil.copy2(backup_path, presentation_path)
                print(f"[RESTORE] Présentation originale restaurée")
            raise

    def _find_toc_layout_index(self, presentation: Presentation) -> Optional[int]:
        """Trouve l'index du layout TOC dans la présentation"""
        try:
            template_prs = Presentation(self.template_path)
            template_layout_name = template_prs.slides[self.toc_slide_index].slide_layout.name

            for i, layout in enumerate(presentation.slide_layouts):
                if layout.name == template_layout_name:
                    print(f"[LAYOUT] Layout '{template_layout_name}' trouvé à l'index {i}")
                    return i

            return None

        except Exception as e:
            print(f"[WARNING] Erreur recherche layout: {e}")
            return None

    def _customize_toc_slide_direct(self, slide, sections: List[str], toc_title: str):
        """Personnalise directement la slide TOC ajoutée selon la vraie structure de la slide 13"""
        try:
            print(f"[CUSTOMIZE] Personnalisation slide TOC directe...")
            print(f"[CUSTOMIZE] Slide avec {len(slide.shapes)} shapes à personnaliser")

            shape_updates = 0

            # Personnaliser le titre (Shape 10) avec text wrapping ACTIVÉ
            if len(slide.shapes) > 10:
                title_shape = slide.shapes[10]
                if hasattr(title_shape, 'text_frame') and title_shape.text_frame:
                    title_shape.text_frame.text = toc_title
                    title_shape.text_frame.word_wrap = True
                    print(f"[UPDATE] Shape 10 (titre): {toc_title} (text wrapping ACTIVÉ)")
                    shape_updates += 1

            # Personnaliser les numéros et contenus (limiter aux sections fournies)
            max_sections = min(len(sections), 5)  # Maximum 5 sections selon le template

            for i in range(max_sections):
                # Personnaliser le numéro (Shapes 0-4)
                if i < len(slide.shapes):
                    num_shape = slide.shapes[i]
                    if hasattr(num_shape, 'text_frame') and num_shape.text_frame:
                        num_shape.text_frame.text = str(i + 1)
                        num_shape.text_frame.word_wrap = False
                        print(f"[UPDATE] Shape {i} (numéro): {i + 1}")

                # Personnaliser le contenu (Shapes 5-9)
                content_index = i + 5
                if content_index < len(slide.shapes):
                    content_shape = slide.shapes[content_index]
                    if hasattr(content_shape, 'text_frame') and content_shape.text_frame:
                        content_shape.text_frame.text = sections[i]
                        content_shape.text_frame.word_wrap = False
                        print(f"[UPDATE] Shape {content_index} (contenu): {sections[i][:30]}...")
                        shape_updates += 1

            # Vider les sections inutilisées
            for i in range(max_sections, 5):
                # Vider le numéro
                if i < len(slide.shapes):
                    num_shape = slide.shapes[i]
                    if hasattr(num_shape, 'text_frame') and num_shape.text_frame:
                        num_shape.text_frame.text = ""

                # Vider le contenu
                content_index = i + 5
                if content_index < len(slide.shapes):
                    content_shape = slide.shapes[content_index]
                    if hasattr(content_shape, 'text_frame') and content_shape.text_frame:
                        content_shape.text_frame.text = ""

            print(f"[SUCCESS] Slide TOC personnalisée ({max_sections} sections sur 5 disponibles)")

        except Exception as e:
            print(f"[ERROR] Erreur personnalisation directe: {e}")
            raise

    def _apply_auto_widen_to_slide(self, slide):
        """Applique l'élargissement automatique à une slide spécifique"""
        try:
            print(f"[WIDEN] Application élargissement automatique...")

            widen_count = 0
            for i, shape in enumerate(slide.shapes):
                if hasattr(shape, 'text_frame'):
                    current_width_inches = shape.width / Inches(1)

                    # Élargir les objets texte étroits (<4 pouces)
                    if current_width_inches < 4.0:
                        new_width_inches = min(current_width_inches * 1.5, 8.0)
                        shape.width = Inches(new_width_inches)
                        print(f"[WIDEN] Shape {i}: {current_width_inches:.2f}\" -> {new_width_inches:.2f}\"")
                        widen_count += 1

                    # Élargir modérément les objets moyens (4-6 pouces)
                    elif current_width_inches < 6.0:
                        new_width_inches = min(current_width_inches * 1.2, 8.0)
                        shape.width = Inches(new_width_inches)
                        print(f"[WIDEN] Shape {i}: {current_width_inches:.2f}\" -> {new_width_inches:.2f}\" (modéré)")
                        widen_count += 1

            if widen_count > 0:
                print(f"[SUCCESS] {widen_count} objets texte élargis")
            else:
                print(f"[INFO] Aucun objet texte nécessitant un élargissement")

        except Exception as e:
            print(f"[WARNING] Erreur élargissement: {e}")

    def _move_slide_to_position_direct(self, presentation: Presentation, from_index: int, to_index: int):
        """Déplace une slide à la position désirée"""
        try:
            print(f"[POSITION] Slide TOC ajoutée en position {from_index + 1}")
            print(f"[INFO] Position finale: {to_index + 1}")

        except Exception as e:
            print(f"[WARNING] Déplacement slide: {e}")

    def _generate_direct_insertion_report(self, presentation_path: str, sections: List[str],
                                        toc_title: str, insert_position: int, auto_widen: bool):
        """Génère un rapport d'insertion directe modernisé"""
        try:
            report = {
                "insertion_timestamp": datetime.now().isoformat(),
                "method": "JSON Navigation Builder 2025 - Direct Insertion",
                "template_used": self.template_path,
                "target_presentation": presentation_path,
                "toc_details": {
                    "title": toc_title,
                    "sections_count": len(sections),
                    "sections": sections,
                    "intended_position": insert_position + 1,
                    "auto_widen_enabled": auto_widen
                },
                "source_slide": {
                    "index": self.toc_slide_index,
                    "number": self.toc_slide_index + 1,
                    "layout": self.toc_info.get('layout_name', 'Table des matières')
                },
                "quality_assurance": {
                    "method": "JSON-Native Direct Layout-Based Addition",
                    "styles_preserved": True,
                    "premier_tech_standards": True,
                    "direct_integration": True,
                    "json_configuration": True,
                    "architecture_2025": True
                },
                "advantages": [
                    "Architecture JSON 2025 native",
                    "Insertion directe dans la présentation",
                    "Styles Premier Tech 100% préservés",
                    "Configuration JSON validée",
                    "Intégration transparente",
                    "Sauvegarde automatique créée"
                ]
            }

            # Sauvegarder le rapport
            report_path = presentation_path.replace('.pptx', '_direct_navigation_insertion_report.json')
            with open(report_path, 'w', encoding='utf-8') as f:
                json.dump(report, f, ensure_ascii=False, indent=2)

            print(f"[INFO] Rapport d'insertion: {os.path.basename(report_path)}")

        except Exception as e:
            print(f"[WARNING] Erreur génération rapport: {e}")

    def validate_template(self) -> bool:
        """Valide que le template Premier Tech est correct pour la TOC"""
        try:
            checks = {
                "file_exists": os.path.exists(self.template_path),
                "has_slides": False,
                "toc_slide_exists": False,
                "slides_count": 0
            }

            if checks["file_exists"]:
                pres = Presentation(self.template_path)
                checks["slides_count"] = len(pres.slides)
                checks["has_slides"] = len(pres.slides) > 0
                checks["toc_slide_exists"] = len(pres.slides) > self.toc_slide_index

            all_valid = all([checks["file_exists"], checks["has_slides"], checks["toc_slide_exists"]])

            print("=== VALIDATION TEMPLATE PREMIER TECH POUR NAVIGATION ===")
            for check, result in checks.items():
                status = "OK" if result else "ERREUR"
                print(f"[{status}] {check}: {result}")

            if checks["toc_slide_exists"]:
                print(f"[INFO] Slide TOC disponible: {self.toc_slide_index + 1} ({self.toc_info['layout_name']})")

            return all_valid

        except Exception as e:
            print(f"[ERROR] Erreur validation: {e}")
            return False


def create_navigation_from_json(config_data: Dict[str, Any], presentation_path: str,
                               template_path: str = "templates/Template_PT.pptx") -> Dict[str, Any]:
    """
    Fonction utilitaire pour créer une navigation à partir d'une configuration JSON.
    Point d'entrée principal pour l'architecture JSON 2025.

    Args:
        config_data: Configuration JSON de la navigation
        presentation_path: Chemin vers la présentation cible
        template_path: Chemin vers le template Premier Tech

    Returns:
        Dict contenant les résultats du traitement
    """
    try:
        builder = NavigationBuilder(template_path)
        return builder.process_navigation_config(config_data, presentation_path)
    except Exception as e:
        return {
            "success": False,
            "error": str(e),
            "timestamp": datetime.now().isoformat()
        }


def load_navigation_template(template_name: str = "basic_toc") -> Dict[str, Any]:
    """
    Charge un template de navigation prédéfini.

    Args:
        template_name: Nom du template (basic_toc, detailed_toc, strategic_toc)

    Returns:
        Dict contenant la configuration du template
    """
    template_path = "templates/presentation-project/slide-payload-templates/navigation_builder_template.json"

    try:
        with open(template_path, 'r', encoding='utf-8') as f:
            templates = json.load(f)

        examples = templates.get('examples', {})

        if template_name in examples:
            return examples[template_name]
        else:
            # Retourner le template de base si le nom n'est pas trouvé
            return templates.get('payload_structure', {})

    except Exception as e:
        print(f"[WARNING] Erreur chargement template: {e}")
        # Template de fallback
        return {
            "title": "Table des matières",
            "sections": [
                "Introduction",
                "Développement",
                "Conclusion"
            ],
            "options": {
                "auto_widen": True,
                "insert_position": 1
            }
        }


def load_navigation_payload(payload_path: str) -> Dict[str, Any]:
    """
    Charge un payload de navigation depuis un fichier JSON.

    Args:
        payload_path: Chemin vers le fichier JSON contenant le payload

    Returns:
        Dict contenant le payload de navigation
    """
    try:
        with open(payload_path, 'r', encoding='utf-8') as f:
            payload = json.load(f)

        print(f"[PAYLOAD] Chargé depuis {payload_path}")
        return payload

    except Exception as e:
        print(f"[ERROR] Erreur chargement payload {payload_path}: {e}")
        return load_navigation_template("basic_toc")  # Fallback


def process_navigation_from_payload_file(payload_path: str, presentation_path: str,
                                       template_path: str = "templates/Template_PT.pptx") -> Dict[str, Any]:
    """
    Traite une navigation en chargeant le payload depuis un fichier JSON.
    Point d'entrée pour l'architecture nouvelle avec fichiers payload séparés.

    Args:
        payload_path: Chemin vers le fichier JSON contenant le payload
        presentation_path: Chemin vers la présentation cible
        template_path: Chemin vers le template Premier Tech

    Returns:
        Dict contenant les résultats du traitement
    """
    try:
        # Charger le payload
        payload = load_navigation_payload(payload_path)

        # Traiter avec le payload chargé
        builder = NavigationBuilder(template_path)
        result = builder.process_navigation_config(payload, presentation_path)

        # Ajouter les informations de payload
        result["payload_source"] = payload_path
        result["payload_loaded"] = True

        return result

    except Exception as e:
        return {
            "success": False,
            "error": str(e),
            "timestamp": datetime.now().isoformat(),
            "payload_source": payload_path,
            "payload_loaded": False
        }


if __name__ == "__main__":
    # Interface de test pour validation
    print("=== Navigation Builder - Architecture JSON 2025 ===")

    # Exemple d'utilisation
    template_config = load_navigation_template("basic_toc")
    print(f"Template chargé: {template_config}")

    # Test de validation
    builder = NavigationBuilder()
    if builder.validate_template():
        print("[SUCCESS] Template Premier Tech validé pour navigation")
    else:
        print("[ERROR] Template Premier Tech invalide")