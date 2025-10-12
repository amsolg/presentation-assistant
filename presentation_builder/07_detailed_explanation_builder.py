#!/usr/bin/env python3
"""
Detailed Explanation Builder - Création de slides d'explications détaillées Premier Tech
Utilise les slides 35-44 du template Premier Tech pour créer des explications approfondies.
Script spécialisé pour le besoin "Explications détaillées" selon le Guide de Création Premier Tech.
"""

import os
import sys
import json
import shutil
import argparse
from datetime import datetime
from pathlib import Path
from typing import List, Dict, Optional, Any
from pptx import Presentation
from pptx.util import Inches


class DetailedExplanationBuilder:
    """
    Classe pour construire des slides d'explications détaillées Premier Tech.
    Utilise les slides 35-44 du template pour créer des explications approfondies.
    Script spécialisé pour le besoin "Explications détaillées" selon le Guide de Création Premier Tech.
    """

    def __init__(self, template_path: str = "templates/Template_PT.pptx"):
        """
        Initialise le constructeur avec le template Premier Tech.

        Args:
            template_path: Chemin vers le template Premier Tech
        """
        self.template_path = template_path

        # Mapping des slides disponibles pour les explications détaillées (slides 35-44)
        self.explanation_slides = {
            34: {  # Slide 35 (index 34) - 4 énoncés & Mots-clés
                "name": "4 énoncés & Mots-clés",
                "usage": "Explication détaillée avec 4 points clés",
                "audience": "Managers, Executives",
                "style": "four_points",
                "elements": 4,
                "explanation_type": "structured_points"
            },
            # TEMPORAIREMENT COMMENTÉ - Gestion d'images à implémenter plus tard
            # 35: {  # Slide 36 (index 35) - 2 énoncés avec sous-titres et image
            #     "name": "2 énoncés avec sous-titres et image",
            #     "usage": "Explication illustrée avec détails",
            #     "audience": "Toutes audiences",
            #     "style": "illustrated_detailed",
            #     "elements": 2,
            #     "explanation_type": "visual_explanation"
            # },
            # 36: {  # Slide 37 (index 36) - Énoncé avec titre et image
            #     "name": "Énoncé avec titre et image",
            #     "usage": "Concept principal illustré",
            #     "audience": "Toutes audiences",
            #     "style": "concept_visual",
            #     "elements": 1,
            #     "explanation_type": "single_concept"
            # },
            # 37: {  # Slide 38 (index 37) - Liste avec titre et image
            #     "name": "Liste avec titre et image",
            #     "usage": "Énumération illustrée",
            #     "audience": "Toutes audiences",
            #     "style": "list_visual",
            #     "elements": "multiple",
            #     "explanation_type": "illustrated_list"
            # },
            38: {  # Slide 39 (index 38) - 2 énoncés avec sous-titres et ligne bleue
                "name": "2 énoncés avec sous-titres et ligne bleue",
                "usage": "Comparaison détaillée avec ligne bleue",
                "audience": "Managers, Analystes",
                "style": "dual_detailed_blue",
                "elements": 2,
                "explanation_type": "dual_comparison"
            },
            39: {  # Slide 40 (index 39) - 2 énoncés avec sous-titres et ligne grise
                "name": "2 énoncés avec sous-titres et ligne grise",
                "usage": "Comparaison détaillée avec ligne grise",
                "audience": "Managers, Analystes",
                "style": "dual_detailed_grey",
                "elements": 2,
                "explanation_type": "dual_comparison"
            },
            40: {  # Slide 41 (index 40) - 2 énoncés avec titre et ligne bleue
                "name": "2 énoncés avec titre et ligne bleue",
                "usage": "Comparaison avec titre et ligne bleue",
                "audience": "Toutes audiences",
                "style": "dual_titled_blue",
                "elements": 2,
                "explanation_type": "titled_comparison"
            },
            41: {  # Slide 42 (index 41) - 2 énoncés avec titre et ligne grise
                "name": "2 énoncés avec titre et ligne grise",
                "usage": "Comparaison avec titre et ligne grise",
                "audience": "Toutes audiences",
                "style": "dual_titled_grey",
                "elements": 2,
                "explanation_type": "titled_comparison"
            },
            42: {  # Slide 43 (index 42) - 2 listes avec sous-titres et ligne bleue
                "name": "2 listes avec sous-titres et ligne bleue",
                "usage": "Listes structurées avec ligne bleue",
                "audience": "Managers, Analystes",
                "style": "dual_lists_blue",
                "elements": 2,
                "explanation_type": "dual_lists"
            },
            43: {  # Slide 44 (index 43) - 2 listes avec sous-titres et ligne grise
                "name": "2 listes avec sous-titres et ligne grise",
                "usage": "Listes structurées avec ligne grise",
                "audience": "Managers, Analystes",
                "style": "dual_lists_grey",
                "elements": 2,
                "explanation_type": "dual_lists"
            }
        }

        if not os.path.exists(template_path):
            raise FileNotFoundError(f"Template Premier Tech non trouvé: {template_path}")

        # Analyser la structure des slides d'explications détaillées
        self._analyze_explanation_structures()

    def _analyze_explanation_structures(self):
        """Analyse la structure des slides d'explications détaillées disponibles"""
        try:
            pres = Presentation(self.template_path)

            self.explanation_info = {}
            for slide_index, slide_data in self.explanation_slides.items():
                if len(pres.slides) > slide_index:
                    explanation_slide = pres.slides[slide_index]

                    self.explanation_info[slide_index] = {
                        'layout_name': explanation_slide.slide_layout.name,
                        'shape_count': len(explanation_slide.shapes),
                        'slide_index': slide_index,
                        'slide_number': slide_index + 1,
                        'style': slide_data['style'],
                        'usage': slide_data['usage'],
                        'elements': slide_data['elements'],
                        'explanation_type': slide_data['explanation_type']
                    }

            print(f"[INFO] {len(self.explanation_info)} slides d'explications détaillées analysées")
            for idx, info in self.explanation_info.items():
                print(f"[INFO] Slide {info['slide_number']}: {info['layout_name']} ({info['style']})")

        except Exception as e:
            raise Exception(f"Erreur analyse templates explications détaillées: {e}")

    def create_explanation_slide(self,
                               content: str,
                               subtitle: Optional[str] = None,
                               additional_content: Optional[List[str]] = None,
                               title: Optional[str] = None,
                               explanation_style: str = "concept_visual",
                               output_path: Optional[str] = None,
                               auto_widen: bool = True) -> str:
        """
        Crée une slide d'explication détaillée en clonant la slide appropriée du template.

        Args:
            content: Contenu principal de l'explication
            subtitle: Sous-titre ou description (optionnel)
            additional_content: Contenu additionnel (liste de points, etc.)
            title: Titre de la slide (optionnel)
            explanation_style: Style d'explication (voir styles disponibles)
            output_path: Chemin de sortie (optionnel)
            auto_widen: Active l'élargissement automatique des objets texte

        Returns:
            str: Chemin vers le fichier créé
        """
        try:
            # Déterminer la slide à utiliser selon le style
            slide_index = self._get_slide_index_for_style(explanation_style)

            if slide_index is None:
                available_styles = list(self.explanation_slides.values())
                available_style_names = [s['style'] for s in available_styles]
                raise ValueError(f"Style '{explanation_style}' non reconnu. Styles disponibles: {', '.join(available_style_names)}")

            # Vérifier la cohérence du contenu avec le style
            elements_count = self.explanation_slides[slide_index]['elements']
            has_subtitle = explanation_style in ['illustrated_detailed', 'concept_visual']

            if has_subtitle and not subtitle:
                print(f"[WARNING] Le style '{explanation_style}' supporte les sous-titres mais aucun n'est fourni")
            elif not has_subtitle and subtitle:
                print(f"[WARNING] Le style '{explanation_style}' ne supporte pas les sous-titres, il sera ignoré")

            # Vérifier que le contenu additionnel est cohérent avec le style
            if elements_count == "multiple" and not additional_content:
                print(f"[WARNING] Le style '{explanation_style}' supporte du contenu additionnel (liste de points)")
            elif elements_count != "multiple" and additional_content:
                print(f"[WARNING] Le style '{explanation_style}' ne supporte pas de contenu multiple, il sera ignoré")

            # Générer le chemin de sortie si non fourni
            if not output_path:
                output_path = self._generate_explanation_output_path(content, explanation_style)

            # Créer le dossier parent si nécessaire
            os.makedirs(os.path.dirname(output_path), exist_ok=True)

            print(f"[INFO] Création explication détaillée avec slide {slide_index + 1} du template ({explanation_style})")

            # ÉTAPE 1: Cloner la slide d'explication du template avec préservation complète des styles
            success = self._clone_explanation_slide(slide_index, output_path)
            if not success:
                raise Exception(f"Échec du clonage de la slide d'explication {slide_index + 1}")

            print(f"[SUCCESS] Slide d'explication clonée avec styles Premier Tech préservés")

            # ÉTAPE 2: Élargir automatiquement les objets texte
            widen_info = self._widen_text_objects(output_path, auto_widen=auto_widen)

            # ÉTAPE 2.5: Gérer le text wrapping selon le style
            self._configure_text_wrapping(output_path, explanation_style)

            # ÉTAPE 3: Personnaliser le contenu en préservant les styles
            self._customize_explanation_content(output_path, content, subtitle,
                                              additional_content, title, explanation_style)

            print(f"[SUCCESS] Slide d'explication détaillée créée: {output_path}")

            return output_path

        except Exception as e:
            print(f"[ERROR] Erreur création explication détaillée: {e}")
            raise

    def _get_slide_index_for_style(self, style: str) -> Optional[int]:
        """Retourne l'index de slide approprié pour le style demandé"""
        style_mapping = {
            "four_points": 34,           # Slide 35 - 4 énoncés & Mots-clés
            "illustrated_detailed": 35,   # Slide 36 - 2 énoncés avec sous-titres et image
            "concept_visual": 36,        # Slide 37 - Énoncé avec titre et image
            "list_visual": 37,           # Slide 38 - Liste avec titre et image
            "dual_detailed_blue": 38,    # Slide 39 - 2 énoncés avec sous-titres et ligne bleue
            "dual_detailed_grey": 39,    # Slide 40 - 2 énoncés avec sous-titres et ligne grise
            "dual_titled_blue": 40,      # Slide 41 - 2 énoncés avec titre et ligne bleue
            "dual_titled_grey": 41,      # Slide 42 - 2 énoncés avec titre et ligne grise
            "dual_lists_blue": 42,       # Slide 43 - 2 listes avec sous-titres et ligne bleue
            "dual_lists_grey": 43        # Slide 44 - 2 listes avec sous-titres et ligne grise
        }
        return style_mapping.get(style)

    def _clone_explanation_slide(self, slide_index: int, output_file: str) -> bool:
        """
        Clone la slide d'explication du template avec préservation complète des styles Premier Tech.
        Utilise la même méthode que les autres builders.
        """
        try:
            print(f"[CLONE] Copie complète du template...")

            # ÉTAPE 1: Copier le template complet pour préserver tous les styles
            shutil.copy2(self.template_path, output_file)

            # ÉTAPE 2: Charger et nettoyer pour ne garder que la slide comparaison désirée
            prs = Presentation(output_file)

            if slide_index >= len(prs.slides):
                print(f"[ERROR] Slide {slide_index + 1} n'existe pas dans le template")
                return False

            print(f"[CLONE] Suppression des slides non désirées (garder seulement slide {slide_index + 1})...")

            # ÉTAPE 3: Identifier toutes les slides à supprimer
            slides_to_remove = []
            for i in range(len(prs.slides)):
                if i != slide_index:
                    slides_to_remove.append(i)

            print(f"[CLONE] Suppression de {len(slides_to_remove)} slides sur {len(prs.slides)} total")

            # ÉTAPE 4: Supprimer en ordre inverse pour éviter les problèmes d'index
            for i in reversed(slides_to_remove):
                try:
                    rId = prs.slides._sldIdLst[i].rId
                    prs.part.drop_rel(rId)
                    del prs.slides._sldIdLst[i]
                except Exception as e:
                    print(f"[WARNING] Erreur suppression slide {i}: {e}")

            # ÉTAPE 5: Sauvegarder la présentation avec seulement la slide comparaison clonée
            prs.save(output_file)

            print(f"[SUCCESS] Slide comparaison {slide_index + 1} clonée avec styles Premier Tech intacts")
            return True

        except Exception as e:
            print(f"[ERROR] Erreur clonage slide comparaison {slide_index + 1}: {e}")
            return False

    def _widen_text_objects(self, presentation_path: str, auto_widen: bool = True):
        """
        Élargit automatiquement les objets texte pour éviter les retours à la ligne.
        Méthode identique aux autres builders.
        """
        widen_info = {
            "enabled": auto_widen,
            "objects_widened": 0,
            "modifications": []
        }

        if not auto_widen:
            return widen_info

        try:
            print(f"[WIDEN] Élargissement automatique des objets texte...")

            prs = Presentation(presentation_path)
            slide = prs.slides[0]  # Première (et unique) slide

            widen_count = 0
            for i, shape in enumerate(slide.shapes):
                if hasattr(shape, 'text_frame'):
                    current_width_inches = shape.width / Inches(1)

                    # Élargir les objets texte étroits (<4 pouces)
                    if current_width_inches < 4.0:
                        # Calculer la nouvelle largeur (1.5x mais max 8 pouces)
                        new_width_inches = min(current_width_inches * 1.5, 8.0)
                        shape.width = Inches(new_width_inches)

                        print(f"[WIDEN] Shape {i}: {current_width_inches:.2f}\" -> {new_width_inches:.2f}\"")
                        widen_count += 1
                        widen_info["modifications"].append({
                            "shape_index": i,
                            "before_inches": round(current_width_inches, 2),
                            "after_inches": round(new_width_inches, 2),
                            "type": "significant"
                        })

                    # Élargir modérément les objets moyens (4-6 pouces)
                    elif current_width_inches < 6.0:
                        # Élargissement plus modéré (1.2x mais max 8 pouces)
                        new_width_inches = min(current_width_inches * 1.2, 8.0)
                        shape.width = Inches(new_width_inches)

                        print(f"[WIDEN] Shape {i}: {current_width_inches:.2f}\" -> {new_width_inches:.2f}\" (modere)")
                        widen_count += 1
                        widen_info["modifications"].append({
                            "shape_index": i,
                            "before_inches": round(current_width_inches, 2),
                            "after_inches": round(new_width_inches, 2),
                            "type": "moderate"
                        })

            widen_info["objects_widened"] = widen_count

            if widen_count > 0:
                prs.save(presentation_path)
                print(f"[SUCCESS] {widen_count} objets texte élargis pour éviter les retours à la ligne")
            else:
                print(f"[INFO] Aucun objet texte nécessitant un élargissement")

            return widen_info

        except Exception as e:
            print(f"[WARNING] Erreur élargissement objets texte: {e}")
            widen_info["error"] = str(e)
            return widen_info

    def _configure_text_wrapping(self, presentation_path: str, explanation_style: str):
        """
        Configure le renvoi à la ligne selon le style d'explication.
        """
        try:
            print(f"[WRAP] Configuration text wrapping pour style {explanation_style}...")

            prs = Presentation(presentation_path)
            slide = prs.slides[0]  # Première (et unique) slide

            # Activer le wrapping pour tous les styles
            enable_wrapping = True  # Toujours activé maintenant

            wrap_count = 0
            for i, shape in enumerate(slide.shapes):
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    shape.text_frame.word_wrap = enable_wrapping
                    print(f"[WRAP] Shape {i}: Word wrap {'activé' if enable_wrapping else 'désactivé'}")
                    wrap_count += 1

            if wrap_count > 0:
                prs.save(presentation_path)
                print(f"[SUCCESS] Text wrapping configuré sur {wrap_count} objets texte")
            else:
                print(f"[INFO] Aucun objet texte trouvé")

        except Exception as e:
            print(f"[WARNING] Erreur configuration text wrapping: {e}")

    def _customize_explanation_content(self, presentation_path: str, content: str,
                                     subtitle: Optional[str], additional_content: Optional[List[str]],
                                     title: Optional[str], explanation_style: str):
        """
        Personnalise le contenu de la slide d'explication en préservant les styles Premier Tech.
        REMPLACE le contenu sans modifier les styles.
        """
        try:
            print(f"[CUSTOMIZE] Personnalisation du contenu d'explication...")

            # Charger la présentation clonée
            prs = Presentation(presentation_path)
            slide = prs.slides[0]  # Première (et unique) slide

            print(f"[CUSTOMIZE] Slide d'explication avec {len(slide.shapes)} shapes à traiter")
            print(f"[CUSTOMIZE] Style: {explanation_style}")

            updated_count = 0

            # Préparer le contenu selon le style
            all_content = [content]
            if additional_content:
                all_content.extend(additional_content)

            # Détection automatique des éléments à personnaliser
            title_updated = False
            content_updated = False
            content_index = 0

            for i, shape in enumerate(slide.shapes):
                try:
                    if hasattr(shape, 'text_frame') and shape.text_frame:
                        current_text = shape.text_frame.text.strip()
                        print(f"[DEBUG] Shape {i}: '{current_text}' (longueur: {len(current_text)})")

                        # Personnalisation du titre principal
                        if title and not title_updated and self._is_title_shape(current_text, i, shape):
                            shape.text_frame.text = title
                            print(f"[UPDATE] Shape {i}: Titre - {title}")
                            updated_count += 1
                            title_updated = True
                            continue

                        # Personnalisation du sous-titre
                        if subtitle and self._is_subtitle_shape(current_text, i):
                            shape.text_frame.text = subtitle
                            print(f"[UPDATE] Shape {i}: Sous-titre - {subtitle}")
                            updated_count += 1
                            continue

                        # Personnalisation du contenu principal et additionnel
                        if self._is_content_shape(current_text, i) and content_index < len(all_content):
                            shape.text_frame.text = all_content[content_index]
                            print(f"[UPDATE] Shape {i}: Contenu {content_index + 1} - {all_content[content_index]}")
                            updated_count += 1
                            content_index += 1

                except Exception as e:
                    print(f"[WARNING] Erreur personnalisation shape {i}: {e}")

            # Si pas assez de contenu personnalisé, utiliser approche par position
            if updated_count == 0:
                print(f"[FALLBACK] Détection automatique échouée, passage en mode position")
                self._customize_by_position_explanation(slide, title, subtitle, content, additional_content, explanation_style)
                updated_count = 1 + (1 if title else 0) + (1 if subtitle else 0) + len(additional_content or [])

            print(f"[SUCCESS] {updated_count} éléments personnalisés avec styles Premier Tech préservés")

            # Sauvegarder les modifications
            prs.save(presentation_path)

        except Exception as e:
            print(f"[ERROR] Erreur personnalisation explication: {e}")
            raise

    def _is_title_shape(self, text: str, index: int, shape=None) -> bool:
        """Détermine si un shape est le titre principal pour explications détaillées"""
        title_indicators = ["titre", "title", "heading", "header"]

        # 1. Vérifier les indicateurs textuels explicites
        has_title_indicator = any(indicator in text.lower() for indicator in title_indicators)
        if has_title_indicator:
            return True

        # 2. Si on a accès au shape, utiliser la position Y comme critère principal
        if shape and hasattr(shape, 'top'):
            # Le titre est généralement dans la partie supérieure (Y < 150)
            is_top_positioned = shape.top < 150
            if is_top_positioned and len(text) > 3:
                return True

        # 3. Fallback sur l'ancienne logique seulement si pas de position disponible
        return index == 0 and len(text) > 5

    def _is_subtitle_shape(self, text: str, index: int) -> bool:
        """Détermine si un shape est un sous-titre"""
        subtitle_indicators = ["sous-titre", "subtitle", "description", "explication"]
        return (any(indicator in text.lower() for indicator in subtitle_indicators) or
                (len(text) > 20 and len(text) < 100))

    def _is_content_shape(self, text: str, index: int) -> bool:
        """Détermine si un shape contient du contenu à personnaliser"""
        # Éviter les shapes de métadonnées ou de template
        metadata_indicators = ["premier tech", "©", "copyright", "date", "2024", "2025"]
        if any(indicator in text.lower() for indicator in metadata_indicators):
            return False
        # C'est probablement du contenu si c'est du texte non vide
        return len(text.strip()) > 2

    def insert_explanation_into_existing_presentation(self,
                                                    presentation_path: str,
                                                    content: str,
                                                    subtitle: Optional[str] = None,
                                                    additional_content: Optional[List[str]] = None,
                                                    title: Optional[str] = None,
                                                    explanation_style: str = "concept_visual",
                                                    insert_position: Optional[int] = None) -> str:
        """
        Insère une slide d'explication détaillée directement dans une présentation existante.

        Args:
            presentation_path: Chemin vers la présentation existante
            content: Contenu principal de l'explication
            subtitle: Sous-titre ou description (optionnel)
            additional_content: Contenu additionnel (liste de points, etc.)
            title: Titre de la slide (optionnel)
            explanation_style: Style d'explication (voir styles disponibles)
            insert_position: Position d'insertion (défaut: fin)

        Returns:
            str: Chemin vers le fichier modifié
        """
        try:
            print(f"[INSERT] Insertion directe explication dans: {os.path.basename(presentation_path)}")
            print(f"[INSERT] Style: {explanation_style}")

            # ÉTAPE 1: Créer une copie de sauvegarde
            backup_path = presentation_path.replace('.pptx', '_backup_before_explanation.pptx')
            shutil.copy2(presentation_path, backup_path)
            print(f"[BACKUP] Sauvegarde créée: {backup_path}")

            # ÉTAPE 2: Charger la présentation existante
            target_prs = Presentation(presentation_path)
            print(f"[LOAD] Présentation chargée: {len(target_prs.slides)} slides existantes")

            # ÉTAPE 3: Déterminer la slide source à utiliser
            source_slide_index = self._get_slide_index_for_style(explanation_style)
            if source_slide_index is None:
                available_styles = list(self.explanation_slides.values())
                available_style_names = [s['style'] for s in available_styles]
                raise ValueError(f"Style '{explanation_style}' non reconnu. Styles disponibles: {', '.join(available_style_names)}")

            # ÉTAPE 4: Vérifier que le layout existe dans la présentation cible
            explanation_layout_index = self._find_explanation_layout_index(target_prs, explanation_style)
            if explanation_layout_index is None:
                raise Exception(f"Layout pour style '{explanation_style}' non trouvé dans la présentation cible")

            # ÉTAPE 5: Ajouter la slide avec le bon layout
            explanation_layout = target_prs.slide_layouts[explanation_layout_index]
            new_slide = target_prs.slides.add_slide(explanation_layout)
            print(f"[ADD] Slide explication ajoutée avec layout: {explanation_layout.name}")

            # ÉTAPE 6: Personnaliser le contenu
            self._customize_explanation_slide_direct(new_slide, content, subtitle, additional_content, title, explanation_style)

            # ÉTAPE 7: Repositionner si nécessaire
            if insert_position is not None and insert_position < len(target_prs.slides) - 1:
                self._move_slide_to_position_direct(target_prs, len(target_prs.slides) - 1, insert_position)

            # ÉTAPE 8: Sauvegarder
            target_prs.save(presentation_path)

            print(f"[SUCCESS] Slide d'explication insérée dans: {presentation_path}")
            return presentation_path

        except Exception as e:
            print(f"[ERROR] Erreur insertion explication: {e}")
            raise

    def _customize_by_position_explanation(self, slide, title: Optional[str], subtitle: Optional[str],
                                         content: str, additional_content: Optional[List[str]],
                                         explanation_style: str):
        """Personnalise par position quand la détection automatique échoue"""
        content_items = [content]
        if additional_content:
            content_items.extend(additional_content)

        item_index = 0
        title_set = False
        subtitle_set = False

        for i, shape in enumerate(slide.shapes):
            if hasattr(shape, 'text_frame') and shape.text_frame:
                current_text = shape.text_frame.text.strip()

                # Premier shape non vide = titre
                if title and not title_set and len(current_text) > 0:
                    shape.text_frame.text = title
                    print(f"[POSITION] Shape {i}: Titre - {title}")
                    title_set = True
                    continue

                # Deuxième shape = sous-titre si disponible
                if subtitle and not subtitle_set and title_set and len(current_text) > 0:
                    shape.text_frame.text = subtitle
                    print(f"[POSITION] Shape {i}: Sous-titre - {subtitle}")
                    subtitle_set = True
                    continue

                # Shapes suivants = contenu
                if item_index < len(content_items) and len(current_text) > 0:
                    shape.text_frame.text = content_items[item_index]
                    print(f"[POSITION] Shape {i}: Contenu {item_index + 1} - {content_items[item_index]}")
                    item_index += 1

    def _find_explanation_layout_index(self, presentation: Presentation, explanation_style: str) -> Optional[int]:
        """Trouve l'index du layout d'explication dans la présentation"""
        try:
            # Obtenir le slide index pour ce style
            source_slide_index = self._get_slide_index_for_style(explanation_style)
            if source_slide_index is None:
                return None

            # Charger le template et obtenir le nom du layout
            template_prs = Presentation(self.template_path)
            template_layout_name = template_prs.slides[source_slide_index].slide_layout.name

            # Chercher ce layout dans la présentation cible
            for i, layout in enumerate(presentation.slide_layouts):
                if layout.name == template_layout_name:
                    print(f"[LAYOUT] Layout '{template_layout_name}' trouvé à l'index {i}")
                    return i

            print(f"[WARNING] Layout '{template_layout_name}' non trouvé dans la présentation")
            return None

        except Exception as e:
            print(f"[ERROR] Erreur recherche layout: {e}")
            return None

    def _customize_explanation_slide_direct(self, slide, content: str, subtitle: Optional[str],
                                          additional_content: Optional[List[str]], title: Optional[str],
                                          explanation_style: str):
        """
        Personnalise directement une slide d'explication dans une présentation existante.
        NOUVELLE VERSION: Utilise le mapping précis des shapes comme les scripts 02-06.
        """
        try:
            print(f"[CUSTOMIZE] Personnalisation améliorée slide d'explication ({explanation_style})")

            # Obtenir l'index de la slide source pour ce style
            source_slide_index = self._get_slide_index_for_style(explanation_style)
            if source_slide_index is None:
                print(f"[ERROR] Index de slide non trouvé pour style {explanation_style}")
                return

            # Obtenir le mapping précis pour cette slide
            shape_mapping = self._get_precise_shape_mapping(source_slide_index)
            if not shape_mapping:
                print(f"[FALLBACK] Mapping non trouvé pour slide {source_slide_index}, utilisation de l'ancienne méthode")
                return self._customize_explanation_slide_direct_fallback(slide, content, subtitle, additional_content, title, explanation_style)

            print(f"[MAPPING] Utilisation du pattern '{shape_mapping['pattern']}' pour slide {source_slide_index + 1}")

            # Convertir les shapes en liste pour accès par index
            shapes_list = list(slide.shapes)
            print(f"[SHAPES] {len(shapes_list)} shapes disponibles dans la slide")

            updates_made = 0

            # 1. Personnaliser le titre principal
            if title and shape_mapping.get("title_shapes"):
                for title_idx in shape_mapping["title_shapes"]:
                    if title_idx < len(shapes_list):
                        shape = shapes_list[title_idx]
                        if hasattr(shape, 'text_frame') and shape.text_frame:
                            original_text = shape.text_frame.text
                            shape.text_frame.text = title
                            shape.text_frame.word_wrap = True
                            print(f"[UPDATE] TITRE (index {title_idx}): '{original_text}' -> '{title}'")
                            updates_made += 1

            # 2. Pour slide 35 (four_points): Gérer les titres et contenus des points
            if shape_mapping["pattern"] == "four_points_with_subtitles":
                # Préparer le contenu principal et additionnel - assurer 4 points minimum
                all_content = [content]
                if additional_content:
                    all_content.extend(additional_content)

                # Si pas assez de contenu pour 4 points, dupliquer intelligemment
                while len(all_content) < 4:
                    if len(all_content) == 1:
                        # Diviser le contenu principal en aspects différents
                        base_content = all_content[0]
                        all_content.extend([
                            f"Architecture {base_content.split()[0] if base_content.split() else 'robuste'}",
                            f"Performance {base_content.split()[1] if len(base_content.split()) > 1 else 'optimisée'}",
                            f"Sécurité {base_content.split()[2] if len(base_content.split()) > 2 else 'renforcée'}"
                        ])
                    else:
                        # Ajouter des points génériques mais pertinents
                        all_content.append(f"Innovation continue pour point {len(all_content) + 1}")

                # Garder seulement les 4 premiers si plus de contenu
                all_content = all_content[:4]
                print(f"[INFO] Contenu préparé pour 4 points: {[c[:30]+'...' if len(c) > 30 else c for c in all_content]}")

                # Personnaliser les titres des points (mots-clés courts)
                point_title_indices = shape_mapping.get("point_title_shapes", [])
                for i, point_idx in enumerate(point_title_indices):
                    if i < len(all_content) and point_idx < len(shapes_list):
                        shape = shapes_list[point_idx]
                        if hasattr(shape, 'text_frame') and shape.text_frame:
                            # Générer un mot-clé pertinent pour chaque point
                            content_words = all_content[i].split()
                            if len(content_words) >= 2:
                                # Prendre le mot le plus significatif (éviter les mots courants)
                                significant_words = [w for w in content_words if len(w) > 4 and w.lower() not in ['avec', 'pour', 'dans', 'comme', 'cette', 'leurs']]
                                if significant_words:
                                    point_title_text = significant_words[0].capitalize()
                                else:
                                    point_title_text = " ".join(content_words[0:2]).title()
                            else:
                                point_title_text = all_content[i].split()[0].capitalize() if content_words else f"Point {i+1}"

                            original_text = shape.text_frame.text
                            shape.text_frame.text = point_title_text
                            shape.text_frame.word_wrap = True
                            print(f"[UPDATE] POINT {i+1} TITRE (index {point_idx}): '{original_text}' -> '{point_title_text}'")
                            updates_made += 1

                # Personnaliser les contenus des points
                point_content_indices = shape_mapping.get("point_content_shapes", [])
                for i, content_idx in enumerate(point_content_indices):
                    if i < len(all_content) and content_idx < len(shapes_list):
                        shape = shapes_list[content_idx]
                        if hasattr(shape, 'text_frame') and shape.text_frame:
                            original_text = shape.text_frame.text
                            shape.text_frame.text = all_content[i]
                            shape.text_frame.word_wrap = True
                            print(f"[UPDATE] POINT {i+1} CONTENU (index {content_idx}): '{original_text}' -> '{all_content[i][:50]}...'")
                            updates_made += 1

            # TEMPORAIREMENT COMMENTÉ - Gestion d'images à implémenter plus tard
            # # 3. Pour slide 36 (two_points_illustrated): Gérer titre, sous-titre principal et 2 points
            # elif shape_mapping["pattern"] == "two_points_illustrated":
            #     # Préparer le contenu pour 2 points
            #     all_content = [content]
            #     if additional_content:
            #         all_content.extend(additional_content)

            #     # S'assurer qu'on a exactement 2 points de contenu
            #     while len(all_content) < 2:
            #         all_content.append(f"Point complémentaire {len(all_content)}")
            #     all_content = all_content[:2]  # Garder seulement 2 points

            #     # 1. Gérer le sous-titre principal
            #     if subtitle and shape_mapping.get("main_subtitle_shape"):
            #         main_subtitle_idx = shape_mapping["main_subtitle_shape"][0]
            #         if main_subtitle_idx < len(shapes_list):
            #             shape = shapes_list[main_subtitle_idx]
            #             if hasattr(shape, 'text_frame') and shape.text_frame:
            #                 original_text = shape.text_frame.text
            #                 shape.text_frame.text = subtitle
            #                 shape.text_frame.word_wrap = True
            #                 print(f"[UPDATE] SOUS-TITRE PRINCIPAL (index {main_subtitle_idx}): '{original_text}' -> '{subtitle}'")
            #                 updates_made += 1

            #     # 2. Gérer les contenus des 2 points
            #     content_indices = shape_mapping.get("content_shapes", [])
            #     for i, content_idx in enumerate(content_indices):
            #         if i < len(all_content) and content_idx < len(shapes_list):
            #             shape = shapes_list[content_idx]
            #             if hasattr(shape, 'text_frame') and shape.text_frame:
            #                 original_text = shape.text_frame.text
            #                 shape.text_frame.text = all_content[i]
            #                 shape.text_frame.word_wrap = True
            #                 print(f"[UPDATE] CONTENU POINT {i+1} (index {content_idx}): '{original_text}' -> '{all_content[i][:50]}...'")
            #                 updates_made += 1

            #     # 3. Gérer les sous-titres des 2 points
            #     subtitle_indices = shape_mapping.get("subtitle_shapes", [])
            #     for i, subtitle_idx in enumerate(subtitle_indices):
            #         if i < len(all_content) and subtitle_idx < len(shapes_list):
            #             shape = shapes_list[subtitle_idx]
            #             if hasattr(shape, 'text_frame') and shape.text_frame:
            #                 # Générer un sous-titre descriptif pour chaque point
            #                 point_subtitle = f"Sous-titre"
            #                 words = all_content[i].split()
            #                 if len(words) > 1:
            #                     # Utiliser les 2 premiers mots comme sous-titre
            #                     point_subtitle = " ".join(words[:2]).title()

            #                 original_text = shape.text_frame.text
            #                 shape.text_frame.text = point_subtitle
            #                 shape.text_frame.word_wrap = True
            #                 print(f"[UPDATE] SOUS-TITRE POINT {i+1} (index {subtitle_idx}): '{original_text}' -> '{point_subtitle}'")
            #                 updates_made += 1

            # 3. Pour slide 39 (dual_detailed_blue): Gérer les 2 énoncés avec sous-titres
            elif shape_mapping["pattern"] == "dual_detailed_blue":
                # Préparer le contenu pour 2 côtés (gauche et droite)
                left_content = content  # Contenu principal va à gauche
                right_content = additional_content[0] if additional_content else "Contenu complémentaire"

                # Générer les sous-titres courts et percutants
                def generate_short_subtitle(content_text, default_fallback):
                    """Génère un sous-titre court (max 2-3 mots) basé sur le contenu"""
                    if not content_text:
                        return default_fallback

                    # Mots-clés technologiques importants
                    tech_keywords = {
                        'cloud': 'Cloud',
                        'microservices': 'Microservices',
                        'moderne': 'Moderne',
                        'scalable': 'Scalable',
                        'architecture': 'Architecture',
                        'monolithique': 'Monolithique',
                        'traditionnel': 'Traditionnel',
                        'legacy': 'Legacy',
                        'limitations': 'Limité',
                        'performance': 'Performance',
                        'sécurité': 'Sécurité',
                        'innovation': 'Innovation'
                    }

                    # Extraire les premiers mots-clés pertinents
                    words = content_text.lower().split()
                    found_keywords = []

                    for word in words:
                        for key, display in tech_keywords.items():
                            if key in word and display not in found_keywords:
                                found_keywords.append(display)
                                if len(found_keywords) >= 2:  # Max 2 mots-clés
                                    break
                        if len(found_keywords) >= 2:
                            break

                    if found_keywords:
                        return " ".join(found_keywords)
                    else:
                        # Utiliser les 2 premiers mots du contenu, nettoyés
                        first_words = " ".join(words[:2]).title()
                        return first_words if len(first_words) <= 15 else default_fallback

                # Générer les sous-titres courts basés sur le contenu
                left_subtitle = generate_short_subtitle(left_content, "Solution Moderne")
                right_subtitle = generate_short_subtitle(right_content, "Approche Legacy")

                print(f"[INFO] Configuration dual_detailed_blue - Gauche: '{left_subtitle}' / Droite: '{right_subtitle}'")

                # Mettre à jour le sous-titre gauche
                if shape_mapping.get("left_subtitle_shapes"):
                    for subtitle_idx in shape_mapping["left_subtitle_shapes"]:
                        if subtitle_idx < len(shapes_list):
                            shape = shapes_list[subtitle_idx]
                            if hasattr(shape, 'text_frame') and shape.text_frame:
                                original_text = shape.text_frame.text
                                shape.text_frame.text = left_subtitle
                                shape.text_frame.word_wrap = True
                                print(f"[UPDATE] SOUS-TITRE GAUCHE (index {subtitle_idx}): '{original_text}' -> '{left_subtitle}'")
                                updates_made += 1

                # Mettre à jour le contenu gauche
                if shape_mapping.get("left_content_shapes"):
                    for content_idx in shape_mapping["left_content_shapes"]:
                        if content_idx < len(shapes_list):
                            shape = shapes_list[content_idx]
                            if hasattr(shape, 'text_frame') and shape.text_frame:
                                original_text = shape.text_frame.text
                                shape.text_frame.text = left_content
                                shape.text_frame.word_wrap = True
                                print(f"[UPDATE] CONTENU GAUCHE (index {content_idx}): '{original_text}' -> '{left_content[:50]}...'")
                                updates_made += 1

                # Mettre à jour le sous-titre droite
                if shape_mapping.get("right_subtitle_shapes"):
                    for subtitle_idx in shape_mapping["right_subtitle_shapes"]:
                        if subtitle_idx < len(shapes_list):
                            shape = shapes_list[subtitle_idx]
                            if hasattr(shape, 'text_frame') and shape.text_frame:
                                original_text = shape.text_frame.text
                                shape.text_frame.text = right_subtitle
                                shape.text_frame.word_wrap = True
                                print(f"[UPDATE] SOUS-TITRE DROITE (index {subtitle_idx}): '{original_text}' -> '{right_subtitle}'")
                                updates_made += 1

                # Mettre à jour le contenu droite
                if shape_mapping.get("right_content_shapes"):
                    for content_idx in shape_mapping["right_content_shapes"]:
                        if content_idx < len(shapes_list):
                            shape = shapes_list[content_idx]
                            if hasattr(shape, 'text_frame') and shape.text_frame:
                                original_text = shape.text_frame.text
                                shape.text_frame.text = right_content
                                shape.text_frame.word_wrap = True
                                print(f"[UPDATE] CONTENU DROITE (index {content_idx}): '{original_text}' -> '{right_content[:50]}...'")
                                updates_made += 1

            # 4. Pour slide 40 (dual_detailed_grey): Gérer les 2 énoncés avec sous-titres (identique à dual_detailed_blue)
            elif shape_mapping["pattern"] == "dual_detailed_grey":
                # Préparer le contenu pour 2 côtés (gauche et droite)
                left_content = content  # Contenu principal va à gauche
                right_content = additional_content[0] if additional_content else "Contenu complémentaire"

                # Générer les sous-titres courts et percutants (réutiliser la fonction définie plus haut)
                def generate_short_subtitle(content_text, default_fallback):
                    """Génère un sous-titre court (max 2-3 mots) basé sur le contenu"""
                    if not content_text:
                        return default_fallback

                    # Mots-clés technologiques importants
                    tech_keywords = {
                        'cloud': 'Cloud',
                        'microservices': 'Microservices',
                        'moderne': 'Moderne',
                        'scalable': 'Scalable',
                        'architecture': 'Architecture',
                        'monolithique': 'Monolithique',
                        'traditionnel': 'Traditionnel',
                        'legacy': 'Legacy',
                        'limitations': 'Limité',
                        'performance': 'Performance',
                        'sécurité': 'Sécurité',
                        'innovation': 'Innovation',
                        'on-premise': 'On-Premise'
                    }

                    # Extraire les premiers mots-clés pertinents
                    words = content_text.lower().split()
                    found_keywords = []

                    for word in words:
                        for key, display in tech_keywords.items():
                            if key in word and display not in found_keywords:
                                found_keywords.append(display)
                                if len(found_keywords) >= 2:  # Max 2 mots-clés
                                    break
                        if len(found_keywords) >= 2:
                            break

                    if found_keywords:
                        return " ".join(found_keywords)
                    else:
                        # Utiliser les 2 premiers mots du contenu, nettoyés
                        first_words = " ".join(words[:2]).title()
                        return first_words if len(first_words) <= 15 else default_fallback

                # Générer les sous-titres courts basés sur le contenu
                left_subtitle = generate_short_subtitle(left_content, "Solution On-Premise")
                right_subtitle = generate_short_subtitle(right_content, "Solution Cloud")

                print(f"[INFO] Configuration dual_detailed_grey - Gauche: '{left_subtitle}' / Droite: '{right_subtitle}'")

                # Mettre à jour le sous-titre gauche
                if shape_mapping.get("left_subtitle_shapes"):
                    for subtitle_idx in shape_mapping["left_subtitle_shapes"]:
                        if subtitle_idx < len(shapes_list):
                            shape = shapes_list[subtitle_idx]
                            if hasattr(shape, 'text_frame') and shape.text_frame:
                                original_text = shape.text_frame.text
                                shape.text_frame.text = left_subtitle
                                shape.text_frame.word_wrap = True
                                print(f"[UPDATE] SOUS-TITRE GAUCHE (index {subtitle_idx}): '{original_text}' -> '{left_subtitle}'")
                                updates_made += 1

                # Mettre à jour le contenu gauche
                if shape_mapping.get("left_content_shapes"):
                    for content_idx in shape_mapping["left_content_shapes"]:
                        if content_idx < len(shapes_list):
                            shape = shapes_list[content_idx]
                            if hasattr(shape, 'text_frame') and shape.text_frame:
                                original_text = shape.text_frame.text
                                shape.text_frame.text = left_content
                                shape.text_frame.word_wrap = True
                                print(f"[UPDATE] CONTENU GAUCHE (index {content_idx}): '{original_text}' -> '{left_content[:50]}...'")
                                updates_made += 1

                # Mettre à jour le sous-titre droite
                if shape_mapping.get("right_subtitle_shapes"):
                    for subtitle_idx in shape_mapping["right_subtitle_shapes"]:
                        if subtitle_idx < len(shapes_list):
                            shape = shapes_list[subtitle_idx]
                            if hasattr(shape, 'text_frame') and shape.text_frame:
                                original_text = shape.text_frame.text
                                shape.text_frame.text = right_subtitle
                                shape.text_frame.word_wrap = True
                                print(f"[UPDATE] SOUS-TITRE DROITE (index {subtitle_idx}): '{original_text}' -> '{right_subtitle}'")
                                updates_made += 1

                # Mettre à jour le contenu droite
                if shape_mapping.get("right_content_shapes"):
                    for content_idx in shape_mapping["right_content_shapes"]:
                        if content_idx < len(shapes_list):
                            shape = shapes_list[content_idx]
                            if hasattr(shape, 'text_frame') and shape.text_frame:
                                original_text = shape.text_frame.text
                                shape.text_frame.text = right_content
                                shape.text_frame.word_wrap = True
                                print(f"[UPDATE] CONTENU DROITE (index {content_idx}): '{original_text}' -> '{right_content[:50]}...'")
                                updates_made += 1

            # 5. Pour slide 41 (dual_titled_blue): Gérer les 2 énoncés avec titre global
            elif shape_mapping["pattern"] == "dual_titled_blue":
                # Préparer le contenu pour 2 côtés (gauche et droite)
                left_content = content  # Contenu principal va à gauche
                right_content = additional_content[0] if additional_content else "Contenu complémentaire"

                print(f"[INFO] Configuration dual_titled_blue - Gauche: '{left_content[:30]}...' / Droite: '{right_content[:30]}...'")

                # Le titre est déjà géré dans le bloc général (ligne 669-678), pas besoin de le refaire ici

                # Mettre à jour le label central (rectangle du milieu) avec un thème ou catégorie
                if shape_mapping.get("central_label_shapes"):
                    # Déterminer le texte du label central basé sur le contexte
                    if subtitle:
                        central_label = subtitle  # Utiliser le sous-titre si fourni
                    else:
                        # Générer automatiquement un label central pertinent
                        central_label = "Comparaison"  # Par défaut pour dual_titled_blue

                    for label_idx in shape_mapping["central_label_shapes"]:
                        if label_idx < len(shapes_list):
                            shape = shapes_list[label_idx]
                            if hasattr(shape, 'text_frame') and shape.text_frame:
                                original_text = shape.text_frame.text
                                shape.text_frame.text = central_label
                                shape.text_frame.word_wrap = True
                                print(f"[UPDATE] LABEL CENTRAL (index {label_idx}): '{original_text}' -> '{central_label}'")
                                updates_made += 1

                # Mettre à jour le contenu gauche
                if shape_mapping.get("left_content_shapes"):
                    for content_idx in shape_mapping["left_content_shapes"]:
                        if content_idx < len(shapes_list):
                            shape = shapes_list[content_idx]
                            if hasattr(shape, 'text_frame') and shape.text_frame:
                                original_text = shape.text_frame.text
                                shape.text_frame.text = left_content
                                shape.text_frame.word_wrap = True
                                print(f"[UPDATE] CONTENU GAUCHE (index {content_idx}): '{original_text}' -> '{left_content[:50]}...'")
                                updates_made += 1

                # Mettre à jour le contenu droite
                if shape_mapping.get("right_content_shapes"):
                    for content_idx in shape_mapping["right_content_shapes"]:
                        if content_idx < len(shapes_list):
                            shape = shapes_list[content_idx]
                            if hasattr(shape, 'text_frame') and shape.text_frame:
                                original_text = shape.text_frame.text
                                shape.text_frame.text = right_content
                                shape.text_frame.word_wrap = True
                                print(f"[UPDATE] CONTENU DROITE (index {content_idx}): '{original_text}' -> '{right_content[:50]}...'")
                                updates_made += 1

                # Mettre à jour le sous-titre global si fourni
                if subtitle and shape_mapping.get("subtitle_shapes"):
                    for subtitle_idx in shape_mapping["subtitle_shapes"]:
                        if subtitle_idx < len(shapes_list):
                            shape = shapes_list[subtitle_idx]
                            if hasattr(shape, 'text_frame') and shape.text_frame:
                                original_text = shape.text_frame.text
                                shape.text_frame.text = subtitle
                                shape.text_frame.word_wrap = True
                                print(f"[UPDATE] SOUS-TITRE GLOBAL (index {subtitle_idx}): '{original_text}' -> '{subtitle}'")
                                updates_made += 1

            # 6. Pour slide 42 (dual_titled_grey): Gérer les 2 énoncés avec titre global (identique à dual_titled_blue)
            elif shape_mapping["pattern"] == "dual_titled_grey":
                # Préparer le contenu pour 2 côtés (gauche et droite)
                left_content = content  # Contenu principal va à gauche
                right_content = additional_content[0] if additional_content else "Contenu complémentaire"

                print(f"[INFO] Configuration dual_titled_grey - Gauche: '{left_content[:30]}...' / Droite: '{right_content[:30]}...'")

                # Le titre est déjà géré dans le bloc général (ligne 669-678), pas besoin de le refaire ici

                # Mettre à jour le contenu gauche
                if shape_mapping.get("left_content_shapes"):
                    for content_idx in shape_mapping["left_content_shapes"]:
                        if content_idx < len(shapes_list):
                            shape = shapes_list[content_idx]
                            if hasattr(shape, 'text_frame') and shape.text_frame:
                                original_text = shape.text_frame.text
                                shape.text_frame.text = left_content
                                shape.text_frame.word_wrap = True
                                print(f"[UPDATE] CONTENU GAUCHE (index {content_idx}): '{original_text}' -> '{left_content[:50]}...'")
                                updates_made += 1

                # Mettre à jour le contenu droite
                if shape_mapping.get("right_content_shapes"):
                    for content_idx in shape_mapping["right_content_shapes"]:
                        if content_idx < len(shapes_list):
                            shape = shapes_list[content_idx]
                            if hasattr(shape, 'text_frame') and shape.text_frame:
                                original_text = shape.text_frame.text
                                shape.text_frame.text = right_content
                                shape.text_frame.word_wrap = True
                                print(f"[UPDATE] CONTENU DROITE (index {content_idx}): '{original_text}' -> '{right_content[:50]}...'")
                                updates_made += 1

                # Mettre à jour le sous-titre global si fourni
                if subtitle and shape_mapping.get("subtitle_shapes"):
                    for subtitle_idx in shape_mapping["subtitle_shapes"]:
                        if subtitle_idx < len(shapes_list):
                            shape = shapes_list[subtitle_idx]
                            if hasattr(shape, 'text_frame') and shape.text_frame:
                                original_text = shape.text_frame.text
                                shape.text_frame.text = subtitle
                                shape.text_frame.word_wrap = True
                                print(f"[UPDATE] SOUS-TITRE GLOBAL (index {subtitle_idx}): '{original_text}' -> '{subtitle}'")
                                updates_made += 1

            # 7. Pour slides 43/44 (dual_lists_blue/grey): Gérer les 2 listes avec sous-titres
            elif shape_mapping["pattern"] in ["dual_lists_blue", "dual_lists_grey"]:
                # Préparer le contenu pour 2 listes (gauche et droite)
                left_content = content  # Contenu principal va à gauche
                right_content = additional_content[0] if additional_content else "Contenu liste complémentaire"

                # Fonction pour convertir en format liste
                def format_as_list(text_content, side_hint=""):
                    """Convertit le texte en format liste avec puces"""
                    if not text_content:
                        return "• Point principal\n• Point secondaire\n• Point tertiaire"

                    # Si déjà formaté comme une liste, conserver
                    if '\n' in text_content and any(line.strip().startswith(('•', '-', '*')) for line in text_content.split('\n')):
                        return text_content

                    # Sinon, créer des points de liste intelligents
                    words = text_content.split()
                    if len(words) <= 3:
                        # Texte court: créer des points descriptifs spécifiques
                        base_text = text_content.lower()

                        # Mapping spécifique pour architectures communes
                        if "monolithique" in base_text:
                            return "• Application unique et centralisée\n• Déploiement simplifié\n• Base de code unifiée"
                        elif "microservice" in base_text:
                            return "• Services indépendants et découplés\n• Déploiement par service\n• Technologies hétérogènes"
                        elif "architecture" in base_text:
                            # Générique pour architecture
                            base_word = words[-1] if len(words) > 1 else words[0] if words else "Solution"
                            return f"• {base_word.capitalize()} robuste\n• {base_word.capitalize()} évolutive\n• {base_word.capitalize()} performante"
                        else:
                            # Fallback générique avec plus de variété
                            base_word = words[0] if words else "Solution"
                            if side_hint == "left":
                                return f"• {base_word} traditionnel\n• {base_word} éprouvé\n• {base_word} stable"
                            else:
                                return f"• {base_word} moderne\n• {base_word} agile\n• {base_word} innovant"
                    else:
                        # Texte long: diviser en points logiques
                        points = []
                        chunk_size = max(2, len(words) // 3)
                        for i in range(0, min(len(words), 9), chunk_size):
                            chunk = ' '.join(words[i:i+chunk_size])
                            if chunk:
                                points.append(f"• {chunk}")
                        return '\n'.join(points[:3])  # Max 3 points

                # Générer les sous-titres courts pour les listes
                def generate_list_subtitle(content_text, default_fallback):
                    """Génère un sous-titre court pour une liste"""
                    if not content_text:
                        return default_fallback

                    # Mots-clés pour listes
                    list_keywords = {
                        'avantage': 'Avantages',
                        'bénéfice': 'Bénéfices',
                        'fonctionnalité': 'Fonctionnalités',
                        'caractéristique': 'Caractéristiques',
                        'limitation': 'Limitations',
                        'challenge': 'Défis',
                        'problème': 'Problèmes',
                        'solution': 'Solutions',
                        'innovation': 'Innovations',
                        'architecture': 'Architecture',
                        'microservice': 'Microservices',
                        'monolithique': 'Monolithique'
                    }

                    # Chercher des mots-clés dans le contenu
                    words = content_text.lower().split()
                    for word in words:
                        for key, display in list_keywords.items():
                            if key in word:
                                return display

                    # Utiliser le premier mot, capitalisé
                    if words:
                        return words[0].capitalize()
                    else:
                        return default_fallback

                # Formater le contenu en listes avec indications de côté
                left_list = format_as_list(left_content, "left")
                right_list = format_as_list(right_content, "right")

                # Générer les sous-titres
                left_subtitle = generate_list_subtitle(left_content, "Points Clés")
                right_subtitle = generate_list_subtitle(right_content, "Éléments Complémentaires")

                print(f"[INFO] Configuration {shape_mapping['pattern']} - Gauche: '{left_subtitle}' / Droite: '{right_subtitle}'")

                # Mettre à jour le sous-titre gauche
                if shape_mapping.get("left_subtitle_shapes"):
                    for subtitle_idx in shape_mapping["left_subtitle_shapes"]:
                        if subtitle_idx < len(shapes_list):
                            shape = shapes_list[subtitle_idx]
                            if hasattr(shape, 'text_frame') and shape.text_frame:
                                original_text = shape.text_frame.text
                                shape.text_frame.text = left_subtitle
                                shape.text_frame.word_wrap = True
                                print(f"[UPDATE] SOUS-TITRE GAUCHE (index {subtitle_idx}): '{original_text}' -> '{left_subtitle}'")
                                updates_made += 1

                # Mettre à jour le contenu liste gauche
                if shape_mapping.get("left_content_shapes"):
                    for content_idx in shape_mapping["left_content_shapes"]:
                        if content_idx < len(shapes_list):
                            shape = shapes_list[content_idx]
                            if hasattr(shape, 'text_frame') and shape.text_frame:
                                original_text = shape.text_frame.text
                                shape.text_frame.text = left_list
                                shape.text_frame.word_wrap = True
                                print(f"[UPDATE] LISTE GAUCHE (index {content_idx}): '{original_text}' -> '{left_list[:50]}...'")
                                updates_made += 1

                # Mettre à jour le sous-titre droite
                if shape_mapping.get("right_subtitle_shapes"):
                    for subtitle_idx in shape_mapping["right_subtitle_shapes"]:
                        if subtitle_idx < len(shapes_list):
                            shape = shapes_list[subtitle_idx]
                            if hasattr(shape, 'text_frame') and shape.text_frame:
                                original_text = shape.text_frame.text
                                shape.text_frame.text = right_subtitle
                                shape.text_frame.word_wrap = True
                                print(f"[UPDATE] SOUS-TITRE DROITE (index {subtitle_idx}): '{original_text}' -> '{right_subtitle}'")
                                updates_made += 1

                # Mettre à jour le contenu liste droite
                if shape_mapping.get("right_content_shapes"):
                    for content_idx in shape_mapping["right_content_shapes"]:
                        if content_idx < len(shapes_list):
                            shape = shapes_list[content_idx]
                            if hasattr(shape, 'text_frame') and shape.text_frame:
                                original_text = shape.text_frame.text
                                shape.text_frame.text = right_list
                                shape.text_frame.word_wrap = True
                                print(f"[UPDATE] LISTE DROITE (index {content_idx}): '{original_text}' -> '{right_list[:50]}...'")
                                updates_made += 1

            # 8. Pour autres patterns: gérer contenus simples
            else:
                all_content = [content]
                if additional_content:
                    all_content.extend(additional_content)

                content_indices = shape_mapping.get("content_shapes", [])
                for i, content_idx in enumerate(content_indices):
                    if i < len(all_content) and content_idx < len(shapes_list):
                        shape = shapes_list[content_idx]
                        if hasattr(shape, 'text_frame') and shape.text_frame:
                            original_text = shape.text_frame.text
                            shape.text_frame.text = all_content[i]
                            shape.text_frame.word_wrap = True
                            print(f"[UPDATE] CONTENU {i+1} (index {content_idx}): '{original_text}' -> '{all_content[i][:50]}...'")
                            updates_made += 1

                # Gérer les sous-titres si disponibles
                if subtitle and shape_mapping.get("subtitle_shapes"):
                    subtitle_indices = shape_mapping["subtitle_shapes"]
                    for subtitle_idx in subtitle_indices[:1]:  # Prendre le premier sous-titre disponible
                        if subtitle_idx < len(shapes_list):
                            shape = shapes_list[subtitle_idx]
                            if hasattr(shape, 'text_frame') and shape.text_frame:
                                original_text = shape.text_frame.text
                                shape.text_frame.text = subtitle
                                shape.text_frame.word_wrap = True
                                print(f"[UPDATE] SOUS-TITRE (index {subtitle_idx}): '{original_text}' -> '{subtitle}'")
                                updates_made += 1

            # 5. Vider les placeholders indésirables
            if shape_mapping.get("placeholder_shapes"):
                for placeholder_idx in shape_mapping["placeholder_shapes"]:
                    if placeholder_idx < len(shapes_list):
                        shape = shapes_list[placeholder_idx]
                        if hasattr(shape, 'text_frame') and shape.text_frame:
                            original_text = shape.text_frame.text
                            shape.text_frame.text = ""  # Vider le placeholder
                            print(f"[CLEAR] PLACEHOLDER (index {placeholder_idx}): '{original_text}' -> ''")
                            updates_made += 1

            # 6. Vérification finale - s'assurer que tous les points ont du contenu
            if shape_mapping["pattern"] == "four_points_with_subtitles":
                expected_points = 4
                actual_content_updates = len([u for u in range(updates_made) if "CONTENU" in str(u)])
                if actual_content_updates < expected_points:
                    print(f"[WARNING] Seulement {actual_content_updates}/{expected_points} points de contenu mis à jour")

            print(f"[SUCCESS] Slide d'explication personnalisée avec mapping précis: {updates_made} éléments mis à jour")

        except Exception as e:
            print(f"[ERROR] Erreur personnalisation slide directe: {e}")
            raise

    def _customize_explanation_slide_direct_fallback(self, slide, content: str, subtitle: Optional[str],
                                                   additional_content: Optional[List[str]], title: Optional[str],
                                                   explanation_style: str):
        """Méthode de fallback pour la personnalisation (ancienne logique)"""
        print(f"[FALLBACK] Utilisation de la méthode de personnalisation de secours")

        # Préparer le contenu - assurer suffisamment de contenu pour le style
        all_content = [content]
        if additional_content:
            all_content.extend(additional_content)

        # Pour four_points, assurer 4 éléments de contenu
        if explanation_style == "four_points" and len(all_content) < 4:
            base_content = content
            content_keywords = base_content.split()[:3] if len(base_content.split()) >= 3 else ["Innovation", "Performance", "Sécurité"]
            while len(all_content) < 4:
                keyword_idx = (len(all_content) - 1) % len(content_keywords)
                all_content.append(f"{content_keywords[keyword_idx]} avancée pour optimisation continue")

        # Utiliser la même logique que l'ancienne méthode
        content_items = all_content
        item_index = 0
        title_set = False
        subtitle_set = False

        for i, shape in enumerate(slide.shapes):
            if hasattr(shape, 'text_frame') and shape.text_frame:
                current_text = shape.text_frame.text.strip()

                # Premier shape non vide = titre
                if title and not title_set and len(current_text) > 0:
                    shape.text_frame.text = title
                    print(f"[FALLBACK] Shape {i}: Titre - {title}")
                    title_set = True
                    continue

                # Deuxième shape = sous-titre si disponible
                if subtitle and not subtitle_set and title_set and len(current_text) > 0:
                    shape.text_frame.text = subtitle
                    print(f"[FALLBACK] Shape {i}: Sous-titre - {subtitle}")
                    subtitle_set = True
                    continue

                # Shapes suivants = contenu
                if item_index < len(content_items) and len(current_text) > 0:
                    shape.text_frame.text = content_items[item_index]
                    print(f"[FALLBACK] Shape {i}: Contenu {item_index + 1} - {content_items[item_index]}")
                    item_index += 1

        total_updated = item_index + (1 if title_set else 0) + (1 if subtitle_set else 0)
        print(f"[FALLBACK] Personnalisation terminée: {total_updated} éléments mis à jour")

        # Vérification spéciale pour four_points
        if explanation_style == "four_points" and item_index < 4:
            print(f"[FALLBACK-WARNING] Style four_points nécessite 4 points, seulement {item_index} assignés")

    def _get_precise_shape_mapping(self, slide_index: int) -> Optional[Dict[str, Any]]:
        """
        Retourne le mapping précis des shapes pour chaque slide d'explication.
        Basé sur l'analyse réelle de la structure des templates Premier Tech.
        """
        try:
            # Mapping précis pour slide 35 (index 34) - 4 énoncés & Mots-clés
            if slide_index == 34:  # Slide 35
                return {
                    "pattern": "four_points_with_subtitles",
                    "description": "4 énoncés avec titres et contenus séparés",
                    "title_shapes": [0],  # Shape 0: Titre principal
                    "point_title_shapes": [1, 3, 5, 7],  # Shapes 1,3,5,7: Titres des 4 points (mots-clés)
                    "point_content_shapes": [2, 4, 6, 8],  # Shapes 2,4,6,8: Contenus des 4 points
                    "page_number_shapes": [9],  # Shape 9: Numéro de page (à ignorer)
                    "total_shapes": 10
                }

            # TEMPORAIREMENT COMMENTÉ - Gestion d'images à implémenter plus tard
            # # Mapping précis pour slide 36 (index 35) - 2 énoncés avec sous-titres et image
            # elif slide_index == 35:  # Slide 36
            #     return {
            #         "pattern": "two_points_illustrated",
            #         "description": "2 énoncés avec sous-titres et espace pour image",
            #         "title_shapes": [0],  # Shape 0: Titre principal (position problématique corrigée)
            #         "subtitle_shapes": [3, 5],  # Shapes 3,5: Sous-titres des 2 points (actuellement vides)
            #         "content_shapes": [2, 4],  # Shapes 2,4: Contenus des 2 points
            #         "main_subtitle_shape": [1],  # Shape 1: Sous-titre principal
            #         "image_placeholder_shapes": [],  # Pas de placeholder dans les shapes analysées
            #         "total_shapes": 6
            #     }

            # # Mapping précis pour slide 37 (index 36) - Énoncé avec titre et image
            # elif slide_index == 36:  # Slide 37
            #     return {
            #         "pattern": "single_concept_visual",
            #         "description": "Concept unique avec titre et image",
            #         "title_shapes": [0],  # Titre principal
            #         "content_shapes": [1],  # Contenu principal
            #         "image_placeholder_shapes": [2],  # Espace pour image
            #         "page_number_shapes": [3],  # Numéro de page
            #         "total_shapes": 4
            #     }

            # Mapping précis pour slide 39 (index 38) - 2 énoncés avec sous-titres et ligne bleue
            elif slide_index == 38:  # Slide 39
                return {
                    "pattern": "dual_detailed_blue",
                    "description": "2 énoncés avec sous-titres et ligne bleue",
                    "title_shapes": [0],  # Shape 0 = Titre principal (position la plus haute)
                    "left_subtitle_shapes": [2],  # Shape 2 = Sous-titre gauche (gauche, position basse = sous-titre)
                    "left_content_shapes": [1],  # Shape 1 = Contenu gauche (gauche, position haute = contenu principal)
                    "right_subtitle_shapes": [4],  # Shape 4 = Sous-titre droite (droite, position basse = sous-titre)
                    "right_content_shapes": [3],  # Shape 3 = Contenu droite (droite, position haute = contenu principal)
                    "page_number_shapes": [5] if "page_number_shapes" in str(slide_index) else [],  # Shape 5 = Numéro de page (optionnel)
                    "total_shapes": 5
                }

            # Mapping précis pour slide 40 (index 39) - 2 énoncés avec sous-titres et ligne grise
            elif slide_index == 39:  # Slide 40
                return {
                    "pattern": "dual_detailed_grey",
                    "description": "2 énoncés avec sous-titres et ligne grise",
                    "title_shapes": [0],  # Shape 0 = Titre principal (shape_1/ID:2)
                    "left_subtitle_shapes": [2],  # Shape 2 = Sous-titre gauche (shape_3/ID:4)
                    "left_content_shapes": [1],  # Shape 1 = Contenu gauche (shape_2/ID:3)
                    "right_subtitle_shapes": [4],  # Shape 4 = Sous-titre droite (shape_5/ID:8)
                    "right_content_shapes": [3],  # Shape 3 = Contenu droite (shape_4/ID:7)
                    "page_number_shapes": [5],  # Shape 5 = Numéro de page (shape_6/ID:6)
                    "total_shapes": 6
                }

            # Mapping précis pour slide 41 (index 40) - 2 énoncés avec titre et ligne bleue
            elif slide_index == 40:  # Slide 41
                return {
                    "pattern": "dual_titled_blue",
                    "description": "2 énoncés avec titre et ligne bleue",
                    "title_shapes": [2],  # Shape 2 = Titre principal (shape_3/ID:4, position la plus haute: 121.51)
                    "central_label_shapes": [0],  # Shape 0 = Rectangle central pour catégorie/thème (shape_1/ID:2, position: 176.1)
                    "left_content_shapes": [1],  # Shape 1 = Contenu gauche (shape_2/ID:3, position basse gauche: 292.85)
                    "right_content_shapes": [3],  # Shape 3 = Contenu droite (shape_4/ID:5, position basse droite: 292.85)
                    "page_number_shapes": [4],  # Shape 4 = Numéro de page (shape_5/ID:6)
                    "total_shapes": 5
                }

            # Mapping précis pour slide 42 (index 41) - 2 énoncés avec titre et ligne grise
            elif slide_index == 41:  # Slide 42
                return {
                    "pattern": "dual_titled_grey",
                    "description": "2 énoncés avec titre et ligne grise",
                    "title_shapes": [2],  # Shape 2 = Titre principal (shape_3/ID:4, position la plus haute: 121.51)
                    "subtitle_shapes": [0],  # Shape 0 = Sous-titre global (shape_1/ID:2, position moyenne: 176.1)
                    "left_content_shapes": [1],  # Shape 1 = Contenu gauche (shape_2/ID:3, position basse gauche: 292.85)
                    "right_content_shapes": [3],  # Shape 3 = Contenu droite (shape_4/ID:5, position basse droite: 292.85)
                    "page_number_shapes": [4],  # Shape 4 = Numéro de page (shape_5/ID:6)
                    "total_shapes": 5
                }

            # Mapping précis pour slide 43 (index 42) - 2 listes avec sous-titres et ligne bleue
            elif slide_index == 42:  # Slide 43
                return {
                    "pattern": "dual_lists_blue",
                    "description": "2 listes avec sous-titres et ligne bleue",
                    "title_shapes": [0],  # Shape 0 = Titre principal (shape_1/ID:2)
                    "left_subtitle_shapes": [2],  # Shape 2 = Sous-titre gauche (shape_3/ID:4)
                    "left_content_shapes": [1],  # Shape 1 = Contenu liste gauche (shape_2/ID:3)
                    "right_subtitle_shapes": [3],  # Shape 3 = Sous-titre droite (shape_4/ID:5)
                    "right_content_shapes": [4],  # Shape 4 = Contenu liste droite (shape_5/ID:10)
                    "page_number_shapes": [5],  # Shape 5 = Numéro de page (shape_6/ID:11)
                    "total_shapes": 6
                }

            # Mapping précis pour slide 44 (index 43) - 2 listes avec sous-titres et ligne grise
            elif slide_index == 43:  # Slide 44
                return {
                    "pattern": "dual_lists_grey",
                    "description": "2 listes avec sous-titres et ligne grise",
                    "title_shapes": [0],  # Shape 0 = Titre principal (shape_1/ID:2)
                    "left_subtitle_shapes": [2],  # Shape 2 = Sous-titre gauche (shape_3/ID:4)
                    "left_content_shapes": [1],  # Shape 1 = Contenu liste gauche (shape_2/ID:3)
                    "right_subtitle_shapes": [3],  # Shape 3 = Sous-titre droite (shape_4/ID:5)
                    "right_content_shapes": [4],  # Shape 4 = Contenu liste droite (shape_5/ID:10)
                    "page_number_shapes": [5],  # Shape 5 = Numéro de page (shape_6/ID:11)
                    "total_shapes": 6
                }

            # Mapping précis pour autres slides 38 (index 37) - Formats divers
            elif slide_index == 37:  # Slide 38
                # Mapping générique pour cette slide (à personnaliser selon besoin)
                return {
                    "pattern": "dual_content",
                    "description": "Contenu dual avec ligne de séparation",
                    "title_shapes": [0],  # Titre principal
                    "content_shapes": [1, 2],  # 2 contenus principaux
                    "subtitle_shapes": [3, 4],  # Sous-titres si applicable
                    "page_number_shapes": [-1],  # Dernière shape généralement
                    "total_shapes": 5
                }

            else:
                print(f"[WARNING] Aucun mapping précis disponible pour slide index {slide_index}")
                return None

        except Exception as e:
            print(f"[ERROR] Erreur lors de la récupération du mapping pour slide {slide_index}: {e}")
            return None

    def _generate_explanation_output_path(self, content: str, explanation_style: str) -> str:
        """Génère le chemin de sortie pour la slide d'explication"""

        # Nettoyer le contenu pour le nom de fichier
        clean_content = "".join(c for c in content if c.isalnum() or c in (' ', '-', '_')).rstrip()
        clean_content = clean_content.replace(' ', '_').lower()[:50]  # Limiter la longueur

        # Timestamp pour l'unicité
        timestamp = datetime.now().strftime("%Y%m%d_%H%M")

        # Nom du fichier
        filename = f"{timestamp}_explanation_{explanation_style}_{clean_content}.pptx"

        # Dossier de destination
        base_dir = "presentations"
        explanation_dir = os.path.join(base_dir, f"explanation_{timestamp}")

        return os.path.join(explanation_dir, "explanation", filename)

    def _move_slide_to_position_direct(self, presentation, source_index, target_index):
        """Déplace une slide vers une nouvelle position"""
        try:
            if target_index >= len(presentation.slides) or target_index < 0:
                print(f"[WARNING] Position {target_index} invalide, slide laissée en position actuelle")
                return

            # Cette fonctionnalité nécessite une manipulation avancée des slides
            print(f"[INFO] Déplacement de slide {source_index} vers position {target_index}")
            # Implementation simplifiée: la slide reste à la fin
            print(f"[INFO] Slide ajoutée en fin de présentation")

        except Exception as e:
            print(f"[WARNING] Erreur déplacement slide: {e}")

    def debug_shape_content_assignment(self, presentation_path: str, slide_index: int = 0):
        """
        Méthode de debug pour analyser l'assignation du contenu dans une slide.
        Utile pour identifier les problèmes d'assignation comme dans l'image fournie.
        """
        try:
            print(f"[DEBUG] Analyse de l'assignation du contenu pour slide {slide_index + 1}")

            prs = Presentation(presentation_path)
            if slide_index >= len(prs.slides):
                print(f"[DEBUG-ERROR] Slide {slide_index + 1} n'existe pas")
                return

            slide = prs.slides[slide_index]
            print(f"[DEBUG] Slide analysée: {len(slide.shapes)} shapes trouvées")

            for i, shape in enumerate(slide.shapes):
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    content = shape.text_frame.text.strip()
                    print(f"[DEBUG] Shape {i}: '{content}' (longueur: {len(content)} chars)")

                    # Analyser la position et taille
                    try:
                        left_inches = shape.left / Inches(1)
                        top_inches = shape.top / Inches(1)
                        width_inches = shape.width / Inches(1)
                        height_inches = shape.height / Inches(1)
                        print(f"[DEBUG]   Position: ({left_inches:.1f}\", {top_inches:.1f}\") Taille: {width_inches:.1f}\"x{height_inches:.1f}\"")
                    except:
                        print(f"[DEBUG]   Position/Taille: Non disponible")

                else:
                    print(f"[DEBUG] Shape {i}: Non-texte (type: {type(shape).__name__})")

            return True

        except Exception as e:
            print(f"[DEBUG-ERROR] Erreur analyse: {e}")
            return False


    def validate_template(self) -> bool:
        """Valide que le template Premier Tech est correct pour les explications détaillées"""
        try:
            checks = {
                "file_exists": os.path.exists(self.template_path),
                "has_slides": False,
                "explanation_slides_exist": False,
                "slides_count": 0,
                "available_styles": []
            }

            if checks["file_exists"]:
                pres = Presentation(self.template_path)
                checks["slides_count"] = len(pres.slides)
                checks["has_slides"] = len(pres.slides) > 0

                # Vérifier que toutes les slides d'explication existent
                available_styles = []
                for slide_index in self.explanation_slides.keys():
                    if len(pres.slides) > slide_index:
                        style = self.explanation_slides[slide_index]['style']
                        available_styles.append(style)

                checks["available_styles"] = available_styles
                checks["explanation_slides_exist"] = len(available_styles) == len(self.explanation_slides)

            all_valid = all([checks["file_exists"], checks["has_slides"], checks["explanation_slides_exist"]])

            print("=== VALIDATION TEMPLATE PREMIER TECH POUR EXPLICATIONS DÉTAILLÉES ===")
            for check, result in checks.items():
                if check == "available_styles":
                    print(f"[INFO] Styles disponibles: {', '.join(result)}")
                else:
                    status = "OK" if result else "ERREUR"
                    print(f"[{status}] {check}: {result}")

            if checks["explanation_slides_exist"]:
                print(f"[INFO] {len(checks['available_styles'])} styles d'explication disponibles:")
                for slide_index, slide_data in self.explanation_slides.items():
                    if slide_index in [idx for idx in self.explanation_slides.keys() if len(pres.slides) > idx]:
                        print(f"  - {slide_data['style']}: Slide {slide_index + 1} ({slide_data['usage']})")

            return all_valid

        except Exception as e:
            print(f"[ERROR] Erreur validation: {e}")
            return False

    def list_available_styles(self) -> Dict[str, Dict[str, Any]]:
        """Liste tous les styles d'explication disponibles"""
        return {
            slide_data['style']: {
                "slide_number": slide_index + 1,
                "name": slide_data['name'],
                "usage": slide_data['usage'],
                "audience": slide_data['audience'],
                "elements": slide_data['elements'],
                "explanation_type": slide_data['explanation_type']
            }
            for slide_index, slide_data in self.explanation_slides.items()
        }


def main():
    """Interface en ligne de commande"""

    parser = argparse.ArgumentParser(
        description='Construction de slides d\'explications détaillées Premier Tech (slides 35-44)'
    )

    parser.add_argument('content', nargs='?', help='Contenu principal de l\'explication')
    parser.add_argument('--subtitle', help='Sous-titre ou description')
    parser.add_argument('--additional', nargs='*', help='Contenu additionnel (liste de points)')
    parser.add_argument('--title', help='Titre de la slide')
    parser.add_argument('--style',
                       choices=['four_points', 'dual_detailed_blue', 'dual_detailed_grey', 'dual_titled_blue',
                               'dual_titled_grey', 'dual_lists_blue', 'dual_lists_grey'],
                       default='four_points',
                       help='Style d\'explication (four_points=slide35, dual_detailed_blue=slide39, etc.)')
    parser.add_argument('--output', help='Chemin de sortie spécifique')
    parser.add_argument('--template', default='templates/Template_PT.pptx',
                       help='Chemin vers le template Premier Tech')
    parser.add_argument('--no-widen', action='store_true',
                       help='Désactiver l\'élargissement automatique des objets texte')
    parser.add_argument('--validate', action='store_true',
                       help='Valider le template seulement')
    parser.add_argument('--list-styles', action='store_true',
                       help='Lister les styles disponibles')
    parser.add_argument('--insert-into',
                       help='Insérer dans une présentation existante (chemin)')
    parser.add_argument('--position', type=int,
                       help='Position d\'insertion dans la présentation (défaut: fin)')
    parser.add_argument('--debug-assignment',
                       help='Debug l\'assignation du contenu dans un fichier existant')
    parser.add_argument('--debug-slide-index', type=int, default=0,
                       help='Index de la slide à analyser pour debug (défaut: 0)')

    args = parser.parse_args()

    try:
        # Initialiser le constructeur
        builder = DetailedExplanationBuilder(args.template)

        # Mode validation
        if args.validate:
            is_valid = builder.validate_template()
            sys.exit(0 if is_valid else 1)

        # Mode liste des styles
        if args.list_styles:
            styles = builder.list_available_styles()
            print("=== STYLES D'EXPLICATIONS DÉTAILLÉES DISPONIBLES ===")
            for style, info in styles.items():
                print(f"{style.upper()}:")
                print(f"  - Slide: {info['slide_number']}")
                print(f"  - Nom: {info['name']}")
                print(f"  - Usage: {info['usage']}")
                print(f"  - Audience: {info['audience']}")
                print(f"  - Éléments: {info['elements']}")
                print(f"  - Type: {info['explanation_type']}")
                print()
            sys.exit(0)

        # Mode debug de l'assignation
        if args.debug_assignment:
            if not os.path.exists(args.debug_assignment):
                print(f"ERREUR: Fichier non trouvé: {args.debug_assignment}")
                sys.exit(1)

            success = builder.debug_shape_content_assignment(args.debug_assignment, args.debug_slide_index)
            sys.exit(0 if success else 1)

        # Vérifier que le contenu requis est fourni (sauf en mode validation/list)
        if not args.content:
            print("ERREUR: Le contenu principal de l'explication est requis")
            print("Utilisez --validate ou --list-styles pour les modes sans paramètres")
            sys.exit(1)

        # Validation cohérence style et contenu
        subtitle_provided = bool(args.subtitle)
        if args.style in ['dual_detailed_blue', 'dual_detailed_grey'] and not subtitle_provided:
            print(f"WARNING: Le style '{args.style}' est plus efficace avec un sous-titre (--subtitle)")
        elif args.style not in ['dual_detailed_blue', 'dual_detailed_grey'] and subtitle_provided:
            print(f"WARNING: Le style '{args.style}' ne supporte pas les sous-titres, il sera ignoré")

        # Préparer les variables
        content = args.content
        subtitle = args.subtitle
        additional = args.additional
        title = args.title

        # Mode insertion dans présentation existante
        if args.insert_into:
            output_path = builder.insert_explanation_into_existing_presentation(
                presentation_path=args.insert_into,
                content=content,
                subtitle=subtitle,
                additional_content=additional,
                title=title,
                explanation_style=args.style,
                insert_position=args.position
            )
        else:
            # Créer une slide d'explication individuelle
            output_path = builder.create_explanation_slide(
                content=content,
                subtitle=subtitle,
                additional_content=additional,
                title=title,
                explanation_style=args.style,
                output_path=args.output,
                auto_widen=not args.no_widen
            )
            print(f"\nSUCCES: Slide d'explication détaillée créée: {output_path}")

        print(f"Style utilisé: {args.style}")
        print(f"Contenu: {args.content}")
        if args.subtitle:
            print(f"Sous-titre: {args.subtitle}")
        if args.additional:
            print(f"Contenu additionnel: {', '.join(args.additional)}")
        if args.title:
            print(f"Titre: {args.title}")

    except Exception as e:
        print(f"\nERREUR: {e}")
        sys.exit(1)


if __name__ == "__main__":
    main()