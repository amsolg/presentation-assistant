#!/usr/bin/env python3
"""
Enhanced Charts Builder - Création avancée de slides graphiques Premier Tech
Amélioration du script 09 avec support CSV, multi-séries, et standards visuels Premier Tech.

Nouvelles fonctionnalités:
- Import automatique depuis CSV/Excel
- Support multi-séries pour comparaisons complexes
- Configuration JSON pour paramètres avancés
- Validation et enrichissement automatique des données
- Export des données de graphiques
- Respect strict des standards visuels Premier Tech
"""

import os
import sys
import json
import shutil
import argparse
import csv
from datetime import datetime
from pathlib import Path
from typing import List, Dict, Optional, Any, Tuple, Union
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.dml.color import RGBColor

# Import optionnel de pandas pour fonctionnalités avancées
try:
    import pandas as pd
    PANDAS_AVAILABLE = True
except ImportError:
    PANDAS_AVAILABLE = False
    print("[INFO] pandas non installé - fonctionnalités CSV avancées limitées")

# Standards visuels Premier Tech avec centrage
# Dimensions slide: 1280×720 pixels (13.333×7.5 inches)
PT_CHART_STANDARDS = {
    'COLUMN_CLUSTERED': {
        'width': Inches(8.5),   # Largeur raisonnable pour le graphique
        'height': Inches(4.5),  # Hauteur proportionnelle
        'left': Inches(2.4),    # Centré: (13.333 - 8.5) / 2
        'top': Inches(1.8)      # Position verticale avec espace pour titre
    },
    'LINE': {
        'width': Inches(8.5),
        'height': Inches(4.5),
        'left': Inches(2.4),
        'top': Inches(1.8)
    },
    'PIE': {
        'width': Inches(6.0),   # Plus petit pour graphique en secteurs
        'height': Inches(4.5),
        'left': Inches(3.7),    # Centré: (13.333 - 6.0) / 2
        'top': Inches(1.8)
    },
    'BAR_CLUSTERED': {
        'width': Inches(8.5),
        'height': Inches(4.5),
        'left': Inches(2.4),
        'top': Inches(1.8)
    },
    'COMPACT': {
        'width': Inches(7.0),   # Version compacte
        'height': Inches(4.0),
        'left': Inches(3.2),    # Centré: (13.333 - 7.0) / 2
        'top': Inches(2.0)
    }
}

# Palette de couleurs Premier Tech - Basée sur les standards officiels
# Couleurs principales
PT_COLOR_PALETTE = [
    RGBColor(65, 182, 230),   # #41B6E6 - Bleu Premier Tech (Pantone 298)
    RGBColor(138, 141, 143),  # #8A8D8F - Gris (Pantone 877)
    RGBColor(4, 14, 30),      # #040E1E - Bleu Noir PPT (fond)
    RGBColor(0, 119, 200),    # #0077C8 - Bleu vif (Pantone 3005)
    RGBColor(84, 88, 91),     # #54585B - Gris foncé (Pantone 425)
    RGBColor(189, 189, 189),  # #BDBDBD - Gris clair (Pantone 877 - 50%)
    RGBColor(255, 255, 255),  # #FFFFFF - Blanc (texte)
    RGBColor(0, 63, 127),     # #003F7F - Bleu foncé (pour variation)
]

# Palette étendue pour graphiques avec nombreuses séries
PT_COLOR_PALETTE_EXTENDED = [
    RGBColor(65, 182, 230),   # Bleu Premier Tech
    RGBColor(0, 119, 200),    # Bleu vif
    RGBColor(138, 141, 143),  # Gris moyen
    RGBColor(84, 88, 91),     # Gris foncé
    RGBColor(189, 189, 189),  # Gris clair
    RGBColor(0, 63, 127),     # Bleu foncé
    RGBColor(102, 178, 255),  # Bleu clair
    RGBColor(51, 153, 255),   # Bleu moyen clair
]

# Configuration des styles de graphiques
CHART_STYLE_CONFIG = {
    'column_clustered': {
        'colors': PT_COLOR_PALETTE_EXTENDED[:6],  # Utilise une diversité de couleurs
        'has_legend': True,
        'has_data_labels': True,
        'value_axis_format': '0',
        'chart_type': XL_CHART_TYPE.COLUMN_CLUSTERED
    },
    'line_chart': {
        'colors': PT_COLOR_PALETTE_EXTENDED[:6],  # Utilise une diversité de couleurs
        'has_legend': True,
        'has_data_labels': True,
        'value_axis_format': '0',
        'chart_type': XL_CHART_TYPE.LINE
    },
    'pie_chart': {
        'colors': PT_COLOR_PALETTE_EXTENDED,  # Toutes les couleurs pour les secteurs
        'has_legend': True,
        'has_data_labels': True,
        'data_label_format': 'percentage',
        'chart_type': XL_CHART_TYPE.PIE
        # Pas de value_axis_format car les graphiques en secteurs n'ont pas d'axe de valeurs
    },
    'bar_clustered': {
        'colors': PT_COLOR_PALETTE_EXTENDED[:6],  # Utilise une diversité de couleurs
        'has_legend': True,
        'has_data_labels': True,
        'value_axis_format': '0',
        'chart_type': XL_CHART_TYPE.BAR_CLUSTERED
    },
    'column_compact': {
        'colors': PT_COLOR_PALETTE[:3],  # Moins de couleurs pour compact
        'has_legend': False,
        'has_data_labels': True,
        'value_axis_format': '0',
        'chart_type': XL_CHART_TYPE.COLUMN_CLUSTERED
    },
    'bar_compact': {
        'colors': PT_COLOR_PALETTE[:3],  # Moins de couleurs pour compact
        'has_legend': False,
        'has_data_labels': True,
        'value_axis_format': '0',
        'chart_type': XL_CHART_TYPE.BAR_CLUSTERED
    }
}


class EnhancedChartsBuilder:
    """
    Classe améliorée pour construction de slides graphiques Premier Tech.
    Support CSV, multi-séries, et standards visuels stricts.
    """

    def __init__(self, template_path: str = "../templates/Template_PT.pptx"):
        """
        Initialise le constructeur avec le template Premier Tech.

        Args:
            template_path: Chemin vers le template Premier Tech
        """
        self.template_path = template_path

        # Mapping EXACT des slides 46-51 selon template_analysis_output/
        self.chart_slides = {
            45: {  # Slide 46 (index 45)
                "name": "Graphique en colonnes groupées",
                "usage": "Comparaisons catégorielles multiples",
                "audience": "Managers, Analystes",
                "style": "column_clustered",
                "chart_type": "COLUMN_CLUSTERED",
                "powerpoint_chart_type": 51,
                "description": "Graphiques en colonnes pour comparaisons multiples"
            },
            46: {  # Slide 47 (index 46)
                "name": "Graphique linéaire",
                "usage": "Évolutions temporelles et tendances",
                "audience": "Analystes, Spécialistes",
                "style": "line_chart",
                "chart_type": "LINE",
                "powerpoint_chart_type": 4,
                "description": "Graphiques linéaires pour tendances temporelles"
            },
            47: {  # Slide 48 (index 47)
                "name": "Graphique en secteurs",
                "usage": "Répartitions et proportions",
                "audience": "Leaders, Managers",
                "style": "pie_chart",
                "chart_type": "PIE",
                "powerpoint_chart_type": 5,
                "description": "Graphiques en secteurs pour proportions"
            },
            48: {  # Slide 49 (index 48)
                "name": "Graphique en barres horizontales",
                "usage": "Comparaisons horizontales de catégories",
                "audience": "Managers, Analystes",
                "style": "bar_clustered",
                "chart_type": "BAR_CLUSTERED",
                "powerpoint_chart_type": 57,
                "description": "Graphiques en barres horizontales pour comparaisons"
            },
            49: {  # Slide 50 (index 49)
                "name": "Graphique en colonnes compact",
                "usage": "Visualisations condensées multiples",
                "audience": "Spécialistes, Techniques",
                "style": "column_compact",
                "chart_type": "COLUMN_CLUSTERED",
                "powerpoint_chart_type": 51,
                "description": "Version compacte des colonnes groupées"
            },
            50: {  # Slide 51 (index 50)
                "name": "Graphique en barres compact",
                "usage": "Comparaisons condensées horizontales",
                "audience": "Analystes avancés",
                "style": "bar_compact",
                "chart_type": "BAR_CLUSTERED",
                "powerpoint_chart_type": 57,
                "description": "Version compacte des barres horizontales"
            }
        }

        if not os.path.exists(self.template_path):
            raise FileNotFoundError(f"Template Premier Tech non trouvé: {self.template_path}")

        # Valider le template
        self._validate_template()

    def _validate_template(self):
        """Valide que le template contient toutes les slides nécessaires"""
        try:
            pres = Presentation(self.template_path)
            max_slide_needed = max(self.chart_slides.keys()) + 1

            if len(pres.slides) < max_slide_needed:
                raise ValueError(f"Template ne contient que {len(pres.slides)} slides, "
                               f"minimum {max_slide_needed} requis pour les graphiques")

            print(f"[OK] Template validé: {len(pres.slides)} slides disponibles")

        except Exception as e:
            raise ValueError(f"Erreur validation template: {e}")

    def load_data_from_csv(self, csv_path: str, has_header: bool = True) -> Dict[str, Any]:
        """
        Charge les données depuis un fichier CSV.

        Args:
            csv_path: Chemin vers le fichier CSV
            has_header: Si le CSV a une ligne d'en-tête

        Returns:
            Dict contenant les labels, valeurs et séries optionnelles
        """
        if not os.path.exists(csv_path):
            raise FileNotFoundError(f"Fichier CSV non trouvé: {csv_path}")

        print(f"[CSV] Chargement des données depuis: {csv_path}")

        if PANDAS_AVAILABLE:
            # Utilisation de pandas pour plus de flexibilité
            df = pd.read_csv(csv_path)

            # Détecter le format du CSV
            if len(df.columns) == 2:
                # Format simple: Label, Valeur
                labels = df.iloc[:, 0].tolist()
                values = df.iloc[:, 1].tolist()
                series_data = None
                print(f"[CSV] Format simple détecté: {len(labels)} entrées")
            else:
                # Format multi-séries: Label, Série1, Série2, ...
                labels = df.iloc[:, 0].tolist()
                series_data = {}
                for col in df.columns[1:]:
                    series_data[col] = df[col].tolist()
                values = None
                print(f"[CSV] Format multi-séries détecté: {len(series_data)} séries")

            return {
                'labels': labels,
                'values': values,
                'series_data': series_data,
                'dataframe': df
            }
        else:
            # Fallback sans pandas
            with open(csv_path, 'r', encoding='utf-8') as f:
                reader = csv.reader(f)
                rows = list(reader)

                if has_header:
                    headers = rows[0]
                    data_rows = rows[1:]
                else:
                    data_rows = rows
                    headers = None

                if len(data_rows[0]) == 2:
                    # Format simple
                    labels = [row[0] for row in data_rows]
                    values = [row[1] for row in data_rows]
                    return {
                        'labels': labels,
                        'values': values,
                        'series_data': None
                    }
                else:
                    # Format multi-séries
                    labels = [row[0] for row in data_rows]
                    series_data = {}
                    for i in range(1, len(data_rows[0])):
                        series_name = headers[i] if headers else f"Série {i}"
                        series_data[series_name] = [row[i] for row in data_rows]
                    return {
                        'labels': labels,
                        'values': None,
                        'series_data': series_data
                    }

    def load_chart_config(self, json_path: str) -> Dict[str, Any]:
        """
        Charge une configuration de graphique depuis un fichier JSON.

        Args:
            json_path: Chemin vers le fichier JSON

        Returns:
            Dict contenant la configuration complète
        """
        if not os.path.exists(json_path):
            raise FileNotFoundError(f"Fichier JSON non trouvé: {json_path}")

        print(f"[JSON] Chargement de la configuration depuis: {json_path}")

        with open(json_path, 'r', encoding='utf-8') as f:
            config = json.load(f)

        # Valider la configuration
        required_fields = ['title', 'style']
        for field in required_fields:
            if field not in config:
                raise ValueError(f"Champ requis manquant dans JSON: {field}")

        print(f"[JSON] Configuration chargée: {config.get('title')}")
        return config

    def validate_and_enrich_data(
        self,
        labels: List[str],
        values: List[Union[str, float]],
        chart_type: str
    ) -> Tuple[List[str], List[float]]:
        """
        Valide et enrichit les données pour le type de graphique.

        Args:
            labels: Labels des données
            values: Valeurs (peuvent être des strings avec %)
            chart_type: Type de graphique

        Returns:
            Tuple (labels validés, valeurs numériques)
        """
        if not labels or not values:
            raise ValueError("Labels et valeurs sont requis")

        if len(labels) != len(values):
            # Ajuster à la longueur minimale
            min_len = min(len(labels), len(values))
            labels = labels[:min_len]
            values = values[:min_len]
            print(f"[WARNING] Ajustement des données à {min_len} entrées")

        # Conversion des valeurs en nombres
        numeric_values = []
        for val in values:
            if isinstance(val, str):
                # Nettoyer les valeurs string
                clean_val = val.replace('%', '').replace('€', '').replace('$', '')
                clean_val = clean_val.replace(',', '.').replace(' ', '')
                try:
                    numeric_values.append(float(clean_val))
                except ValueError:
                    print(f"[WARNING] Impossible de convertir '{val}', utilisation de 0")
                    numeric_values.append(0.0)
            else:
                numeric_values.append(float(val))

        # Enrichissement selon le type de graphique
        if chart_type == 'PIE':
            # Pour un graphique en secteurs, normaliser à 100% si nécessaire
            total = sum(numeric_values)
            if total > 0 and abs(total - 100) > 0.01:
                print(f"[INFO] Normalisation des valeurs pour graphique en secteurs (total: {total:.1f}%)")
                numeric_values = [v / total * 100 for v in numeric_values]

            # Limiter à 8 catégories max pour la lisibilité
            if len(labels) > 8:
                print(f"[INFO] Limitation à 8 catégories pour graphique en secteurs")
                # Regrouper les petites valeurs
                sorted_data = sorted(zip(labels, numeric_values), key=lambda x: x[1], reverse=True)
                top_7 = sorted_data[:7]
                other_value = sum(v for _, v in sorted_data[7:])
                labels = [l for l, _ in top_7] + ["Autres"]
                numeric_values = [v for _, v in top_7] + [other_value]

        elif chart_type in ['COLUMN_CLUSTERED', 'BAR_CLUSTERED']:
            # Pour les graphiques en barres/colonnes, limiter à 12 catégories
            if len(labels) > 12:
                print(f"[INFO] Limitation à 12 catégories pour graphique en barres")
                labels = labels[:12]
                numeric_values = numeric_values[:12]

        elif chart_type == 'LINE':
            # Pour les graphiques linéaires, s'assurer qu'on a au moins 3 points
            if len(labels) < 3:
                print(f"[WARNING] Graphique linéaire avec moins de 3 points")

        print(f"[VALIDATE] {len(labels)} catégories, valeurs de {min(numeric_values):.1f} à {max(numeric_values):.1f}")
        return labels, numeric_values

    def create_chart_data(
        self,
        labels: List[str],
        values: Optional[List[float]] = None,
        series_data: Optional[Dict[str, List[float]]] = None,
        chart_type: str = 'COLUMN_CLUSTERED',
        series_title: str = None
    ) -> CategoryChartData:
        """
        Crée les données de graphique pour python-pptx.

        Args:
            labels: Catégories du graphique
            values: Valeurs simples (pour une seule série)
            series_data: Données multi-séries
            chart_type: Type de graphique

        Returns:
            CategoryChartData pour python-pptx
        """
        chart_data = CategoryChartData()

        if series_data:
            # Multi-séries
            # Utiliser les labels originaux pour l'instant
            chart_data.categories = labels
            for series_name, series_values in series_data.items():
                # Valider les valeurs de la série
                _, numeric_values = self.validate_and_enrich_data(
                    labels, series_values, chart_type
                )
                chart_data.add_series(series_name, numeric_values)
            print(f"[CHART_DATA] {len(series_data)} séries ajoutées")
        elif values:
            # Série unique avec titre personnalisable
            # Pour PIE charts, validate_and_enrich_data peut modifier les labels
            validated_labels, numeric_values = self.validate_and_enrich_data(
                labels, values, chart_type
            )
            # Utilise les labels validés (potentiellement modifiés pour PIE)
            chart_data.categories = validated_labels

            # Utilise le titre fourni ou un titre par défaut
            series_name = series_title if series_title else 'Valeurs'
            chart_data.add_series(series_name, numeric_values)
            print(f"[CHART_DATA] 1 série ajoutée: '{series_name}'")
        else:
            raise ValueError("Ni values ni series_data fournis")

        return chart_data

    def apply_chart_formatting(self, chart, style: str, series_title: str = None):
        """
        Applique le formatage Premier Tech au graphique.

        Args:
            chart: Objet chart python-pptx
            style: Style du graphique
            series_title: Titre personnalisé pour la série (optionnel)
        """
        config = CHART_STYLE_CONFIG.get(style, {})

        # Appliquer les couleurs avec diversité
        if 'colors' in config and hasattr(chart, 'series'):
            colors = config['colors']
            for i, series in enumerate(chart.series):
                # Utilise une rotation de couleurs pour plus de diversité
                color_index = i % len(colors)
                fill = series.format.fill
                fill.solid()
                fill.fore_color.rgb = colors[color_index]

        # Configurer la légende
        if config.get('has_legend', True):
            chart.has_legend = True
            if chart.has_legend and hasattr(chart, 'legend'):
                try:
                    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
                    chart.legend.include_in_layout = False
                except Exception:
                    # Certains types de graphiques (comme PIE) peuvent ne pas supporter toutes les propriétés
                    pass

        # Configurer les étiquettes de données pour les graphiques en secteurs
        if config.get('has_data_labels', False):
            plot = chart.plots[0]
            plot.has_data_labels = True

            if style == 'pie_chart':
                # Pour les graphiques en secteurs
                data_labels = plot.data_labels
                data_labels.show_percentage = True
                data_labels.show_category_name = True
                data_labels.show_value = False

                # Désactiver le text wrapping pour éviter la coupure des mots
                try:
                    # Position des étiquettes à l'extérieur pour plus de place
                    from pptx.enum.chart import XL_DATA_LABEL_POSITION
                    data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END
                except:
                    pass

        # Format de l'axe des valeurs
        if 'value_axis_format' in config:
            try:
                chart.value_axis.number_format = config['value_axis_format']
            except (AttributeError, ValueError):
                # Les graphiques en secteurs n'ont pas d'axe de valeurs
                pass

        # Personnaliser le titre de la série si fourni
        if series_title and hasattr(chart, 'series') and len(chart.series) > 0:
            try:
                chart.series[0].name = series_title
            except:
                pass

        print(f"[FORMAT] Style '{style}' appliqué au graphique avec palette de couleurs diversifiée")

    def export_chart_data(self, presentation_path: str, output_csv: str, slide_index: int = -1):
        """
        Exporte les données d'un graphique vers CSV.

        Args:
            presentation_path: Chemin vers la présentation
            output_csv: Chemin de sortie CSV
            slide_index: Index de la slide (défaut: dernière)
        """
        pres = Presentation(presentation_path)
        slide = pres.slides[slide_index]

        data_rows = []
        chart_found = False

        for shape in slide.shapes:
            if hasattr(shape, 'chart'):
                chart_found = True
                chart = shape.chart

                # Extraire les catégories
                categories = []
                try:
                    for category in chart.plots[0].categories:
                        categories.append(str(category))
                except:
                    categories = [f"Cat{i+1}" for i in range(5)]

                # Extraire les séries
                for series in chart.series:
                    series_name = series.name or "Série"
                    for i, value in enumerate(series.values):
                        if i < len(categories):
                            data_rows.append({
                                'Catégorie': categories[i],
                                'Série': series_name,
                                'Valeur': value
                            })

        if not chart_found:
            print(f"[WARNING] Aucun graphique trouvé dans la slide {slide_index + 1}")
            return

        # Écrire le CSV
        if PANDAS_AVAILABLE:
            df = pd.DataFrame(data_rows)
            df.to_csv(output_csv, index=False, encoding='utf-8')
        else:
            with open(output_csv, 'w', newline='', encoding='utf-8') as f:
                if data_rows:
                    writer = csv.DictWriter(f, fieldnames=data_rows[0].keys())
                    writer.writeheader()
                    writer.writerows(data_rows)

        print(f"[EXPORT] Données exportées vers: {output_csv}")

    def insert_enhanced_chart_slide(
        self,
        target_presentation_path: str,
        title: str,
        style: str = "column_clustered",
        csv_path: Optional[str] = None,
        json_config: Optional[str] = None,
        data_labels: Optional[List[str]] = None,
        data_values: Optional[List[str]] = None,
        series_data: Optional[Dict[str, List]] = None,
        insights: Optional[str] = None,
        position: Optional[int] = None,
        series_title: Optional[str] = None
    ) -> Dict[str, Any]:
        """
        Insère une slide graphique améliorée dans une présentation.

        Args:
            target_presentation_path: Chemin vers la présentation cible
            title: Titre de la slide
            style: Style de graphique
            csv_path: Chemin vers fichier CSV (optionnel)
            json_config: Chemin vers configuration JSON (optionnel)
            data_labels: Labels directs
            data_values: Valeurs directes
            series_data: Données multi-séries
            insights: Points clés
            position: Position d'insertion

        Returns:
            Dict avec informations sur l'insertion
        """
        if not os.path.exists(target_presentation_path):
            raise FileNotFoundError(f"Présentation cible non trouvée: {target_presentation_path}")

        print(f"[INSERT] Insertion slide graphique améliorée dans: {target_presentation_path}")

        # Créer une sauvegarde
        backup_path = target_presentation_path.replace('.pptx', '_backup_before_enhanced_chart.pptx')
        shutil.copy2(target_presentation_path, backup_path)
        print(f"[BACKUP] Sauvegarde créée: {backup_path}")

        try:
            # Charger configuration JSON si fournie
            if json_config:
                config = self.load_chart_config(json_config)
                title = config.get('title', title)
                style = config.get('style', style)
                insights = config.get('insights', insights)
                if 'data' in config:
                    data_labels = config['data'].get('labels', data_labels)
                    data_values = config['data'].get('values', data_values)
                    series_data = config['data'].get('series', series_data)

            # Charger données CSV si fournies
            if csv_path:
                csv_data = self.load_data_from_csv(csv_path)
                if csv_data['series_data']:
                    series_data = csv_data['series_data']
                    data_labels = csv_data['labels']
                    data_values = None
                else:
                    data_labels = csv_data['labels']
                    data_values = csv_data['values']
                    series_data = None

            # Valider le style
            available_styles = [info['style'] for info in self.chart_slides.values()]
            if style not in available_styles:
                raise ValueError(f"Style '{style}' non supporté. Styles disponibles: {available_styles}")

            # Trouver la slide template correspondante
            slide_index = None
            slide_info = None
            for idx, info in self.chart_slides.items():
                if info['style'] == style:
                    slide_index = idx
                    slide_info = info
                    break

            print(f"[CREATE] Création graphique amélioré style '{style}'")

            # Charger les présentations
            template_pres = Presentation(self.template_path)
            target_pres = Presentation(target_presentation_path)

            # Déterminer la position
            if position is None:
                position = len(target_pres.slides)

            # Obtenir le layout
            source_slide = template_pres.slides[slide_index]
            layout_name = source_slide.slide_layout.name

            # Chercher le layout dans la présentation cible
            target_layout = None
            for layout in target_pres.slide_layouts:
                if layout.name == layout_name:
                    target_layout = layout
                    break

            if target_layout is None:
                # Utiliser le premier layout disponible comme fallback
                target_layout = target_pres.slide_layouts[0]
                print(f"[WARNING] Layout '{layout_name}' non trouvé, utilisation du layout par défaut")

            # Ajouter la slide
            new_slide = target_pres.slides.add_slide(target_layout)

            # Ajouter le titre
            if new_slide.shapes.title:
                new_slide.shapes.title.text = title

            # Créer les données du graphique
            chart_type = slide_info['chart_type']

            if not data_labels:
                # Données par défaut si aucune fournie
                data_labels = ["T1", "T2", "T3", "T4"]
                data_values = [2500000, 3200000, 2800000, 3700000]
                print("[INFO] Utilisation de données d'exemple")

            # Créer l'objet de données avec titre personnalisable
            chart_data = self.create_chart_data(
                data_labels, data_values, series_data, chart_type, series_title
            )

            # Obtenir les dimensions standards Premier Tech
            if chart_type in ['COLUMN_CLUSTERED', 'BAR_CLUSTERED']:
                if 'compact' in style:
                    dimensions = PT_CHART_STANDARDS['COMPACT']
                else:
                    dimensions = PT_CHART_STANDARDS.get(chart_type, PT_CHART_STANDARDS['COLUMN_CLUSTERED'])
            else:
                dimensions = PT_CHART_STANDARDS.get(chart_type, PT_CHART_STANDARDS['COLUMN_CLUSTERED'])

            # Ajouter le graphique avec dimensions Premier Tech
            xl_chart_type = CHART_STYLE_CONFIG[style]['chart_type']
            chart_shape = new_slide.shapes.add_chart(
                xl_chart_type,
                dimensions['left'],
                dimensions['top'],
                dimensions['width'],
                dimensions['height'],
                chart_data
            )

            # Appliquer le formatage Premier Tech avec titre de série
            self.apply_chart_formatting(chart_shape.chart, style, series_title)

            # Ajouter les insights si fournis
            if insights:
                # Ajouter une zone de texte pour les insights
                left = Inches(1)
                top = dimensions['top'] + dimensions['height'] + Inches(0.2)
                width = Inches(10)
                height = Inches(1)

                textbox = new_slide.shapes.add_textbox(left, top, width, height)
                text_frame = textbox.text_frame
                text_frame.text = f"Points clés: {insights}"

                # Formater le texte en blanc pour meilleure visibilité
                paragraph = text_frame.paragraphs[0]
                paragraph.font.size = Pt(12)
                paragraph.font.color.rgb = RGBColor(255, 255, 255)  # Blanc

            # Repositionner si nécessaire
            if position < len(target_pres.slides) - 1:
                xml_slides = target_pres.slides._sldIdLst
                xml_slides.insert(position, xml_slides.pop())

            # Sauvegarder
            target_pres.save(target_presentation_path)

            print(f"[SUCCESS] Slide graphique améliorée insérée à la position {position + 1}")

            # Créer rapport détaillé
            report_data = {
                "insertion_info": {
                    "timestamp": datetime.now().isoformat(),
                    "script": "09_charts_builder.py",
                    "target_presentation": target_presentation_path,
                    "backup_created": backup_path
                },
                "slide_details": {
                    "title": title,
                    "style": style,
                    "chart_type": chart_type,
                    "data_source": "CSV" if csv_path else "JSON" if json_config else "Direct",
                    "has_multi_series": series_data is not None,
                    "num_categories": len(data_labels) if data_labels else 0,
                    "has_insights": insights is not None
                },
                "standards_applied": {
                    "dimensions": f"{dimensions['width'].inches:.1f}x{dimensions['height'].inches:.1f} inches",
                    "colors": "Palette Premier Tech appliquée",
                    "formatting": "Standards visuels respectés"
                },
                "position": position + 1,
                "total_slides": len(target_pres.slides)
            }

            # Sauvegarder le rapport
            report_path = target_presentation_path.replace('.pptx', '_enhanced_chart_report.json')
            with open(report_path, 'w', encoding='utf-8') as f:
                json.dump(report_data, f, ensure_ascii=False, indent=2)

            print(f"[REPORT] Rapport d'insertion: {report_path}")

            return {
                "presentation_path": target_presentation_path,
                "backup_path": backup_path,
                "report_path": report_path,
                "style_used": style,
                "position": position + 1,
                "total_slides": len(target_pres.slides)
            }

        except Exception as e:
            print(f"[ERROR] Erreur lors de l'insertion: {e}")
            # Restaurer la sauvegarde
            if os.path.exists(backup_path):
                shutil.copy2(backup_path, target_presentation_path)
                print(f"[RESTORE] Présentation originale restaurée")
            raise


def main():
    """Fonction principale avec interface CLI améliorée"""
    parser = argparse.ArgumentParser(
        description="Enhanced Charts Builder - Création avancée de graphiques Premier Tech",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Exemples d'usage:

  # Import depuis CSV simple (2 colonnes: Label, Valeur)
  python 09_charts_builder.py "Ventes Q4" --insert-into presentation.pptx \\
    --csv data/ventes_q4.csv --style column_clustered

  # Import CSV multi-séries
  python 09_charts_builder.py "Comparaison Régions" --insert-into presentation.pptx \\
    --csv data/regions_comparison.csv --style line_chart

  # Configuration JSON complète
  python 09_charts_builder.py --insert-into presentation.pptx \\
    --json-config charts/config_budget.json

  # Données directes avec insights
  python 09_charts_builder.py "Performance KPI" --insert-into presentation.pptx \\
    --style pie_chart --data-labels "Atteint" "En cours" "Retard" \\
    --data-values "65" "25" "10" --insights "65% des objectifs atteints"

  # Export des données d'un graphique existant
  python 09_charts_builder.py --export-from presentation.pptx \\
    --export-csv exported_data.csv --slide-index -1

Formats CSV supportés:

  Format simple (2 colonnes):
    Catégorie,Valeur
    T1,2500000
    T2,3200000
    T3,2800000
    T4,3700000

  Format multi-séries:
    Trimestre,Europe,Amérique,Asie
    T1,1200000,800000,500000
    T2,1500000,1000000,700000
    T3,1300000,900000,600000
    T4,1800000,1200000,700000

Format JSON:
  {
    "title": "Titre du graphique",
    "style": "column_clustered",
    "data": {
      "labels": ["T1", "T2", "T3", "T4"],
      "values": [2500000, 3200000, 2800000, 3700000]
    },
    "insights": "Croissance continue avec pic en T4"
  }

Styles disponibles:
  - column_clustered : Colonnes groupées (défaut)
  - line_chart       : Graphique linéaire
  - pie_chart        : Graphique en secteurs
  - bar_clustered    : Barres horizontales
  - column_compact   : Colonnes compactes
  - bar_compact      : Barres compactes
        """
    )

    # Arguments principaux
    parser.add_argument("title", nargs="?", help="Titre de la slide graphique")

    # Import de données
    parser.add_argument("--csv", help="Importer données depuis fichier CSV")
    parser.add_argument("--json-config", help="Configuration complète depuis JSON")

    # Données directes
    parser.add_argument("--data-labels", nargs="*", help="Labels des données")
    parser.add_argument("--data-values", nargs="*", help="Valeurs des données")
    parser.add_argument("--insights", help="Points clés à retenir")
    parser.add_argument("--series-title", help="Titre personnalisé pour la série de données")

    # Options de style
    parser.add_argument("--style", choices=[
        "column_clustered", "line_chart", "pie_chart",
        "bar_clustered", "column_compact", "bar_compact"
    ], default="column_clustered", help="Style de graphique")

    # Options de fichier
    parser.add_argument("--template", default="../templates/Template_PT.pptx",
                        help="Chemin vers le template Premier Tech")
    parser.add_argument("--insert-into", help="Présentation cible")
    parser.add_argument("--position", type=int, help="Position d'insertion")

    # Export
    parser.add_argument("--export-from", help="Exporter données depuis présentation")
    parser.add_argument("--export-csv", help="Fichier CSV de sortie pour export")
    parser.add_argument("--slide-index", type=int, default=-1,
                        help="Index de la slide à exporter (défaut: dernière)")

    # Options utilitaires
    parser.add_argument("--validate", action="store_true", help="Valider le template")
    parser.add_argument("--list-styles", action="store_true", help="Lister les styles")
    parser.add_argument("--create-sample-csv", help="Créer un CSV d'exemple")
    parser.add_argument("--create-sample-json", help="Créer un JSON d'exemple")

    args = parser.parse_args()

    try:
        # Initialiser le constructeur
        builder = EnhancedChartsBuilder(args.template)

        # Mode validation
        if args.validate:
            print("[OK] Template validé avec succès")
            return

        # Mode liste des styles
        if args.list_styles:
            print("\n[STYLES] Styles de graphiques disponibles:\n")
            for info in builder.chart_slides.values():
                print(f"- {info['style']:<15} : {info['name']}")
                print(f"  Usage: {info['usage']}")
                print(f"  Type: {info['chart_type']}")
                print(f"  Audience: {info['audience']}\n")
            return

        # Créer CSV d'exemple
        if args.create_sample_csv:
            sample_data = """Catégorie,Valeur
T1 2024,2500000
T2 2024,3200000
T3 2024,2800000
T4 2024,3700000"""
            with open(args.create_sample_csv, 'w', encoding='utf-8') as f:
                f.write(sample_data)
            print(f"[SAMPLE] CSV d'exemple créé: {args.create_sample_csv}")
            return

        # Créer JSON d'exemple
        if args.create_sample_json:
            sample_config = {
                "title": "Exemple de Configuration Graphique",
                "style": "column_clustered",
                "data": {
                    "labels": ["T1 2024", "T2 2024", "T3 2024", "T4 2024"],
                    "values": [2500000, 3200000, 2800000, 3700000]
                },
                "insights": "Croissance continue avec pic en T4",
                "formatting": {
                    "show_legend": True,
                    "show_data_labels": True
                }
            }
            with open(args.create_sample_json, 'w', encoding='utf-8') as f:
                json.dump(sample_config, f, ensure_ascii=False, indent=2)
            print(f"[SAMPLE] JSON d'exemple créé: {args.create_sample_json}")
            return

        # Mode export
        if args.export_from:
            if not args.export_csv:
                parser.error("--export-csv requis avec --export-from")
            builder.export_chart_data(args.export_from, args.export_csv, args.slide_index)
            return

        # Mode insertion
        if not args.insert_into:
            parser.error("--insert-into est requis pour insérer une slide")

        if not args.title and not args.json_config:
            parser.error("Le titre est requis (sauf avec --json-config)")

        # Insertion avec toutes les options
        result = builder.insert_enhanced_chart_slide(
            args.insert_into,
            args.title or "Graphique",
            args.style,
            args.csv,
            args.json_config,
            args.data_labels,
            args.data_values,
            None,  # series_data géré via CSV
            args.insights,
            args.position,
            args.series_title
        )

        print(f"\n[SUCCESS] Slide graphique améliorée insérée!")
        print(f"[FILE] Présentation: {result['presentation_path']}")
        print(f"[POSITION] Position: {result['position']}/{result['total_slides']}")

    except FileNotFoundError as e:
        print(f"[ERROR] Fichier non trouvé: {e}")
        sys.exit(1)
    except ValueError as e:
        print(f"[ERROR] Erreur de validation: {e}")
        sys.exit(1)
    except Exception as e:
        print(f"[ERROR] Erreur inattendue: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)


if __name__ == "__main__":
    main()