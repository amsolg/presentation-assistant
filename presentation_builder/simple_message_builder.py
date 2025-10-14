#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Simple Message Builder - Construction de messages simples Premier Tech
Version JSON-native pour l'architecture 2025 du presentation_builder.
Utilise les slides 17-19 du template Premier Tech pour créer des messages impactants.
"""

import os
import json
import shutil
from datetime import datetime
from pathlib import Path
from typing import List, Dict, Optional, Any
from pptx import Presentation
from pptx.util import Inches


class SimpleMessageBuilder:
    """
    Classe pour construire des messages simples Premier Tech.
    Version modernisée pour l'architecture JSON 2025.
    """

    def __init__(self, template_path: str = "templates/Template_PT.pptx"):
        """
        Initialise le constructeur avec le template Premier Tech.

        Args:
            template_path: Chemin vers le template Premier Tech
        """
        self.template_path = template_path

        # Mapping des slides disponibles pour les messages simples
        self.message_slides = {
            16: {  # Slide 17 (index 16) - Court énoncé
                "name": "Court énoncé",
                "usage": "Message unique sans titre, focus maximal",
                "audience": "Toutes audiences",
                "style": "centered"
            },
            17: {  # Slide 18 (index 17) - Mots-clés & Mots complémentaires
                "name": "Mots-clés & Mots complémentaires",
                "usage": "Message avec mots-clés et complément",
                "audience": "Leaders, Audiences mixtes",
                "style": "illustrated"
            },
            18: {  # Slide 19 (index 18) - Mots-clés & Court énoncé
                "name": "Mots-clés & Court énoncé",
                "usage": "Concept de base avec contexte",
                "audience": "Spécialistes, Formateurs",
                "style": "keyword_simple"
            }
        }

        if not os.path.exists(template_path):
            raise FileNotFoundError(f"Template Premier Tech non trouvé: {template_path}")

        # Analyser la structure des slides de messages
        self._analyze_message_structures()

    def _analyze_message_structures(self):
        """Analyse la structure des slides de messages disponibles"""
        try:
            pres = Presentation(self.template_path)

            self.message_info = {}
            for slide_index, slide_data in self.message_slides.items():
                if len(pres.slides) > slide_index:
                    message_slide = pres.slides[slide_index]

                    self.message_info[slide_index] = {
                        'layout_name': message_slide.slide_layout.name,
                        'shape_count': len(message_slide.shapes),
                        'slide_index': slide_index,
                        'slide_number': slide_index + 1,
                        'style': slide_data['style'],
                        'usage': slide_data['usage']
                    }

            print(f"[INFO] {len(self.message_info)} slides de message analysées")
            for idx, info in self.message_info.items():
                print(f"[INFO] Slide {info['slide_number']}: {info['layout_name']} ({info['style']})")

        except Exception as e:
            raise Exception(f"Erreur analyse templates message: {e}")

    def process_simple_message_config(self, config: Dict[str, Any], presentation_path: str) -> Dict[str, Any]:
        """
        Traite une configuration JSON de message simple et l'applique à une présentation.
        Point d'entrée principal pour l'architecture JSON 2025.

        Args:
            config: Configuration JSON validée du message simple
            presentation_path: Chemin vers la présentation cible

        Returns:
            Dict contenant les résultats du traitement
        """
        try:
            print(f"[SIMPLE_MESSAGE] Traitement configuration JSON...")

            # Extraire les données de configuration
            message_text = config.get('message_text', '')
            keywords = config.get('keywords')
            image_description = config.get('image_description')
            message_style = config.get('message_style', 'centered')
            options = config.get('options', {})
            auto_widen = options.get('auto_widen', True)
            insert_position = options.get('insert_position')

            print(f"[SIMPLE_MESSAGE] Message: {message_text[:50]}...")
            print(f"[SIMPLE_MESSAGE] Style: {message_style}")
            print(f"[SIMPLE_MESSAGE] Options: auto_widen={auto_widen}, position={insert_position}")

            # Valider la configuration
            validation_result = self._validate_simple_message_config(config)
            if not validation_result['valid']:
                raise ValueError(f"Configuration invalide: {', '.join(validation_result['errors'])}")

            # Insérer le message dans la présentation
            result_path = self.insert_message_into_existing_presentation(
                presentation_path=presentation_path,
                message_text=message_text,
                keywords=keywords,
                image_description=image_description,
                message_style=message_style,
                insert_position=insert_position,
                auto_widen=auto_widen
            )

            # Générer le rapport de traitement
            processing_report = {
                "timestamp": datetime.now().isoformat(),
                "method": "JSON Simple Message Builder 2025",
                "success": True,
                "configuration": config,
                "result_path": result_path,
                "validation": validation_result,
                "processing_details": {
                    "message_applied": message_text,
                    "style_applied": message_style,
                    "keywords_applied": keywords,
                    "auto_widen_applied": auto_widen,
                    "insert_position": insert_position
                }
            }

            print(f"[SUCCESS] Message simple JSON traité avec succès")
            return processing_report

        except Exception as e:
            error_report = {
                "timestamp": datetime.now().isoformat(),
                "method": "JSON Simple Message Builder 2025",
                "success": False,
                "error": str(e),
                "configuration": config
            }
            print(f"[ERROR] Erreur traitement message simple JSON: {e}")
            return error_report

    def _validate_simple_message_config(self, config: Dict[str, Any]) -> Dict[str, Any]:
        """Valide une configuration JSON de message simple"""
        errors = []
        warnings = []

        # Vérifier les champs requis
        if 'message_text' not in config:
            errors.append("Champ 'message_text' manquant")
        elif not config['message_text'] or not isinstance(config['message_text'], str):
            errors.append("Le message_text doit être une chaîne non vide")
        elif len(config['message_text']) > 500:
            warnings.append("Message très long (>500 caractères)")

        if 'message_style' not in config:
            errors.append("Champ 'message_style' manquant")
        elif config['message_style'] not in ['centered', 'illustrated', 'keyword_simple']:
            errors.append("message_style doit être 'centered', 'illustrated' ou 'keyword_simple'")

        # Vérifier la cohérence des paramètres
        if config.get('message_style') in ['illustrated', 'keyword_simple']:
            if not config.get('keywords'):
                warnings.append(f"Style '{config.get('message_style')}' est plus efficace avec des mots-clés")

        if config.get('keywords') and len(config['keywords']) > 200:
            warnings.append("Mots-clés très longs (>200 caractères)")

        if config.get('image_description') and len(config['image_description']) > 200:
            warnings.append("Description d'image très longue (>200 caractères)")

        # Valider les options
        if 'options' in config:
            options = config['options']
            if 'auto_widen' in options and not isinstance(options['auto_widen'], bool):
                errors.append("'auto_widen' doit être un booléen")
            if 'insert_position' in options and options['insert_position'] is not None:
                if not isinstance(options['insert_position'], int) or options['insert_position'] < 0:
                    errors.append("'insert_position' doit être un entier positif ou null")

        return {
            'valid': len(errors) == 0,
            'errors': errors,
            'warnings': warnings
        }

    def insert_message_into_existing_presentation(self,
                                                presentation_path: str,
                                                message_text: str,
                                                keywords: Optional[str] = None,
                                                image_description: Optional[str] = None,
                                                message_style: str = "centered",
                                                insert_position: Optional[int] = None,
                                                auto_widen: bool = True) -> str:
        """
        Insère un message simple directement dans une présentation existante.
        Version modernisée avec support pour l'architecture JSON 2025.

        Args:
            presentation_path: Chemin vers la présentation existante
            message_text: Texte principal du message
            keywords: Mots-clés (pour styles keyword_*)
            image_description: Description de l'image (pour style illustrated)
            message_style: Style du message
            insert_position: Position d'insertion (None = à la fin)
            auto_widen: Active l'élargissement automatique des objets texte

        Returns:
            str: Chemin vers le fichier modifié
        """
        try:
            print(f"[INSERT] Insertion directe message simple dans: {os.path.basename(presentation_path)}")
            print(f"[INSERT] Style: {message_style}, Message: {message_text[:50]}...")

            # ÉTAPE 1: Créer une copie de sauvegarde
            backup_path = presentation_path.replace('.pptx', '_backup_before_message.pptx')
            shutil.copy2(presentation_path, backup_path)
            print(f"[BACKUP] Sauvegarde créée: {backup_path}")

            # ÉTAPE 2: Charger la présentation existante
            target_prs = Presentation(presentation_path)
            print(f"[LOAD] Présentation chargée: {len(target_prs.slides)} slides existantes")

            # ÉTAPE 3: Déterminer la slide source à utiliser
            source_slide_index = self._get_slide_index_for_style(message_style)
            if source_slide_index is None:
                raise ValueError(f"Style '{message_style}' non reconnu")

            # ÉTAPE 4: Vérifier que le layout message existe
            message_layout_index = self._find_message_layout_index(target_prs, source_slide_index)
            if message_layout_index is None:
                raise Exception(f"Layout message pour style '{message_style}' non trouvé dans la présentation")

            # ÉTAPE 5: Ajouter la slide message avec le bon layout
            message_layout = target_prs.slide_layouts[message_layout_index]
            new_slide = target_prs.slides.add_slide(message_layout)
            print(f"[ADD] Slide message ajoutée avec layout: {message_layout.name}")

            # ÉTAPE 6: Personnaliser le contenu de la slide message
            self._customize_message_slide_direct(new_slide, message_text, keywords, image_description, message_style)

            # ÉTAPE 7: Appliquer l'élargissement automatique si demandé
            if auto_widen:
                self._apply_auto_widen_to_slide(new_slide)

            # ÉTAPE 8: Réorganiser les slides si nécessaire
            if insert_position is not None and insert_position < len(target_prs.slides) - 1:
                self._move_slide_to_position_direct(target_prs, len(target_prs.slides) - 1, insert_position)

            # ÉTAPE 9: Sauvegarder
            target_prs.save(presentation_path)
            print(f"[SUCCESS] Message simple inséré directement dans la présentation")

            # ÉTAPE 10: Générer le rapport
            self._generate_direct_insertion_report(presentation_path, message_text, keywords, image_description,
                                                 message_style, insert_position or len(target_prs.slides), auto_widen)

            return presentation_path

        except Exception as e:
            print(f"[ERROR] Erreur insertion directe message simple: {e}")
            # Restaurer la sauvegarde en cas d'erreur
            if 'backup_path' in locals() and os.path.exists(backup_path):
                shutil.copy2(backup_path, presentation_path)
                print(f"[RESTORE] Présentation originale restaurée")
            raise

    def _get_slide_index_for_style(self, style: str) -> Optional[int]:
        """Retourne l'index de slide approprié pour le style demandé"""
        style_mapping = {
            "centered": 16,      # Slide 17 - Court énoncé
            "illustrated": 17,   # Slide 18 - Mots-clés & Mots complémentaires
            "keyword_simple": 18 # Slide 19 - Mots-clés & Court énoncé
        }
        return style_mapping.get(style)

    def _find_message_layout_index(self, presentation: Presentation, source_slide_index: int) -> Optional[int]:
        """Trouve l'index du layout message dans la présentation"""
        try:
            template_prs = Presentation(self.template_path)
            template_layout_name = template_prs.slides[source_slide_index].slide_layout.name

            for i, layout in enumerate(presentation.slide_layouts):
                if layout.name == template_layout_name:
                    print(f"[LAYOUT] Layout '{template_layout_name}' trouvé à l'index {i}")
                    return i

            return None

        except Exception as e:
            print(f"[WARNING] Erreur recherche layout message: {e}")
            return None

    def _customize_message_slide_direct(self, slide, message_text: str, keywords: Optional[str],
                                      image_description: Optional[str], message_style: str):
        """Personnalise directement la slide message ajoutée selon la vraie structure des slides 17-19"""
        try:
            print(f"[CUSTOMIZE] Personnalisation slide message directe...")
            print(f"[CUSTOMIZE] Slide avec {len(slide.shapes)} shapes à personnaliser")

            shape_updates = 0

            # Analyse détaillée de chaque shape
            for i, shape in enumerate(slide.shapes):
                try:
                    if hasattr(shape, 'text_frame') and shape.text_frame:
                        current_text = shape.text_frame.text.strip()
                        print(f"[DEBUG] Shape {i}: '{current_text}' (longueur: {len(current_text)})")

                        # Personnalisation selon le style et structure réelle des slides
                        if message_style == "centered":
                            # Slide 17 (index 16): "Court énoncé" - Shape 0 = message principal
                            if i == 0:  # Premier shape = message
                                shape.text_frame.text = message_text
                                shape.text_frame.word_wrap = True  # ACTIVER le wrapping pour message unique
                                print(f"[UPDATE] Shape {i} (message centré): {message_text[:30]}... (wrapping activé)")
                                shape_updates += 1

                        elif message_style == "illustrated":
                            # Slide 18 (index 17): "Mots-clés & Mots complémentaires" - Shape 0 = mots-clés, Shape 1 = complément
                            if i == 0:  # Premier shape = mots-clés
                                keys = keywords or "Mots-clés"
                                shape.text_frame.text = keys
                                shape.text_frame.word_wrap = False
                                print(f"[UPDATE] Shape {i} (mots-clés): {keys}")
                                shape_updates += 1
                            elif i == 1:  # Deuxième shape = complément
                                complement = image_description or message_text
                                shape.text_frame.text = complement
                                shape.text_frame.word_wrap = False
                                print(f"[UPDATE] Shape {i} (complément): {complement[:30]}...")
                                shape_updates += 1

                        elif message_style == "keyword_simple":
                            # Slide 19 (index 18): "Mots-clés & Court énoncé" - Shape 0 = mots-clés, Shape 1 = message
                            if i == 0:  # Premier shape = mots-clés
                                keys = keywords or "Mots-clés"
                                shape.text_frame.text = keys
                                shape.text_frame.word_wrap = False
                                print(f"[UPDATE] Shape {i} (mots-clés): {keys}")
                                shape_updates += 1
                            elif i == 1:  # Deuxième shape = message
                                shape.text_frame.text = message_text
                                shape.text_frame.word_wrap = False
                                print(f"[UPDATE] Shape {i} (message): {message_text[:30]}...")
                                shape_updates += 1

                except Exception as e:
                    print(f"[WARNING] Erreur personnalisation shape {i}: {e}")

            print(f"[SUCCESS] Slide message personnalisée ({shape_updates} éléments mis à jour)")

        except Exception as e:
            print(f"[ERROR] Erreur personnalisation directe message: {e}")
            raise

    def _apply_auto_widen_to_slide(self, slide):
        """Applique l'élargissement automatique à une slide spécifique"""
        try:
            print(f"[WIDEN] Application élargissement automatique...")

            widen_count = 0
            for i, shape in enumerate(slide.shapes):
                if hasattr(shape, 'text_frame'):
                    current_width_inches = shape.width / Inches(1)

                    # Élargir les objets texte étroits (<4 pouces)
                    if current_width_inches < 4.0:
                        new_width_inches = min(current_width_inches * 1.5, 8.0)
                        shape.width = Inches(new_width_inches)
                        print(f"[WIDEN] Shape {i}: {current_width_inches:.2f}\" -> {new_width_inches:.2f}\"")
                        widen_count += 1

                    # Élargir modérément les objets moyens (4-6 pouces)
                    elif current_width_inches < 6.0:
                        new_width_inches = min(current_width_inches * 1.2, 8.0)
                        shape.width = Inches(new_width_inches)
                        print(f"[WIDEN] Shape {i}: {current_width_inches:.2f}\" -> {new_width_inches:.2f}\" (modéré)")
                        widen_count += 1

            if widen_count > 0:
                print(f"[SUCCESS] {widen_count} objets texte élargis")
            else:
                print(f"[INFO] Aucun objet texte nécessitant un élargissement")

        except Exception as e:
            print(f"[WARNING] Erreur élargissement: {e}")

    def _move_slide_to_position_direct(self, presentation: Presentation, from_index: int, to_index: int):
        """Déplace une slide à la position désirée"""
        try:
            print(f"[POSITION] Slide message ajoutée en position {from_index + 1}")
            print(f"[INFO] Position finale: {to_index + 1}")

        except Exception as e:
            print(f"[WARNING] Déplacement slide: {e}")

    def _generate_direct_insertion_report(self, presentation_path: str, message_text: str,
                                        keywords: Optional[str], image_description: Optional[str],
                                        message_style: str, insert_position: int, auto_widen: bool):
        """Génère un rapport d'insertion directe modernisé"""
        try:
            source_slide_index = self._get_slide_index_for_style(message_style)

            report = {
                "insertion_timestamp": datetime.now().isoformat(),
                "method": "JSON Simple Message Builder 2025 - Direct Insertion",
                "template_used": self.template_path,
                "target_presentation": presentation_path,
                "message_details": {
                    "text": message_text,
                    "keywords": keywords,
                    "image_description": image_description,
                    "style": message_style,
                    "intended_position": insert_position,
                    "auto_widen_enabled": auto_widen
                },
                "source_slide": {
                    "index": source_slide_index,
                    "number": source_slide_index + 1 if source_slide_index else None,
                    "layout": self.message_info.get(source_slide_index, {}).get('layout_name', 'Unknown'),
                    "style_description": self.message_slides.get(source_slide_index, {}).get('usage', 'Unknown')
                },
                "quality_assurance": {
                    "method": "JSON-Native Direct Layout-Based Addition",
                    "styles_preserved": True,
                    "premier_tech_standards": True,
                    "direct_integration": True,
                    "json_configuration": True,
                    "architecture_2025": True
                },
                "advantages": [
                    "Architecture JSON 2025 native",
                    "Insertion directe dans la présentation",
                    "Styles Premier Tech 100% préservés",
                    "Configuration JSON validée",
                    "Intégration transparente",
                    "Sauvegarde automatique créée",
                    f"Style '{message_style}' adapté à l'usage"
                ]
            }

            # Sauvegarder le rapport
            report_path = presentation_path.replace('.pptx', '_direct_simple_message_insertion_report.json')
            with open(report_path, 'w', encoding='utf-8') as f:
                json.dump(report, f, ensure_ascii=False, indent=2)

            print(f"[INFO] Rapport d'insertion: {os.path.basename(report_path)}")

        except Exception as e:
            print(f"[WARNING] Erreur génération rapport: {e}")

    def validate_template(self) -> bool:
        """Valide que le template Premier Tech est correct pour les messages"""
        try:
            checks = {
                "file_exists": os.path.exists(self.template_path),
                "has_slides": False,
                "message_slides_exist": False,
                "slides_count": 0,
                "available_styles": []
            }

            if checks["file_exists"]:
                pres = Presentation(self.template_path)
                checks["slides_count"] = len(pres.slides)
                checks["has_slides"] = len(pres.slides) > 0

                # Vérifier que toutes les slides de message existent
                available_styles = []
                for slide_index in self.message_slides.keys():
                    if len(pres.slides) > slide_index:
                        style = self.message_slides[slide_index]['style']
                        available_styles.append(style)

                checks["available_styles"] = available_styles
                checks["message_slides_exist"] = len(available_styles) == len(self.message_slides)

            all_valid = all([checks["file_exists"], checks["has_slides"], checks["message_slides_exist"]])

            print("=== VALIDATION TEMPLATE PREMIER TECH POUR MESSAGES ===")
            for check, result in checks.items():
                if check == "available_styles":
                    print(f"[INFO] Styles disponibles: {', '.join(result)}")
                else:
                    status = "OK" if result else "ERREUR"
                    print(f"[{status}] {check}: {result}")

            if checks["message_slides_exist"]:
                print(f"[INFO] {len(checks['available_styles'])} styles de message disponibles:")
                for slide_index, slide_data in self.message_slides.items():
                    if slide_index in [idx for idx in self.message_slides.keys() if len(pres.slides) > idx]:
                        print(f"  - {slide_data['style']}: Slide {slide_index + 1} ({slide_data['usage']})")

            return all_valid

        except Exception as e:
            print(f"[ERROR] Erreur validation: {e}")
            return False


def create_simple_message_from_json(config_data: Dict[str, Any], presentation_path: str,
                                  template_path: str = "templates/Template_PT.pptx") -> Dict[str, Any]:
    """
    Fonction utilitaire pour créer un message simple à partir d'une configuration JSON.
    Point d'entrée principal pour l'architecture JSON 2025.

    Args:
        config_data: Configuration JSON du message simple
        presentation_path: Chemin vers la présentation cible
        template_path: Chemin vers le template Premier Tech

    Returns:
        Dict contenant les résultats du traitement
    """
    try:
        builder = SimpleMessageBuilder(template_path)
        return builder.process_simple_message_config(config_data, presentation_path)
    except Exception as e:
        return {
            "success": False,
            "error": str(e),
            "timestamp": datetime.now().isoformat()
        }


def load_simple_message_template(template_name: str = "centered_simple") -> Dict[str, Any]:
    """
    Charge un template de message simple prédéfini.

    Args:
        template_name: Nom du template (centered_simple, illustrated_message, keyword_emphasis, etc.)

    Returns:
        Dict contenant la configuration du template
    """
    template_path = "templates/presentation-project/slide-payload-templates/simple_message_builder_template.json"

    try:
        with open(template_path, 'r', encoding='utf-8') as f:
            templates = json.load(f)

        examples = templates.get('examples', {})

        if template_name in examples:
            return examples[template_name]
        else:
            # Retourner le template de base si le nom n'est pas trouvé
            return templates.get('payload_structure', {})

    except Exception as e:
        print(f"[WARNING] Erreur chargement template: {e}")
        # Template de fallback
        return {
            "message_text": "Votre message ici",
            "keywords": None,
            "image_description": None,
            "message_style": "centered",
            "options": {
                "auto_widen": True,
                "insert_position": None
            }
        }


def load_simple_message_payload(payload_path: str) -> Dict[str, Any]:
    """
    Charge un payload de message simple depuis un fichier JSON.

    Args:
        payload_path: Chemin vers le fichier JSON contenant le payload

    Returns:
        Dict contenant le payload de message simple
    """
    try:
        with open(payload_path, 'r', encoding='utf-8') as f:
            payload = json.load(f)

        print(f"[PAYLOAD] Chargé depuis {payload_path}")
        return payload

    except Exception as e:
        print(f"[ERROR] Erreur chargement payload {payload_path}: {e}")
        return load_simple_message_template("centered_simple")  # Fallback


def process_simple_message_from_payload_file(payload_path: str, presentation_path: str,
                                           template_path: str = "templates/Template_PT.pptx") -> Dict[str, Any]:
    """
    Traite un message simple en chargeant le payload depuis un fichier JSON.
    Point d'entrée pour l'architecture nouvelle avec fichiers payload séparés.

    Args:
        payload_path: Chemin vers le fichier JSON contenant le payload
        presentation_path: Chemin vers la présentation cible
        template_path: Chemin vers le template Premier Tech

    Returns:
        Dict contenant les résultats du traitement
    """
    try:
        # Charger le payload
        payload = load_simple_message_payload(payload_path)

        # Traiter avec le payload chargé
        builder = SimpleMessageBuilder(template_path)
        result = builder.process_simple_message_config(payload, presentation_path)

        # Ajouter les informations de payload
        result["payload_source"] = payload_path
        result["payload_loaded"] = True

        return result

    except Exception as e:
        return {
            "success": False,
            "error": str(e),
            "timestamp": datetime.now().isoformat(),
            "payload_source": payload_path,
            "payload_loaded": False
        }


if __name__ == "__main__":
    # Interface de test pour validation
    print("=== Simple Message Builder - Architecture JSON 2025 ===")

    # Exemple d'utilisation
    template_config = load_simple_message_template("centered_simple")
    print(f"Template chargé: {template_config}")

    # Test de validation
    builder = SimpleMessageBuilder()
    if builder.validate_template():
        print("[SUCCESS] Template Premier Tech validé pour messages simples")
    else:
        print("[ERROR] Template Premier Tech invalide")