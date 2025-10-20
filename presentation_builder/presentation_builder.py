#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Presentation Builder 3 - Architecture slide-structure JSON
Script principal qui coordonne la création de présentations complètes
à partir d'un fichier JSON de configuration avec structures slide-structure.

Architecture:
- Crée une présentation avec les slides définies dans le JSON (slide-structure)
- Utilise une fonction de personnalisation universelle pour toutes les slides
- Plus de slides obligatoires titre/fermeture
- Output unique : fichier .pptx (pas de rapports/backups)
"""

import os
import sys
import json
import argparse
from datetime import datetime
from pathlib import Path
from typing import Dict, List, Any, Optional
import shutil
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor


class PresentationBuilder:
    """
    Orchestrateur principal pour construire des présentations à partir de JSON.

    Workflow:
    1. Parse le JSON de configuration avec nouveau schéma slide-structure
    2. Crée une présentation avec la première slide du JSON
    3. Ajoute toutes les autres slides définies dans le JSON
    4. Utilise une fonction de personnalisation universelle pour toutes les slides
    """

    def __init__(self):
        """Initialise l'orchestrateur avec les chemins de base"""
        self.script_dir = Path(__file__).parent
        self.template_path = self.script_dir.parent / "templates" / "Template_PT.pptx"
        self.slide_structures_path = self.script_dir.parent / "templates" / "presentation-project" / "slide-structure"
        self.premier_tech_enums_path = self.script_dir.parent / "templates" / "presentation-project" / "premier_tech_schema_enums.json"

        # Vérifier l'existence du template
        if not self.template_path.exists():
            raise FileNotFoundError(f"Template Premier Tech non trouvé: {self.template_path}")

        # Vérifier l'existence des structures de slides
        if not self.slide_structures_path.exists():
            raise FileNotFoundError(f"Dossier slide-structure non trouvé: {self.slide_structures_path}")

        # Charger les enums Premier Tech pour validation
        self.premier_tech_enums = self._load_premier_tech_enums()

        print(f"[INIT] Template Premier Tech: {self.template_path}")
        print(f"[INIT] Structures slides: {self.slide_structures_path}")
        print(f"[INIT] Premier Tech enums loaded: {len(self.premier_tech_enums.get('enums', {}))} categories")

    def _load_premier_tech_enums(self) -> Dict[str, Any]:
        """
        Charge les enums Premier Tech pour validation des propriétés.

        Returns:
            Dict: Enums Premier Tech ou dict vide si fichier non trouvé
        """
        try:
            if not self.premier_tech_enums_path.exists():
                print(f"[WARNING] Fichier enums Premier Tech non trouvé: {self.premier_tech_enums_path}")
                return {}

            with open(self.premier_tech_enums_path, 'r', encoding='utf-8') as f:
                enums_data = json.load(f)

            print(f"[INIT] Premier Tech enums chargés: {enums_data.get('total_slides_analyzed', 0)} slides analysées")
            return enums_data

        except Exception as e:
            print(f"[WARNING] Erreur chargement enums Premier Tech: {e}")
            return {}

    def _validate_property_value(self, property_name: str, value: Any) -> bool:
        """
        Valide une valeur de propriété contre les enums Premier Tech.

        Args:
            property_name: Nom de la propriété à valider
            value: Valeur à valider

        Returns:
            bool: True si valide, False sinon
        """
        if not self.premier_tech_enums or 'enums' not in self.premier_tech_enums:
            return True  # Pas de validation si enums non disponibles

        enums = self.premier_tech_enums['enums']

        if property_name not in enums:
            return True  # Propriété non validée

        allowed_values = enums[property_name].get('enum', [])

        if value not in allowed_values:
            print(f"[WARNING] Valeur '{value}' non valide pour '{property_name}'. Valeurs autorisées: {allowed_values}")
            return False

        return True

    def load_presentation_config(self, json_path: str) -> Dict[str, Any]:
        """
        Charge et valide le fichier JSON de configuration selon le nouveau schéma.

        Args:
            json_path: Chemin vers le fichier JSON

        Returns:
            Dict: Configuration de la présentation

        Raises:
            ValueError: Si le JSON est invalide
            FileNotFoundError: Si le fichier n'existe pas
        """
        try:
            with open(json_path, 'r', encoding='utf-8') as f:
                config = json.load(f)

            # Validation des champs requis selon le nouveau schéma
            required_fields = ["presentation_name", "subject", "audience", "slides", "output_path"]
            for field in required_fields:
                if field not in config:
                    raise ValueError(f"Champ requis manquant: {field}")

            # Validation des slides
            if not isinstance(config["slides"], list) or len(config["slides"]) == 0:
                raise ValueError("Le tableau 'slides' ne peut pas être vide")

            for i, slide in enumerate(config["slides"]):
                if "slide_number" not in slide:
                    raise ValueError(f"Slide {i+1}: 'slide_number' manquant")
                if not isinstance(slide["slide_number"], int) or slide["slide_number"] < 1 or slide["slide_number"] > 57:
                    raise ValueError(f"Slide {i+1}: 'slide_number' doit être entre 1 et 57")

            # Validation output_path
            if not config["output_path"].endswith(".pptx"):
                raise ValueError("output_path doit se terminer par .pptx")

            # Traitement du chemin de base selon is_test
            is_test = config.get("is_test", False)
            if is_test:
                # Remplacer "presentations/" par "tests/" dans le chemin de sortie
                output_path = config["output_path"]
                if output_path.startswith("presentations/"):
                    config["output_path"] = output_path.replace("presentations/", "tests/", 1)
                    print(f"[CONFIG] Mode test activé - Redirection vers dossier 'tests'")

            print(f"[CONFIG] Présentation: {config['presentation_name']}")
            print(f"[CONFIG] Sujet: {config['subject']}")
            print(f"[CONFIG] Audience: {config['audience']}")
            print(f"[CONFIG] Mode test: {is_test}")
            print(f"[CONFIG] Output: {config['output_path']}")
            print(f"[CONFIG] Slides à créer: {len(config['slides'])}")

            return config

        except json.JSONDecodeError as e:
            raise ValueError(f"JSON invalide: {e}")
        except FileNotFoundError:
            raise FileNotFoundError(f"Fichier de configuration non trouvé: {json_path}")

    def load_slide_structure(self, slide_number: int) -> Dict[str, Any]:
        """
        Charge la structure d'une slide depuis les fichiers slide-structure.

        Args:
            slide_number: Numéro de la slide (1-57)

        Returns:
            Dict: Structure de la slide

        Raises:
            FileNotFoundError: Si le fichier de structure n'existe pas
        """
        structure_file = self.slide_structures_path / f"slide_{slide_number}.json"

        if not structure_file.exists():
            raise FileNotFoundError(f"Structure slide {slide_number} non trouvée: {structure_file}")

        try:
            with open(structure_file, 'r', encoding='utf-8') as f:
                structure = json.load(f)

            print(f"[STRUCTURE] Slide {slide_number}: {structure.get('layout_name', 'Unknown')} ({structure.get('total_shapes', 0)} shapes)")
            return structure

        except json.JSONDecodeError as e:
            raise ValueError(f"Structure slide {slide_number} invalide: {e}")

    def _customize_slide_universal(self, slide, slide_config: Dict[str, Any]) -> bool:
        """
        Fonction de personnalisation universelle pour toutes les slides.
        Utilise la configuration slide-structure pour personnaliser n'importe quelle slide.

        Args:
            slide: Objet slide PowerPoint
            slide_config: Configuration de la slide avec shapes à personnaliser

        Returns:
            bool: True si succès, False sinon
        """
        try:
            print(f"[CUSTOMIZE] Personnalisation universelle de la slide...")

            # Extraire les shapes à personnaliser depuis la configuration
            shapes_to_customize = slide_config.get("shapes", [])
            if not shapes_to_customize:
                print(f"[INFO] Aucune personnalisation définie pour cette slide")
                return True

            print(f"[CUSTOMIZE] {len(shapes_to_customize)} shapes à personnaliser")

            customized_count = 0
            for shape_config in shapes_to_customize:
                shape_id = shape_config.get("shape_id")
                if shape_id is None:
                    print(f"[WARNING] Shape sans shape_id, ignoré")
                    continue

                # Trouver le shape correspondant dans la slide
                target_shape = self._find_shape_by_id(slide, shape_id)
                if target_shape is None:
                    print(f"[WARNING] Shape ID {shape_id} non trouvé dans la slide")
                    continue

                # Appliquer la personnalisation
                if self._apply_shape_customization(target_shape, shape_config):
                    customized_count += 1

            print(f"[SUCCESS] {customized_count}/{len(shapes_to_customize)} shapes personnalisés")
            return customized_count > 0

        except Exception as e:
            print(f"[ERROR] Erreur personnalisation universelle: {e}")
            return False

    def _find_shape_by_id(self, slide, shape_id: int):
        """
        Trouve un shape dans une slide par son ID selon la structure slide-structure.

        Args:
            slide: Objet slide PowerPoint
            shape_id: ID du shape à trouver

        Returns:
            Shape PowerPoint ou None si non trouvé
        """
        try:
            # Dans les structures slide-structure, les shape_id correspondent aux index des shapes
            # shape_id 1 = index 0, shape_id 2 = index 1, etc.
            shape_index = shape_id - 1

            if 0 <= shape_index < len(slide.shapes):
                return slide.shapes[shape_index]
            else:
                print(f"[WARNING] Shape ID {shape_id} (index {shape_index}) hors limites (0-{len(slide.shapes)-1})")
                return None

        except Exception as e:
            print(f"[ERROR] Erreur recherche shape ID {shape_id}: {e}")
            return None

    def _apply_shape_customization(self, shape, shape_config: Dict[str, Any]) -> bool:
        """
        Applique la personnalisation COMPLETE à un shape selon la configuration slide-structure.
        Support de TOUTES les propriétés Premier Tech.

        Args:
            shape: Objet shape PowerPoint
            shape_config: Configuration du shape à appliquer

        Returns:
            bool: True si succès, False sinon
        """
        try:
            shape_id = shape_config.get('shape_id', 'unknown')
            print(f"[CUSTOMIZE] Application personnalisation complète shape {shape_id}")

            success = True

            # 1. Propriétés géométriques (position et dimensions)
            if not self._apply_geometry_properties(shape, shape_config):
                success = False

            # 2. Propriétés de texte de base (existant amélioré)
            if not self._apply_text_properties(shape, shape_config):
                success = False

            # 3. Propriétés de formatage avancées (marges, alignement vertical)
            if not self._apply_advanced_formatting(shape, shape_config):
                success = False

            # 4. Propriétés PowerPoint spécifiques (autofit, placeholders)
            if not self._apply_powerpoint_properties(shape, shape_config):
                success = False

            if success:
                print(f"[SUCCESS] Shape {shape_id} personnalisé avec succès")
            else:
                print(f"[WARNING] Shape {shape_id} partiellement personnalisé (certaines propriétés ont échoué)")

            return success

        except Exception as e:
            print(f"[ERROR] Erreur application personnalisation complète shape {shape_config.get('shape_id')}: {e}")
            return False

    def _apply_geometry_properties(self, shape, shape_config: Dict[str, Any]) -> bool:
        """
        Applique les propriétés géométriques (position et dimensions).

        Args:
            shape: Objet shape PowerPoint
            shape_config: Configuration du shape

        Returns:
            bool: True si succès, False sinon
        """
        try:
            # Position
            left = shape_config.get("left")
            if left is not None:
                self._validate_property_value("position", left)
                shape.left = Inches(left)

            top = shape_config.get("top")
            if top is not None:
                self._validate_property_value("position", top)
                shape.top = Inches(top)

            # Dimensions
            width = shape_config.get("width")
            if width is not None:
                self._validate_property_value("dimension", width)
                shape.width = Inches(width)

            height = shape_config.get("height")
            if height is not None:
                self._validate_property_value("dimension", height)
                shape.height = Inches(height)

            return True

        except Exception as e:
            print(f"[ERROR] Erreur application propriétés géométriques: {e}")
            return False

    def _apply_text_properties(self, shape, shape_config: Dict[str, Any]) -> bool:
        """
        Applique les propriétés de texte (version améliorée de l'existant).

        Args:
            shape: Objet shape PowerPoint
            shape_config: Configuration du shape

        Returns:
            bool: True si succès, False sinon
        """
        try:
            # Vérifier si le shape a un text_frame
            if not hasattr(shape, 'text_frame') or not shape.text_frame:
                # Pas d'erreur, certains shapes n'ont pas de texte
                return True

            # Appliquer le texte
            text = shape_config.get("text")
            if text is not None:
                shape.text_frame.text = text
                print(f"[UPDATE] Texte appliqué: '{text[:50]}{'...' if len(text) > 50 else ''}'")

            # Appliquer le formatage si le text_frame a des paragraphes
            if shape.text_frame.paragraphs:
                paragraph = shape.text_frame.paragraphs[0]

                # Créer ou récupérer le run
                if paragraph.runs:
                    run = paragraph.runs[0]
                else:
                    run = paragraph.add_run()

                # Police
                font_name = shape_config.get("font_name")
                if font_name:
                    self._validate_property_value("font_name", font_name)
                    run.font.name = font_name

                # Taille
                font_size = shape_config.get("font_size")
                if font_size:
                    self._validate_property_value("font_size", font_size)
                    run.font.size = Pt(font_size)

                # Gras
                bold = shape_config.get("bold")
                if bold is not None:
                    run.font.bold = bold

                # Couleur
                color = shape_config.get("color")
                if color and color.startswith("#"):
                    self._validate_property_value("color", color)
                    hex_color = color.lstrip("#")
                    rgb = tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
                    run.font.color.rgb = RGBColor(*rgb)

                # Alignement horizontal
                alignment = shape_config.get("alignment")
                if alignment:
                    self._validate_property_value("alignment", alignment)
                    if alignment == "LEFT":
                        paragraph.alignment = PP_ALIGN.LEFT
                    elif alignment == "CENTER":
                        paragraph.alignment = PP_ALIGN.CENTER
                    elif alignment == "RIGHT":
                        paragraph.alignment = PP_ALIGN.RIGHT

            return True

        except Exception as e:
            print(f"[ERROR] Erreur application propriétés texte: {e}")
            return False

    def _apply_advanced_formatting(self, shape, shape_config: Dict[str, Any]) -> bool:
        """
        Applique les propriétés de formatage avancées (marges, alignement vertical).

        Args:
            shape: Objet shape PowerPoint
            shape_config: Configuration du shape

        Returns:
            bool: True si succès, False sinon
        """
        try:
            if not hasattr(shape, 'text_frame') or not shape.text_frame:
                return True

            text_frame = shape.text_frame

            # Marges
            margin_left = shape_config.get("margin_left")
            if margin_left is not None:
                self._validate_property_value("margin_values", margin_left)
                text_frame.margin_left = Pt(margin_left)

            margin_right = shape_config.get("margin_right")
            if margin_right is not None:
                self._validate_property_value("margin_values", margin_right)
                text_frame.margin_right = Pt(margin_right)

            margin_top = shape_config.get("margin_top")
            if margin_top is not None:
                self._validate_property_value("margin_values", margin_top)
                text_frame.margin_top = Pt(margin_top)

            margin_bottom = shape_config.get("margin_bottom")
            if margin_bottom is not None:
                self._validate_property_value("margin_values", margin_bottom)
                text_frame.margin_bottom = Pt(margin_bottom)

            # Alignement vertical
            vertical_alignment = shape_config.get("vertical_alignment")
            if vertical_alignment:
                self._validate_property_value("vertical_alignment", vertical_alignment)
                if vertical_alignment == "TOP":
                    text_frame.vertical_anchor = MSO_ANCHOR.TOP
                elif vertical_alignment == "MIDDLE":
                    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                elif vertical_alignment == "BOTTOM":
                    text_frame.vertical_anchor = MSO_ANCHOR.BOTTOM

            # Text wrapping
            text_wrapping = shape_config.get("text_wrapping")
            if text_wrapping:
                self._validate_property_value("text_wrapping", text_wrapping)
                text_frame.word_wrap = (text_wrapping == "square")

            return True

        except Exception as e:
            print(f"[ERROR] Erreur application formatage avancé: {e}")
            return False

    def _apply_powerpoint_properties(self, shape, shape_config: Dict[str, Any]) -> bool:
        """
        Applique les propriétés PowerPoint spécifiques (autofit, placeholders).

        Args:
            shape: Objet shape PowerPoint
            shape_config: Configuration du shape

        Returns:
            bool: True si succès, False sinon
        """
        try:
            if not hasattr(shape, 'text_frame') or not shape.text_frame:
                return True

            text_frame = shape.text_frame

            # Autofit
            autofit_type = shape_config.get("autofit_type")
            if autofit_type:
                self._validate_property_value("autofit_type", autofit_type)
                if autofit_type == "none":
                    text_frame.auto_size = MSO_AUTO_SIZE.NONE
                elif autofit_type == "normal":
                    text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

            # Font scale (pour autofit)
            font_scale = shape_config.get("font_scale")
            if font_scale is not None:
                # python-pptx ne supporte pas directement font_scale
                # On applique une approximation via la taille de police
                if hasattr(text_frame, 'paragraphs') and text_frame.paragraphs:
                    for paragraph in text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.font.size:
                                new_size = int(run.font.size.pt * (font_scale / 100))
                                run.font.size = Pt(new_size)

            # Line spacing reduction (pour autofit)
            line_spacing_reduction = shape_config.get("line_spacing_reduction")
            if line_spacing_reduction is not None:
                if hasattr(text_frame, 'paragraphs') and text_frame.paragraphs:
                    for paragraph in text_frame.paragraphs:
                        # Ajuster l'espacement des lignes
                        paragraph.line_spacing = 1.0 - (line_spacing_reduction / 100)

            # Propriétés placeholder
            placeholder_type = shape_config.get("placeholder_type")
            if placeholder_type:
                self._validate_property_value("placeholder_type", placeholder_type)
                # Note: La modification du type de placeholder nécessite une approche plus complexe
                # car elle implique la structure XML sous-jacente

            return True

        except Exception as e:
            print(f"[ERROR] Erreur application propriétés PowerPoint: {e}")
            return False

    def create_presentation_from_first_slide(self, config: Dict[str, Any]) -> bool:
        """
        Crée une nouvelle présentation en utilisant la première slide du JSON.
        Remplace la logique des slides obligatoires.

        Args:
            config: Configuration de la présentation

        Returns:
            bool: True si succès, False sinon
        """
        try:
            output_path = config["output_path"]
            first_slide_config = config["slides"][0]
            slide_number = first_slide_config["slide_number"]

            print(f"[CREATE] Création présentation avec slide {slide_number} comme première slide")

            # Créer le dossier parent
            os.makedirs(os.path.dirname(output_path), exist_ok=True)

            # Copier le template complet
            shutil.copy2(str(self.template_path), output_path)
            print(f"[CREATE] Template copié vers {output_path}")

            # Charger la présentation et nettoyer
            prs = Presentation(output_path)
            slide_index = slide_number - 1  # Convertir en index (slide 11 = index 10)

            if slide_index >= len(prs.slides):
                raise ValueError(f"Slide {slide_number} n'existe pas dans le template")

            print(f"[CREATE] Suppression des autres slides, conservation de la slide {slide_number}")

            # Supprimer toutes les slides sauf celle désirée
            slides_to_remove = []
            for i in range(len(prs.slides)):
                if i != slide_index:
                    slides_to_remove.append(i)

            # Supprimer en ordre inverse
            for i in reversed(slides_to_remove):
                try:
                    rId = prs.slides._sldIdLst[i].rId
                    prs.part.drop_rel(rId)
                    del prs.slides._sldIdLst[i]
                except Exception as e:
                    print(f"[WARNING] Erreur suppression slide {i}: {e}")

            # Personnaliser la première slide
            remaining_slide = prs.slides[0]  # Il ne reste qu'une slide
            self._customize_slide_universal(remaining_slide, first_slide_config)

            # Sauvegarder
            prs.save(output_path)
            print(f"[SUCCESS] Présentation créée avec slide {slide_number} personnalisée")

            return True

        except Exception as e:
            print(f"[ERROR] Erreur création présentation: {e}")
            return False

    def add_slides_to_presentation(self, config: Dict[str, Any]) -> bool:
        """
        Ajoute les slides supplémentaires à une présentation existante.

        Args:
            config: Configuration de la présentation

        Returns:
            bool: True si succès, False sinon
        """
        try:
            output_path = config["output_path"]
            slides_to_add = config["slides"][1:]  # Toutes sauf la première

            if not slides_to_add:
                print(f"[INFO] Aucune slide supplémentaire à ajouter")
                return True

            print(f"[ADD] Ajout de {len(slides_to_add)} slides supplémentaires...")

            # Charger la présentation existante
            target_prs = Presentation(output_path)
            template_prs = Presentation(self.template_path)

            for i, slide_config in enumerate(slides_to_add):
                slide_number = slide_config["slide_number"]
                slide_index = slide_number - 1

                print(f"[ADD] Ajout slide {slide_number} (position {i+2})...")

                # Trouver le layout approprié
                layout_index = self._find_layout_index(target_prs, template_prs, slide_index)
                if layout_index is None:
                    print(f"[ERROR] Layout pour slide {slide_number} non trouvé")
                    continue

                # Ajouter la slide avec le bon layout
                layout = target_prs.slide_layouts[layout_index]
                new_slide = target_prs.slides.add_slide(layout)

                # Personnaliser la slide
                self._customize_slide_universal(new_slide, slide_config)

            # Sauvegarder
            target_prs.save(output_path)
            print(f"[SUCCESS] {len(slides_to_add)} slides ajoutées")

            return True

        except Exception as e:
            print(f"[ERROR] Erreur ajout slides: {e}")
            return False

    def _find_layout_index(self, target_prs, template_prs, slide_index: int):
        """
        Trouve l'index du layout dans la présentation cible.

        Args:
            target_prs: Présentation cible
            template_prs: Template source
            slide_index: Index de la slide dans le template

        Returns:
            int: Index du layout ou None si non trouvé
        """
        try:
            template_layout_name = template_prs.slides[slide_index].slide_layout.name

            for i, layout in enumerate(target_prs.slide_layouts):
                if layout.name == template_layout_name:
                    return i

            print(f"[WARNING] Layout '{template_layout_name}' non trouvé")
            return None

        except Exception as e:
            print(f"[ERROR] Erreur recherche layout: {e}")
            return None

    def build_presentation(self, json_path: str) -> str:
        """
        Construit une présentation complète à partir du JSON selon la nouvelle architecture.

        Args:
            json_path: Chemin vers le fichier JSON de configuration

        Returns:
            str: Chemin vers la présentation créée

        Raises:
            Exception: Si la construction échoue
        """
        try:
            print(f"=== PRESENTATION BUILDER v3 - Démarrage ===")
            print(f"Configuration: {json_path}")

            # 1. Charger la configuration
            config = self.load_presentation_config(json_path)

            # 2. Créer la présentation avec la première slide
            if not self.create_presentation_from_first_slide(config):
                raise Exception("Échec création présentation avec première slide")

            # 3. Ajouter les slides supplémentaires si nécessaire
            if not self.add_slides_to_presentation(config):
                print(f"[WARNING] Certaines slides supplémentaires ont échoué")

            # 4. Vérifier le succès
            output_path = config["output_path"]
            success = os.path.exists(output_path)

            if success:
                print(f"=== SUCCESS: Présentation créée ===")
                print(f"Fichier: {output_path}")
                return output_path
            else:
                raise Exception("Construction échouée")

        except Exception as e:
            print(f"=== ERROR: Construction échouée ===")
            print(f"Erreur: {e}")
            raise


def main():
    """Interface en ligne de commande"""
    parser = argparse.ArgumentParser(
        description='Construction de présentations Premier Tech à partir de JSON - Version 3'
    )

    parser.add_argument('json_file', help='Fichier JSON de configuration de la présentation')
    parser.add_argument('--validate', action='store_true', help='Valider seulement le JSON')

    args = parser.parse_args()

    try:
        builder = PresentationBuilder()

        if args.validate:
            config = builder.load_presentation_config(args.json_file)
            print(f"JSON valide: {args.json_file}")
            sys.exit(0)

        output_path = builder.build_presentation(args.json_file)
        print(f"\nSUCCES: {output_path}")

    except Exception as e:
        print(f"\nERREUR: {e}")
        sys.exit(1)


if __name__ == "__main__":
    main()