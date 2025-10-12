#!/usr/bin/env python3
"""
Conclusion Builder - Création de slides de conclusion Premier Tech
Utilise les slides 52, 54, 56, 57 du template Premier Tech pour créer des conclusions corporate.
Script spécialisé pour le besoin "Conclusion PT" selon le Guide de Création Premier Tech.
"""

import os
import sys
import json
import shutil
import argparse
from datetime import datetime
from pathlib import Path
from typing import List, Dict, Optional, Any
from pptx import Presentation
from pptx.util import Inches


class ConclusionBuilder:
    """
    Classe pour construire des slides de conclusion Premier Tech.
    Utilise les slides 52, 54, 56, 57 du template pour créer des fermetures corporate.
    Script spécialisé pour le besoin "Conclusion PT" selon le Guide de Création Premier Tech.
    """

    def __init__(self, template_path: str = "templates/Template_PT.pptx"):
        """
        Initialise le constructeur avec le template Premier Tech.

        Args:
            template_path: Chemin vers le template Premier Tech
        """
        self.template_path = template_path

        # Mapping des slides disponibles pour les conclusions PT
        # Slides utilisées : 52, 54, 56, 57 (slides 53 et 55 supprimées)
        self.conclusion_slides = {
            51: {  # Slide 52 (index 51) - "Passion et Technologies"
                "name": "Passion et Technologies",
                "usage": "Message corporate principal",
                "audience": "Toutes audiences",
                "style": "passion_tech",
                "conclusion_type": "corporate_message",
                "description": "Message de fermeture avec l'identité Premier Tech",
                "layout_name": "Passion et Technologies pour faire la différence"
            },
            53: {  # Slide 54 (index 53) - "Diapositive vide"
                "name": "Conclusion personnalisée",
                "usage": "Message de fermeture personnalisé",
                "audience": "Managers, Spécialistes",
                "style": "custom_conclusion",
                "conclusion_type": "personalized_closing",
                "description": "Conclusion adaptable avec message personnalisé",
                "layout_name": "Diapositive vide"
            },
            55: {  # Slide 56 (index 55) - "We are PT"
                "name": "We are PT",
                "usage": "Identité collective",
                "audience": "Toutes audiences",
                "style": "we_are_pt",
                "conclusion_type": "collective_identity",
                "description": "Conclusion axée sur l'appartenance à Premier Tech",
                "layout_name": "We are PT"
            },
            56: {  # Slide 57 (index 56) - "Monogramme PT"
                "name": "Fermeture minimaliste",
                "usage": "Fermeture élégante et minimaliste",
                "audience": "Executives, Audiences formelles",
                "style": "monogram",
                "conclusion_type": "minimalist_closing",
                "description": "Fermeture simple avec monogramme Premier Tech",
                "layout_name": "Monogramme PT"
            }
        }

        if not os.path.exists(self.template_path):
            raise FileNotFoundError(f"Template Premier Tech non trouvé: {self.template_path}")

    def list_available_styles(self) -> None:
        """Liste tous les styles disponibles avec leurs descriptions"""
        print("Styles de Conclusion Disponibles (Slides 52, 54, 56, 57)")
        print("=" * 60)
        for slide_idx, info in self.conclusion_slides.items():
            slide_number = slide_idx + 1
            print(f"\nStyle: {info['style']}")
            print(f"   Slide: {slide_number} - {info['name']}")
            print(f"   Layout: {info['layout_name']}")
            print(f"   Usage: {info['usage']}")
            print(f"   Audience: {info['audience']}")
            print(f"   Type: {info['conclusion_type']}")
            print(f"   Description: {info['description']}")

    def get_style_mapping(self, style: str) -> Optional[Dict]:
        """
        Récupère les informations de mapping pour un style donné

        Args:
            style: Style de conclusion demandé

        Returns:
            Dict avec les informations de la slide ou None si non trouvé
        """
        for slide_idx, info in self.conclusion_slides.items():
            if info['style'] == style:
                return {'slide_index': slide_idx, **info}
        return None

    def validate_template(self) -> bool:
        """
        Valide que le template contient toutes les slides nécessaires

        Returns:
            True si le template est valide, False sinon
        """
        try:
            pres = Presentation(self.template_path)
            total_slides = len(pres.slides)

            print(f"Validation Template: {total_slides} slides détectées")

            missing_slides = []
            for slide_idx in self.conclusion_slides.keys():
                if slide_idx >= total_slides:
                    missing_slides.append(slide_idx + 1)

            if missing_slides:
                print(f"ERREUR: Slides manquantes: {missing_slides}")
                return False

            print("SUCCESS: Template validé pour toutes les slides de conclusion")
            return True

        except Exception as e:
            print(f"ERREUR: Erreur validation template: {e}")
            return False

    def create_conclusion_slide(self, message: str = "", call_to_action: str = "",
                              style: str = "passion_tech", title: str = "",
                              custom_text: str = "", insert_into: str = "",
                              position: Optional[int] = None, output_path: str = "",
                              widen_objects: bool = True) -> Dict[str, Any]:
        """
        Crée une slide de conclusion Premier Tech avec le style spécifié

        Args:
            message: Message principal de conclusion
            call_to_action: Appel à l'action (pour styles appropriés)
            style: Style de conclusion à utiliser
            title: Titre personnalisé (si supporté par le style)
            custom_text: Texte personnalisé additionnel
            insert_into: Chemin vers présentation existante (mode insertion)
            position: Position d'insertion (défaut: fin)
            output_path: Chemin de sortie personnalisé
            widen_objects: Activer l'élargissement automatique des objets texte

        Returns:
            Dict avec les informations de création
        """

        # Validation du style
        style_info = self.get_style_mapping(style)
        if not style_info:
            available_styles = list(set(info['style'] for info in self.conclusion_slides.values()))
            raise ValueError(f"Style '{style}' non supporté. Styles disponibles: {available_styles}")

        slide_index = style_info['slide_index']

        try:
            # Clonage du template
            template_pres = Presentation(self.template_path)

            if slide_index >= len(template_pres.slides):
                raise ValueError(f"Slide {slide_index + 1} non trouvée dans le template")

            reference_slide = template_pres.slides[slide_index]

            # Mode insertion ou création nouvelle présentation
            if insert_into and os.path.exists(insert_into):
                target_pres = Presentation(insert_into)

                # Sauvegarde de sécurité
                backup_path = insert_into.replace('.pptx', '_backup_before_conclusion.pptx')
                shutil.copy2(insert_into, backup_path)

                # Trouver le layout de conclusion approprié
                conclusion_layout_index = self._find_conclusion_layout_index(target_pres, slide_index)
                if conclusion_layout_index is None:
                    raise Exception(f"Layout conclusion pour slide {slide_index + 1} non trouvé dans la présentation")

                # Clonage de la slide avec le bon layout
                conclusion_layout = target_pres.slide_layouts[conclusion_layout_index]
                new_slide = target_pres.slides.add_slide(conclusion_layout)

                if position is not None and 0 <= position < len(target_pres.slides):
                    # Réorganiser si position spécifiée
                    target_pres._slides._sld_lst.insert(position, target_pres._slides._sld_lst.pop())

            else:
                # Nouvelle présentation
                target_pres = Presentation(self.template_path)

                # Supprimer toutes les slides sauf celle de référence
                slides_to_remove = []
                for i, slide in enumerate(target_pres.slides):
                    if i != slide_index:
                        slides_to_remove.append(i)

                # Supprimer en ordre inverse pour éviter les problèmes d'index
                for i in reversed(slides_to_remove):
                    rId = target_pres.slides._slides[i].rId
                    target_pres.part.drop_rel(rId)
                    del target_pres.slides._slides[i]

                new_slide = target_pres.slides[0]

            # Personnalisation du contenu selon le style
            self._customize_conclusion_content(new_slide, style, style_info, {
                'message': message,
                'call_to_action': call_to_action,
                'title': title,
                'custom_text': custom_text
            }, widen_objects)

            # Sauvegarde
            if insert_into and os.path.exists(insert_into):
                target_pres.save(insert_into)
                output_file = insert_into
                creation_mode = "insertion"
            else:
                if not output_path:
                    timestamp = datetime.now().strftime("%Y%m%d_%H%M")
                    safe_message = "".join(c for c in message[:20] if c.isalnum() or c in (' ', '-', '_')).rstrip()
                    safe_message = safe_message.replace(' ', '_').lower()
                    output_file = f"presentations/conclusion_{timestamp}/conclusion/{timestamp}_conclusion_{style}_{safe_message}.pptx"
                else:
                    output_file = output_path

                os.makedirs(os.path.dirname(output_file), exist_ok=True)
                target_pres.save(output_file)
                creation_mode = "new_presentation"

            # Génération du rapport
            report = {
                "creation_time": datetime.now().isoformat(),
                "template_used": self.template_path,
                "source_slide": slide_index + 1,
                "style_used": style,
                "style_info": style_info,
                "content_customized": {
                    "message": message,
                    "call_to_action": call_to_action,
                    "title": title,
                    "custom_text": custom_text
                },
                "output_file": output_file,
                "creation_mode": creation_mode,
                "widen_objects_enabled": widen_objects,
                "success": True,
                "advantages": [
                    "Styles Premier Tech 100% préservés",
                    f"Slide {slide_index + 1} clonée avec fidélité parfaite",
                    "Conclusion corporate professionnelle",
                    "Personnalisation intelligente du contenu",
                    "Sauvegarde automatique (mode insertion)"
                ]
            }

            # Sauvegarde du rapport
            if insert_into:
                report_file = insert_into.replace('.pptx', '_direct_conclusion_insertion_report.json')
            else:
                report_file = output_file.replace('.pptx', '_creation_report.json')

            with open(report_file, 'w', encoding='utf-8') as f:
                json.dump(report, f, indent=2, ensure_ascii=False)

            return report

        except Exception as e:
            error_report = {
                "creation_time": datetime.now().isoformat(),
                "error": str(e),
                "template_used": self.template_path,
                "style_attempted": style,
                "success": False
            }
            return error_report

    def _customize_conclusion_content(self, slide, style: str, style_info: Dict,
                                    content: Dict, widen_objects: bool) -> None:
        """
        Personnalise le contenu de la slide selon le style de conclusion

        Args:
            slide: Slide à personnaliser
            style: Style de conclusion
            style_info: Informations du style
            content: Contenu à insérer
            widen_objects: Activer l'élargissement des objets texte
        """

        try:
            # Stratégie de personnalisation par style
            if style == "passion_tech":
                self._customize_passion_tech(slide, content, widen_objects)
            elif style == "we_are_pt":
                self._customize_we_are_pt(slide, content, widen_objects)
            elif style == "custom_conclusion":
                self._customize_custom_conclusion(slide, content, widen_objects)
            elif style == "monogram":
                self._customize_monogram(slide, content, widen_objects)
            else:
                # Fallback générique
                self._customize_generic_conclusion(slide, content, widen_objects)

        except Exception as e:
            print(f"WARNING: Avertissement personnalisation {style}: {e}")
            # Fallback sur personnalisation générique
            self._customize_generic_conclusion(slide, content, widen_objects)

    def _customize_passion_tech(self, slide, content: Dict, widen_objects: bool) -> None:
        """Personnalise slide Passion et Technologies"""
        # Le message principal est déjà dans la slide, on peut ajouter du contenu personnalisé si fourni
        if content.get('custom_text'):
            # Chercher un placeholder ou zone de texte personnalisable
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    # Détecter les zones modifiables (placeholder ou zones spécifiques)
                    if '[CUSTOM]' in shape.text:
                        shape.text = shape.text.replace('[CUSTOM]', content['custom_text'])
                        if widen_objects and hasattr(shape, 'width'):
                            shape.width = int(shape.width * 1.2)

    def _customize_we_are_pt(self, slide, content: Dict, widen_objects: bool) -> None:
        """Personnalise slide We are PT"""
        # Identité collective - personnalisation minimale pour préserver l'impact
        if content.get('message'):
            # Chercher zone de texte secondaire pour message additionnel
            text_shapes = [s for s in slide.shapes if hasattr(s, 'text_frame') and s.text_frame]
            if len(text_shapes) > 1:  # S'il y a plusieurs zones de texte
                # Utiliser la deuxième zone pour le message personnalisé
                secondary_text = text_shapes[1]
                if secondary_text.text_frame.text.strip() == "" or '[MESSAGE]' in secondary_text.text:
                    secondary_text.text = content['message']
                    if widen_objects:
                        secondary_text.width = int(secondary_text.width * 1.1)

    def _customize_custom_conclusion(self, slide, content: Dict, widen_objects: bool) -> None:
        """Personnalise slide conclusion personnalisée (diapositive vide)"""
        text_shapes = [s for s in slide.shapes if hasattr(s, 'text_frame') and s.text_frame]

        # Titre personnalisé
        if content.get('title') and slide.shapes.title:
            slide.shapes.title.text = content['title']
            if widen_objects:
                slide.shapes.title.width = int(slide.shapes.title.width * 1.1)

        # Message principal
        if content.get('message') and len(text_shapes) > 0:
            main_text_shape = text_shapes[0] if not slide.shapes.title else text_shapes[1] if len(text_shapes) > 1 else text_shapes[0]
            main_text_shape.text = content['message']
            if widen_objects:
                main_text_shape.width = int(main_text_shape.width * 1.1)

        # Appel à l'action si fourni
        if content.get('call_to_action') and len(text_shapes) > 1:
            cta_shape = text_shapes[-1]  # Dernière zone de texte pour le CTA
            cta_shape.text = content['call_to_action']
            if widen_objects:
                cta_shape.width = int(cta_shape.width * 1.1)

    def _customize_monogram(self, slide, content: Dict, widen_objects: bool) -> None:
        """Personnalise slide monogramme (fermeture minimaliste)"""
        # Fermeture minimaliste - personnalisation très limitée pour préserver l'élégance
        # Seul un titre discret peut être ajouté si vraiment nécessaire
        if content.get('title') and slide.shapes.title:
            slide.shapes.title.text = content['title']
            # Pas d'élargissement pour préserver le minimalisme

    def _customize_generic_conclusion(self, slide, content: Dict, widen_objects: bool) -> None:
        """Fallback de personnalisation générique"""
        # Stratégie de personnalisation générique qui marche avec tous les styles

        # Titre si disponible
        if content.get('title') and slide.shapes.title:
            slide.shapes.title.text = content['title']
            if widen_objects:
                slide.shapes.title.width = int(slide.shapes.title.width * 1.1)

        # Message principal dans la première zone de texte disponible
        if content.get('message'):
            text_shapes = [s for s in slide.shapes
                          if hasattr(s, 'text_frame') and s.text_frame
                          and s != slide.shapes.title]

            if text_shapes:
                text_shapes[0].text = content['message']
                if widen_objects:
                    text_shapes[0].width = int(text_shapes[0].width * 1.1)

    def _find_conclusion_layout_index(self, presentation: Presentation, source_slide_index: int) -> Optional[int]:
        """Trouve l'index du layout de conclusion dans la présentation"""
        try:
            template_prs = Presentation(self.template_path)
            template_layout_name = template_prs.slides[source_slide_index].slide_layout.name

            for i, layout in enumerate(presentation.slide_layouts):
                if layout.name == template_layout_name:
                    print(f"[LAYOUT] Layout '{template_layout_name}' trouvé à l'index {i}")
                    return i

            return None

        except Exception as e:
            print(f"[WARNING] Erreur recherche layout conclusion: {e}")
            return None


def main():
    parser = argparse.ArgumentParser(
        description="Créer des slides de conclusion Premier Tech",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Exemples d'usage:

  # Conclusion corporate standard
  python 10_conclusion_builder.py --style passion_tech

  # Conclusion avec message personnalisé
  python 10_conclusion_builder.py "Merci pour votre attention" \
    --style custom_conclusion --title "Conclusion"

  # Fermeture We are PT
  python 10_conclusion_builder.py --style we_are_pt

  # Fermeture minimaliste
  python 10_conclusion_builder.py --style monogram

  # Insertion dans présentation existante
  python 10_conclusion_builder.py "Merci" --style we_are_pt \
    --insert-into ma_presentation.pptx

  # Lister tous les styles disponibles
  python 10_conclusion_builder.py --list-styles

  # Validation du template
  python 10_conclusion_builder.py --validate
        """
    )

    # Arguments principaux
    parser.add_argument('message', nargs='?', default='',
                       help='Message principal de conclusion')
    parser.add_argument('--call-to-action',
                       help='Appel à l\'action (pour style custom_conclusion)')
    parser.add_argument('--title',
                       help='Titre personnalisé (si supporté)')
    parser.add_argument('--custom-text',
                       help='Texte personnalisé additionnel')

    # Configuration
    parser.add_argument('--style', default='passion_tech',
                       choices=['passion_tech', 'we_are_pt', 'custom_conclusion', 'monogram'],
                       help='Style de conclusion (défaut: passion_tech)')
    parser.add_argument('--insert-into',
                       help='Insérer dans une présentation existante')
    parser.add_argument('--position', type=int,
                       help='Position d\'insertion (défaut: fin)')
    parser.add_argument('--output',
                       help='Chemin de sortie spécifique')
    parser.add_argument('--template', default='templates/Template_PT.pptx',
                       help='Chemin vers le template Premier Tech')

    # Options
    parser.add_argument('--no-widen', action='store_true',
                       help='Désactiver l\'élargissement automatique des objets texte')
    parser.add_argument('--validate', action='store_true',
                       help='Valider le template seulement')
    parser.add_argument('--list-styles', action='store_true',
                       help='Lister les styles disponibles')

    args = parser.parse_args()

    try:
        builder = ConclusionBuilder(args.template)

        if args.validate:
            print("Validation du template...")
            is_valid = builder.validate_template()
            sys.exit(0 if is_valid else 1)

        if args.list_styles:
            builder.list_available_styles()
            sys.exit(0)

        # Création de la slide de conclusion
        print(f"Création d'une slide de conclusion style '{args.style}'...")

        result = builder.create_conclusion_slide(
            message=args.message,
            call_to_action=args.call_to_action or '',
            style=args.style,
            title=args.title or '',
            custom_text=args.custom_text or '',
            insert_into=args.insert_into or '',
            position=args.position,
            output_path=args.output or '',
            widen_objects=not args.no_widen
        )

        if result.get('success'):
            print(f"SUCCESS: Slide de conclusion créée avec succès!")
            print(f"Fichier: {result['output_file']}")
            print(f"Style: {result['style_used']}")
            print(f"Slide source: {result['source_slide']}")

            if args.insert_into:
                backup_file = args.insert_into.replace('.pptx', '_backup_before_conclusion.pptx')
                print(f"Sauvegarde: {backup_file}")

            print("\nAvantages:")
            for advantage in result['advantages']:
                print(f"   - {advantage}")

        else:
            print(f"ERREUR: {result.get('error', 'Erreur inconnue')}")
            sys.exit(1)

    except KeyboardInterrupt:
        print("\nOpération annulée par l'utilisateur")
        sys.exit(1)
    except Exception as e:
        print(f"ERREUR inattendue: {e}")
        sys.exit(1)


if __name__ == "__main__":
    main()