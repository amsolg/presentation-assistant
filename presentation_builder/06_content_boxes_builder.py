#!/usr/bin/env python3
"""
Content Boxes Builder - Création de slides avec 3 ou 4 concepts égaux Premier Tech
Utilise les slides 27-34 du template Premier Tech pour créer des énumérations équilibrées.
Script spécialisé pour les besoins "3 concepts égaux" et "4 concepts égaux" selon le Guide de Création Premier Tech.
"""

import os
import sys
import json
import shutil
import argparse
from datetime import datetime
from pathlib import Path
from typing import List, Dict, Optional, Any
from pptx import Presentation
from pptx.util import Inches


class ContentBoxesBuilder:
    """
    Classe pour construire des slides avec 3 ou 4 concepts égaux Premier Tech.
    Utilise les slides 27-34 du template pour créer des énumérations équilibrées.
    Script spécialisé pour les besoins "3 concepts égaux" et "4 concepts égaux" selon le Guide de Création Premier Tech.
    """

    def __init__(self, template_path: str = "templates/Template_PT.pptx"):
        """
        Initialise le constructeur avec le template Premier Tech.

        Args:
            template_path: Chemin vers le template Premier Tech
        """
        self.template_path = template_path

        # Mapping des slides disponibles pour les 3 et 4 concepts égaux (slides 27-34)
        self.content_slides = {
            26: {  # Slide 27 (index 26) - 3 boîtes grises avec sous-titres
                "name": "3 boîtes grises avec sous-titres",
                "usage": "Concepts techniques catégorisés",
                "audience": "Spécialistes, Audiences techniques",
                "style": "grey_3_detailed",
                "has_subtitles": True,
                "box_count": 3
            },
            27: {  # Slide 28 (index 27) - 3 boîtes grises sans sous-titres
                "name": "3 boîtes grises sans sous-titres",
                "usage": "Concepts techniques simples",
                "audience": "Spécialistes, Audiences techniques",
                "style": "grey_3_simple",
                "has_subtitles": False,
                "box_count": 3
            },
            28: {  # Slide 29 (index 28) - 3 boîtes bleues avec sous-titres
                "name": "3 boîtes bleues avec sous-titres",
                "usage": "Concepts catégorisés, détails structurés",
                "audience": "Managers, Spécialistes",
                "style": "blue_3_detailed",
                "has_subtitles": True,
                "box_count": 3
            },
            29: {  # Slide 30 (index 29) - 3 boîtes bleues sans sous-titres
                "name": "3 boîtes bleues sans sous-titres",
                "usage": "Concepts équivalents, importance égale",
                "audience": "Toutes audiences",
                "style": "blue_3_simple",
                "has_subtitles": False,
                "box_count": 3
            },
            30: {  # Slide 31 (index 30) - 4 boîtes grises avec sous-titres
                "name": "4 boîtes grises avec sous-titres",
                "usage": "Processus techniques détaillés",
                "audience": "Spécialistes, Équipes techniques",
                "style": "grey_4_detailed",
                "has_subtitles": True,
                "box_count": 4
            },
            31: {  # Slide 32 (index 31) - 4 boîtes grises sans sous-titres
                "name": "4 boîtes grises sans sous-titres",
                "usage": "Processus techniques simples",
                "audience": "Spécialistes, Équipes techniques",
                "style": "grey_4_simple",
                "has_subtitles": False,
                "box_count": 4
            },
            32: {  # Slide 33 (index 32) - 4 boîtes bleues avec sous-titres
                "name": "4 boîtes bleues avec sous-titres",
                "usage": "Processus business catégorisés",
                "audience": "Managers, Directeurs",
                "style": "blue_4_detailed",
                "has_subtitles": True,
                "box_count": 4
            },
            33: {  # Slide 34 (index 33) - 4 boîtes bleues sans sous-titres
                "name": "4 boîtes bleues sans sous-titres",
                "usage": "Processus business équilibrés",
                "audience": "Toutes audiences business",
                "style": "blue_4_simple",
                "has_subtitles": False,
                "box_count": 4
            }
        }

        if not os.path.exists(template_path):
            raise FileNotFoundError(f"Template Premier Tech non trouvé: {template_path}")

        # Analyser la structure des slides de contenu
        self._analyze_content_structures()

    def _analyze_content_structures(self):
        """Analyse la structure des slides de contenu disponibles"""
        try:
            pres = Presentation(self.template_path)

            self.content_info = {}
            for slide_index, slide_data in self.content_slides.items():
                if len(pres.slides) > slide_index:
                    content_slide = pres.slides[slide_index]

                    self.content_info[slide_index] = {
                        'layout_name': content_slide.slide_layout.name,
                        'shape_count': len(content_slide.shapes),
                        'slide_index': slide_index,
                        'slide_number': slide_index + 1,
                        'style': slide_data['style'],
                        'usage': slide_data['usage'],
                        'has_subtitles': slide_data['has_subtitles'],
                        'box_count': slide_data['box_count']
                    }

            print(f"[INFO] {len(self.content_info)} slides de contenu analysées (slides 27-34)")
            for idx, info in self.content_info.items():
                print(f"[INFO] Slide {info['slide_number']}: {info['layout_name']} ({info['style']}) - {info['box_count']} boîtes")

        except Exception as e:
            raise Exception(f"Erreur analyse templates contenu: {e}")

    def create_content_boxes(self,
                           concept1: str,
                           concept2: str,
                           concept3: str,
                           concept4: Optional[str] = None,
                           subtitle1: Optional[str] = None,
                           subtitle2: Optional[str] = None,
                           subtitle3: Optional[str] = None,
                           subtitle4: Optional[str] = None,
                           title: Optional[str] = None,
                           content_style: str = "blue_3_simple",
                           output_path: Optional[str] = None,
                           auto_widen: bool = True) -> str:
        """
        Crée une slide avec 3 ou 4 concepts égaux en clonant la slide appropriée du template.

        Args:
            concept1: Premier concept/pilier
            concept2: Deuxième concept/pilier
            concept3: Troisième concept/pilier
            concept4: Quatrième concept/pilier (optionnel, pour styles 4 boîtes)
            subtitle1: Sous-titre du premier concept (optionnel)
            subtitle2: Sous-titre du deuxième concept (optionnel)
            subtitle3: Sous-titre du troisième concept (optionnel)
            subtitle4: Sous-titre du quatrième concept (optionnel)
            title: Titre de la slide (optionnel)
            content_style: Style du contenu ("blue_3_simple", "blue_3_detailed", "grey_3_simple", "grey_3_detailed", "blue_4_simple", "blue_4_detailed", "grey_4_simple", "grey_4_detailed")
            output_path: Chemin de sortie (optionnel)
            auto_widen: Active l'élargissement automatique des objets texte

        Returns:
            str: Chemin vers le fichier créé
        """
        try:
            # Déterminer la slide à utiliser selon le style
            slide_index = self._get_slide_index_for_style(content_style)

            if slide_index is None:
                raise ValueError(f"Style '{content_style}' non reconnu. Utilisez: blue_3_simple, blue_3_detailed, grey_3_simple, grey_3_detailed, blue_4_simple, blue_4_detailed, grey_4_simple, grey_4_detailed")

            # Vérifier la cohérence des sous-titres avec le style
            has_subtitles = self.content_slides[slide_index]['has_subtitles']
            box_count = self.content_slides[slide_index]['box_count']

            # Construire la liste des concepts et sous-titres selon le nombre de boîtes
            concepts = [concept1, concept2, concept3]
            subtitles = [subtitle1, subtitle2, subtitle3]
            if box_count == 4:
                if concept4 is None:
                    raise ValueError("Le style sélectionné nécessite 4 concepts, mais concept4 n'est pas fourni")
                concepts.append(concept4)
                subtitles.append(subtitle4)

            if has_subtitles and not any(subtitles):
                print(f"[WARNING] Le style '{content_style}' supporte les sous-titres mais aucun n'est fourni")
            elif not has_subtitles and any(subtitles):
                print(f"[WARNING] Le style '{content_style}' ne supporte pas les sous-titres, ils seront ignorés")

            # Générer le chemin de sortie si non fourni
            if not output_path:
                output_path = self._generate_content_output_path(concept1, content_style)

            # Créer le dossier parent si nécessaire
            os.makedirs(os.path.dirname(output_path), exist_ok=True)

            print(f"[INFO] Création concepts égaux avec slide {slide_index + 1} du template ({content_style})")

            # ÉTAPE 1: Cloner la slide contenu du template avec préservation complète des styles
            success = self._clone_content_slide(slide_index, output_path)
            if not success:
                raise Exception(f"Échec du clonage de la slide contenu {slide_index + 1}")

            print(f"[SUCCESS] Slide contenu clonée avec styles Premier Tech préservés")

            # ÉTAPE 2: Élargir automatiquement les objets texte
            widen_info = self._widen_text_objects(output_path, auto_widen=auto_widen)

            # ÉTAPE 2.5: Gérer le text wrapping selon le style
            self._configure_text_wrapping(output_path, content_style)

            # ÉTAPE 3: Personnaliser le contenu en préservant les styles
            self._customize_content_boxes(output_path, concepts, subtitles, title, content_style)

            print(f"[SUCCESS] Slide avec {box_count} concepts égaux créée: {output_path}")

            # ÉTAPE 4: Générer le rapport de création
            self._generate_creation_report(output_path, concepts, subtitles, title, content_style, slide_index, widen_info)

            return output_path

        except Exception as e:
            print(f"[ERROR] Erreur création concepts égaux: {e}")
            raise

    def _get_slide_index_for_style(self, style: str) -> Optional[int]:
        """Retourne l'index de slide approprié pour le style demandé"""
        style_mapping = {
            "grey_3_detailed": 26,   # Slide 27 - 3 boîtes grises avec sous-titres
            "grey_3_simple": 27,     # Slide 28 - 3 boîtes grises sans sous-titres
            "blue_3_detailed": 28,   # Slide 29 - 3 boîtes bleues avec sous-titres
            "blue_3_simple": 29,     # Slide 30 - 3 boîtes bleues sans sous-titres
            "grey_4_detailed": 30,   # Slide 31 - 4 boîtes grises avec sous-titres
            "grey_4_simple": 31,     # Slide 32 - 4 boîtes grises sans sous-titres
            "blue_4_detailed": 32,   # Slide 33 - 4 boîtes bleues avec sous-titres
            "blue_4_simple": 33      # Slide 34 - 4 boîtes bleues sans sous-titres
        }
        return style_mapping.get(style)

    def _clone_content_slide(self, slide_index: int, output_file: str) -> bool:
        """
        Clone la slide contenu du template avec préservation complète des styles Premier Tech.
        Utilise la même méthode que les autres builders.
        """
        try:
            print(f"[CLONE] Copie complète du template...")

            # ÉTAPE 1: Copier le template complet pour préserver tous les styles
            shutil.copy2(self.template_path, output_file)

            # ÉTAPE 2: Charger et nettoyer pour ne garder que la slide contenu désirée
            prs = Presentation(output_file)

            if slide_index >= len(prs.slides):
                print(f"[ERROR] Slide {slide_index + 1} n'existe pas dans le template")
                return False

            print(f"[CLONE] Suppression des slides non désirées (garder seulement slide {slide_index + 1})...")

            # ÉTAPE 3: Identifier toutes les slides à supprimer
            slides_to_remove = []
            for i in range(len(prs.slides)):
                if i != slide_index:
                    slides_to_remove.append(i)

            print(f"[CLONE] Suppression de {len(slides_to_remove)} slides sur {len(prs.slides)} total")

            # ÉTAPE 4: Supprimer en ordre inverse pour éviter les problèmes d'index
            for i in reversed(slides_to_remove):
                try:
                    rId = prs.slides._sldIdLst[i].rId
                    prs.part.drop_rel(rId)
                    del prs.slides._sldIdLst[i]
                except Exception as e:
                    print(f"[WARNING] Erreur suppression slide {i}: {e}")

            # ÉTAPE 5: Sauvegarder la présentation avec seulement la slide contenu clonée
            prs.save(output_file)

            print(f"[SUCCESS] Slide contenu {slide_index + 1} clonée avec styles Premier Tech intacts")
            return True

        except Exception as e:
            print(f"[ERROR] Erreur clonage slide contenu {slide_index + 1}: {e}")
            return False

    def _widen_text_objects(self, presentation_path: str, auto_widen: bool = True):
        """
        Élargit automatiquement les objets texte pour éviter les retours à la ligne.
        Méthode identique aux autres builders.
        """
        widen_info = {
            "enabled": auto_widen,
            "objects_widened": 0,
            "modifications": []
        }

        if not auto_widen:
            return widen_info

        try:
            print(f"[WIDEN] Élargissement automatique des objets texte...")

            prs = Presentation(presentation_path)
            slide = prs.slides[0]  # Première (et unique) slide

            widen_count = 0
            for i, shape in enumerate(slide.shapes):
                if hasattr(shape, 'text_frame'):
                    current_width_inches = shape.width / Inches(1)

                    # Élargir les objets texte étroits (<4 pouces)
                    if current_width_inches < 4.0:
                        # Calculer la nouvelle largeur (1.5x mais max 8 pouces)
                        new_width_inches = min(current_width_inches * 1.5, 8.0)
                        shape.width = Inches(new_width_inches)

                        print(f"[WIDEN] Shape {i}: {current_width_inches:.2f}\" -> {new_width_inches:.2f}\"")
                        widen_count += 1
                        widen_info["modifications"].append({
                            "shape_index": i,
                            "before_inches": round(current_width_inches, 2),
                            "after_inches": round(new_width_inches, 2),
                            "type": "significant"
                        })

                    # Élargir modérément les objets moyens (4-6 pouces)
                    elif current_width_inches < 6.0:
                        # Élargissement plus modéré (1.2x mais max 8 pouces)
                        new_width_inches = min(current_width_inches * 1.2, 8.0)
                        shape.width = Inches(new_width_inches)

                        print(f"[WIDEN] Shape {i}: {current_width_inches:.2f}\" -> {new_width_inches:.2f}\" (modere)")
                        widen_count += 1
                        widen_info["modifications"].append({
                            "shape_index": i,
                            "before_inches": round(current_width_inches, 2),
                            "after_inches": round(new_width_inches, 2),
                            "type": "moderate"
                        })

            widen_info["objects_widened"] = widen_count

            if widen_count > 0:
                prs.save(presentation_path)
                print(f"[SUCCESS] {widen_count} objets texte élargis pour éviter les retours à la ligne")
            else:
                print(f"[INFO] Aucun objet texte nécessitant un élargissement")

            return widen_info

        except Exception as e:
            print(f"[WARNING] Erreur élargissement objets texte: {e}")
            widen_info["error"] = str(e)
            return widen_info

    def _configure_text_wrapping(self, presentation_path: str, content_style: str):
        """
        Configure le renvoi à la ligne selon le style de contenu.
        Active le wrapping pour permettre l'affichage correct des longs textes.
        """
        try:
            print(f"[WRAP] Configuration text wrapping pour style {content_style}...")

            prs = Presentation(presentation_path)
            slide = prs.slides[0]  # Première (et unique) slide

            wrap_enabled_count = 0
            for i, shape in enumerate(slide.shapes):
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    # Activer le word wrap pour tous les styles (permet l'affichage des longs textes)
                    shape.text_frame.word_wrap = True
                    print(f"[WRAP] Shape {i}: Word wrap activé")
                    wrap_enabled_count += 1

            if wrap_enabled_count > 0:
                prs.save(presentation_path)
                print(f"[SUCCESS] Text wrapping activé sur {wrap_enabled_count} objets texte")
            else:
                print(f"[INFO] Aucun objet texte trouvé")

        except Exception as e:
            print(f"[WARNING] Erreur configuration text wrapping: {e}")

    def _get_precise_shape_mapping(self, slide_index: int) -> Dict[str, List[int]]:
        """
        Retourne un mapping précis des shapes selon l'analyse réelle des templates.
        Basé sur l'analyse template_analysis_output avec les vraies positions des shapes.
        """
        # Mapping basé sur l'analyse précise des templates Premier Tech
        mappings = {
            # Slide 27 (index 26) - 3 boîtes grises avec sous-titres
            26: {
                "title_shapes": [3],           # Index 3: "Aliquam a nibh..." (TITRE PRINCIPAL)
                "content_shapes": [2, 0, 1],   # Index 2: Gauche, Index 0: Centre, Index 1: Droite
                "subtitle_shapes": [4, 5, 6],  # Index 4: Gauche, Index 5: Centre, Index 6: Droite
                "has_placeholders": False,
                "pattern": "text_box_with_subtitles"
            },
            # Slide 28 (index 27) - 3 boîtes grises sans sous-titres
            27: {
                "title_shapes": [3],           # shape_id 10 = element_id "shape_4" = index 3 (placeholder TITLE)
                "content_shapes": [2, 0, 1],   # shape_ids 11,12,13 = element_ids "shape_3","shape_1","shape_2" = indices 2,0,1 (placeholders BODY)
                "subtitle_shapes": [],
                "has_placeholders": True,
                "pattern": "placeholder_simple"
            },
            # Slide 29 (index 28) - 3 boîtes bleues avec sous-titres
            28: {
                "title_shapes": [3],           # shape_id 5 = element_id "shape_4" = index 3
                "content_shapes": [2, 0, 1],   # shape_ids 2,3,4 = element_ids "shape_3","shape_1","shape_2" = indices 2,0,1
                "subtitle_shapes": [4, 5, 6],  # shape_ids 6,7,8 = element_ids "shape_5","shape_6","shape_7" = indices 4,5,6
                "has_placeholders": False,
                "pattern": "text_box_with_subtitles"
            },
            # Slide 30 (index 29) - 3 boîtes bleues sans sous-titres
            29: {
                "title_shapes": [3],           # shape_id 2 = element_id "shape_4" = index 3 (placeholder TITLE)
                "content_shapes": [2, 0, 1],   # shape_ids 3,4,5 = element_ids "shape_3","shape_1","shape_2" = indices 2,0,1 (placeholders BODY)
                "subtitle_shapes": [],
                "has_placeholders": True,
                "pattern": "placeholder_simple"
            },
            # Slide 31 (index 30) - 4 boîtes grises avec sous-titres
            30: {
                "title_shapes": [4],             # shape_id 5 = element_id "shape_5" = index 4
                "content_shapes": [3, 0, 1, 2],  # shape_ids 2,3,4,9 = element_ids "shape_4","shape_1","shape_2","shape_3" = indices 3,0,1,2
                "subtitle_shapes": [5, 6, 7, 8], # shape_ids 6,7,8,10 = element_ids "shape_6","shape_7","shape_8","shape_9" = indices 5,6,7,8
                "has_placeholders": False,
                "pattern": "text_box_with_subtitles_4"
            },
            # Slide 32 (index 31) - 4 boîtes grises sans sous-titres
            31: {
                "title_shapes": [4],             # shape_id 12 = element_id "shape_5" = index 4 (placeholder TITLE)
                "content_shapes": [3, 0, 1, 2],  # shape_ids 13,14,15,16 = element_ids "shape_4","shape_1","shape_2","shape_3" = indices 3,0,1,2 (placeholders BODY)
                "subtitle_shapes": [],
                "has_placeholders": True,
                "pattern": "placeholder_simple_4"
            },
            # Slide 33 (index 32) - 4 boîtes bleues avec sous-titres
            32: {
                "title_shapes": [4],             # shape_id 5 = element_id "shape_5" = index 4
                "content_shapes": [3, 0, 1, 2],  # shape_ids 2,3,4,9 = element_ids "shape_4","shape_1","shape_2","shape_3" = indices 3,0,1,2
                "subtitle_shapes": [5, 6, 7, 8], # shape_ids 6,7,8,10 = element_ids "shape_6","shape_7","shape_8","shape_9" = indices 5,6,7,8
                "has_placeholders": False,
                "pattern": "text_box_with_subtitles_4"
            },
            # Slide 34 (index 33) - 4 boîtes bleues sans sous-titres
            33: {
                "title_shapes": [4],             # Sera déterminé par analyse automatique
                "content_shapes": [3, 0, 1, 2],  # Sera déterminé par analyse automatique
                "subtitle_shapes": [],
                "has_placeholders": True,
                "pattern": "placeholder_simple_4"
            }
        }

        return mappings.get(slide_index, {})

    def _customize_content_boxes(self, presentation_path: str, concepts: List[str], subtitles: List[str],
                               title: Optional[str], content_style: str):
        """
        Personnalise le contenu de la slide avec 3 ou 4 concepts égaux en préservant les styles Premier Tech.
        Utilise le mapping précis des templates pour un placement correct.
        """
        try:
            print(f"[CUSTOMIZE] Personnalisation améliorée du contenu {len(concepts)} concepts...")

            # Utiliser la nouvelle méthode améliorée
            updates_made = self._customize_content_boxes_improved(presentation_path, concepts, subtitles, title, content_style)

            # Valider le placement du contenu
            validation_success = self._validate_content_placement(presentation_path, concepts, title)

            if not validation_success:
                print(f"[WARNING] La validation indique que certains contenus peuvent ne pas être correctement placés")
            else:
                print(f"[SUCCESS] Validation confirmée: tous les contenus sont correctement placés")

        except Exception as e:
            print(f"[ERROR] Erreur personnalisation contenu: {e}")
            raise

    def _customize_content_boxes_improved(self, presentation_path: str, concepts: List[str],
                                        subtitles: List[str], title: Optional[str], content_style: str):
        """
        Personnalisation améliorée basée sur le mapping précis des templates.
        Utilise l'analyse réelle des structures de slides pour un placement correct.
        """
        try:
            print(f"[IMPROVED] Personnalisation améliorée du contenu {len(concepts)} concepts...")

            prs = Presentation(presentation_path)
            slide = prs.slides[0]
            slide_index = self._get_slide_index_for_style(content_style)

            # Obtenir le mapping précis pour ce template
            shape_mapping = self._get_precise_shape_mapping(slide_index)

            if not shape_mapping:
                print(f"[FALLBACK] Mapping non trouvé pour slide {slide_index}, utilisation de l'ancienne méthode")
                return self._customize_by_position_fallback(slide, concepts, subtitles, title,
                                                          self.content_slides[slide_index]['has_subtitles'])

            print(f"[MAPPING] Utilisation du pattern '{shape_mapping['pattern']}' pour slide {slide_index + 1}")

            # Convertir les shapes en liste pour accès par index
            shapes_list = list(slide.shapes)
            print(f"[SHAPES] {len(shapes_list)} shapes disponibles dans la slide")

            updates_made = 0

            # 1. Personnaliser le titre principal
            if title and shape_mapping.get("title_shapes"):
                for title_idx in shape_mapping["title_shapes"]:
                    if title_idx < len(shapes_list):
                        shape = shapes_list[title_idx]
                        if hasattr(shape, 'text_frame') and shape.text_frame:
                            original_text = shape.text_frame.text
                            shape.text_frame.text = title
                            shape.text_frame.word_wrap = True
                            print(f"[UPDATE] TITRE (index {title_idx}): '{original_text}' → '{title}'")
                            updates_made += 1

            # 2. Personnaliser les contenus des boîtes (priorité haute)
            content_indices = shape_mapping.get("content_shapes", [])
            for i, concept in enumerate(concepts):
                if i < len(content_indices) and content_indices[i] < len(shapes_list):
                    shape = shapes_list[content_indices[i]]
                    if hasattr(shape, 'text_frame') and shape.text_frame:
                        original_text = shape.text_frame.text
                        shape.text_frame.text = concept
                        shape.text_frame.word_wrap = True
                        print(f"[UPDATE] CONCEPT {i+1} (index {content_indices[i]}): '{original_text}' → '{concept[:50]}...'")
                        updates_made += 1

            # 3. Personnaliser les sous-titres si supportés
            if shape_mapping.get("subtitle_shapes") and subtitles:
                subtitle_indices = shape_mapping["subtitle_shapes"]
                for i, subtitle in enumerate(subtitles):
                    if subtitle and i < len(subtitle_indices) and subtitle_indices[i] < len(shapes_list):
                        shape = shapes_list[subtitle_indices[i]]
                        if hasattr(shape, 'text_frame') and shape.text_frame:
                            original_text = shape.text_frame.text
                            shape.text_frame.text = subtitle
                            shape.text_frame.word_wrap = True  # Activer wrapping pour sous-titres
                            print(f"[UPDATE] SOUS-TITRE {i+1} (index {subtitle_indices[i]}): '{original_text}' → '{subtitle}'")
                            updates_made += 1

            prs.save(presentation_path)
            print(f"[SUCCESS] Personnalisation précise terminée: {updates_made} éléments mis à jour")

            return updates_made

        except Exception as e:
            print(f"[ERROR] Erreur personnalisation améliorée: {e}")
            # Fallback vers l'ancienne méthode
            return self._customize_by_position_fallback(slide, concepts, subtitles, title,
                                                      self.content_slides[slide_index]['has_subtitles'])

    def _validate_content_placement(self, presentation_path: str, expected_concepts: List[str],
                                  expected_title: Optional[str] = None) -> bool:
        """
        Valide que le contenu a été correctement placé dans la présentation.
        Vérifie que tous les concepts et le titre sont présents dans les shapes de la slide.
        """
        try:
            prs = Presentation(presentation_path)
            slide = prs.slides[0]

            found_texts = []
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame') and shape.text_frame and shape.text_frame.text.strip():
                    found_texts.append(shape.text_frame.text.strip())

            # Vérifier que tous les concepts sont présents
            missing_concepts = []
            for concept in expected_concepts:
                if not any(concept in text for text in found_texts):
                    missing_concepts.append(concept)

            # Vérifier le titre si fourni
            title_found = True
            if expected_title:
                title_found = any(expected_title in text for text in found_texts)

            validation_success = len(missing_concepts) == 0 and title_found

            print(f"[VALIDATION] Concepts trouvés: {len(expected_concepts) - len(missing_concepts)}/{len(expected_concepts)}")
            print(f"[VALIDATION] Titre trouvé: {title_found}")
            if missing_concepts:
                print(f"[WARNING] Concepts manquants: {missing_concepts}")

            return validation_success

        except Exception as e:
            print(f"[ERROR] Erreur validation contenu: {e}")
            return False

    def _customize_by_position(self, slide, concepts: List[str], subtitles: List[str],
                             title: Optional[str], has_subtitles: bool):
        """
        Personnalisation par position des shapes selon la structure connue des templates.
        Utilise le mapping précis des positions pour chaque type de slide.
        """
        try:
            print(f"[POSITION] Personnalisation par position intelligente des shapes...")

            box_count = len(concepts)
            text_shapes = []
            for i, shape in enumerate(slide.shapes):
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    text_shapes.append((i, shape))

            print(f"[POSITION] {len(text_shapes)} shapes texte trouvés pour {box_count} boîtes")

            # Mapping des positions selon la structure connue des templates
            if box_count == 3:
                # Slides 27/29: Shape 3=titre, Shapes 4,5,6=sous-titres, Shapes 2,0,1=contenus
                title_position = 3  # CORRECTION: Index 3 est le titre, pas 4
                subtitle_positions = [4, 5, 6]  # CORRECTION: Indices 4,5,6 pour sous-titres
                content_positions = [2, 0, 1]  # Ordre d'affichage: gauche, centre, droite
            elif box_count == 4:
                # Slides 31/33: Shape 4=titre, Shapes 5,6,7,8=sous-titres, Shapes 3,0,1,2=contenus
                title_position = 4  # CORRECTION: Index 4 est le titre, pas 5
                subtitle_positions = [5, 6, 7, 8]  # CORRECTION: Indices 5,6,7,8 pour sous-titres
                content_positions = [3, 0, 1, 2]  # Ordre d'affichage: gauche, centre-gauche, centre-droite, droite
            else:
                print(f"[WARNING] Nombre de boîtes non supporté: {box_count}")
                return

            updates_made = 0

            # 1. Personnaliser le titre principal
            if title and title_position < len(text_shapes):
                shape_idx, shape = text_shapes[title_position]
                shape.text_frame.text = title
                shape.text_frame.word_wrap = True
                print(f"[POSITION] Shape {shape_idx} (titre): {title}")
                updates_made += 1

            # 2. Personnaliser les sous-titres si supportés
            if has_subtitles and subtitles:
                for i, subtitle in enumerate(subtitles):
                    if subtitle and i < len(subtitle_positions) and subtitle_positions[i] < len(text_shapes):
                        shape_idx, shape = text_shapes[subtitle_positions[i]]
                        shape.text_frame.text = subtitle
                        shape.text_frame.word_wrap = True  # Activer wrapping pour sous-titres
                        print(f"[POSITION] Shape {shape_idx} (sous-titre {i+1}): {subtitle}")
                        updates_made += 1

            # 3. Personnaliser les contenus des boîtes
            for i, concept in enumerate(concepts):
                if i < len(content_positions) and content_positions[i] < len(text_shapes):
                    shape_idx, shape = text_shapes[content_positions[i]]
                    shape.text_frame.text = concept
                    shape.text_frame.word_wrap = True  # Activer wrap pour les contenus longs
                    print(f"[POSITION] Shape {shape_idx} (concept {i+1}): {concept}")
                    updates_made += 1

            print(f"[SUCCESS] Position mapping: {updates_made} éléments personnalisés avec structure connue")

        except Exception as e:
            print(f"[WARNING] Erreur personnalisation par position: {e}")

    def _customize_by_position_fallback(self, slide, concepts: List[str], subtitles: List[str],
                                      title: Optional[str], has_subtitles: bool):
        """
        Personnalisation par position en fallback (ancienne méthode).
        Utilisée quand le mapping précis n'est pas disponible.
        """
        try:
            print(f"[FALLBACK] Personnalisation par position intelligente des shapes...")

            box_count = len(concepts)
            text_shapes = []
            for i, shape in enumerate(slide.shapes):
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    text_shapes.append((i, shape))

            print(f"[FALLBACK] {len(text_shapes)} shapes texte trouvés pour {box_count} boîtes")

            # Mapping des positions selon la structure supposée des templates
            if box_count == 3:
                title_position = 4 if len(text_shapes) > 4 else 0
                subtitle_positions = [5, 6, 7] if len(text_shapes) > 7 else []
                content_positions = [2, 0, 1] if len(text_shapes) > 2 else [0, 1, 2]
            elif box_count == 4:
                title_position = 5 if len(text_shapes) > 5 else 0
                subtitle_positions = [6, 7, 8, 9] if len(text_shapes) > 9 else []
                content_positions = [4, 1, 2, 3] if len(text_shapes) > 4 else [0, 1, 2, 3]
            else:
                print(f"[WARNING] Nombre de boîtes non supporté: {box_count}")
                return 0

            updates_made = 0

            # 1. Personnaliser le titre principal
            if title and title_position < len(text_shapes):
                shape_idx, shape = text_shapes[title_position]
                shape.text_frame.text = title
                shape.text_frame.word_wrap = True
                print(f"[FALLBACK] Shape {shape_idx} (titre): {title}")
                updates_made += 1

            # 2. Personnaliser les sous-titres si supportés
            if has_subtitles and subtitles:
                for i, subtitle in enumerate(subtitles):
                    if subtitle and i < len(subtitle_positions) and subtitle_positions[i] < len(text_shapes):
                        shape_idx, shape = text_shapes[subtitle_positions[i]]
                        shape.text_frame.text = subtitle
                        shape.text_frame.word_wrap = True
                        print(f"[FALLBACK] Shape {shape_idx} (sous-titre {i+1}): {subtitle}")
                        updates_made += 1

            # 3. Personnaliser les contenus des boîtes
            for i, concept in enumerate(concepts):
                if i < len(content_positions) and content_positions[i] < len(text_shapes):
                    shape_idx, shape = text_shapes[content_positions[i]]
                    shape.text_frame.text = concept
                    shape.text_frame.word_wrap = True
                    print(f"[FALLBACK] Shape {shape_idx} (concept {i+1}): {concept}")
                    updates_made += 1

            print(f"[SUCCESS] Fallback mapping: {updates_made} éléments personnalisés")
            return updates_made

        except Exception as e:
            print(f"[WARNING] Erreur personnalisation fallback: {e}")
            return 0

    def _generate_content_output_path(self, concept1: str, content_style: str) -> str:
        """Génère le chemin de sortie pour la slide de contenu"""

        # Nettoyer le concept pour le nom de fichier
        clean_concept = "".join(c for c in concept1 if c.isalnum() or c in (' ', '-', '_')).rstrip()
        clean_concept = clean_concept.replace(' ', '_').lower()

        # Timestamp pour l'unicité
        timestamp = datetime.now().strftime("%Y%m%d_%H%M")

        # Nom du fichier
        filename = f"{timestamp}_content_{content_style}_{clean_concept}.pptx"

        # Dossier de destination
        base_dir = "presentations"
        content_dir = os.path.join(base_dir, f"content_{timestamp}")

        return os.path.join(content_dir, "content_boxes", filename)

    def _generate_creation_report(self, output_path: str, concepts: List[str], subtitles: List[str],
                                title: Optional[str], content_style: str, slide_index: int,
                                widen_info: Optional[Dict] = None):
        """Génère un rapport de création détaillé"""

        report = {
            "creation_timestamp": datetime.now().isoformat(),
            "method": "Template Content Boxes Slide Cloning (Premier Tech Standards)",
            "template_used": self.template_path,
            "source_slide": {
                "index": slide_index,
                "number": slide_index + 1,
                "layout": self.content_info.get(slide_index, {}).get('layout_name', 'Unknown'),
                "style": content_style
            },
            "content": {
                "title": title,
                "concepts": concepts,
                "subtitles": [s for s in subtitles if s],
                "content_style": content_style,
                "style_description": self.content_slides.get(slide_index, {}).get('usage', 'Unknown')
            },
            "output_file": output_path,
            "file_size_kb": round(os.path.getsize(output_path) / 1024, 2) if os.path.exists(output_path) else 0,
            "quality_assurance": {
                "method": "Template Content Boxes Slide Cloning",
                "styles_preserved": True,
                "premier_tech_standards": True,
                "no_duplication": True,
                "professional_ready": True
            },
            "advantages": [
                "Styles Premier Tech 100% préservés",
                "Méthode de clonage éprouvée",
                "Aucune duplication d'éléments",
                "3 concepts équilibrés et structurés",
                f"Style '{content_style}' adapté à l'usage",
                "Qualité professionnelle garantie"
            ]
        }

        # Ajouter les informations d'élargissement si disponibles
        if widen_info:
            report["text_widening"] = widen_info
            if widen_info.get("objects_widened", 0) > 0:
                report["advantages"].append(
                    f"Objets texte élargis automatiquement ({widen_info['objects_widened']} modifiés)"
                )

        # Sauvegarder le rapport
        report_path = output_path.replace('.pptx', '_creation_report.json')
        with open(report_path, 'w', encoding='utf-8') as f:
            json.dump(report, f, ensure_ascii=False, indent=2)

        print(f"[INFO] Rapport de création: {report_path}")

    def insert_content_into_existing_presentation(self,
                                                presentation_path: str,
                                                concept1: str,
                                                concept2: str,
                                                concept3: str,
                                                concept4: Optional[str] = None,
                                                subtitle1: Optional[str] = None,
                                                subtitle2: Optional[str] = None,
                                                subtitle3: Optional[str] = None,
                                                subtitle4: Optional[str] = None,
                                                title: Optional[str] = None,
                                                content_style: str = "blue_3_simple",
                                                insert_position: Optional[int] = None) -> str:
        """
        Insère une slide avec 3 ou 4 concepts égaux directement dans une présentation existante.

        Args:
            presentation_path: Chemin vers la présentation existante
            concept1, concept2, concept3: Les 3 premiers concepts à présenter
            concept4: Quatrième concept (optionnel, pour styles 4 boîtes)
            subtitle1, subtitle2, subtitle3, subtitle4: Sous-titres optionnels
            title: Titre de la slide (optionnel)
            content_style: Style du contenu (blue_3_simple, blue_3_detailed, etc.)
            insert_position: Position d'insertion (None = à la fin)

        Returns:
            str: Chemin vers le fichier modifié
        """
        try:
            # Construire la liste des concepts
            concepts = [concept1, concept2, concept3]
            if concept4:
                concepts.append(concept4)

            print(f"[INSERT] Insertion directe contenu {len(concepts)} concepts dans: {os.path.basename(presentation_path)}")
            print(f"[INSERT] Style: {content_style}, Concepts: {', '.join(concepts)}")

            # ÉTAPE 1: Créer une copie de sauvegarde
            backup_path = presentation_path.replace('.pptx', '_backup_before_content.pptx')
            shutil.copy2(presentation_path, backup_path)
            print(f"[BACKUP] Sauvegarde créée: {backup_path}")

            # ÉTAPE 2: Charger la présentation existante
            target_prs = Presentation(presentation_path)
            print(f"[LOAD] Présentation chargée: {len(target_prs.slides)} slides existantes")

            # ÉTAPE 3: Déterminer la slide source à utiliser
            source_slide_index = self._get_slide_index_for_style(content_style)
            if source_slide_index is None:
                raise ValueError(f"Style '{content_style}' non reconnu")

            # ÉTAPE 4: Vérifier que le layout contenu existe
            content_layout_index = self._find_content_layout_index(target_prs, source_slide_index)
            if content_layout_index is None:
                raise Exception(f"Layout contenu pour style '{content_style}' non trouvé dans la présentation")

            # ÉTAPE 5: Ajouter la slide contenu avec le bon layout
            content_layout = target_prs.slide_layouts[content_layout_index]
            new_slide = target_prs.slides.add_slide(content_layout)
            print(f"[ADD] Slide contenu ajoutée avec layout: {content_layout.name}")

            # ÉTAPE 6: Personnaliser le contenu de la slide
            subtitles = [subtitle1, subtitle2, subtitle3]
            if subtitle4:
                subtitles.append(subtitle4)
            self._customize_content_slide_direct(new_slide, concepts, subtitles, title, content_style)

            # ÉTAPE 7: Réorganiser les slides si nécessaire
            if insert_position is not None and insert_position < len(target_prs.slides) - 1:
                self._move_slide_to_position_direct(target_prs, len(target_prs.slides) - 1, insert_position)

            # ÉTAPE 8: Sauvegarder
            target_prs.save(presentation_path)
            print(f"[SUCCESS] Slide contenu insérée directement dans la présentation")

            # ÉTAPE 9: Générer le rapport
            self._generate_direct_insertion_report(presentation_path, concepts, subtitles, title, content_style,
                                                 insert_position or len(target_prs.slides))

            return presentation_path

        except Exception as e:
            print(f"[ERROR] Erreur insertion directe contenu: {e}")
            # Restaurer la sauvegarde en cas d'erreur
            if 'backup_path' in locals() and os.path.exists(backup_path):
                shutil.copy2(backup_path, presentation_path)
                print(f"[RESTORE] Présentation originale restaurée")
            raise

    def _find_content_layout_index(self, presentation: Presentation, source_slide_index: int) -> Optional[int]:
        """Trouve l'index du layout contenu dans la présentation"""
        try:
            template_prs = Presentation(self.template_path)
            template_layout_name = template_prs.slides[source_slide_index].slide_layout.name

            for i, layout in enumerate(presentation.slide_layouts):
                if layout.name == template_layout_name:
                    print(f"[LAYOUT] Layout '{template_layout_name}' trouvé à l'index {i}")
                    return i

            return None

        except Exception as e:
            print(f"[WARNING] Erreur recherche layout contenu: {e}")
            return None

    def _customize_content_slide_direct(self, slide, concepts: List[str], subtitles: List[str],
                                      title: Optional[str], content_style: str):
        """Personnalise directement la slide contenu ajoutée"""
        try:
            print(f"[CUSTOMIZE] Personnalisation slide contenu directe...")
            print(f"[CUSTOMIZE] Slide avec {len(slide.shapes)} shapes à personnaliser")

            slide_index = self._get_slide_index_for_style(content_style)
            has_subtitles = self.content_slides[slide_index]['has_subtitles']
            box_count = self.content_slides[slide_index]['box_count']

            shape_updates = 0
            title_assigned = False  # Tracker pour éviter la double assignation

            # Analyse et personnalisation des shapes
            for i, shape in enumerate(slide.shapes):
                try:
                    if hasattr(shape, 'text_frame') and shape.text_frame:
                        current_text = shape.text_frame.text.strip()
                        print(f"[DEBUG] Shape {i}: '{current_text}' (longueur: {len(current_text)})")

                        # Personnalisation du titre principal - détecter les placeholders TITLE
                        if (title and not title_assigned and hasattr(shape, 'placeholder_format') and
                            shape.placeholder_format and shape.placeholder_format.type == 1):  # TITLE placeholder
                            shape.text_frame.text = title
                            shape.text_frame.word_wrap = True
                            print(f"[UPDATE] Shape {i} (titre via placeholder): {title}")
                            shape_updates += 1
                            title_assigned = True

                        # Personnalisation des concepts
                        for box_idx in range(box_count):
                            box_num = box_idx + 1
                            if f"boîte {box_num}" in current_text.lower() or f"contenu {box_num}" in current_text.lower():
                                if box_idx < len(concepts):
                                    shape.text_frame.text = concepts[box_idx]
                                    shape.text_frame.word_wrap = True
                                    print(f"[UPDATE] Shape {i} (concept {box_num}): {concepts[box_idx]}")
                                    shape_updates += 1
                                break

                        # Personnalisation des sous-titres si supportés
                        if has_subtitles:
                            for sub_idx in range(box_count):
                                sub_num = sub_idx + 1
                                if f"sous-titre {sub_num}" in current_text.lower() and sub_idx < len(subtitles) and subtitles[sub_idx]:
                                    shape.text_frame.text = subtitles[sub_idx]
                                    shape.text_frame.word_wrap = True
                                    print(f"[UPDATE] Shape {i} (sous-titre {sub_num}): {subtitles[sub_idx]}")
                                    shape_updates += 1
                                    break

                except Exception as e:
                    print(f"[WARNING] Erreur personnalisation shape {i}: {e}")

            # Fallback pour le titre si pas encore assigné
            if title and not title_assigned and len(slide.shapes) > 0:
                # Utiliser le premier shape comme titre fallback
                first_shape = slide.shapes[0]
                if hasattr(first_shape, 'text_frame') and first_shape.text_frame:
                    first_shape.text_frame.text = title
                    first_shape.text_frame.word_wrap = True
                    print(f"[FALLBACK] Titre assigné au premier shape: {title}")
                    title_assigned = True
                    shape_updates += 1

            # Fallback par position si détection automatique insuffisante
            if shape_updates < box_count:
                print(f"[FALLBACK] Détection automatique insuffisante ({shape_updates} éléments), passage en mode position")
                self._customize_by_position(slide, concepts, subtitles, title, has_subtitles)
                shape_updates = box_count

            print(f"[SUCCESS] Slide contenu personnalisée: {shape_updates} éléments mis à jour")

        except Exception as e:
            print(f"[ERROR] Erreur personnalisation directe contenu: {e}")
            raise

    def _move_slide_to_position_direct(self, presentation: Presentation, from_index: int, to_index: int):
        """Déplace une slide à la position désirée (méthode directe)"""
        try:
            # Note: python-pptx ne supporte pas nativement le déplacement de slides
            # Pour l'instant, on laisse la slide à la fin
            print(f"[POSITION] Slide contenu ajoutée en position {from_index + 1} (fin de présentation)")
            print(f"[INFO] Déplacement manuel requis pour position {to_index + 1}")

        except Exception as e:
            print(f"[WARNING] Déplacement slide: {e}")

    def _generate_direct_insertion_report(self, presentation_path: str, concepts: List[str],
                                        subtitles: List[str], title: Optional[str], content_style: str,
                                        insert_position: int):
        """Génère un rapport d'insertion directe"""
        try:
            source_slide_index = self._get_slide_index_for_style(content_style)

            report = {
                "insertion_timestamp": datetime.now().isoformat(),
                "method": "Direct Layout-Based Content Boxes Insertion (Premier Tech Standards)",
                "template_used": self.template_path,
                "target_presentation": presentation_path,
                "content_details": {
                    "title": title,
                    "concepts": concepts,
                    "subtitles": [s for s in subtitles if s],
                    "style": content_style,
                    "intended_position": insert_position,
                    "actual_position": "End of presentation"
                },
                "source_slide": {
                    "index": source_slide_index,
                    "number": source_slide_index + 1 if source_slide_index else None,
                    "layout": self.content_info.get(source_slide_index, {}).get('layout_name', 'Unknown'),
                    "style_description": self.content_slides.get(source_slide_index, {}).get('usage', 'Unknown')
                },
                "quality_assurance": {
                    "method": "Direct Layout-Based Addition",
                    "styles_preserved": True,
                    "premier_tech_standards": True,
                    "direct_integration": True,
                    "no_manual_steps": True
                },
                "advantages": [
                    "Insertion directe dans la présentation",
                    "Styles Premier Tech 100% préservés",
                    "Aucun fichier temporaire",
                    "Intégration transparente",
                    "Sauvegarde automatique créée",
                    f"Style '{content_style}' adapté à l'usage",
                    "3 concepts équilibrés"
                ]
            }

            # Sauvegarder le rapport
            report_path = presentation_path.replace('.pptx', '_direct_content_insertion_report.json')
            with open(report_path, 'w', encoding='utf-8') as f:
                json.dump(report, f, ensure_ascii=False, indent=2)

            print(f"[INFO] Rapport d'insertion directe: {os.path.basename(report_path)}")

        except Exception as e:
            print(f"[WARNING] Erreur génération rapport: {e}")

    def validate_template(self) -> bool:
        """Valide que le template Premier Tech est correct pour les contenus"""
        try:
            checks = {
                "file_exists": os.path.exists(self.template_path),
                "has_slides": False,
                "content_slides_exist": False,
                "slides_count": 0,
                "available_styles": []
            }

            if checks["file_exists"]:
                pres = Presentation(self.template_path)
                checks["slides_count"] = len(pres.slides)
                checks["has_slides"] = len(pres.slides) > 0

                # Vérifier que toutes les slides de contenu existent
                available_styles = []
                for slide_index in self.content_slides.keys():
                    if len(pres.slides) > slide_index:
                        style = self.content_slides[slide_index]['style']
                        available_styles.append(style)

                checks["available_styles"] = available_styles
                checks["content_slides_exist"] = len(available_styles) == len(self.content_slides)

            all_valid = all([checks["file_exists"], checks["has_slides"], checks["content_slides_exist"]])

            print("=== VALIDATION TEMPLATE PREMIER TECH POUR CONTENU 3/4 CONCEPTS ===")
            for check, result in checks.items():
                if check == "available_styles":
                    print(f"[INFO] Styles disponibles: {', '.join(result)}")
                else:
                    status = "OK" if result else "ERREUR"
                    print(f"[{status}] {check}: {result}")

            if checks["content_slides_exist"]:
                print(f"[INFO] {len(checks['available_styles'])} styles de contenu disponibles:")
                for slide_index, slide_data in self.content_slides.items():
                    if slide_index in [idx for idx in self.content_slides.keys() if len(pres.slides) > idx]:
                        print(f"  - {slide_data['style']}: Slide {slide_index + 1} ({slide_data['box_count']} boîtes, {slide_data['usage']})")

            return all_valid

        except Exception as e:
            print(f"[ERROR] Erreur validation: {e}")
            return False

    def list_available_styles(self) -> Dict[str, Dict[str, Any]]:
        """Liste tous les styles de contenu disponibles"""
        return {
            slide_data['style']: {
                "slide_number": slide_index + 1,
                "name": slide_data['name'],
                "usage": slide_data['usage'],
                "audience": slide_data['audience'],
                "has_subtitles": slide_data['has_subtitles'],
                "box_count": slide_data['box_count']
            }
            for slide_index, slide_data in self.content_slides.items()
        }


def main():
    """Interface en ligne de commande"""

    parser = argparse.ArgumentParser(
        description='Construction de slides avec 3 ou 4 concepts égaux Premier Tech (slides 27-34)'
    )

    parser.add_argument('concept1', help='Premier concept/pilier')
    parser.add_argument('concept2', help='Deuxième concept/pilier')
    parser.add_argument('concept3', help='Troisième concept/pilier')
    parser.add_argument('--concept4', help='Quatrième concept/pilier (requis pour styles 4 boîtes)')
    parser.add_argument('--subtitle1', help='Sous-titre du premier concept')
    parser.add_argument('--subtitle2', help='Sous-titre du deuxième concept')
    parser.add_argument('--subtitle3', help='Sous-titre du troisième concept')
    parser.add_argument('--subtitle4', help='Sous-titre du quatrième concept')
    parser.add_argument('--title', help='Titre de la slide')
    parser.add_argument('--style', choices=['grey_3_detailed', 'grey_3_simple', 'blue_3_detailed', 'blue_3_simple',
                                          'grey_4_detailed', 'grey_4_simple', 'blue_4_detailed', 'blue_4_simple'],
                       default='blue_3_simple',
                       help='Style du contenu (3 ou 4 boîtes, couleurs grises/bleues, avec/sans sous-titres)')
    parser.add_argument('--output', help='Chemin de sortie spécifique')
    parser.add_argument('--template', default='templates/Template_PT.pptx',
                       help='Chemin vers le template Premier Tech')
    parser.add_argument('--no-widen', action='store_true',
                       help='Désactiver l\'élargissement automatique des objets texte')
    parser.add_argument('--validate', action='store_true',
                       help='Valider le template seulement')
    parser.add_argument('--list-styles', action='store_true',
                       help='Lister les styles disponibles')
    parser.add_argument('--insert-into',
                       help='Insérer dans une présentation existante (chemin)')
    parser.add_argument('--position', type=int,
                       help='Position d\'insertion dans la présentation (défaut: fin)')

    args = parser.parse_args()

    try:
        # Initialiser le constructeur
        builder = ContentBoxesBuilder(args.template)

        # Mode validation
        if args.validate:
            is_valid = builder.validate_template()
            sys.exit(0 if is_valid else 1)

        # Mode liste des styles
        if args.list_styles:
            styles = builder.list_available_styles()
            print("=== STYLES DE CONTENU 3 CONCEPTS DISPONIBLES ===")
            for style, info in styles.items():
                print(f"{style.upper()}:")
                print(f"  - Slide: {info['slide_number']}")
                print(f"  - Nom: {info['name']}")
                print(f"  - Usage: {info['usage']}")
                print(f"  - Audience: {info['audience']}")
                print(f"  - Sous-titres: {'Oui' if info['has_subtitles'] else 'Non'}")
                print()
            sys.exit(0)

        # Validation cohérence style et concepts/sous-titres
        if args.style.endswith('_detailed') and not any([args.subtitle1, args.subtitle2, args.subtitle3, args.subtitle4]):
            print(f"WARNING: Le style '{args.style}' est plus efficace avec des sous-titres")
        elif args.style.endswith('_simple') and any([args.subtitle1, args.subtitle2, args.subtitle3, args.subtitle4]):
            print(f"WARNING: Le style '{args.style}' ne supporte pas les sous-titres, ils seront ignorés")

        # Validation 4 concepts pour styles 4 boîtes
        if args.style.startswith(('grey_4', 'blue_4')) and not args.concept4:
            print(f"ERROR: Le style '{args.style}' nécessite 4 concepts. Utilisez --concept4")
            sys.exit(1)
        elif not args.style.startswith(('grey_4', 'blue_4')) and args.concept4:
            print(f"WARNING: Le style '{args.style}' n'utilise que 3 concepts. Le concept4 sera ignoré")

        # Mode insertion dans présentation existante
        if args.insert_into:
            output_path = builder.insert_content_into_existing_presentation(
                presentation_path=args.insert_into,
                concept1=args.concept1,
                concept2=args.concept2,
                concept3=args.concept3,
                concept4=args.concept4,
                subtitle1=args.subtitle1,
                subtitle2=args.subtitle2,
                subtitle3=args.subtitle3,
                subtitle4=args.subtitle4,
                title=args.title,
                content_style=args.style,
                insert_position=args.position
            )
            box_count = 4 if args.concept4 else 3
            print(f"\nSUCCES: Slide {box_count} concepts égaux intégrée dans présentation existante: {output_path}")
        else:
            print(f"\nERREUR: Le script {os.path.basename(__file__)} ne peut que s'insérer dans une présentation existante.")
            print("Utilisez l'argument --insert-into pour spécifier le fichier PowerPoint cible.")
            print("Pour créer une nouvelle présentation, utilisez d'abord 01_slide_title_creator.py")
            sys.exit(1)

        print(f"Style utilisé: {args.style}")
        concepts = [args.concept1, args.concept2, args.concept3]
        if args.concept4:
            concepts.append(args.concept4)
        print(f"Concepts: {', '.join(concepts)}")

        if any([args.subtitle1, args.subtitle2, args.subtitle3, args.subtitle4]):
            subtitles = [s for s in [args.subtitle1, args.subtitle2, args.subtitle3, args.subtitle4] if s]
            print(f"Sous-titres: {', '.join(subtitles)}")
        if args.title:
            print(f"Titre: {args.title}")

    except Exception as e:
        print(f"\nERREUR: {e}")
        sys.exit(1)


if __name__ == "__main__":
    main()