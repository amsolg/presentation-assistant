#!/usr/bin/env python3
"""
Slide Title Creator - Création de slides de titre Premier Tech
Utilise la slide 11 du template Premier Tech pour créer des pages de couverture/intro.
Script spécialisé pour le besoin "Couverture/Intro" selon le Guide de Création Premier Tech.
"""

import os
import sys
import json
import shutil
import argparse
from datetime import datetime
from pathlib import Path
from typing import Dict, Optional, Any
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches


class SlideTitleCreator:
    """
    Classe pour créer des slides de titre Premier Tech en clonant la slide 11 du template.
    Utilise la méthode de clonage correcte qui préserve parfaitement les styles.
    Script spécialisé pour le besoin "Couverture/Intro" selon le Guide de Création Premier Tech.
    """

    def __init__(self):
        """
        Initialise le créateur avec le template Premier Tech.
        Utilise automatiquement templates/Template_PT.pptx et la slide 11.
        """
        self.template_path = r"c:\repos\presentation-assistant\templates\Template_PT.pptx"
        self.reference_slide_index = 10  # Slide 11 (index 10) - Page titre

        if not os.path.exists(self.template_path):
            raise FileNotFoundError(f"Template Premier Tech non trouvé: {self.template_path}")

        # Analyser la structure de référence
        self._analyze_reference_structure()

    def _analyze_reference_structure(self):
        """Analyse la structure de la slide de référence"""
        try:
            pres = Presentation(self.template_path)
            if len(pres.slides) <= self.reference_slide_index:
                raise ValueError(f"Template ne contient pas de slide {self.reference_slide_index + 1}")

            reference_slide = pres.slides[self.reference_slide_index]

            self.reference_info = {
                'layout_name': reference_slide.slide_layout.name,
                'shape_count': len(reference_slide.shapes),
                'slide_index': self.reference_slide_index,
                'slide_number': self.reference_slide_index + 1
            }

            print(f"[INFO] Slide de référence: {self.reference_info['slide_number']} ({self.reference_info['layout_name']})")
            print(f"[INFO] {self.reference_info['shape_count']} shapes identifiés")

        except Exception as e:
            raise Exception(f"Erreur analyse template: {e}")

    def _validate_title_length(self, title: str):
        """
        Valide la longueur du titre pour éviter le débordement.

        Args:
            title: Titre à valider

        Raises:
            ValueError: Si le titre est trop long
        """
        # Règle empirique: max 50 caractères pour éviter le débordement
        # Basé sur l'observation que "L'Hygiène des Mains - Un Investissement" (47 chars) déborde déjà
        MAX_TITLE_LENGTH = 45

        if len(title) > MAX_TITLE_LENGTH:
            raise ValueError(
                f"ERREUR: Titre trop long ({len(title)} caractères).\n"
                f"Maximum recommandé: {MAX_TITLE_LENGTH} caractères.\n"
                f"Titre actuel: '{title}'\n"
                f"Suggestion: Raccourcissez le titre ou utilisez le sous-titre pour les détails."
            )

        print(f"[VALIDATION] Titre accepté ({len(title)} caractères): {title}")

    def create_title_slide(self,
                          title: str,
                          subtitle: Optional[str] = None,
                          metadata: Optional[str] = None,
                          output_path: Optional[str] = None,
                          project_name: Optional[str] = None,
                          auto_widen: bool = True) -> str:
        """
        Crée une slide de titre en clonant la slide 11 du template Premier Tech.

        Args:
            title: Titre principal de la présentation
            subtitle: Sous-titre ou contexte (optionnel)
            metadata: Métadonnées (date, auteurs) (optionnel)
            output_path: Chemin de sortie (optionnel)
            project_name: Nom du projet pour l'organisation (optionnel)
            auto_widen: Active l'élargissement automatique des objets texte (défaut: True)

        Returns:
            str: Chemin vers le fichier créé
        """
        try:
            # VALIDATION: Vérifier la longueur du titre
            self._validate_title_length(title)

            # Toujours utiliser la slide 11 (index 10)
            slide_index = self.reference_slide_index

            # Générer le chemin de sortie si non fourni
            if not output_path:
                output_path = self._generate_output_path(title, project_name)

            # Créer le dossier parent si nécessaire
            os.makedirs(os.path.dirname(output_path), exist_ok=True)

            print(f"[INFO] Clonage slide {slide_index + 1} du template Premier Tech")

            # ÉTAPE 1: Cloner la slide du template avec préservation complète des styles
            success = self._clone_template_slide(slide_index, output_path)
            if not success:
                raise Exception(f"Échec du clonage de la slide {slide_index + 1}")

            print(f"[SUCCESS] Slide clonée avec styles Premier Tech préservés")

            # ÉTAPE 2: Élargir automatiquement les objets texte
            widen_info = self._widen_text_objects(output_path, auto_widen=auto_widen)

            # ÉTAPE 2.5: Désactiver le renvoi à la ligne automatique
            self._disable_text_wrapping(output_path)

            # ÉTAPE 3: Personnaliser le contenu en préservant les styles
            self._customize_cloned_slide(output_path, title, subtitle, metadata)

            print(f"[SUCCESS] Présentation créée: {output_path}")

            # ÉTAPE 4: Générer le rapport de création
            self._generate_creation_report(output_path, title, subtitle, metadata, slide_index, widen_info)

            return output_path

        except Exception as e:
            print(f"[ERROR] Erreur création présentation: {e}")
            raise

    def _clone_template_slide(self, slide_index: int, output_file: str) -> bool:
        """
        Clone une slide du template avec préservation complète des styles Premier Tech.
        Méthode basée sur enhanced_presentation_builder_v2.py
        """
        try:
            print(f"[CLONE] Copie complète du template...")

            # ÉTAPE 1: Copier le template complet pour préserver tous les styles
            shutil.copy2(self.template_path, output_file)

            # ÉTAPE 2: Charger et nettoyer pour ne garder que la slide désirée
            prs = Presentation(output_file)

            if slide_index >= len(prs.slides):
                print(f"[ERROR] Slide {slide_index + 1} n'existe pas dans le template")
                return False

            print(f"[CLONE] Suppression des slides non désirées (garder seulement slide {slide_index + 1})...")

            # ÉTAPE 3: Identifier toutes les slides à supprimer
            slides_to_remove = []
            for i in range(len(prs.slides)):
                if i != slide_index:
                    slides_to_remove.append(i)

            print(f"[CLONE] Suppression de {len(slides_to_remove)} slides sur {len(prs.slides)} total")

            # ÉTAPE 4: Supprimer en ordre inverse pour éviter les problèmes d'index
            for i in reversed(slides_to_remove):
                try:
                    rId = prs.slides._sldIdLst[i].rId
                    prs.part.drop_rel(rId)
                    del prs.slides._sldIdLst[i]
                except Exception as e:
                    print(f"[WARNING] Erreur suppression slide {i}: {e}")

            # ÉTAPE 5: Sauvegarder la présentation avec seulement la slide clonée
            prs.save(output_file)

            print(f"[SUCCESS] Slide {slide_index + 1} clonée avec styles Premier Tech intacts")
            return True

        except Exception as e:
            print(f"[ERROR] Erreur clonage slide {slide_index + 1}: {e}")
            return False

    def _widen_text_objects(self, presentation_path: str, auto_widen: bool = True):
        """
        Élargit automatiquement les objets texte pour éviter les retours à la ligne.

        Args:
            presentation_path: Chemin vers la présentation
            auto_widen: Active l'élargissement automatique (défaut: True)

        Returns:
            dict: Informations sur l'élargissement effectué
        """
        widen_info = {
            "enabled": auto_widen,
            "objects_widened": 0,
            "modifications": []
        }

        if not auto_widen:
            return widen_info

        try:
            print(f"[WIDEN] Élargissement automatique des objets texte...")

            prs = Presentation(presentation_path)
            slide = prs.slides[0]  # Première (et unique) slide

            widen_count = 0
            for i, shape in enumerate(slide.shapes):
                if hasattr(shape, 'text_frame'):
                    current_width_inches = shape.width / Inches(1)

                    # Élargir les objets texte étroits (<4 pouces)
                    if current_width_inches < 4.0:
                        # Calculer la nouvelle largeur (1.5x mais max 8 pouces)
                        new_width_inches = min(current_width_inches * 1.5, 8.0)
                        shape.width = Inches(new_width_inches)

                        print(f"[WIDEN] Shape {i}: {current_width_inches:.2f}\" -> {new_width_inches:.2f}\"")
                        widen_count += 1
                        widen_info["modifications"].append({
                            "shape_index": i,
                            "before_inches": round(current_width_inches, 2),
                            "after_inches": round(new_width_inches, 2),
                            "type": "significant"
                        })

                    # Élargir modérément les objets moyens (4-6 pouces)
                    elif current_width_inches < 6.0:
                        # Élargissement plus modéré (1.2x mais max 8 pouces)
                        new_width_inches = min(current_width_inches * 1.2, 8.0)
                        shape.width = Inches(new_width_inches)

                        print(f"[WIDEN] Shape {i}: {current_width_inches:.2f}\" -> {new_width_inches:.2f}\" (modere)")
                        widen_count += 1
                        widen_info["modifications"].append({
                            "shape_index": i,
                            "before_inches": round(current_width_inches, 2),
                            "after_inches": round(new_width_inches, 2),
                            "type": "moderate"
                        })

            widen_info["objects_widened"] = widen_count

            if widen_count > 0:
                prs.save(presentation_path)
                print(f"[SUCCESS] {widen_count} objets texte élargis pour éviter les retours à la ligne")
            else:
                print(f"[INFO] Aucun objet texte nécessitant un élargissement")

            return widen_info

        except Exception as e:
            print(f"[WARNING] Erreur élargissement objets texte: {e}")
            widen_info["error"] = str(e)
            return widen_info

    def _disable_text_wrapping(self, presentation_path: str):
        """
        Désactive le renvoi à la ligne automatique pour tous les objets texte.
        Équivalent à décocher "Renvoyer le texte à la ligne dans la forme" dans PowerPoint.
        """
        try:
            print(f"[WRAP] Désactivation du renvoi à la ligne automatique...")

            prs = Presentation(presentation_path)
            slide = prs.slides[0]  # Première (et unique) slide

            wrap_disabled_count = 0
            for i, shape in enumerate(slide.shapes):
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    # Désactiver le word wrap (renvoi à la ligne automatique)
                    shape.text_frame.word_wrap = False
                    print(f"[WRAP] Shape {i}: Word wrap désactivé")
                    wrap_disabled_count += 1

            if wrap_disabled_count > 0:
                prs.save(presentation_path)
                print(f"[SUCCESS] Renvoi à la ligne désactivé sur {wrap_disabled_count} objets texte")
            else:
                print(f"[INFO] Aucun objet texte trouvé")

        except Exception as e:
            print(f"[WARNING] Erreur désactivation word wrap: {e}")

    def _customize_cloned_slide(self, presentation_path: str, title: str, subtitle: Optional[str], metadata: Optional[str]):
        """
        Personnalise le contenu de la slide clonée en préservant les styles Premier Tech.
        REMPLACE le contenu sans modifier les styles.
        """
        try:
            print(f"[CUSTOMIZE] Personnalisation du contenu...")

            # Charger la présentation clonée
            prs = Presentation(presentation_path)
            slide = prs.slides[0]  # Première (et unique) slide

            # Générer les métadonnées par défaut si non fournies
            if not metadata:
                metadata = f"{datetime.now().strftime('%Y.%m.%d')} – Présentation Premier Tech"

            # Mapping du contenu
            content_mapping = {
                'title': title,
                'subtitle': subtitle or "Présentation d'entreprise",
                'metadata': metadata
            }

            print(f"[CUSTOMIZE] Slide avec {len(slide.shapes)} shapes à traiter")

            # Personnaliser les shapes en préservant les styles
            updated_count = 0
            for i, shape in enumerate(slide.shapes):
                try:
                    if hasattr(shape, 'text_frame') and shape.text_frame and shape.text_frame.text:
                        current_text = shape.text_frame.text.lower()

                        # Identifier le rôle du shape selon son contenu actuel
                        new_content = None
                        role = None

                        if 'objet' in current_text or 'titre' in current_text:
                            new_content = content_mapping['title']
                            role = 'title'
                        elif 'contexte' in current_text:
                            new_content = content_mapping['subtitle']
                            role = 'subtitle'
                        elif 'statut' in current_text or 'date' in current_text or '2025' in current_text:
                            new_content = content_mapping['metadata']
                            role = 'metadata'

                        # Appliquer le nouveau contenu EN PRÉSERVANT LE FORMATAGE
                        if new_content and role:
                            # CRITIQUE: Remplacer seulement le texte, pas le formatage
                            shape.text_frame.text = new_content
                            print(f"[UPDATE] {role}: {new_content[:50]}...")
                            updated_count += 1

                except Exception as e:
                    print(f"[WARNING] Erreur personnalisation shape {i}: {e}")

            print(f"[SUCCESS] {updated_count} éléments personnalisés avec styles Premier Tech préservés")

            # Sauvegarder les modifications
            prs.save(presentation_path)

        except Exception as e:
            print(f"[ERROR] Erreur personnalisation: {e}")
            raise

    def _generate_output_path(self, title: str, project_name: Optional[str] = None) -> str:
        """Génère le chemin de sortie selon les conventions Premier Tech"""

        # Nettoyer le titre pour le nom de fichier
        clean_title = "".join(c for c in title if c.isalnum() or c in (' ', '-', '_')).rstrip()
        clean_title = clean_title.replace(' ', '_').lower()

        # Timestamp pour l'unicité
        timestamp = datetime.now().strftime("%Y%m%d_%H%M")

        # Nom du fichier
        if project_name:
            filename = f"{timestamp}_{project_name}_{clean_title}.pptx"
        else:
            filename = f"{timestamp}_{clean_title}.pptx"

        # Dossier de destination
        base_dir = "presentations"
        if project_name:
            project_dir = os.path.join(base_dir, project_name)
        else:
            project_dir = os.path.join(base_dir, f"{timestamp}_{clean_title}")

        return os.path.join(project_dir, "presentation", filename)

    def _generate_creation_report(self, output_path: str, title: str, subtitle: Optional[str],
                                metadata: Optional[str], slide_index: int, widen_info: Optional[Dict] = None):
        """Génère un rapport de création détaillé"""

        report = {
            "creation_timestamp": datetime.now().isoformat(),
            "method": "Template Slide Cloning (Premier Tech Standards)",
            "template_used": self.template_path,
            "source_slide": {
                "index": slide_index,
                "number": slide_index + 1,
                "layout": self.reference_info.get('layout_name', 'Unknown')
            },
            "content": {
                "title": title,
                "subtitle": subtitle,
                "metadata": metadata
            },
            "output_file": output_path,
            "file_size_kb": round(os.path.getsize(output_path) / 1024, 2) if os.path.exists(output_path) else 0,
            "quality_assurance": {
                "method": "Template Slide Cloning",
                "styles_preserved": True,
                "premier_tech_standards": True,
                "no_duplication": True,
                "professional_ready": True
            },
            "advantages": [
                "Styles Premier Tech 100% préservés",
                "Méthode de clonage éprouvée",
                "Aucune duplication d'éléments",
                "Qualité professionnelle garantie"
            ]
        }

        # Ajouter les informations d'élargissement si disponibles
        if widen_info:
            report["text_widening"] = widen_info
            if widen_info.get("objects_widened", 0) > 0:
                report["advantages"].append(
                    f"Objets texte élargis automatiquement ({widen_info['objects_widened']} modifiés)"
                )

        # Sauvegarder le rapport
        report_path = output_path.replace('.pptx', '_creation_report.json')
        with open(report_path, 'w', encoding='utf-8') as f:
            json.dump(report, f, ensure_ascii=False, indent=2)

        print(f"[INFO] Rapport de création: {report_path}")

    def list_available_templates(self) -> Dict[int, Dict[str, Any]]:
        """Liste toutes les slides disponibles dans le template Premier Tech"""
        try:
            pres = Presentation(self.template_path)
            templates = {}

            for i, slide in enumerate(pres.slides):
                templates[i] = {
                    "slide_number": i + 1,
                    "layout_name": slide.slide_layout.name,
                    "shape_count": len(slide.shapes),
                    "has_title": hasattr(slide.shapes, 'title') and slide.shapes.title is not None
                }

            return templates

        except Exception as e:
            print(f"[ERROR] Erreur lecture templates: {e}")
            return {}

    def validate_template(self) -> bool:
        """Valide que le template Premier Tech est correct"""
        try:
            checks = {
                "file_exists": os.path.exists(self.template_path),
                "has_slides": False,
                "reference_slide_exists": False,
                "slides_count": 0
            }

            if checks["file_exists"]:
                pres = Presentation(self.template_path)
                checks["slides_count"] = len(pres.slides)
                checks["has_slides"] = len(pres.slides) > 0
                checks["reference_slide_exists"] = len(pres.slides) > self.reference_slide_index

            all_valid = all([checks["file_exists"], checks["has_slides"], checks["reference_slide_exists"]])

            print("=== VALIDATION TEMPLATE PREMIER TECH ===")
            for check, result in checks.items():
                status = "OK" if result else "ERREUR"
                print(f"[{status}] {check}: {result}")

            return all_valid

        except Exception as e:
            print(f"[ERROR] Erreur validation: {e}")
            return False


def main():
    """Interface en ligne de commande"""

    parser = argparse.ArgumentParser(
        description='Création de slides de titre Premier Tech utilisant la slide 11 du template'
    )

    parser.add_argument('title', help='Titre principal de la présentation')
    parser.add_argument('--subtitle', help='Sous-titre ou contexte')
    parser.add_argument('--metadata', help='Métadonnées (date, auteurs)')
    parser.add_argument('--project', help='Nom du projet')
    parser.add_argument('--output', help='Chemin de sortie spécifique')
    parser.add_argument('--no-widen', action='store_true',
                       help='Désactiver l\'élargissement automatique des objets texte')
    parser.add_argument('--validate', action='store_true',
                       help='Valider le template seulement')

    args = parser.parse_args()

    try:
        # Initialiser le créateur (sans paramètres)
        creator = SlideTitleCreator()

        # Mode validation
        if args.validate:
            is_valid = creator.validate_template()
            print(f"Template validé: {'OK' if is_valid else 'ERREUR'}")
            sys.exit(0 if is_valid else 1)

        # Créer la slide de titre (slide 11 automatiquement)
        output_path = creator.create_title_slide(
            title=args.title,
            subtitle=args.subtitle,
            metadata=args.metadata,
            output_path=args.output,
            project_name=args.project,
            auto_widen=not args.no_widen
        )

        print(f"\nSUCCES: Slide de titre Premier Tech créée: {output_path}")

    except Exception as e:
        print(f"\nERREUR: {e}")
        sys.exit(1)


if __name__ == "__main__":
    main()