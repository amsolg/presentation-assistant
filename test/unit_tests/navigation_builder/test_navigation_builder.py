#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Tests unitaires pour Navigation Builder - Architecture JSON 2025
Teste le nouveau script navigation_builder.py avec différentes configurations JSON.
"""

import os
import sys
import json
import unittest
from pathlib import Path
from datetime import datetime

# Ajouter le chemin vers presentation_builder
sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..', '..', '..', 'presentation_builder'))

try:
    from navigation_builder import NavigationBuilder, create_navigation_from_json, load_navigation_template
    from presentation_builder import PresentationBuilder

except ImportError as e:
    print(f"[ERROR] Impossible d'importer les modules requis: {e}")
    sys.exit(1)


class TestNavigationBuilder(unittest.TestCase):
    """Tests pour le NavigationBuilder avec architecture JSON 2025"""

    def setUp(self):
        """Configuration initiale pour chaque test"""
        self.template_path = os.path.join('..', '..', '..', 'templates', 'Template_PT.pptx')
        self.test_subject = "unit-tests"
        self.test_audience = "navigation_builder_tester"

        # Créer le dossier de sortie
        self.output_dir = os.path.join('output')
        os.makedirs(self.output_dir, exist_ok=True)

        # Timestamp pour l'unicité des fichiers
        self.timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")

    def _create_presentation_from_schema(self, test_name: str) -> str:
        """Crée une présentation complète à partir du schéma JSON"""
        try:
            # Créer le PresentationBuilder
            builder = PresentationBuilder()

            # Changer vers le répertoire du test pour les chemins relatifs des payloads
            original_cwd = os.getcwd()
            test_dir = os.path.join(os.getcwd(), test_name)
            os.chdir(test_dir)

            try:
                # Construire la présentation complète avec le chemin vers le schéma JSON
                result_path = builder.build_presentation("presentation_schema.json")

                self.assertTrue(result_path, f"Échec construction présentation pour {test_name}")
                self.assertTrue(os.path.exists(result_path), f"Fichier de présentation non créé: {result_path}")

                print(f"[TEST] Présentation créée: {os.path.basename(result_path)}")
                return result_path

            finally:
                os.chdir(original_cwd)

        except Exception as e:
            self.fail(f"Erreur création présentation {test_name}: {e}")

    def test_basic_toc_navigation(self):
        """Test avec présentation complète basique via schéma JSON"""
        print(f"\n=== TEST: Navigation TOC Basique (Schéma Complet) ===")

        try:
            # Créer la présentation complète depuis le schéma
            presentation_path = self._create_presentation_from_schema("basic_toc")

            # Vérifier que le fichier existe et contient les bonnes slides
            self.assertTrue(os.path.exists(presentation_path), "Fichier présentation non trouvé")

            # Analyser la présentation pour vérifier qu'elle contient les slides attendues
            from pptx import Presentation
            pres = Presentation(presentation_path)

            # Doit contenir au moins 3 slides : titre + navigation + fermeture
            self.assertGreaterEqual(len(pres.slides), 3, "Présentation doit contenir au moins 3 slides")

            print(f"[SUCCESS] Présentation basique créée avec {len(pres.slides)} slides")

        except Exception as e:
            self.fail(f"Erreur test navigation basique: {e}")

    def test_detailed_toc_navigation(self):
        """Test avec présentation complète détaillée via schéma JSON"""
        print(f"\n=== TEST: Navigation TOC Détaillée (Schéma Complet) ===")

        try:
            # Créer la présentation complète depuis le schéma
            presentation_path = self._create_presentation_from_schema("detailed_toc")

            # Vérifier que le fichier existe
            self.assertTrue(os.path.exists(presentation_path), "Fichier présentation non trouvé")

            # Analyser la présentation
            from pptx import Presentation
            pres = Presentation(presentation_path)

            # Doit contenir au moins 3 slides : titre + navigation + fermeture
            self.assertGreaterEqual(len(pres.slides), 3, "Présentation doit contenir au moins 3 slides")

            print(f"[SUCCESS] Présentation détaillée créée avec {len(pres.slides)} slides")

        except Exception as e:
            self.fail(f"Erreur test navigation détaillée: {e}")

    def test_strategic_toc_navigation(self):
        """Test avec présentation complète stratégique via schéma JSON"""
        print(f"\n=== TEST: Navigation TOC Stratégique (Schéma Complet) ===")

        try:
            # Créer la présentation complète depuis le schéma
            presentation_path = self._create_presentation_from_schema("strategic_toc")

            # Vérifier que le fichier existe
            self.assertTrue(os.path.exists(presentation_path), "Fichier présentation non trouvé")

            # Analyser la présentation
            from pptx import Presentation
            pres = Presentation(presentation_path)

            # Doit contenir au moins 3 slides : titre + navigation + fermeture
            self.assertGreaterEqual(len(pres.slides), 3, "Présentation doit contenir au moins 3 slides")

            print(f"[SUCCESS] Présentation stratégique créée avec {len(pres.slides)} slides")

        except Exception as e:
            self.fail(f"Erreur test navigation stratégique: {e}")

    def test_template_loading(self):
        """Test du chargement des templates prédéfinis"""
        print(f"\n=== TEST: Chargement Templates ===")

        # Tester les différents templates
        templates_to_test = ["basic_toc", "detailed_toc", "strategic_toc", "nonexistent"]

        for template_name in templates_to_test:
            try:
                config = load_navigation_template(template_name)

                # Vérifications de structure
                self.assertIn('title', config, f"Template {template_name}: 'title' manquant")
                self.assertIn('sections', config, f"Template {template_name}: 'sections' manquant")
                self.assertIn('options', config, f"Template {template_name}: 'options' manquant")

                self.assertIsInstance(config['sections'], list, f"Template {template_name}: 'sections' doit être une liste")
                self.assertGreater(len(config['sections']), 0, f"Template {template_name}: sections vides")

                print(f"[SUCCESS] Template '{template_name}' chargé: {len(config['sections'])} sections")

            except Exception as e:
                if template_name == "nonexistent":
                    print(f"[EXPECTED] Template inexistant géré: {template_name}")
                else:
                    self.fail(f"Erreur chargement template {template_name}: {e}")

    def test_configuration_validation(self):
        """Test de validation des configurations JSON"""
        print(f"\n=== TEST: Validation Configurations ===")

        builder = NavigationBuilder(self.template_path)

        # Configurations de test (valides et invalides)
        test_configs = [
            {
                "name": "config_valide",
                "config": {
                    "title": "Test Valide",
                    "sections": ["Section 1", "Section 2"],
                    "options": {"auto_widen": True}
                },
                "should_be_valid": True
            },
            {
                "name": "titre_manquant",
                "config": {
                    "sections": ["Section 1", "Section 2"]
                },
                "should_be_valid": False
            },
            {
                "name": "sections_manquantes",
                "config": {
                    "title": "Test"
                },
                "should_be_valid": False
            },
            {
                "name": "sections_vides",
                "config": {
                    "title": "Test",
                    "sections": []
                },
                "should_be_valid": False
            },
            {
                "name": "titre_trop_long",
                "config": {
                    "title": "T" * 150,  # Titre de 150 caractères
                    "sections": ["Section 1"]
                },
                "should_be_valid": True  # Valide mais avec warning
            }
        ]

        for test_case in test_configs:
            try:
                validation = builder._validate_navigation_config(test_case["config"])

                if test_case["should_be_valid"]:
                    self.assertTrue(validation['valid'],
                                  f"Config '{test_case['name']}' devrait être valide: {validation['errors']}")
                else:
                    self.assertFalse(validation['valid'],
                                   f"Config '{test_case['name']}' devrait être invalide")

                print(f"[SUCCESS] Validation '{test_case['name']}': "
                      f"valide={validation['valid']}, erreurs={len(validation['errors'])}")

            except Exception as e:
                self.fail(f"Erreur validation config '{test_case['name']}': {e}")

    def test_template_validation(self):
        """Test de validation du template Premier Tech"""
        print(f"\n=== TEST: Validation Template PT ===")

        try:
            builder = NavigationBuilder(self.template_path)
            is_valid = builder.validate_template()

            self.assertTrue(is_valid, "Template Premier Tech devrait être valide")

            # Vérifier les informations du template
            self.assertEqual(builder.toc_slide_index, 12, "Index slide TOC incorrect")
            self.assertIn('layout_name', builder.toc_info, "Info layout manquante")

            print(f"[SUCCESS] Template validé: slide {builder.toc_info['slide_number']}, "
                  f"layout '{builder.toc_info['layout_name']}'")

        except Exception as e:
            self.fail(f"Erreur validation template: {e}")

    def tearDown(self):
        """Nettoyage après chaque test"""
        pass  # Conserver les fichiers pour inspection


def run_navigation_tests():
    """Lance tous les tests de navigation"""
    print("=== TESTS UNITAIRES NAVIGATION BUILDER - ARCHITECTURE JSON 2025 ===")
    print(f"Timestamp: {datetime.now().isoformat()}")
    print(f"Sujet: unit-tests | Audience: navigation_builder_tester")

    # Configuration du test runner
    loader = unittest.TestLoader()
    suite = loader.loadTestsFromTestCase(TestNavigationBuilder)

    # Exécuter les tests avec reporting détaillé
    runner = unittest.TextTestRunner(verbosity=2, buffer=False)
    result = runner.run(suite)

    # Rapport final
    print(f"\n=== RAPPORT FINAL ===")
    print(f"Tests exécutés: {result.testsRun}")
    print(f"Échecs: {len(result.failures)}")
    print(f"Erreurs: {len(result.errors)}")

    if result.failures:
        print("\nÉCHECS:")
        for test, traceback in result.failures:
            print(f"- {test}: {traceback}")

    if result.errors:
        print("\nERREURS:")
        for test, traceback in result.errors:
            print(f"- {test}: {traceback}")

    # Générer rapport JSON
    report_path = os.path.join('output', f'test_navigation_builder_{datetime.now().strftime("%Y%m%d_%H%M%S")}_report.json')
    test_report = {
        "timestamp": datetime.now().isoformat(),
        "test_suite": "NavigationBuilder JSON Architecture 2025",
        "subject": "unit-tests",
        "audience": "navigation_builder_tester",
        "results": {
            "tests_run": result.testsRun,
            "failures": len(result.failures),
            "errors": len(result.errors),
            "success_rate": (result.testsRun - len(result.failures) - len(result.errors)) / result.testsRun * 100 if result.testsRun > 0 else 0
        },
        "details": {
            "failures": [{"test": str(test), "error": traceback} for test, traceback in result.failures],
            "errors": [{"test": str(test), "error": traceback} for test, traceback in result.errors]
        }
    }

    with open(report_path, 'w', encoding='utf-8') as f:
        json.dump(test_report, f, ensure_ascii=False, indent=2)

    print(f"\nRapport JSON généré: {report_path}")

    return result.wasSuccessful()


if __name__ == "__main__":
    success = run_navigation_tests()
    sys.exit(0 if success else 1)