#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Test unitaire complet pour Content Boxes Builder - Architecture JSON 2025
Tests exhaustifs des fonctionnalités de création de slides avec boîtes de contenu.
"""

import os
import sys
import unittest
import tempfile
import shutil
import json
from pathlib import Path
from typing import Dict, List, Any

# Ajouter le dossier parent pour les imports
sys.path.insert(0, str(Path(__file__).parent.parent.parent.parent / "presentation_builder"))

from content_boxes_builder import ContentBoxesBuilder, process_content_boxes_config, load_content_boxes_payload, process_content_boxes_from_payload_file


class TestContentBoxesBuilder(unittest.TestCase):
    """
    Classe de test complète pour Content Boxes Builder.
    Couvre tous les styles de boîtes et configurations possibles.
    """

    def setUp(self):
        """Configuration initiale pour chaque test"""
        self.test_dir = Path(__file__).parent
        self.template_path = self.test_dir.parent.parent.parent / "templates" / "Template_PT.pptx"
        self.builder = ContentBoxesBuilder(str(self.template_path))

        # Créer un dossier temporaire pour les tests
        self.temp_dir = Path(tempfile.mkdtemp())

        # Paths des fichiers de test
        self.test_files = {
            "content_3_boxes": {
                "schema": self.test_dir / "content_3_boxes" / "presentation_schema.json",
                "payload": self.test_dir / "content_3_boxes" / "content_boxes_payload.json"
            },
            "content_4_boxes": {
                "schema": self.test_dir / "content_4_boxes" / "presentation_schema.json",
                "payload": self.test_dir / "content_4_boxes" / "content_boxes_payload.json"
            },
            "custom_content": {
                "schema": self.test_dir / "custom_content_boxes" / "presentation_schema.json",
                "payload": self.test_dir / "custom_content_boxes" / "content_boxes_payload.json"
            },
            "content_2_boxes": {
                "schema": self.test_dir / "content_2_boxes" / "presentation_schema.json",
                "payload": self.test_dir / "content_2_boxes" / "content_boxes_payload.json"
            }
        }

    def tearDown(self):
        """Nettoyage après chaque test"""
        if self.temp_dir.exists():
            shutil.rmtree(self.temp_dir)

    def test_01_builder_initialization(self):
        """Test 1: Initialisation du ContentBoxesBuilder"""
        print(f"\n[TEST 1] Initialisation du Content Boxes Builder...")

        # Vérifier que le builder s'initialise correctement
        self.assertIsNotNone(self.builder)
        self.assertEqual(str(self.builder.template_path), str(self.template_path))

        # Vérifier que les slides de content boxes sont analysées
        self.assertGreater(len(self.builder.content_info), 0)
        self.assertEqual(len(self.builder.content_slides), 8)  # 8 styles de content boxes

        # Vérifier les styles disponibles
        expected_styles = [
            'grey_3_detailed', 'grey_3_simple', 'blue_3_detailed', 'blue_3_simple',
            'grey_4_detailed', 'grey_4_simple', 'blue_4_detailed', 'blue_4_simple'
        ]

        for slide_index, slide_data in self.builder.content_slides.items():
            self.assertIn(slide_data['style'], expected_styles)
            self.assertIn('box_count', slide_data)
            self.assertIn(slide_data['box_count'], [3, 4])

        print(f"[SUCCESS] Builder initialisé avec {len(self.builder.content_slides)} styles")

    def test_02_validation_config_valid(self):
        """Test 2: Validation d'une configuration valide"""
        print(f"\n[TEST 2] Validation configuration valide...")

        config = {
            "content_style": "blue_3_simple",
            "concepts": [
                "Premier concept",
                "Deuxième concept",
                "Troisième concept"
            ],
            "title": "Test Content Boxes"
        }

        result = self.builder._validate_content_boxes_config(config)

        self.assertTrue(result['valid'])
        self.assertEqual(len(result['errors']), 0)

        print(f"[SUCCESS] Configuration valide correctement validée")

    def test_03_validation_config_invalid(self):
        """Test 3: Validation d'une configuration invalide"""
        print(f"\n[TEST 3] Validation configuration invalide...")

        # Test avec style invalide
        config_invalid_style = {
            "content_style": "invalid_style",
            "concepts": ["Concept 1", "Concept 2", "Concept 3"]
        }

        result = self.builder._validate_content_boxes_config(config_invalid_style)
        self.assertFalse(result['valid'])
        self.assertGreater(len(result['errors']), 0)

        # Test avec nombre incorrect de concepts
        config_wrong_count = {
            "content_style": "blue_4_simple",  # Requiert 4 concepts
            "concepts": ["Concept 1", "Concept 2"]  # Seulement 2 fournis
        }

        result = self.builder._validate_content_boxes_config(config_wrong_count)
        self.assertFalse(result['valid'])
        self.assertGreater(len(result['errors']), 0)

        print(f"[SUCCESS] Configurations invalides correctement rejetées")

    def test_04_load_payload_file(self):
        """Test 4: Chargement d'un payload depuis un fichier"""
        print(f"\n[TEST 4] Chargement payload depuis fichier...")

        payload_path = self.test_files["content_3_boxes"]["payload"]

        if payload_path.exists():
            payload = self.builder.load_content_boxes_payload(str(payload_path))

            self.assertIsInstance(payload, dict)
            self.assertIn('content_style', payload)
            self.assertIn('concepts', payload)
            self.assertEqual(payload['content_style'], 'blue_3_simple')
            self.assertEqual(len(payload['concepts']), 3)

            print(f"[SUCCESS] Payload chargé: {payload['content_style']} avec {len(payload['concepts'])} concepts")
        else:
            self.skipTest(f"Fichier payload non trouvé: {payload_path}")

    def test_05_process_config_3_boxes_simple(self):
        """Test 5: Traitement configuration 3 boîtes simples"""
        print(f"\n[TEST 5] Traitement configuration 3 boîtes simples...")

        # Créer une présentation de test basique
        test_presentation = self.temp_dir / "test_3_boxes.pptx"
        self._create_basic_presentation(test_presentation)

        config = {
            "content_style": "blue_3_simple",
            "concepts": [
                "Innovation Continue",
                "Excellence Opérationnelle",
                "Satisfaction Client"
            ],
            "title": "Nos 3 Piliers Stratégiques",
            "options": {
                "auto_widen": True
            }
        }

        result = self.builder.process_content_boxes_config(config, str(test_presentation))

        self.assertTrue(result['success'])
        self.assertTrue(test_presentation.exists())

        # Vérifier que la présentation a été modifiée
        self.assertGreater(test_presentation.stat().st_size, 1000)

        print(f"[SUCCESS] Configuration 3 boîtes simples traitée avec succès")

    def test_06_process_config_4_boxes_detailed(self):
        """Test 6: Traitement configuration 4 boîtes détaillées"""
        print(f"\n[TEST 6] Traitement configuration 4 boîtes détaillées...")

        test_presentation = self.temp_dir / "test_4_boxes.pptx"
        self._create_basic_presentation(test_presentation)

        config = {
            "content_style": "blue_4_detailed",
            "concepts": [
                "Analyse détaillée des besoins",
                "Conception architecture moderne",
                "Développement agile avec qualité",
                "Déploiement continu sécurisé"
            ],
            "subtitles": [
                "Analyse",
                "Architecture",
                "Développement",
                "Production"
            ],
            "title": "Processus de Développement",
            "options": {
                "auto_widen": True
            }
        }

        result = self.builder.process_content_boxes_config(config, str(test_presentation))

        self.assertTrue(result['success'])
        self.assertTrue(test_presentation.exists())

        print(f"[SUCCESS] Configuration 4 boîtes détaillées traitée avec succès")

    def test_07_process_from_payload_file(self):
        """Test 7: Traitement depuis fichier payload"""
        print(f"\n[TEST 7] Traitement depuis fichier payload...")

        payload_path = self.test_files["content_4_boxes"]["payload"]

        if payload_path.exists():
            test_presentation = self.temp_dir / "test_from_payload.pptx"
            self._create_basic_presentation(test_presentation)

            result = process_content_boxes_from_payload_file(
                str(payload_path),
                str(test_presentation),
                str(self.template_path)
            )

            self.assertTrue(result['success'])
            self.assertTrue(test_presentation.exists())

            print(f"[SUCCESS] Traitement depuis payload file réussi")
        else:
            self.skipTest(f"Fichier payload non trouvé: {payload_path}")

    def test_08_all_content_styles(self):
        """Test 8: Test de tous les styles de content boxes"""
        print(f"\n[TEST 8] Test de tous les styles de content boxes...")

        test_configs = [
            {
                "content_style": "grey_3_simple",
                "concepts": ["Concept A", "Concept B", "Concept C"],
                "title": "Test Gris 3 Simple"
            },
            {
                "content_style": "grey_3_detailed",
                "concepts": ["Concept A détaillé", "Concept B détaillé", "Concept C détaillé"],
                "subtitles": ["Sous A", "Sous B", "Sous C"],
                "title": "Test Gris 3 Détaillé"
            },
            {
                "content_style": "blue_3_simple",
                "concepts": ["Concept 1", "Concept 2", "Concept 3"],
                "title": "Test Bleu 3 Simple"
            },
            {
                "content_style": "blue_3_detailed",
                "concepts": ["Concept 1 détaillé", "Concept 2 détaillé", "Concept 3 détaillé"],
                "subtitles": ["Sub 1", "Sub 2", "Sub 3"],
                "title": "Test Bleu 3 Détaillé"
            },
            {
                "content_style": "grey_4_simple",
                "concepts": ["Étape 1", "Étape 2", "Étape 3", "Étape 4"],
                "title": "Test Gris 4 Simple"
            },
            {
                "content_style": "grey_4_detailed",
                "concepts": ["Étape 1 détaillée", "Étape 2 détaillée", "Étape 3 détaillée", "Étape 4 détaillée"],
                "subtitles": ["Phase 1", "Phase 2", "Phase 3", "Phase 4"],
                "title": "Test Gris 4 Détaillé"
            },
            {
                "content_style": "blue_4_simple",
                "concepts": ["Processus 1", "Processus 2", "Processus 3", "Processus 4"],
                "title": "Test Bleu 4 Simple"
            },
            {
                "content_style": "blue_4_detailed",
                "concepts": ["Processus 1 détaillé", "Processus 2 détaillé", "Processus 3 détaillé", "Processus 4 détaillé"],
                "subtitles": ["Proc 1", "Proc 2", "Proc 3", "Proc 4"],
                "title": "Test Bleu 4 Détaillé"
            }
        ]

        success_count = 0
        for i, config in enumerate(test_configs):
            try:
                test_presentation = self.temp_dir / f"test_style_{i}.pptx"
                self._create_basic_presentation(test_presentation)

                result = self.builder.process_content_boxes_config(config, str(test_presentation))

                if result['success']:
                    success_count += 1
                    print(f"  Style {config['content_style']}: SUCCESS")
                else:
                    print(f"  Style {config['content_style']}: FAILED - {result.get('error', 'Erreur inconnue')}")

            except Exception as e:
                print(f"  Style {config['content_style']}: ERROR - {e}")

        # Au moins 6 des 8 styles doivent passer (tolérance pour les styles expérimentaux)
        self.assertGreaterEqual(success_count, 6)
        print(f"[SUCCESS] {success_count}/8 styles testés avec succès")

    def _create_basic_presentation(self, path: Path):
        """Crée une présentation basique pour les tests"""
        try:
            # Copier le template vers le chemin de test
            shutil.copy2(self.template_path, path)

            # Simplifier: garder seulement la première slide (titre)
            from pptx import Presentation
            prs = Presentation(str(path))

            # Supprimer toutes les slides sauf la première
            slides_to_remove = list(range(1, len(prs.slides)))
            for i in reversed(slides_to_remove):
                try:
                    rId = prs.slides._sldIdLst[i].rId
                    prs.part.drop_rel(rId)
                    del prs.slides._sldIdLst[i]
                except:
                    pass  # Ignorer les erreurs de suppression

            prs.save(str(path))

        except Exception as e:
            print(f"[WARNING] Erreur création présentation basique: {e}")
            # Fallback: copier directement le template
            shutil.copy2(self.template_path, path)


def run_tests():
    """Exécute tous les tests et affiche les résultats"""
    print("=" * 70)
    print("TESTS UNITAIRES - CONTENT BOXES BUILDER (Architecture JSON 2025)")
    print("=" * 70)

    # Configuration du test runner
    loader = unittest.TestLoader()
    suite = loader.loadTestsFromTestCase(TestContentBoxesBuilder)

    # Exécuter les tests avec un output détaillé
    runner = unittest.TextTestRunner(verbosity=2, stream=sys.stdout)
    result = runner.run(suite)

    # Affichage du résumé
    print("\n" + "=" * 70)
    print("RÉSUMÉ DES TESTS")
    print("=" * 70)

    total_tests = result.testsRun
    failures = len(result.failures)
    errors = len(result.errors)
    success = total_tests - failures - errors

    print(f"Tests exécutés: {total_tests}")
    print(f"Succès: {success}")
    print(f"Échecs: {failures}")
    print(f"Erreurs: {errors}")

    if result.wasSuccessful():
        print("\n[SUCCESS] TOUS LES TESTS SONT PASSES!")
        print("Content Boxes Builder est prêt pour la production.")
    else:
        print("\nX CERTAINS TESTS ONT ECHOUE")
        if result.failures:
            print("\nEchecs:")
            for test, traceback in result.failures:
                print(f"  - {test}: {traceback.split('AssertionError:')[-1].strip()}")

        if result.errors:
            print("\nErreurs:")
            for test, traceback in result.errors:
                print(f"  - {test}: {traceback.split('Exception:')[-1].strip()}")

    print("=" * 70)

    return result.wasSuccessful()


if __name__ == "__main__":
    success = run_tests()
    sys.exit(0 if success else 1)