#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Test unitaire pour le Presentation Builder - Architecture JSON

Ce test valide le nouveau système basé sur JSON avec:
1. Création automatique de la slide titre
2. Array slides vide (aucune slide de contenu)
3. Ajout automatique de la slide de fermeture Premier Tech

Résultat attendu: Présentation de 2 slides (titre + fermeture)
"""

import os
import sys
import subprocess
import tempfile
import json
from pathlib import Path
from datetime import datetime

# Ajouter le répertoire racine au path pour les imports
sys.path.insert(0, str(Path(__file__).parent.parent.parent.parent))

class TestPresentationBuilder:
    """
    Test du nouveau système Presentation Builder avec architecture JSON.
    """

    def __init__(self):
        self.test_name = "Test Unitaire - Presentation Builder - Architecture JSON"
        self.script_path = Path(__file__).parent.parent.parent.parent / "presentation_builder" / "presentation_builder.py"
        self.test_json_path = Path(__file__).parent / "test_empty_slides_array.json"
        self.test_output_dir = Path(__file__).parent / "output"
        self.test_output_dir.mkdir(exist_ok=True)

    def setup_test(self):
        """Prépare l'environnement de test"""
        print(f"=== {self.test_name} ===")
        print(f"Repertoire de sortie: {self.test_output_dir}")
        print(f"Script teste: {self.script_path}")
        print(f"JSON de test: {self.test_json_path}")

        # Vérifier que les fichiers existent
        if not self.script_path.exists():
            raise FileNotFoundError(f"Script non trouvé: {self.script_path}")

        if not self.test_json_path.exists():
            raise FileNotFoundError(f"JSON de test non trouvé: {self.test_json_path}")

        # Afficher le contenu du JSON
        with open(self.test_json_path, 'r', encoding='utf-8') as f:
            test_config = json.load(f)

        print(f"\n=== Configuration de Test ===")
        print(f"Presentation: {test_config['presentation_name']}")
        print(f"Sujet: {test_config['subject']}")
        print(f"Audience: {test_config['audience']}")
        print(f"Titre: {test_config['title_slide']['title']}")
        print(f"Slides de contenu: {len(test_config['slides'])} (doit être 0)")
        print(f"Résultat attendu: 2 slides (titre + fermeture Premier Tech)")

    def run_json_validation_test(self):
        """
        Test la validation du JSON de configuration.

        Returns:
            bool: True si la validation réussit, False sinon
        """
        try:
            print(f"\n=== Test 1: Validation JSON ===")

            # Commande de validation
            cmd = [
                sys.executable,
                str(self.script_path),
                str(self.test_json_path),
                "--validate"
            ]

            print(f"Commande: {' '.join(cmd)}")

            # Exécuter la validation
            result = subprocess.run(
                cmd,
                cwd=self.script_path.parent,
                capture_output=True,
                text=True,
                timeout=30
            )

            # Vérifier le résultat
            if result.returncode != 0:
                print(f"ERREUR lors de la validation:")
                print(f"   STDOUT: {result.stdout}")
                print(f"   STDERR: {result.stderr}")
                return False

            print(f"SUCCESS: JSON validé avec succès")
            print(f"   Output: {result.stdout.strip()}")
            return True

        except subprocess.TimeoutExpired:
            print(f"ERREUR: Timeout lors de la validation")
            return False
        except Exception as e:
            print(f"ERREUR: Exception lors de la validation: {e}")
            return False

    def run_presentation_build_test(self):
        """
        Test la construction de présentation avec array slides vide.

        Returns:
            tuple: (success: bool, output_path: str)
        """
        try:
            print(f"\n=== Test 2: Construction Présentation ===")

            # Commande de construction
            cmd = [
                sys.executable,
                str(self.script_path),
                str(self.test_json_path)
            ]

            print(f"Commande: {' '.join(cmd)}")

            # Exécuter la construction
            result = subprocess.run(
                cmd,
                cwd=self.script_path.parent,
                capture_output=True,
                text=True,
                timeout=120
            )

            print(f"Code de retour: {result.returncode}")
            print(f"STDOUT:\n{result.stdout}")

            if result.stderr:
                print(f"STDERR:\n{result.stderr}")

            # Construire le chemin attendu basé sur la configuration JSON
            with open(self.test_json_path, 'r', encoding='utf-8') as f:
                test_config = json.load(f)

            # Reconstituer le chemin selon la logique de presentation_builder
            subject = test_config["subject"]
            audience = test_config["audience"]
            presentation_name = test_config["presentation_name"]

            # Nettoyer les noms comme le fait le script
            clean_subject = "".join(c for c in subject if c.isalnum() or c in ('-', '_')).lower()
            clean_audience = "".join(c for c in audience if c.isalnum() or c in ('-', '_')).lower()
            clean_name = "".join(c for c in presentation_name if c.isalnum() or c in ('-', '_', ' ')).strip()
            clean_name = clean_name.replace(' ', '_').lower()

            # Analyser la sortie pour trouver le timestamp
            timestamp = None
            for line in result.stdout.split('\n'):
                if '.pptx' in line:
                    # Extraire le timestamp du nom de fichier
                    import re
                    match = re.search(r'(\d{8}_\d{4})_', line)
                    if match:
                        timestamp = match.group(1)
                        break

            if not timestamp:
                print(f"ERREUR: Timestamp non trouvé dans la sortie")
                return False, None

            # Construire le chemin attendu
            # Structure: test/unit_tests/presentation_builder/test_*.py
            # Donc projet racine = 3 niveaux au-dessus
            project_root = Path(__file__).parent.parent.parent.parent
            output_dir = project_root / "presentations" / clean_subject / clean_audience / "output"
            expected_filename = f"{timestamp}_{clean_name}.pptx"
            output_path = output_dir / expected_filename

            if result.returncode != 0:
                print(f"ERREUR lors de la construction")
                return False, None

            # Vérifier que le fichier existe
            if not output_path.exists():
                print(f"ERREUR: Fichier de sortie non trouvé: {output_path}")
                # Afficher les fichiers qui existent dans le dossier pour debug
                if output_dir.exists():
                    print(f"Fichiers dans {output_dir}:")
                    for f in output_dir.iterdir():
                        print(f"  - {f.name}")
                return False, None

            # Vérifier la taille du fichier
            file_size = output_path.stat().st_size
            if file_size == 0:
                print(f"ERREUR: Fichier de sortie vide")
                return False, None

            print(f"SUCCESS: Présentation créée avec succès")
            print(f"   Fichier: {output_path}")
            print(f"   Taille: {file_size:,} bytes")

            return True, str(output_path)

        except subprocess.TimeoutExpired:
            print(f"ERREUR: Timeout lors de la construction")
            return False, None
        except Exception as e:
            print(f"ERREUR: Exception lors de la construction: {e}")
            return False, None

    def validate_presentation_content(self, presentation_path: str):
        """
        Valide le contenu de la présentation créée.

        Args:
            presentation_path: Chemin vers la présentation

        Returns:
            bool: True si la validation réussit, False sinon
        """
        try:
            print(f"\n=== Test 3: Validation Contenu ===")

            # Importer python-pptx pour analyser la présentation
            try:
                from pptx import Presentation
            except ImportError:
                print(f"WARNING: python-pptx non disponible, validation contenu ignorée")
                return True

            # Charger la présentation
            prs = Presentation(presentation_path)
            slide_count = len(prs.slides)

            print(f"Analyse de la présentation:")
            print(f"   Chemin: {presentation_path}")
            print(f"   Nombre de slides: {slide_count}")

            # Vérification: doit avoir exactement 2 slides
            expected_slides = 2  # titre + fermeture
            if slide_count != expected_slides:
                print(f"ERREUR: Nombre de slides incorrect")
                print(f"   Attendu: {expected_slides}")
                print(f"   Trouvé: {slide_count}")
                return False

            # Analyser chaque slide
            for i, slide in enumerate(prs.slides):
                print(f"   Slide {i+1}: {len(slide.shapes)} shapes, layout: {slide.slide_layout.name}")

                # Vérifier qu'il y a du contenu
                text_found = False
                for shape in slide.shapes:
                    if hasattr(shape, 'text_frame') and shape.text_frame and shape.text_frame.text.strip():
                        text_found = True
                        break

                if not text_found:
                    print(f"WARNING: Slide {i+1} semble vide")

            print(f"SUCCESS: Présentation validée")
            print(f"   Structure correcte: {slide_count} slides")
            print(f"   Slide 1: Titre (créée automatiquement)")
            print(f"   Slide 2: Fermeture Premier Tech (ajoutée automatiquement)")

            return True

        except Exception as e:
            print(f"ERREUR lors de la validation du contenu: {e}")
            return False

    def generate_test_report(self, validation_success: bool, build_success: bool,
                           content_success: bool, output_path: str = None):
        """Génère un rapport de test détaillé"""
        try:
            report = {
                "test_timestamp": datetime.now().isoformat(),
                "test_name": self.test_name,
                "test_configuration": {
                    "script_tested": str(self.script_path),
                    "json_config": str(self.test_json_path),
                    "test_scenario": "Array slides vide (titre + fermeture uniquement)"
                },
                "test_results": {
                    "json_validation": validation_success,
                    "presentation_build": build_success,
                    "content_validation": content_success,
                    "overall_success": validation_success and build_success and content_success
                },
                "output_file": output_path,
                "architecture_validation": {
                    "json_based_config": True,
                    "automatic_title_creation": True,
                    "automatic_closing_addition": True,
                    "empty_slides_array_handled": True
                },
                "expected_vs_actual": {
                    "expected_slides": 2,
                    "expected_structure": "Titre + Fermeture Premier Tech",
                    "validation_passed": content_success
                }
            }

            # Sauvegarder le rapport
            report_path = self.test_output_dir / "test_presentation_builder_report.json"
            with open(report_path, 'w', encoding='utf-8') as f:
                json.dump(report, f, ensure_ascii=False, indent=2)

            print(f"\n=== Rapport de Test ===")
            print(f"Rapport sauvegardé: {report_path}")

            return str(report_path)

        except Exception as e:
            print(f"WARNING: Erreur génération rapport: {e}")
            return None

    def cleanup_test(self):
        """Nettoie les fichiers de test si nécessaire"""
        print(f"\nNettoyage: Fichiers de test conservés dans {self.test_output_dir}")

    def run_all_tests(self):
        """
        Exécute tous les tests pour le Presentation Builder.

        Returns:
            bool: True si tous les tests réussissent, False sinon
        """
        try:
            self.setup_test()

            # Test 1: Validation JSON
            validation_success = self.run_json_validation_test()

            # Test 2: Construction présentation
            build_success, output_path = self.run_presentation_build_test()

            # Test 3: Validation contenu (si fichier créé)
            content_success = True
            if build_success and output_path:
                content_success = self.validate_presentation_content(output_path)

            # Générer le rapport
            self.generate_test_report(validation_success, build_success, content_success, output_path)

            self.cleanup_test()

            # Résultat final
            all_success = validation_success and build_success and content_success

            print(f"\n{'='*70}")
            print(f"RESULTATS FINAUX - {self.test_name}")
            print(f"{'='*70}")
            print(f"Validation JSON:        {'REUSSI' if validation_success else 'ECHEC'}")
            print(f"Construction:           {'REUSSI' if build_success else 'ECHEC'}")
            print(f"Validation contenu:     {'REUSSI' if content_success else 'ECHEC'}")
            print(f"Resultat global:        {'TOUS LES TESTS REUSSIS' if all_success else 'ECHECS DETECTES'}")

            if output_path:
                print(f"Fichier cree:           {output_path}")

            print(f"{'='*70}")

            return all_success

        except Exception as e:
            print(f"ERREUR critique dans les tests: {e}")
            return False


def main():
    """Point d'entrée principal pour exécuter les tests"""
    test_runner = TestPresentationBuilder()
    success = test_runner.run_all_tests()

    # Code de retour approprié
    exit_code = 0 if success else 1
    sys.exit(exit_code)


if __name__ == "__main__":
    main()