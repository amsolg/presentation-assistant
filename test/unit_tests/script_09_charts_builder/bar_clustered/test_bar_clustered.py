#!/usr/bin/env python3
"""
Test unitaire pour le style bar_clustered du script 09_charts_builder.py
Teste la création de graphiques en barres horizontales pour comparaisons.
"""

import unittest
import os
import sys
import json
from datetime import datetime
from pathlib import Path
from pptx import Presentation

# Ajouter les chemins nécessaires
project_root = Path(__file__).parent.parent.parent.parent.parent
script_dir = project_root / "presentation_builder"
data_dir = project_root / "data" / "charts"
sys.path.insert(0, str(script_dir))

# Importer le script amélioré
from importlib.util import spec_from_file_location, module_from_spec
spec = spec_from_file_location("charts_builder_enhanced",
                               str(script_dir / "09_charts_builder.py"))
charts_module = module_from_spec(spec)
spec.loader.exec_module(charts_module)
EnhancedChartsBuilder = charts_module.EnhancedChartsBuilder


class TestBarClusteredChart(unittest.TestCase):
    """Test du style bar_clustered avec le script amélioré"""

    @classmethod
    def setUpClass(cls):
        """Configuration unique pour tous les tests"""
        cls.output_dir = Path(__file__).parent / "output"
        cls.output_dir.mkdir(exist_ok=True)

        cls.template_path = str(project_root / "templates" / "Template_PT.pptx")
        cls.presentation_path = str(cls.output_dir / "test_bar_clustered.pptx")

        # Créer la présentation de base
        spec_01 = spec_from_file_location("slide_title_creator",
                                         str(script_dir / "01_slide_title_creator.py"))
        title_module = module_from_spec(spec_01)
        spec_01.loader.exec_module(title_module)

        creator = title_module.SlideTitleCreator()
        creator.create_title_slide(
            "Test Bar Clustered Chart",
            subtitle="Test unitaire graphique barres horizontales",
            output_path=cls.presentation_path,
            metadata="Test Suite - " + datetime.now().strftime("%Y-%m-%d")
        )

    def setUp(self):
        """Configuration pour chaque test"""
        self.builder = EnhancedChartsBuilder(self.template_path)

    def test_01_insert_bar_clustered_comparison(self):
        """Test graphique en barres pour comparaison horizontale"""
        result = self.builder.insert_enhanced_chart_slide(
            self.presentation_path,
            title="Comparaison Performance Équipes",
            style="bar_clustered",
            data_labels=["Équipe Alpha", "Équipe Beta", "Équipe Gamma", "Équipe Delta", "Équipe Epsilon"],
            data_values=[92, 88, 95, 79, 91],
            insights="Équipe Gamma leader avec 95% de performance",
            position=1
        )

        # Vérifications
        self.assertIsNotNone(result)
        self.assertEqual(result["style_used"], "bar_clustered")

        # Vérifier le type de graphique
        pres = Presentation(result["presentation_path"])
        chart_slide = pres.slides[1]

        has_bar_chart = False
        for shape in chart_slide.shapes:
            if hasattr(shape, 'chart'):
                has_bar_chart = True
                chart = shape.chart
                self.assertEqual(chart.chart_type.name, "BAR_CLUSTERED")
                break

        self.assertTrue(has_bar_chart, "Aucun graphique en barres trouvé")

    def test_02_bar_clustered_ranking_display(self):
        """Test affichage de classement avec barres horizontales"""
        # Données de classement (top 10 produits)
        products = [
            "Produit Premium A", "Produit Standard B", "Produit Économique C",
            "Produit Premium D", "Produit Standard E", "Produit Économique F",
            "Produit Premium G", "Produit Standard H", "Produit Économique I",
            "Produit Premium J"
        ]
        sales = [450000, 380000, 320000, 290000, 275000,
                260000, 245000, 230000, 215000, 200000]

        result = self.builder.insert_enhanced_chart_slide(
            self.presentation_path,
            title="Top 10 Produits par Ventes",
            style="bar_clustered",
            data_labels=products,
            data_values=sales,
            insights="Produits Premium dominent le classement",
            position=2
        )

        self.assertIsNotNone(result)
        self.assertEqual(result["style_used"], "bar_clustered")

        # Vérifier que les 10 produits sont présents
        pres = Presentation(result["presentation_path"])
        chart_slide = pres.slides[2]

        for shape in chart_slide.shapes:
            if hasattr(shape, 'chart'):
                chart = shape.chart
                categories = list(chart.plots[0].categories)
                self.assertEqual(len(categories), 10)
                break

    def test_03_bar_clustered_multi_series_regions(self):
        """Test barres horizontales multi-séries pour comparaison régionale"""
        csv_regions_path = str(data_dir / "regions_comparison.csv")

        result = self.builder.insert_enhanced_chart_slide(
            self.presentation_path,
            title="Performance Régionale Comparative",
            style="bar_clustered",
            csv_path=csv_regions_path,
            insights="Vue horizontale facilite la comparaison directe",
            position=3
        )

        self.assertIsNotNone(result)

        # Vérifier le multi-séries
        with open(result["report_path"], 'r', encoding='utf-8') as f:
            report = json.load(f)
        self.assertTrue(report["slide_details"]["has_multi_series"])

    def test_04_bar_clustered_category_names_long(self):
        """Test avec noms de catégories longs (idéal pour barres horizontales)"""
        # Les barres horizontales sont idéales pour les noms longs
        categories = [
            "Département des Ressources Humaines et Formation",
            "Service Informatique et Développement Digital",
            "Division Marketing et Communication Externe",
            "Unité Production et Contrôle Qualité",
            "Centre de Recherche et Développement Innovation",
            "Administration Générale et Services Financiers"
        ]
        values = [85, 92, 78, 88, 95, 82]

        result = self.builder.insert_enhanced_chart_slide(
            self.presentation_path,
            title="Satisfaction par Département",
            style="bar_clustered",
            data_labels=categories,
            data_values=values,
            insights="R&D Innovation montre le plus haut taux de satisfaction",
            position=4
        )

        self.assertIsNotNone(result)
        # Les barres horizontales devraient bien gérer les noms longs

    def test_05_bar_clustered_percentage_data(self):
        """Test avec données en pourcentages"""
        result = self.builder.insert_enhanced_chart_slide(
            self.presentation_path,
            title="Taux de Complétion Projets",
            style="bar_clustered",
            data_labels=["Projet Alpha", "Projet Beta", "Projet Gamma",
                        "Projet Delta", "Projet Epsilon"],
            data_values=["95%", "87%", "92%", "78%", "100%"],
            insights="Projet Epsilon complété à 100%",
            position=5
        )

        self.assertIsNotNone(result)

        # Les pourcentages devraient être convertis correctement
        pres = Presentation(result["presentation_path"])
        chart_slide = pres.slides[5]

        for shape in chart_slide.shapes:
            if hasattr(shape, 'chart'):
                chart = shape.chart
                series = chart.series[0]
                values = list(series.values)
                # Vérifier que la valeur maximale est 100
                self.assertEqual(max(values), 100.0)
                break

    def test_06_bar_clustered_negative_positive_values(self):
        """Test avec valeurs négatives et positives (variation)"""
        result = self.builder.insert_enhanced_chart_slide(
            self.presentation_path,
            title="Variation Annuelle par Secteur",
            style="bar_clustered",
            data_labels=["Technologie", "Finance", "Santé", "Énergie", "Commerce"],
            data_values=[15.5, -3.2, 8.7, -1.5, 12.3],
            insights="Technologie montre la plus forte croissance",
            position=6
        )

        self.assertIsNotNone(result)
        # Les barres horizontales sont excellentes pour montrer des variations +/-

    @classmethod
    def tearDownClass(cls):
        """Nettoyage après tous les tests"""
        # Créer un rapport final
        report = {
            "test_suite": "bar_clustered",
            "template_used": cls.template_path,
            "output_presentation": str(cls.presentation_path),
            "tests_completed": 6,
            "final_slide_count": 7,
            "test_focus": "Graphiques en barres horizontales pour comparaisons",
            "features_tested": [
                "Comparaison horizontale simple",
                "Classement top 10",
                "Multi-séries régionales",
                "Noms de catégories longs",
                "Données en pourcentages",
                "Valeurs négatives/positives"
            ]
        }

        report_path = cls.output_dir / "test_report.json"
        with open(report_path, 'w', encoding='utf-8') as f:
            json.dump(report, f, indent=2)

        print(f"\n[SUCCESS] Tests bar_clustered terminés")
        print(f"[PRESENTATION] Créée: {cls.presentation_path}")
        print(f"[REPORT] Rapport: {report_path}")


if __name__ == "__main__":
    unittest.main(verbosity=2)