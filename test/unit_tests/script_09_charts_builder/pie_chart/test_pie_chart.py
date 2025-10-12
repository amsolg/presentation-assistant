#!/usr/bin/env python3
"""
Test unitaire pour le style pie_chart du script 09_charts_builder.py
Teste la création de graphiques en secteurs pour répartitions et proportions.
"""

import unittest
import os
import sys
import json
from datetime import datetime
from pathlib import Path
from pptx import Presentation

# Ajouter les chemins nécessaires
project_root = Path(__file__).parent.parent.parent.parent.parent
script_dir = project_root / "presentation_builder"
data_dir = project_root / "data" / "charts"
sys.path.insert(0, str(script_dir))

# Importer le script amélioré
from importlib.util import spec_from_file_location, module_from_spec
spec = spec_from_file_location("charts_builder_enhanced",
                               str(script_dir / "09_charts_builder.py"))
charts_module = module_from_spec(spec)
spec.loader.exec_module(charts_module)
EnhancedChartsBuilder = charts_module.EnhancedChartsBuilder


class TestPieChart(unittest.TestCase):
    """Test du style pie_chart avec le script amélioré"""

    @classmethod
    def setUpClass(cls):
        """Configuration unique pour tous les tests"""
        cls.output_dir = Path(__file__).parent / "output"
        cls.output_dir.mkdir(exist_ok=True)

        cls.template_path = str(project_root / "templates" / "Template_PT.pptx")
        cls.csv_path = str(data_dir / "budget_repartition.csv")
        cls.presentation_path = str(cls.output_dir / "test_pie_chart.pptx")

        # Créer la présentation de base
        spec_01 = spec_from_file_location("slide_title_creator",
                                         str(script_dir / "01_slide_title_creator.py"))
        title_module = module_from_spec(spec_01)
        spec_01.loader.exec_module(title_module)

        creator = title_module.SlideTitleCreator()
        creator.create_title_slide(
            "Test Pie Chart",
            subtitle="Test unitaire graphique en secteurs",
            output_path=cls.presentation_path,
            metadata="Test Suite - " + datetime.now().strftime("%Y-%m-%d")
        )

    def setUp(self):
        """Configuration pour chaque test"""
        self.builder = EnhancedChartsBuilder(self.template_path)

    def test_01_insert_pie_chart_budget_repartition(self):
        """Test graphique en secteurs pour répartition budgétaire"""
        result = self.builder.insert_enhanced_chart_slide(
            self.presentation_path,
            title="Répartition Budget 2024",
            style="pie_chart",
            csv_path=self.csv_path,
            insights="R&D représente plus d'un tiers du budget total",
            position=1
        )

        # Vérifications
        self.assertIsNotNone(result)
        self.assertEqual(result["style_used"], "pie_chart")

        # Vérifier le type de graphique
        pres = Presentation(result["presentation_path"])
        chart_slide = pres.slides[1]

        has_pie_chart = False
        for shape in chart_slide.shapes:
            if hasattr(shape, 'chart'):
                has_pie_chart = True
                chart = shape.chart
                self.assertEqual(chart.chart_type.name, "PIE")
                break

        self.assertTrue(has_pie_chart, "Aucun graphique en secteurs trouvé")

    def test_02_pie_chart_percentage_normalization(self):
        """Test normalisation automatique à 100%"""
        # Données qui ne totalisent pas 100
        result = self.builder.insert_enhanced_chart_slide(
            self.presentation_path,
            title="Test Normalisation",
            style="pie_chart",
            data_labels=["Catégorie A", "Catégorie B", "Catégorie C"],
            data_values=[30, 45, 60],  # Total = 135, sera normalisé à 100%
            insights="Données normalisées automatiquement",
            position=2
        )

        self.assertIsNotNone(result)

        # Les valeurs devraient être normalisées à 100%
        # A: 30/135 = 22.2%, B: 45/135 = 33.3%, C: 60/135 = 44.4%
        pres = Presentation(result["presentation_path"])
        chart_slide = pres.slides[2]

        for shape in chart_slide.shapes:
            if hasattr(shape, 'chart'):
                chart = shape.chart
                series = chart.series[0]
                values = list(series.values)
                # Vérifier que le total est proche de 100
                total = sum(values)
                self.assertAlmostEqual(total, 100.0, places=1)
                break

    def test_03_pie_chart_with_kpi_data(self):
        """Test avec données KPI depuis CSV"""
        csv_kpi_path = str(data_dir / "kpi_performance.csv")

        result = self.builder.insert_enhanced_chart_slide(
            self.presentation_path,
            title="Performance KPI Q4",
            style="pie_chart",
            csv_path=csv_kpi_path,
            insights="65% des objectifs atteints, excellent résultat",
            position=3
        )

        self.assertIsNotNone(result)
        self.assertEqual(result["style_used"], "pie_chart")

    def test_04_pie_chart_many_categories_grouping(self):
        """Test regroupement automatique pour trop de catégories"""
        # Plus de 8 catégories - devrait être regroupé
        labels = [f"Département {i}" for i in range(1, 11)]
        values = [5, 8, 12, 15, 10, 7, 3, 2, 1, 1]  # Les 3 derniers sont petits

        result = self.builder.insert_enhanced_chart_slide(
            self.presentation_path,
            title="Test Regroupement Catégories",
            style="pie_chart",
            data_labels=labels,
            data_values=values,
            insights="Petites catégories regroupées en 'Autres'",
            position=4
        )

        self.assertIsNotNone(result)

        # Vérifier qu'il y a maintenant 8 catégories max (7 + "Autres")
        pres = Presentation(result["presentation_path"])
        chart_slide = pres.slides[4]

        for shape in chart_slide.shapes:
            if hasattr(shape, 'chart'):
                chart = shape.chart
                # Le nombre de catégories devrait être limité à 8
                categories = chart.plots[0].categories
                self.assertLessEqual(len(list(categories)), 8)
                break

    def test_05_pie_chart_with_json_config(self):
        """Test configuration JSON pour graphique en secteurs"""
        json_config_path = str(data_dir / "config_budget.json")

        result = self.builder.insert_enhanced_chart_slide(
            self.presentation_path,
            title="Override",  # Sera remplacé par config JSON
            style="column_clustered",  # Sera remplacé par config JSON
            json_config=json_config_path,
            position=5
        )

        self.assertIsNotNone(result)
        self.assertEqual(result["style_used"], "pie_chart")

        # Vérifier le titre depuis la config
        pres = Presentation(result["presentation_path"])
        chart_slide = pres.slides[5]
        self.assertEqual(chart_slide.shapes.title.text, "Répartition du Budget 2024")

    def test_06_pie_chart_minimal_data(self):
        """Test avec données minimales (2 catégories)"""
        result = self.builder.insert_enhanced_chart_slide(
            self.presentation_path,
            title="Répartition Simple",
            style="pie_chart",
            data_labels=["Complété", "En cours"],
            data_values=[75, 25],
            insights="Trois quarts du projet complété",
            position=6
        )

        self.assertIsNotNone(result)

        # Un graphique en secteurs devrait fonctionner avec 2 catégories minimum
        pres = Presentation(result["presentation_path"])
        self.assertEqual(len(pres.slides), 7)

    def test_07_pie_chart_custom_series_title(self):
        """Test avec titre de série personnalisé"""
        result = self.builder.insert_enhanced_chart_slide(
            self.presentation_path,
            title="Test Titre Personnalisé",
            style="pie_chart",
            data_labels=["Marketing", "Développement", "Support"],
            data_values=[30, 50, 20],
            series_title="Répartition des Ressources",
            insights="Développement utilise 50% des ressources",
            position=7
        )

        self.assertIsNotNone(result)

        # Vérifier le titre personnalisé
        pres = Presentation(result["presentation_path"])
        chart_slide = pres.slides[7]
        for shape in chart_slide.shapes:
            if hasattr(shape, 'chart'):
                chart = shape.chart
                if len(chart.series) > 0:
                    self.assertEqual(chart.series[0].name, "Répartition des Ressources")
                break

    def test_08_pie_chart_white_insights_text(self):
        """Test que le texte des insights est en blanc"""
        result = self.builder.insert_enhanced_chart_slide(
            self.presentation_path,
            title="Test Couleur Insights",
            style="pie_chart",
            data_labels=["A", "B", "C"],
            data_values=[33, 33, 34],
            insights="Le texte des points clés devrait être en blanc",
            position=8
        )

        self.assertIsNotNone(result)

        # Vérifier que le texte des insights est blanc
        pres = Presentation(result["presentation_path"])
        chart_slide = pres.slides[8]

        # Chercher la textbox avec les insights
        for shape in chart_slide.shapes:
            if hasattr(shape, 'text_frame') and shape.text_frame.text.startswith("Points clés:"):
                paragraph = shape.text_frame.paragraphs[0]
                font_color = paragraph.font.color.rgb
                # Vérifier que c'est blanc (255, 255, 255)
                self.assertEqual(str(font_color), "FFFFFF")
                break

    @classmethod
    def tearDownClass(cls):
        """Nettoyage après tous les tests"""
        # Créer un rapport final
        report = {
            "test_suite": "pie_chart",
            "template_used": cls.template_path,
            "output_presentation": str(cls.presentation_path),
            "tests_completed": 8,
            "final_slide_count": 9,
            "test_focus": "Graphiques en secteurs pour répartitions et proportions",
            "features_tested": [
                "Import CSV budget",
                "Normalisation automatique à 100%",
                "Données KPI",
                "Regroupement catégories (>8)",
                "Configuration JSON",
                "Données minimales",
                "Titre de série personnalisé",
                "Texte des insights en blanc"
            ],
            "new_features": [
                "Diversité de couleurs pour secteurs",
                "Prévention de la coupure de mots dans les labels",
                "Positionnement externe des étiquettes"
            ]
        }

        report_path = cls.output_dir / "test_report.json"
        with open(report_path, 'w', encoding='utf-8') as f:
            json.dump(report, f, indent=2)

        print(f"\n[SUCCESS] Tests pie_chart terminés")
        print(f"[PRESENTATION] Créée: {cls.presentation_path}")
        print(f"[REPORT] Rapport: {report_path}")


if __name__ == "__main__":
    unittest.main(verbosity=2)