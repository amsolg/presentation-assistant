#!/usr/bin/env python3
"""
Test unitaire pour le style column_compact du script 09_charts_builder.py
Teste la création de graphiques en colonnes compactes pour visualisations condensées.
"""

import unittest
import os
import sys
import json
from datetime import datetime
from pathlib import Path
from pptx import Presentation

# Ajouter les chemins nécessaires
project_root = Path(__file__).parent.parent.parent.parent.parent
script_dir = project_root / "presentation_builder"
data_dir = project_root / "data" / "charts"
sys.path.insert(0, str(script_dir))

# Importer le script amélioré
from importlib.util import spec_from_file_location, module_from_spec
spec = spec_from_file_location("charts_builder_enhanced",
                               str(script_dir / "09_charts_builder.py"))
charts_module = module_from_spec(spec)
spec.loader.exec_module(charts_module)
EnhancedChartsBuilder = charts_module.EnhancedChartsBuilder


class TestColumnCompactChart(unittest.TestCase):
    """Test du style column_compact avec le script amélioré"""

    @classmethod
    def setUpClass(cls):
        """Configuration unique pour tous les tests"""
        cls.output_dir = Path(__file__).parent / "output"
        cls.output_dir.mkdir(exist_ok=True)

        cls.template_path = str(project_root / "templates" / "Template_PT.pptx")
        cls.presentation_path = str(cls.output_dir / "test_column_compact.pptx")

        # Créer la présentation de base
        spec_01 = spec_from_file_location("slide_title_creator",
                                         str(script_dir / "01_slide_title_creator.py"))
        title_module = module_from_spec(spec_01)
        spec_01.loader.exec_module(title_module)

        creator = title_module.SlideTitleCreator()
        creator.create_title_slide(
            "Test Column Compact",
            subtitle="Test unitaire graphique colonnes compactes",
            output_path=cls.presentation_path,
            metadata="Test Suite - " + datetime.now().strftime("%Y-%m-%d")
        )

    def setUp(self):
        """Configuration pour chaque test"""
        self.builder = EnhancedChartsBuilder(self.template_path)

    def test_01_column_compact_simple_metrics(self):
        """Test graphique compact pour métriques simples"""
        result = self.builder.insert_enhanced_chart_slide(
            self.presentation_path,
            title="KPI Mensuels Compacts",
            style="column_compact",
            data_labels=["Jan", "Fév", "Mar", "Avr"],
            data_values=[82, 85, 88, 91],
            insights="Progression constante des KPI",
            position=1
        )

        # Vérifications
        self.assertIsNotNone(result)
        self.assertEqual(result["style_used"], "column_compact")

        # Vérifier le type de graphique (column_compact utilise COLUMN_CLUSTERED)
        pres = Presentation(result["presentation_path"])
        chart_slide = pres.slides[1]

        has_chart = False
        for shape in chart_slide.shapes:
            if hasattr(shape, 'chart'):
                has_chart = True
                chart = shape.chart
                self.assertEqual(chart.chart_type.name, "COLUMN_CLUSTERED")
                # Version compacte ne devrait pas avoir de légende
                self.assertFalse(chart.has_legend)
                break

        self.assertTrue(has_chart, "Aucun graphique compact trouvé")

    def test_02_column_compact_dashboard_element(self):
        """Test élément de dashboard compact"""
        # Données typiques pour un élément de dashboard
        result = self.builder.insert_enhanced_chart_slide(
            self.presentation_path,
            title="Ventes Hebdomadaires",
            style="column_compact",
            data_labels=["S1", "S2", "S3", "S4"],
            data_values=[125000, 132000, 128000, 141000],
            insights="Semaine 4 record du mois",
            position=2
        )

        self.assertIsNotNone(result)
        # Les graphiques compacts sont idéaux pour les dashboards

    def test_03_column_compact_limited_space(self):
        """Test avec espace limité (peu de catégories)"""
        result = self.builder.insert_enhanced_chart_slide(
            self.presentation_path,
            title="Focus Trimestriel",
            style="column_compact",
            data_labels=["T1", "T2", "T3"],
            data_values=[2.1, 2.4, 2.8],
            insights="Croissance trimestre après trimestre",
            position=3
        )

        self.assertIsNotNone(result)

        # Vérifier que le graphique est bien compact (dimensions réduites)
        with open(result["report_path"], 'r', encoding='utf-8') as f:
            report = json.load(f)

        # Le rapport devrait indiquer les dimensions compactes
        self.assertIn("7.0x4.0 inches", report["standards_applied"]["dimensions"])

    def test_04_column_compact_comparison_simple(self):
        """Test comparaison simple sans légende"""
        result = self.builder.insert_enhanced_chart_slide(
            self.presentation_path,
            title="Avant vs Après",
            style="column_compact",
            data_labels=["Avant", "Après"],
            data_values=[65, 92],
            insights="+42% d'amélioration",
            position=4
        )

        self.assertIsNotNone(result)

        # Graphique compact idéal pour comparaisons simples
        pres = Presentation(result["presentation_path"])
        chart_slide = pres.slides[4]

        for shape in chart_slide.shapes:
            if hasattr(shape, 'chart'):
                chart = shape.chart
                # Pas de légende nécessaire pour comparaison simple
                self.assertFalse(chart.has_legend)
                break

    def test_05_column_compact_year_over_year(self):
        """Test comparaison année sur année compact"""
        result = self.builder.insert_enhanced_chart_slide(
            self.presentation_path,
            title="Évolution 3 Ans",
            style="column_compact",
            data_labels=["2022", "2023", "2024"],
            data_values=[8.5, 9.2, 10.1],
            insights="Croissance stable sur 3 ans",
            position=5
        )

        self.assertIsNotNone(result)

    def test_06_column_compact_percentage_display(self):
        """Test affichage pourcentages compact"""
        result = self.builder.insert_enhanced_chart_slide(
            self.presentation_path,
            title="Taux de Réussite",
            style="column_compact",
            data_labels=["Objectif 1", "Objectif 2", "Objectif 3", "Objectif 4"],
            data_values=["88%", "92%", "85%", "95%"],
            insights="Tous les objectifs au-dessus de 85%",
            position=6
        )

        self.assertIsNotNone(result)

        # Les pourcentages devraient être convertis et affichés proprement
        pres = Presentation(result["presentation_path"])
        self.assertEqual(len(pres.slides), 7)

    @classmethod
    def tearDownClass(cls):
        """Nettoyage après tous les tests"""
        # Créer un rapport final
        report = {
            "test_suite": "column_compact",
            "template_used": cls.template_path,
            "output_presentation": str(cls.presentation_path),
            "tests_completed": 6,
            "final_slide_count": 7,
            "test_focus": "Graphiques en colonnes compactes pour visualisations condensées",
            "features_tested": [
                "Métriques simples compactes",
                "Élément de dashboard",
                "Espace limité (3 catégories)",
                "Comparaison simple sans légende",
                "Évolution sur 3 ans",
                "Affichage pourcentages"
            ],
            "key_characteristics": [
                "Pas de légende",
                "Dimensions réduites (7.0x4.0 inches)",
                "Idéal pour dashboards",
                "Maximum 4-5 catégories recommandé"
            ]
        }

        report_path = cls.output_dir / "test_report.json"
        with open(report_path, 'w', encoding='utf-8') as f:
            json.dump(report, f, indent=2)

        print(f"\n[SUCCESS] Tests column_compact terminés")
        print(f"[PRESENTATION] Créée: {cls.presentation_path}")
        print(f"[REPORT] Rapport: {report_path}")


if __name__ == "__main__":
    unittest.main(verbosity=2)