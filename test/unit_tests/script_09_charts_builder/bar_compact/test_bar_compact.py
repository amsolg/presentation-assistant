#!/usr/bin/env python3
"""
Test unitaire pour le style bar_compact du script 09_charts_builder.py
Teste la création de graphiques en barres horizontales compactes.
"""

import unittest
import os
import sys
import json
from datetime import datetime
from pathlib import Path
from pptx import Presentation

# Ajouter les chemins nécessaires
project_root = Path(__file__).parent.parent.parent.parent.parent
script_dir = project_root / "presentation_builder"
data_dir = project_root / "data" / "charts"
sys.path.insert(0, str(script_dir))

# Importer le script amélioré
from importlib.util import spec_from_file_location, module_from_spec
spec = spec_from_file_location("charts_builder_enhanced",
                               str(script_dir / "09_charts_builder.py"))
charts_module = module_from_spec(spec)
spec.loader.exec_module(charts_module)
EnhancedChartsBuilder = charts_module.EnhancedChartsBuilder


class TestBarCompactChart(unittest.TestCase):
    """Test du style bar_compact avec le script amélioré"""

    @classmethod
    def setUpClass(cls):
        """Configuration unique pour tous les tests"""
        cls.output_dir = Path(__file__).parent / "output"
        cls.output_dir.mkdir(exist_ok=True)

        cls.template_path = str(project_root / "templates" / "Template_PT.pptx")
        cls.presentation_path = str(cls.output_dir / "test_bar_compact.pptx")

        # Créer la présentation de base
        spec_01 = spec_from_file_location("slide_title_creator",
                                         str(script_dir / "01_slide_title_creator.py"))
        title_module = module_from_spec(spec_01)
        spec_01.loader.exec_module(title_module)

        creator = title_module.SlideTitleCreator()
        creator.create_title_slide(
            "Test Bar Compact",
            subtitle="Test unitaire graphique barres compactes",
            output_path=cls.presentation_path,
            metadata="Test Suite - " + datetime.now().strftime("%Y-%m-%d")
        )

    def setUp(self):
        """Configuration pour chaque test"""
        self.builder = EnhancedChartsBuilder(self.template_path)

    def test_01_bar_compact_quick_comparison(self):
        """Test graphique barres compactes pour comparaison rapide"""
        result = self.builder.insert_enhanced_chart_slide(
            self.presentation_path,
            title="Comparaison Rapide",
            style="bar_compact",
            data_labels=["Option A", "Option B", "Option C"],
            data_values=[75, 82, 68],
            insights="Option B présente le meilleur score",
            position=1
        )

        # Vérifications
        self.assertIsNotNone(result)
        self.assertEqual(result["style_used"], "bar_compact")

        # Vérifier le type de graphique (bar_compact utilise BAR_CLUSTERED)
        pres = Presentation(result["presentation_path"])
        chart_slide = pres.slides[1]

        has_chart = False
        for shape in chart_slide.shapes:
            if hasattr(shape, 'chart'):
                has_chart = True
                chart = shape.chart
                self.assertEqual(chart.chart_type.name, "BAR_CLUSTERED")
                # Version compacte ne devrait pas avoir de légende
                self.assertFalse(chart.has_legend)
                break

        self.assertTrue(has_chart, "Aucun graphique barres compact trouvé")

    def test_02_bar_compact_top_5_ranking(self):
        """Test classement Top 5 compact"""
        result = self.builder.insert_enhanced_chart_slide(
            self.presentation_path,
            title="Top 5 Vendeurs",
            style="bar_compact",
            data_labels=["Marie L.", "Jean P.", "Sophie M.", "Pierre D.", "Anne B."],
            data_values=[145000, 138000, 132000, 128000, 125000],
            insights="Marie L. en tête avec 145K€",
            position=2
        )

        self.assertIsNotNone(result)

        # Bar compact parfait pour petits classements
        pres = Presentation(result["presentation_path"])
        chart_slide = pres.slides[2]

        for shape in chart_slide.shapes:
            if hasattr(shape, 'chart'):
                chart = shape.chart
                categories = list(chart.plots[0].categories)
                self.assertEqual(len(categories), 5)
                break

    def test_03_bar_compact_binary_comparison(self):
        """Test comparaison binaire (Oui/Non, Pour/Contre)"""
        result = self.builder.insert_enhanced_chart_slide(
            self.presentation_path,
            title="Résultat Vote",
            style="bar_compact",
            data_labels=["Pour", "Contre"],
            data_values=[73, 27],
            insights="Majorité claire à 73%",
            position=3
        )

        self.assertIsNotNone(result)

        # Idéal pour comparaisons binaires simples
        pres = Presentation(result["presentation_path"])
        self.assertEqual(len(pres.slides), 4)

    def test_04_bar_compact_satisfaction_scores(self):
        """Test scores de satisfaction compact"""
        result = self.builder.insert_enhanced_chart_slide(
            self.presentation_path,
            title="Satisfaction Client",
            style="bar_compact",
            data_labels=["Très satisfait", "Satisfait", "Neutre", "Insatisfait"],
            data_values=[42, 38, 15, 5],
            insights="80% de satisfaction positive",
            position=4
        )

        self.assertIsNotNone(result)

        # Vérifier les dimensions compactes
        with open(result["report_path"], 'r', encoding='utf-8') as f:
            report = json.load(f)

        self.assertIn("7.0x4.0 inches", report["standards_applied"]["dimensions"])

    def test_05_bar_compact_progress_indicators(self):
        """Test indicateurs de progression compacts"""
        result = self.builder.insert_enhanced_chart_slide(
            self.presentation_path,
            title="État d'Avancement",
            style="bar_compact",
            data_labels=["Phase 1", "Phase 2", "Phase 3", "Phase 4"],
            data_values=["100%", "85%", "45%", "10%"],
            insights="Phase 1 complétée, Phase 2 presque terminée",
            position=5
        )

        self.assertIsNotNone(result)

        # Les pourcentages devraient être bien gérés
        pres = Presentation(result["presentation_path"])
        chart_slide = pres.slides[5]

        for shape in chart_slide.shapes:
            if hasattr(shape, 'chart'):
                chart = shape.chart
                series = chart.series[0]
                values = list(series.values)
                # Phase 1 devrait être à 100%
                self.assertEqual(values[0], 100.0)
                break

    def test_06_bar_compact_minimal_categories(self):
        """Test avec nombre minimal de catégories (2-3)"""
        result = self.builder.insert_enhanced_chart_slide(
            self.presentation_path,
            title="Analyse Simple",
            style="bar_compact",
            data_labels=["Actuel", "Cible", "Écart"],
            data_values=[85, 100, 15],
            insights="15 points d'écart avec la cible",
            position=6
        )

        self.assertIsNotNone(result)

        # Bar compact excellent pour peu de catégories
        pres = Presentation(result["presentation_path"])
        self.assertEqual(len(pres.slides), 7)

    @classmethod
    def tearDownClass(cls):
        """Nettoyage après tous les tests"""
        # Créer un rapport final
        report = {
            "test_suite": "bar_compact",
            "template_used": cls.template_path,
            "output_presentation": str(cls.presentation_path),
            "tests_completed": 6,
            "final_slide_count": 7,
            "test_focus": "Graphiques en barres horizontales compactes",
            "features_tested": [
                "Comparaison rapide 3 options",
                "Top 5 classement compact",
                "Comparaison binaire (Pour/Contre)",
                "Scores de satisfaction",
                "Indicateurs de progression",
                "Catégories minimales (2-3)"
            ],
            "key_characteristics": [
                "Pas de légende",
                "Dimensions réduites (7.0x4.0 inches)",
                "Idéal pour comparaisons simples",
                "Maximum 5 catégories recommandé",
                "Parfait pour dashboards"
            ]
        }

        report_path = cls.output_dir / "test_report.json"
        with open(report_path, 'w', encoding='utf-8') as f:
            json.dump(report, f, indent=2)

        print(f"\n[SUCCESS] Tests bar_compact terminés")
        print(f"[PRESENTATION] Créée: {cls.presentation_path}")
        print(f"[REPORT] Rapport: {report_path}")


if __name__ == "__main__":
    unittest.main(verbosity=2)