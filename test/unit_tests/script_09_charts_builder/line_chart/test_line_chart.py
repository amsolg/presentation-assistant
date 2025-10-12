#!/usr/bin/env python3
"""
Test unitaire pour le style line_chart du script 09_charts_builder.py
Teste la création de graphiques linéaires pour tendances temporelles.
"""

import unittest
import os
import sys
import json
from datetime import datetime
from pathlib import Path
from pptx import Presentation

# Ajouter les chemins nécessaires
project_root = Path(__file__).parent.parent.parent.parent.parent
script_dir = project_root / "presentation_builder"
data_dir = project_root / "data" / "charts"
sys.path.insert(0, str(script_dir))

# Importer le script amélioré
from importlib.util import spec_from_file_location, module_from_spec
spec = spec_from_file_location("charts_builder_enhanced",
                               str(script_dir / "09_charts_builder.py"))
charts_module = module_from_spec(spec)
spec.loader.exec_module(charts_module)
EnhancedChartsBuilder = charts_module.EnhancedChartsBuilder


class TestLineChart(unittest.TestCase):
    """Test du style line_chart avec le script amélioré"""

    @classmethod
    def setUpClass(cls):
        """Configuration unique pour tous les tests"""
        cls.output_dir = Path(__file__).parent / "output"
        cls.output_dir.mkdir(exist_ok=True)

        cls.template_path = str(project_root / "templates" / "Template_PT.pptx")
        cls.csv_path = str(data_dir / "ventes_trimestrielles.csv")
        cls.presentation_path = str(cls.output_dir / "test_line_chart.pptx")

        # Créer la présentation de base avec le script 01
        spec_01 = spec_from_file_location("slide_title_creator",
                                         str(script_dir / "01_slide_title_creator.py"))
        title_module = module_from_spec(spec_01)
        spec_01.loader.exec_module(title_module)

        creator = title_module.SlideTitleCreator()
        creator.create_title_slide(
            "Test Line Chart",
            subtitle="Test unitaire graphique linéaire",
            output_path=cls.presentation_path,
            metadata="Test Suite - " + datetime.now().strftime("%Y-%m-%d")
        )

    def setUp(self):
        """Configuration pour chaque test"""
        self.builder = EnhancedChartsBuilder(self.template_path)

    def test_01_insert_line_chart_temporal_evolution(self):
        """Test graphique linéaire pour évolution temporelle"""
        result = self.builder.insert_enhanced_chart_slide(
            self.presentation_path,
            title="Évolution Temporelle des Ventes",
            style="line_chart",
            csv_path=self.csv_path,
            insights="Tendance haussière sur 8 trimestres",
            position=1
        )

        # Vérifications
        self.assertIsNotNone(result)
        self.assertEqual(result["style_used"], "line_chart")
        self.assertEqual(result["position"], 2)

        # Vérifier le graphique
        pres = Presentation(result["presentation_path"])
        chart_slide = pres.slides[1]

        has_chart = False
        for shape in chart_slide.shapes:
            if hasattr(shape, 'chart'):
                has_chart = True
                chart = shape.chart
                self.assertEqual(chart.chart_type.name, "LINE")
                break

        self.assertTrue(has_chart, "Aucun graphique linéaire trouvé")

    def test_02_line_chart_multi_series_comparison(self):
        """Test graphique linéaire multi-séries pour comparaison régionale"""
        csv_multi_path = str(data_dir / "regions_comparison.csv")

        result = self.builder.insert_enhanced_chart_slide(
            self.presentation_path,
            title="Évolution Comparative par Région",
            style="line_chart",
            csv_path=csv_multi_path,
            insights="Europe montre la croissance la plus stable",
            position=2
        )

        self.assertIsNotNone(result)
        self.assertEqual(result["style_used"], "line_chart")

        # Vérifier le multi-séries
        with open(result["report_path"], 'r', encoding='utf-8') as f:
            report = json.load(f)
        self.assertTrue(report["slide_details"]["has_multi_series"])

    def test_03_line_chart_with_json_config(self):
        """Test avec configuration JSON pour graphique linéaire"""
        # Créer une config JSON temporaire
        json_config = {
            "title": "Tendance Mensuelle KPI",
            "style": "line_chart",
            "data": {
                "labels": ["Jan", "Fév", "Mar", "Avr", "Mai", "Jun"],
                "values": [85, 88, 92, 90, 95, 98]
            },
            "insights": "Amélioration continue du KPI avec légère baisse en avril"
        }

        json_path = str(self.output_dir / "line_config.json")
        with open(json_path, 'w', encoding='utf-8') as f:
            json.dump(json_config, f)

        result = self.builder.insert_enhanced_chart_slide(
            self.presentation_path,
            title="Config JSON Override",  # Sera remplacé par le titre JSON
            style="column_clustered",  # Sera remplacé par le style JSON
            json_config=json_path,
            position=3
        )

        self.assertIsNotNone(result)
        self.assertEqual(result["style_used"], "line_chart")

        # Vérifier le titre
        pres = Presentation(result["presentation_path"])
        chart_slide = pres.slides[3]
        self.assertEqual(chart_slide.shapes.title.text, "Tendance Mensuelle KPI")

    def test_04_line_chart_trend_analysis(self):
        """Test analyse de tendance avec graphique linéaire"""
        # Données montrant une tendance claire
        result = self.builder.insert_enhanced_chart_slide(
            self.presentation_path,
            title="Analyse de Tendance",
            style="line_chart",
            data_labels=["T1", "T2", "T3", "T4", "T5", "T6", "T7", "T8"],
            data_values=[100, 110, 125, 140, 160, 185, 210, 240],
            insights="Croissance exponentielle: +140% sur 8 périodes",
            position=4
        )

        self.assertIsNotNone(result)

        # Vérifier que le graphique a bien 8 points
        pres = Presentation(result["presentation_path"])
        chart_slide = pres.slides[4]

        for shape in chart_slide.shapes:
            if hasattr(shape, 'chart'):
                chart = shape.chart
                # Le graphique devrait avoir une série avec 8 points
                series = chart.series[0]
                self.assertEqual(len(list(series.values)), 8)
                break

    def test_05_line_chart_seasonal_pattern(self):
        """Test motif saisonnier avec graphique linéaire"""
        # Données avec pattern saisonnier (4 ans de données trimestrielles)
        labels = []
        values = []
        base_values = [100, 120, 110, 130]  # Pattern saisonnier

        for year in range(2021, 2025):
            for quarter in range(1, 5):
                labels.append(f"Q{quarter} {year}")
                # Ajouter une croissance annuelle au pattern
                growth_factor = 1 + (year - 2021) * 0.1
                values.append(base_values[quarter-1] * growth_factor)

        result = self.builder.insert_enhanced_chart_slide(
            self.presentation_path,
            title="Analyse Saisonnière",
            style="line_chart",
            data_labels=labels,
            data_values=values,
            insights="Pattern saisonnier avec pic récurrent au Q4",
            position=5
        )

        self.assertIsNotNone(result)
        self.assertEqual(len(labels), 16)  # 4 ans × 4 trimestres

    def test_06_line_chart_minimal_data(self):
        """Test avec données minimales (3 points)"""
        result = self.builder.insert_enhanced_chart_slide(
            self.presentation_path,
            title="Tendance Simple",
            style="line_chart",
            data_labels=["Début", "Milieu", "Fin"],
            data_values=[50, 75, 100],
            insights="Progression linéaire simple",
            position=6
        )

        self.assertIsNotNone(result)

        # Un graphique linéaire devrait fonctionner avec 3 points minimum
        pres = Presentation(result["presentation_path"])
        self.assertEqual(len(pres.slides), 7)  # 1 titre + 6 graphiques

    @classmethod
    def tearDownClass(cls):
        """Nettoyage après tous les tests"""
        # Créer un rapport final
        report = {
            "test_suite": "line_chart",
            "template_used": cls.template_path,
            "output_presentation": str(cls.presentation_path),
            "tests_completed": 6,
            "final_slide_count": 7,
            "test_focus": "Graphiques linéaires pour tendances temporelles"
        }

        report_path = cls.output_dir / "test_report.json"
        with open(report_path, 'w', encoding='utf-8') as f:
            json.dump(report, f, indent=2)

        print(f"\n[SUCCESS] Tests line_chart terminés")
        print(f"[PRESENTATION] Créée: {cls.presentation_path}")
        print(f"[REPORT] Rapport: {report_path}")


if __name__ == "__main__":
    unittest.main(verbosity=2)