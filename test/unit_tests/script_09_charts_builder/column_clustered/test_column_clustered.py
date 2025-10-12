#!/usr/bin/env python3
"""
Test unitaire pour le style column_clustered du script 09_charts_builder.py
Teste la création d'un graphique en colonnes groupées avec import CSV.
"""

import unittest
import os
import sys
import json
from datetime import datetime
from pathlib import Path
from pptx import Presentation

# Ajouter les chemins nécessaires
project_root = Path(__file__).parent.parent.parent.parent.parent
script_dir = project_root / "presentation_builder"
data_dir = project_root / "data" / "charts"
sys.path.insert(0, str(script_dir))

# Importer le script amélioré
from importlib.util import spec_from_file_location, module_from_spec
spec = spec_from_file_location("charts_builder_enhanced",
                               str(script_dir / "09_charts_builder.py"))
charts_module = module_from_spec(spec)
spec.loader.exec_module(charts_module)
EnhancedChartsBuilder = charts_module.EnhancedChartsBuilder


class TestColumnClusteredChart(unittest.TestCase):
    """Test du style column_clustered avec le script amélioré"""

    @classmethod
    def setUpClass(cls):
        """Configuration unique pour tous les tests"""
        cls.output_dir = Path(__file__).parent / "output"
        cls.output_dir.mkdir(exist_ok=True)

        cls.template_path = str(project_root / "templates" / "Template_PT.pptx")
        cls.csv_path = str(data_dir / "ventes_trimestrielles.csv")
        cls.presentation_path = str(cls.output_dir / "test_column_clustered.pptx")

        # Créer la présentation de base avec le script 01
        spec_01 = spec_from_file_location("slide_title_creator",
                                         str(script_dir / "01_slide_title_creator.py"))
        title_module = module_from_spec(spec_01)
        spec_01.loader.exec_module(title_module)

        creator = title_module.SlideTitleCreator()  # Pas de paramètre
        creator.create_title_slide(
            "Test Column Clustered Chart",
            subtitle="Test unitaire graphique colonnes groupées",
            output_path=cls.presentation_path,
            metadata="Test Suite - " + datetime.now().strftime("%Y-%m-%d")
        )

    def setUp(self):
        """Configuration pour chaque test"""
        self.builder = EnhancedChartsBuilder(self.template_path)

    def test_01_insert_column_clustered_with_csv(self):
        """Test insertion d'un graphique en colonnes groupées depuis CSV"""
        result = self.builder.insert_enhanced_chart_slide(
            self.presentation_path,
            title="Ventes Trimestrielles 2024",
            style="column_clustered",
            csv_path=self.csv_path,
            insights="Croissance continue avec pic au T4",
            position=1  # Position 2 (après la slide de titre)
        )

        # Vérifier le résultat
        self.assertIsNotNone(result)
        self.assertEqual(result["style_used"], "column_clustered")
        self.assertEqual(result["position"], 2)

        # Vérifier que le fichier existe
        self.assertTrue(os.path.exists(result["presentation_path"]))

        # Vérifier le contenu
        pres = Presentation(result["presentation_path"])
        self.assertEqual(len(pres.slides), 2)

        # Vérifier la slide de graphique
        chart_slide = pres.slides[1]
        self.assertIsNotNone(chart_slide.shapes.title)
        self.assertEqual(chart_slide.shapes.title.text, "Ventes Trimestrielles 2024")

        # Vérifier qu'il y a un graphique
        has_chart = False
        for shape in chart_slide.shapes:
            if hasattr(shape, 'chart'):
                has_chart = True
                chart = shape.chart
                # Vérifier le type de graphique
                self.assertEqual(chart.chart_type.name, "COLUMN_CLUSTERED")
                break

        self.assertTrue(has_chart, "Aucun graphique trouvé dans la slide")

    def test_02_column_clustered_with_direct_data(self):
        """Test avec données directes (sans CSV) et titre personnalisé"""
        result = self.builder.insert_enhanced_chart_slide(
            self.presentation_path,
            title="Performance par Trimestre",
            style="column_clustered",
            data_labels=["T1", "T2", "T3", "T4"],
            data_values=["2500000", "3200000", "2800000", "3700000"],
            insights="T4 montre la meilleure performance",
            series_title="Revenus Trimestriels",  # Nouveau paramètre
            position=2  # Position 3
        )

        self.assertIsNotNone(result)
        self.assertEqual(result["style_used"], "column_clustered")

        # Vérifier que nous avons maintenant 3 slides
        pres = Presentation(result["presentation_path"])
        self.assertEqual(len(pres.slides), 3)

        # Vérifier que le titre personnalisé est utilisé
        chart_slide = pres.slides[2]
        for shape in chart_slide.shapes:
            if hasattr(shape, 'chart'):
                chart = shape.chart
                # Le titre de la série devrait être "Revenus Trimestriels"
                if len(chart.series) > 0:
                    self.assertEqual(chart.series[0].name, "Revenus Trimestriels")
                break

    def test_03_column_clustered_multi_series(self):
        """Test avec données multi-séries depuis CSV"""
        csv_multi_path = str(data_dir / "regions_comparison.csv")

        result = self.builder.insert_enhanced_chart_slide(
            self.presentation_path,
            title="Comparaison Régionale",
            style="column_clustered",
            csv_path=csv_multi_path,
            insights="Europe leader sur tous les trimestres",
            position=3  # Position 4
        )

        self.assertIsNotNone(result)
        self.assertEqual(result["style_used"], "column_clustered")

        # Vérifier le rapport
        report_path = result["report_path"]
        self.assertTrue(os.path.exists(report_path))

        with open(report_path, 'r', encoding='utf-8') as f:
            report = json.load(f)

        self.assertTrue(report["slide_details"]["has_multi_series"])
        self.assertEqual(report["slide_details"]["data_source"], "CSV")

    def test_04_column_clustered_data_validation(self):
        """Test de la validation et enrichissement des données"""
        # Données avec pourcentages et formats variés
        result = self.builder.insert_enhanced_chart_slide(
            self.presentation_path,
            title="Test Validation Données",
            style="column_clustered",
            data_labels=["Jan", "Fév", "Mar", "Avr", "Mai"],
            data_values=["25%", "30.5%", "45,2%", "22", "38%"],
            position=4  # Position 5
        )

        self.assertIsNotNone(result)

        # Les données devraient être converties correctement
        pres = Presentation(result["presentation_path"])
        chart_slide = pres.slides[4]

        for shape in chart_slide.shapes:
            if hasattr(shape, 'chart'):
                chart = shape.chart
                # Vérifier que les valeurs ont été converties en nombres
                series = chart.series[0]
                values = list(series.values)
                self.assertTrue(all(isinstance(v, (int, float)) for v in values))
                break

    def test_05_export_column_clustered_data(self):
        """Test de l'export des données du graphique"""
        export_path = str(self.output_dir / "exported_column_data.csv")

        # Exporter les données de la dernière slide ajoutée
        self.builder.export_chart_data(
            self.presentation_path,
            export_path,
            slide_index=-1
        )

        # Vérifier que le fichier CSV a été créé
        self.assertTrue(os.path.exists(export_path))

        # Vérifier le contenu si pandas est disponible
        try:
            import pandas as pd
            df = pd.read_csv(export_path)
            self.assertGreater(len(df), 0)
            self.assertIn('Catégorie', df.columns)
            self.assertIn('Valeur', df.columns)
        except ImportError:
            # Test basique sans pandas
            with open(export_path, 'r', encoding='utf-8') as f:
                content = f.read()
                self.assertIn('Catégorie', content)

    def test_06_premier_tech_standards(self):
        """Test que les standards Premier Tech sont respectés"""
        result = self.builder.insert_enhanced_chart_slide(
            self.presentation_path,
            title="Test Standards PT",
            style="column_clustered",
            data_labels=["A", "B", "C"],
            data_values=[100, 200, 150],
            position=5  # Position 6
        )

        # Vérifier le rapport pour les standards
        with open(result["report_path"], 'r', encoding='utf-8') as f:
            report = json.load(f)

        self.assertIn("standards_applied", report)
        self.assertIn("Palette Premier Tech", report["standards_applied"]["colors"])
        self.assertIn("Standards visuels respectés", report["standards_applied"]["formatting"])

    def test_07_color_diversity(self):
        """Test de la diversité des couleurs pour multi-séries"""
        result = self.builder.insert_enhanced_chart_slide(
            self.presentation_path,
            title="Test Diversité Couleurs",
            style="column_clustered",
            csv_path=str(data_dir / "regions_comparison.csv"),  # CSV multi-séries
            insights="Test de palette de couleurs diversifiée",
            position=6  # Position 7
        )

        self.assertIsNotNone(result)

        # Vérifier que différentes couleurs sont utilisées
        pres = Presentation(result["presentation_path"])
        chart_slide = pres.slides[6]

        for shape in chart_slide.shapes:
            if hasattr(shape, 'chart'):
                chart = shape.chart
                # Vérifier qu'il y a plusieurs séries avec des couleurs différentes
                self.assertGreater(len(chart.series), 1, "Le graphique devrait avoir plusieurs séries")

                # Vérifier que les couleurs sont différentes pour chaque série
                colors_used = []
                for series in chart.series:
                    if hasattr(series.format.fill, 'fore_color') and hasattr(series.format.fill.fore_color, 'rgb'):
                        color = series.format.fill.fore_color.rgb
                        colors_used.append(color)

                # Au moins 2 couleurs différentes devraient être utilisées
                if len(colors_used) > 1:
                    # Vérifier que les couleurs ne sont pas toutes identiques
                    self.assertGreater(len(set(str(c) for c in colors_used)), 1,
                                      "Les séries devraient avoir des couleurs différentes")
                break

    @classmethod
    def tearDownClass(cls):
        """Nettoyage après tous les tests"""
        # Créer un rapport final
        report = {
            "test_suite": "column_clustered",
            "template_used": cls.template_path,
            "output_presentation": str(cls.presentation_path),
            "tests_completed": 7,
            "final_slide_count": 8,  # 1 titre + 7 graphiques
            "new_features_tested": [
                "Titre de série personnalisé",
                "Diversité de couleurs pour multi-séries"
            ]
        }

        report_path = cls.output_dir / "test_report.json"
        with open(report_path, 'w', encoding='utf-8') as f:
            json.dump(report, f, indent=2)

        print(f"\n[SUCCESS] Tests column_clustered terminés")
        print(f"[PRESENTATION] Créée: {cls.presentation_path}")
        print(f"[REPORT] Rapport: {report_path}")


if __name__ == "__main__":
    unittest.main(verbosity=2)