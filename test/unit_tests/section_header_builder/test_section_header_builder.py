#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Tests unitaires pour Section Header Builder - Architecture JSON 2025
Teste le nouveau script section_header_builder.py avec différentes configurations JSON.
"""

import os
import sys
import json
import unittest
from pathlib import Path
from datetime import datetime

# Ajouter le chemin vers presentation_builder
sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..', '..', '..', 'presentation_builder'))

try:
    from section_header_builder import SectionHeaderBuilder, create_section_header_from_json, load_section_header_template
    from presentation_builder import PresentationBuilder

except ImportError as e:
    print(f"[ERROR] Impossible d'importer les modules requis: {e}")
    sys.exit(1)


class TestSectionHeaderBuilder(unittest.TestCase):
    """Tests pour le SectionHeaderBuilder avec architecture JSON 2025"""

    def setUp(self):
        """Configuration initiale pour chaque test"""
        self.template_path = os.path.join('..', '..', '..', 'templates', 'Template_PT.pptx')
        self.test_subject = "unit-tests"
        self.test_audience = "section_header_builder_tester"

        # Créer le dossier de sortie
        self.output_dir = os.path.join('output')
        os.makedirs(self.output_dir, exist_ok=True)

        # Timestamp pour l'unicité des fichiers
        self.timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")

    def _create_presentation_from_schema(self, test_name: str) -> str:
        """Crée une présentation complète à partir du schéma JSON"""
        try:
            # Créer le PresentationBuilder
            builder = PresentationBuilder()

            # Changer vers le répertoire du test pour les chemins relatifs des payloads
            original_cwd = os.getcwd()
            test_dir = os.path.join(os.getcwd(), test_name)
            os.chdir(test_dir)

            try:
                # Construire la présentation complète avec le chemin vers le schéma JSON
                result_path = builder.build_presentation("presentation_schema.json")

                self.assertTrue(result_path, f"Échec construction présentation pour {test_name}")
                self.assertTrue(os.path.exists(result_path), f"Fichier de présentation non créé: {result_path}")

                print(f"[TEST] Présentation créée: {os.path.basename(result_path)}")
                return result_path

            finally:
                os.chdir(original_cwd)

        except Exception as e:
            self.fail(f"Erreur création présentation {test_name}: {e}")

    def test_major_section_header(self):
        """Test avec section header majeure via schéma JSON"""
        print(f"\n=== TEST: Section Header Majeure (Schéma Complet) ===")

        try:
            # Créer la présentation complète depuis le schéma
            presentation_path = self._create_presentation_from_schema("major_section")

            # Vérifier que le fichier existe et contient les bonnes slides
            self.assertTrue(os.path.exists(presentation_path), "Fichier présentation non trouvé")

            # Analyser la présentation pour vérifier qu'elle contient les slides attendues
            from pptx import Presentation
            pres = Presentation(presentation_path)

            # Doit contenir au moins 3 slides : titre + section header + fermeture
            self.assertGreaterEqual(len(pres.slides), 3, "Présentation doit contenir au moins 3 slides")

            print(f"[SUCCESS] Présentation avec section majeure créée avec {len(pres.slides)} slides")

        except Exception as e:
            self.fail(f"Erreur test section header majeure: {e}")

    def test_moderate_section_header(self):
        """Test avec section header modérée via schéma JSON"""
        print(f"\n=== TEST: Section Header Modérée (Schéma Complet) ===")

        try:
            # Créer la présentation complète depuis le schéma
            presentation_path = self._create_presentation_from_schema("moderate_section")

            # Vérifier que le fichier existe
            self.assertTrue(os.path.exists(presentation_path), "Fichier présentation non trouvé")

            # Analyser la présentation
            from pptx import Presentation
            pres = Presentation(presentation_path)

            # Doit contenir au moins 3 slides : titre + section header + fermeture
            self.assertGreaterEqual(len(pres.slides), 3, "Présentation doit contenir au moins 3 slides")

            print(f"[SUCCESS] Présentation avec section modérée créée avec {len(pres.slides)} slides")

        except Exception as e:
            self.fail(f"Erreur test section header modérée: {e}")

    def test_custom_section_header(self):
        """Test avec section header personnalisée"""
        print(f"\n=== TEST: Section Header Personnalisée ===")

        try:
            # Créer la présentation complète depuis le schéma
            presentation_path = self._create_presentation_from_schema("custom_section")

            # Vérifier que le fichier existe
            self.assertTrue(os.path.exists(presentation_path), "Fichier présentation non trouvé")

            # Analyser la présentation
            from pptx import Presentation
            pres = Presentation(presentation_path)

            # Doit contenir au moins 3 slides : titre + section header + fermeture
            self.assertGreaterEqual(len(pres.slides), 3, "Présentation doit contenir au moins 3 slides")

            print(f"[SUCCESS] Présentation avec section personnalisée créée avec {len(pres.slides)} slides")

        except Exception as e:
            self.fail(f"Erreur test section header personnalisée: {e}")


    def test_template_loading(self):
        """Test du chargement des templates prédéfinis"""
        print(f"\n=== TEST: Chargement Templates ===")

        # Tester les différents templates
        templates_to_test = ["major_section", "moderate_section", "numbered_major", "custom_section", "nonexistent"]

        for template_name in templates_to_test:
            try:
                config = load_section_header_template(template_name)

                # Vérifications de structure
                self.assertIn('section_title', config, f"Template {template_name}: 'section_title' manquant")
                self.assertIn('header_style', config, f"Template {template_name}: 'header_style' manquant")
                self.assertIn('options', config, f"Template {template_name}: 'options' manquant")

                self.assertIsInstance(config['section_title'], str, f"Template {template_name}: 'section_title' doit être une chaîne")
                self.assertIn(config['header_style'], ['major', 'moderate'], f"Template {template_name}: style invalide")

                print(f"[SUCCESS] Template '{template_name}' chargé: {config['section_title']} ({config['header_style']})")

            except Exception as e:
                if template_name == "nonexistent":
                    print(f"[EXPECTED] Template inexistant géré: {template_name}")
                else:
                    self.fail(f"Erreur chargement template {template_name}: {e}")

    def test_configuration_validation(self):
        """Test de validation des configurations JSON"""
        print(f"\n=== TEST: Validation Configurations ===")

        builder = SectionHeaderBuilder(self.template_path)

        # Configurations de test (valides et invalides)
        test_configs = [
            {
                "name": "config_valide_major",
                "config": {
                    "section_title": "Section Test",
                    "header_style": "major",
                    "options": {"auto_widen": True}
                },
                "should_be_valid": True
            },
            {
                "name": "config_valide_moderate",
                "config": {
                    "section_title": "Section Technique",
                    "header_style": "moderate",
                    "section_number": 2,
                    "options": {"auto_widen": False, "insert_position": 3}
                },
                "should_be_valid": True
            },
            {
                "name": "titre_manquant",
                "config": {
                    "header_style": "major"
                },
                "should_be_valid": False
            },
            {
                "name": "style_manquant",
                "config": {
                    "section_title": "Test"
                },
                "should_be_valid": False
            },
            {
                "name": "style_invalide",
                "config": {
                    "section_title": "Test",
                    "header_style": "invalid_style"
                },
                "should_be_valid": False
            },
            {
                "name": "titre_trop_long",
                "config": {
                    "section_title": "T" * 150,  # Titre de 150 caractères
                    "header_style": "major"
                },
                "should_be_valid": True  # Valide mais avec warning
            },
            {
                "name": "numero_negatif",
                "config": {
                    "section_title": "Test",
                    "header_style": "major",
                    "section_number": -1
                },
                "should_be_valid": False
            }
        ]

        for test_case in test_configs:
            try:
                validation = builder._validate_section_header_config(test_case["config"])

                if test_case["should_be_valid"]:
                    self.assertTrue(validation['valid'],
                                  f"Config '{test_case['name']}' devrait être valide: {validation['errors']}")
                else:
                    self.assertFalse(validation['valid'],
                                   f"Config '{test_case['name']}' devrait être invalide")

                print(f"[SUCCESS] Validation '{test_case['name']}': "
                      f"valide={validation['valid']}, erreurs={len(validation['errors'])}")

            except Exception as e:
                self.fail(f"Erreur validation config '{test_case['name']}': {e}")

    def test_style_mapping(self):
        """Test du mapping des styles vers les slides appropriées"""
        print(f"\n=== TEST: Mapping Styles ===")

        try:
            builder = SectionHeaderBuilder(self.template_path)

            # Test du mapping des styles
            style_tests = [
                {"style": "major", "expected_index": 14},
                {"style": "moderate", "expected_index": 15},
                {"style": "invalid", "expected_index": None}
            ]

            for test in style_tests:
                index = builder._get_slide_index_for_style(test["style"])
                self.assertEqual(index, test["expected_index"],
                               f"Style '{test['style']}' doit mapper vers l'index {test['expected_index']}")

                print(f"[SUCCESS] Style '{test['style']}' -> index {index}")

        except Exception as e:
            self.fail(f"Erreur test mapping styles: {e}")

    def test_available_styles_listing(self):
        """Test de la liste des styles disponibles"""
        print(f"\n=== TEST: Liste Styles Disponibles ===")

        try:
            builder = SectionHeaderBuilder(self.template_path)
            styles = builder.list_available_styles()

            # Vérifier que les styles attendus sont présents
            expected_styles = ["major", "moderate"]
            for style in expected_styles:
                self.assertIn(style, styles, f"Style '{style}' doit être disponible")

                style_info = styles[style]
                self.assertIn('slide_number', style_info, f"Info slide_number manquante pour {style}")
                self.assertIn('name', style_info, f"Info name manquante pour {style}")
                self.assertIn('usage', style_info, f"Info usage manquante pour {style}")
                self.assertIn('audience', style_info, f"Info audience manquante pour {style}")

                print(f"[SUCCESS] Style '{style}': slide {style_info['slide_number']} - {style_info['name']}")

        except Exception as e:
            self.fail(f"Erreur test liste styles: {e}")

    def test_template_validation(self):
        """Test de validation du template Premier Tech"""
        print(f"\n=== TEST: Validation Template PT ===")

        try:
            builder = SectionHeaderBuilder(self.template_path)
            is_valid = builder.validate_template()

            self.assertTrue(is_valid, "Template Premier Tech devrait être valide")

            # Vérifier les informations du template
            self.assertIn(14, builder.section_info, "Info slide 15 (index 14) manquante")
            self.assertIn(15, builder.section_info, "Info slide 16 (index 15) manquante")

            # Vérifier les styles disponibles
            expected_styles = {"major", "moderate"}
            available_styles = {builder.section_info[idx]['style'] for idx in builder.section_info}
            self.assertEqual(available_styles, expected_styles, "Styles disponibles incorrects")

            print(f"[SUCCESS] Template validé avec {len(builder.section_info)} slides de section")
            for idx, info in builder.section_info.items():
                print(f"  - Slide {info['slide_number']}: {info['layout_name']} ({info['style']})")

        except Exception as e:
            self.fail(f"Erreur validation template: {e}")

    def tearDown(self):
        """Nettoyage après chaque test"""
        pass  # Conserver les fichiers pour inspection


def run_section_header_tests():
    """Lance tous les tests de section header"""
    print("=== TESTS UNITAIRES SECTION HEADER BUILDER - ARCHITECTURE JSON 2025 ===")
    print(f"Timestamp: {datetime.now().isoformat()}")
    print(f"Sujet: unit-tests | Audience: section_header_builder_tester")

    # Configuration du test runner
    loader = unittest.TestLoader()
    suite = loader.loadTestsFromTestCase(TestSectionHeaderBuilder)

    # Exécuter les tests avec reporting détaillé
    runner = unittest.TextTestRunner(verbosity=2, buffer=False)
    result = runner.run(suite)

    # Rapport final
    print(f"\n=== RAPPORT FINAL ===")
    print(f"Tests exécutés: {result.testsRun}")
    print(f"Échecs: {len(result.failures)}")
    print(f"Erreurs: {len(result.errors)}")

    if result.failures:
        print("\nÉCHECS:")
        for test, traceback in result.failures:
            print(f"- {test}: {traceback}")

    if result.errors:
        print("\nERREURS:")
        for test, traceback in result.errors:
            print(f"- {test}: {traceback}")

    # Générer rapport JSON
    report_path = os.path.join('output', f'test_section_header_builder_{datetime.now().strftime("%Y%m%d_%H%M%S")}_report.json')
    test_report = {
        "timestamp": datetime.now().isoformat(),
        "test_suite": "SectionHeaderBuilder JSON Architecture 2025",
        "subject": "unit-tests",
        "audience": "section_header_builder_tester",
        "results": {
            "tests_run": result.testsRun,
            "failures": len(result.failures),
            "errors": len(result.errors),
            "success_rate": (result.testsRun - len(result.failures) - len(result.errors)) / result.testsRun * 100 if result.testsRun > 0 else 0
        },
        "details": {
            "failures": [{"test": str(test), "error": traceback} for test, traceback in result.failures],
            "errors": [{"test": str(test), "error": traceback} for test, traceback in result.errors]
        }
    }

    with open(report_path, 'w', encoding='utf-8') as f:
        json.dump(test_report, f, ensure_ascii=False, indent=2)

    print(f"\nRapport JSON généré: {report_path}")

    return result.wasSuccessful()


if __name__ == "__main__":
    success = run_section_header_tests()
    sys.exit(0 if success else 1)