#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Tests unitaires pour detailed_explanation_builder.py
Version JSON-native pour l'architecture 2025 du presentation_builder.
"""

import os
import sys
import json
import unittest
import tempfile
import shutil
from pathlib import Path
from typing import Dict, Any

# Ajouter le chemin des modules du projet
sys.path.insert(0, str(Path(__file__).parent.parent.parent.parent / "presentation_builder"))

from detailed_explanation_builder import DetailedExplanationBuilder, load_detailed_explanation_payload, process_detailed_explanation_from_payload_file
from pptx import Presentation


class TestDetailedExplanationBuilder(unittest.TestCase):
    """Tests pour le Detailed Explanation Builder"""

    def setUp(self):
        """Configuration initiale pour chaque test"""
        self.test_dir = Path(__file__).parent
        self.project_root = self.test_dir.parent.parent.parent
        self.template_path = self.project_root / "templates" / "Template_PT.pptx"

        # Créer un dossier temporaire pour les tests
        self.temp_dir = Path(tempfile.mkdtemp())
        self.temp_presentation = self.temp_dir / "test_presentation.pptx"

        # Créer une présentation de base pour les tests
        if self.template_path.exists():
            shutil.copy2(self.template_path, self.temp_presentation)

        # Instance du builder
        self.builder = DetailedExplanationBuilder(str(self.template_path))

    def tearDown(self):
        """Nettoyage après chaque test"""
        if self.temp_dir.exists():
            shutil.rmtree(self.temp_dir)

    def test_01_builder_initialization(self):
        """Test 1: Vérifier l'initialisation du builder"""
        self.assertIsInstance(self.builder, DetailedExplanationBuilder)
        self.assertTrue(os.path.exists(self.builder.template_path))
        self.assertIsInstance(self.builder.explanation_slides, dict)
        self.assertGreaterEqual(len(self.builder.explanation_slides), 5)

        # Vérifier les styles disponibles
        expected_styles = ["four_points", "dual_detailed_blue", "dual_detailed_grey",
                          "dual_titled_blue", "dual_titled_grey", "dual_lists_blue", "dual_lists_grey"]
        available_styles = [slide_data["style"] for slide_data in self.builder.explanation_slides.values()]

        for style in expected_styles:
            self.assertIn(style, available_styles, f"Style {style} manquant")

    def test_02_load_four_points_payload(self):
        """Test 2: Charger le payload four_points_explanation"""
        payload_path = self.test_dir / "four_points_explanation" / "detailed_explanation_payload.json"

        if payload_path.exists():
            payload = load_detailed_explanation_payload(str(payload_path))

            self.assertIn("content", payload)
            self.assertIn("explanation_style", payload)
            self.assertEqual(payload["explanation_style"], "four_points")
            self.assertIn("additional_content", payload)
            self.assertIsInstance(payload["additional_content"], list)
            self.assertGreaterEqual(len(payload["additional_content"]), 3)
        else:
            self.skipTest(f"Payload non trouvé: {payload_path}")

    def test_03_load_dual_detailed_payload(self):
        """Test 3: Charger le payload dual_detailed_explanation"""
        payload_path = self.test_dir / "dual_detailed_explanation" / "detailed_explanation_payload.json"

        if payload_path.exists():
            payload = load_detailed_explanation_payload(str(payload_path))

            self.assertIn("content", payload)
            self.assertIn("explanation_style", payload)
            self.assertEqual(payload["explanation_style"], "dual_detailed_blue")
            self.assertIn("subtitle", payload)
            self.assertIn("additional_content", payload)
            self.assertEqual(len(payload["additional_content"]), 2)
        else:
            self.skipTest(f"Payload non trouvé: {payload_path}")

    def test_04_validate_payload_configuration(self):
        """Test 4: Validation des configurations de payload"""
        # Test configuration valide
        valid_config = {
            "content": "Contenu de test valide pour validation",
            "explanation_style": "four_points",
            "title": "Test Titre",
            "additional_content": ["Point 1", "Point 2"]
        }

        result = self.builder._validate_detailed_explanation_config(valid_config)
        self.assertTrue(result["valid"])
        self.assertEqual(len(result["errors"]), 0)

        # Test configuration invalide
        invalid_config = {
            "content": "Short",  # Trop court
            "explanation_style": "invalid_style",  # Style inexistant
        }

        result = self.builder._validate_detailed_explanation_config(invalid_config)
        self.assertFalse(result["valid"])
        self.assertGreater(len(result["errors"]), 0)

    def test_05_process_four_points_explanation(self):
        """Test 5: Traitement d'une explication four_points"""
        if not self.temp_presentation.exists():
            self.skipTest("Template de présentation non disponible")

        config = {
            "content": "Test explication avec quatre points structurés pour validation complète",
            "title": "Test Four Points",
            "additional_content": [
                "Premier point détaillé",
                "Deuxième point explicatif",
                "Troisième point technique",
                "Quatrième point stratégique"
            ],
            "explanation_style": "four_points"
        }

        result = self.builder.process_detailed_explanation_config(config, str(self.temp_presentation))

        self.assertTrue(result.get("success", False), f"Erreur: {result.get('error', 'Inconnue')}")
        self.assertTrue(result.get("slide_added", False))
        self.assertEqual(result.get("explanation_style"), "four_points")

    def test_06_process_dual_detailed_explanation(self):
        """Test 6: Traitement d'une explication dual_detailed"""
        if not self.temp_presentation.exists():
            self.skipTest("Template de présentation non disponible")

        config = {
            "content": "Test comparaison détaillée entre deux approches complémentaires",
            "title": "Test Dual Detailed",
            "subtitle": "Comparaison approfondie",
            "additional_content": [
                "Première approche: Méthodologie traditionnelle éprouvée",
                "Seconde approche: Innovation disruptive moderne"
            ],
            "explanation_style": "dual_detailed_grey"
        }

        result = self.builder.process_detailed_explanation_config(config, str(self.temp_presentation))

        self.assertTrue(result.get("success", False), f"Erreur: {result.get('error', 'Inconnue')}")
        self.assertTrue(result.get("slide_added", False))
        self.assertEqual(result.get("explanation_style"), "dual_detailed_grey")

    def test_07_process_dual_lists_explanation(self):
        """Test 7: Traitement d'une explication dual_lists"""
        if not self.temp_presentation.exists():
            self.skipTest("Template de présentation non disponible")

        config = {
            "content": "Test listes duales pour analyse comparative structurée et détaillée",
            "title": "Test Dual Lists",
            "additional_content": [
                "Avantages|Performance élevée|Fiabilité prouvée|Maintenance simplifiée",
                "Défis|Coût initial|Formation requise|Temps d'adaptation"
            ],
            "explanation_style": "dual_lists_blue"
        }

        result = self.builder.process_detailed_explanation_config(config, str(self.temp_presentation))

        self.assertTrue(result.get("success", False), f"Erreur: {result.get('error', 'Inconnue')}")
        self.assertTrue(result.get("slide_added", False))
        self.assertEqual(result.get("explanation_style"), "dual_lists_blue")

    def test_08_process_from_payload_file(self):
        """Test 8: Traitement complet depuis fichier payload"""
        payload_path = self.test_dir / "four_points_explanation" / "detailed_explanation_payload.json"

        if not payload_path.exists() or not self.temp_presentation.exists():
            self.skipTest("Fichiers de test non disponibles")

        result = process_detailed_explanation_from_payload_file(
            str(payload_path),
            str(self.temp_presentation),
            str(self.template_path)
        )

        self.assertTrue(result.get("success", False), f"Erreur: {result.get('error', 'Inconnue')}")
        self.assertTrue(result.get("payload_loaded", False))
        self.assertTrue(result.get("slide_added", False))
        self.assertEqual(result.get("payload_source"), str(payload_path))

    def test_09_validate_presentation_integrity(self):
        """Test 9: Validation de l'intégrité de la présentation après insertion"""
        if not self.temp_presentation.exists():
            self.skipTest("Template de présentation non disponible")

        # Compter les slides avant insertion
        pres_before = Presentation(str(self.temp_presentation))
        slides_before = len(pres_before.slides)

        # Insérer une slide d'explication
        config = {
            "content": "Test intégrité présentation avec validation structure",
            "title": "Test Intégrité",
            "explanation_style": "dual_titled_blue",
            "additional_content": ["Section A: Test", "Section B: Validation"]
        }

        result = self.builder.process_detailed_explanation_config(config, str(self.temp_presentation))
        self.assertTrue(result.get("success", False))

        # Vérifier l'intégrité après insertion
        pres_after = Presentation(str(self.temp_presentation))
        slides_after = len(pres_after.slides)

        self.assertEqual(slides_after, slides_before + 1, "Une slide devrait être ajoutée")

        # Vérifier que la présentation peut être ouverte
        self.assertIsNotNone(pres_after)

    def test_10_error_handling(self):
        """Test 10: Gestion des erreurs"""
        # Test avec fichier inexistant
        result = self.builder.process_detailed_explanation_config(
            {"content": "Test", "explanation_style": "four_points"},
            "/path/inexistant/presentation.pptx"
        )
        self.assertFalse(result.get("success", True))
        self.assertIn("error", result)

        # Test avec configuration invalide
        result = self.builder.process_detailed_explanation_config(
            {"content": "", "explanation_style": "invalid"},
            str(self.temp_presentation)
        )
        self.assertFalse(result.get("success", True))
        self.assertIn("error", result)

    def test_11_all_explanation_styles(self):
        """Test 11: Test de tous les styles d'explication disponibles"""
        if not self.temp_presentation.exists():
            self.skipTest("Template de présentation non disponible")

        test_configs = [
            {"style": "four_points", "additional": ["Point 1", "Point 2", "Point 3", "Point 4"]},
            {"style": "dual_detailed_blue", "additional": ["Section 1", "Section 2"]},
            {"style": "dual_detailed_grey", "additional": ["Aspect A", "Aspect B"]},
            {"style": "dual_titled_blue", "additional": ["Titre 1: Contenu", "Titre 2: Contenu"]},
            {"style": "dual_titled_grey", "additional": ["Cat 1: Info", "Cat 2: Info"]},
            {"style": "dual_lists_blue", "additional": ["Liste A|Item 1|Item 2", "Liste B|Item 3|Item 4"]},
            {"style": "dual_lists_grey", "additional": ["Groupe 1|Elem 1|Elem 2", "Groupe 2|Elem 3|Elem 4"]}
        ]

        success_count = 0
        for i, test_config in enumerate(test_configs):
            config = {
                "content": f"Test style {test_config['style']} avec contenu détaillé pour validation",
                "title": f"Test {test_config['style'].replace('_', ' ').title()}",
                "additional_content": test_config['additional'],
                "explanation_style": test_config['style']
            }

            result = self.builder.process_detailed_explanation_config(config, str(self.temp_presentation))
            if result.get("success", False):
                success_count += 1
            else:
                print(f"[WARNING] Style {test_config['style']} a échoué: {result.get('error', 'Inconnue')}")

        # Au moins 5 styles sur 7 doivent fonctionner
        self.assertGreaterEqual(success_count, 5, f"Au moins 5 styles doivent fonctionner, {success_count} réussis")

    def test_12_content_validation_limits(self):
        """Test 12: Validation des limites de contenu"""
        test_cases = [
            # Contenu trop court
            {"content": "Court", "explanation_style": "four_points", "should_fail": True},
            # Contenu trop long
            {"content": "x" * 600, "explanation_style": "four_points", "should_fail": True},
            # Titre trop long
            {"content": "Contenu valide pour test", "title": "x" * 150, "explanation_style": "four_points", "should_fail": True},
            # Configuration valide
            {"content": "Contenu de longueur appropriée pour validation", "title": "Titre Normal", "explanation_style": "four_points", "should_fail": False}
        ]

        for i, test_case in enumerate(test_cases):
            should_fail = test_case.pop("should_fail")
            result = self.builder._validate_detailed_explanation_config(test_case)

            if should_fail:
                self.assertFalse(result["valid"], f"Test case {i+1} devrait échouer")
            else:
                self.assertTrue(result["valid"], f"Test case {i+1} devrait réussir: {result.get('errors', [])}")


def run_specific_test(test_name: str = None):
    """Exécute un test spécifique ou tous les tests"""
    if test_name:
        suite = unittest.TestSuite()
        suite.addTest(TestDetailedExplanationBuilder(test_name))
        runner = unittest.TextTestRunner(verbosity=2)
        result = runner.run(suite)
    else:
        unittest.main(verbosity=2)

    return result


if __name__ == "__main__":
    print("=== Tests Unitaires Detailed Explanation Builder ===")
    print(f"Dossier de test: {Path(__file__).parent}")
    print(f"Template: {Path(__file__).parent.parent.parent.parent / 'templates' / 'Template_PT.pptx'}")
    print()

    # Vérifier la disponibilité du template
    template_path = Path(__file__).parent.parent.parent.parent / "templates" / "Template_PT.pptx"
    if not template_path.exists():
        print(f"[WARNING] Template non trouvé: {template_path}")
        print("[INFO] Certains tests seront ignorés")
    else:
        print(f"[INFO] Template trouvé: {template_path}")

    print()
    unittest.main(verbosity=2)