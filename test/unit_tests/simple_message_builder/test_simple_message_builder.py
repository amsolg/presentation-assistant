#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Tests unitaires pour Simple Message Builder - Architecture JSON 2025
Teste le nouveau script simple_message_builder.py avec différentes configurations JSON.
"""

import os
import sys
import json
import unittest
from pathlib import Path
from datetime import datetime

# Ajouter le chemin vers presentation_builder
sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..', '..', '..', 'presentation_builder'))

try:
    from simple_message_builder import SimpleMessageBuilder, create_simple_message_from_json, load_simple_message_template
    from presentation_builder import PresentationBuilder

except ImportError as e:
    print(f"[ERROR] Impossible d'importer les modules requis: {e}")
    sys.exit(1)


class TestSimpleMessageBuilder(unittest.TestCase):
    """Tests pour le SimpleMessageBuilder avec architecture JSON 2025"""

    def setUp(self):
        """Configuration initiale pour chaque test"""
        self.template_path = os.path.join('..', '..', '..', 'templates', 'Template_PT.pptx')
        self.test_subject = "unit-tests"
        self.test_audience = "simple_message_builder_tester"

        # Créer le dossier de sortie
        self.output_dir = os.path.join('output')
        os.makedirs(self.output_dir, exist_ok=True)

        # Timestamp pour l'unicité des fichiers
        self.timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")

    def _create_presentation_from_schema(self, test_name: str) -> str:
        """Crée une présentation complète à partir du schéma JSON"""
        try:
            # Créer le PresentationBuilder
            builder = PresentationBuilder()

            # Changer vers le répertoire du test pour les chemins relatifs des payloads
            original_cwd = os.getcwd()
            test_dir = os.path.join(os.getcwd(), test_name)
            os.chdir(test_dir)

            try:
                # La méthode build_presentation attend un chemin JSON, pas un config dict
                output_path = builder.build_presentation("presentation_schema.json")

                self.assertTrue(output_path, f"Échec construction présentation pour {test_name}")
                self.assertTrue(os.path.exists(output_path), f"Fichier de présentation non créé: {output_path}")

                print(f"[TEST] Présentation créée: {os.path.basename(output_path)}")
                return output_path

            finally:
                os.chdir(original_cwd)

        except Exception as e:
            self.fail(f"Erreur création présentation {test_name}: {e}")

    def test_centered_simple_message(self):
        """Test avec message simple centré via schéma JSON"""
        print(f"\n=== TEST: Message Simple Centré (Schéma Complet) ===")

        try:
            # Créer la présentation complète depuis le schéma
            presentation_path = self._create_presentation_from_schema("centered_simple")

            # Vérifier que le fichier existe et contient les bonnes slides
            self.assertTrue(os.path.exists(presentation_path), "Fichier présentation non trouvé")

            # Analyser la présentation pour vérifier qu'elle contient les slides attendues
            from pptx import Presentation
            pres = Presentation(presentation_path)

            # Doit contenir au moins 3 slides : titre + message + fermeture
            self.assertGreaterEqual(len(pres.slides), 3, "Présentation doit contenir au moins 3 slides")

            print(f"[SUCCESS] Présentation message centré créée avec {len(pres.slides)} slides")

        except Exception as e:
            self.fail(f"Erreur test message centré: {e}")

    def test_illustrated_message(self):
        """Test avec message illustré via schéma JSON"""
        print(f"\n=== TEST: Message Illustré (Schéma Complet) ===")

        try:
            # Créer la présentation complète depuis le schéma
            presentation_path = self._create_presentation_from_schema("illustrated_message")

            # Vérifier que le fichier existe
            self.assertTrue(os.path.exists(presentation_path), "Fichier présentation non trouvé")

            # Analyser la présentation
            from pptx import Presentation
            pres = Presentation(presentation_path)

            # Doit contenir au moins 3 slides : titre + message + fermeture
            self.assertGreaterEqual(len(pres.slides), 3, "Présentation doit contenir au moins 3 slides")

            print(f"[SUCCESS] Présentation message illustré créée avec {len(pres.slides)} slides")

        except Exception as e:
            self.fail(f"Erreur test message illustré: {e}")

    def test_keyword_emphasis_message(self):
        """Test avec message à emphase mots-clés via schéma JSON"""
        print(f"\n=== TEST: Message avec Mots-clés (Schéma Complet) ===")

        try:
            # Créer la présentation complète depuis le schéma
            presentation_path = self._create_presentation_from_schema("keyword_emphasis")

            # Vérifier que le fichier existe
            self.assertTrue(os.path.exists(presentation_path), "Fichier présentation non trouvé")

            # Analyser la présentation
            from pptx import Presentation
            pres = Presentation(presentation_path)

            # Doit contenir au moins 3 slides : titre + message + fermeture
            self.assertGreaterEqual(len(pres.slides), 3, "Présentation doit contenir au moins 3 slides")

            print(f"[SUCCESS] Présentation message avec mots-clés créée avec {len(pres.slides)} slides")

        except Exception as e:
            self.fail(f"Erreur test message avec mots-clés: {e}")

    def test_direct_json_message_creation(self):
        """Test de création directe avec configuration JSON"""
        print(f"\n=== TEST: Création Directe JSON ===")

        # Créer une présentation de base d'abord
        base_path = self._create_base_presentation("direct_json")

        # Configuration JSON pour message simple
        message_config = {
            "message_text": "Excellence Premier Tech : innovation et qualité depuis 95 ans",
            "keywords": "Excellence • Innovation • Qualité",
            "image_description": "Histoire et expertise Premier Tech",
            "message_style": "illustrated",
            "options": {
                "auto_widen": True,
                "insert_position": None
            }
        }

        # Appliquer le message
        try:
            builder = SimpleMessageBuilder(self.template_path)
            result = builder.process_simple_message_config(message_config, base_path)

            # Vérifications
            self.assertTrue(result.get('success'), f"Échec traitement: {result.get('error')}")
            self.assertTrue(os.path.exists(base_path), "Fichier résultat non trouvé")
            self.assertEqual(result['processing_details']['style_applied'], 'illustrated')
            self.assertTrue(result['processing_details']['auto_widen_applied'])

            print(f"[SUCCESS] Message JSON direct appliqué avec style illustrated")

        except Exception as e:
            self.fail(f"Erreur test création directe JSON: {e}")

    def test_template_loading(self):
        """Test du chargement des templates prédéfinis"""
        print(f"\n=== TEST: Chargement Templates ===")

        # Tester les différents templates
        templates_to_test = ["centered_simple", "illustrated_message", "keyword_emphasis", "strategic_message", "nonexistent"]

        for template_name in templates_to_test:
            try:
                config = load_simple_message_template(template_name)

                # Vérifications de structure
                self.assertIn('message_text', config, f"Template {template_name}: 'message_text' manquant")
                self.assertIn('message_style', config, f"Template {template_name}: 'message_style' manquant")
                self.assertIn('options', config, f"Template {template_name}: 'options' manquant")

                self.assertIsInstance(config['message_text'], str, f"Template {template_name}: 'message_text' doit être une chaîne")
                self.assertIn(config['message_style'], ['centered', 'illustrated', 'keyword_simple'],
                            f"Template {template_name}: style invalide")

                print(f"[SUCCESS] Template '{template_name}' chargé: style '{config['message_style']}'")

            except Exception as e:
                if template_name == "nonexistent":
                    print(f"[EXPECTED] Template inexistant géré: {template_name}")
                else:
                    self.fail(f"Erreur chargement template {template_name}: {e}")

    def test_configuration_validation(self):
        """Test de validation des configurations JSON"""
        print(f"\n=== TEST: Validation Configurations ===")

        builder = SimpleMessageBuilder(self.template_path)

        # Configurations de test (valides et invalides)
        test_configs = [
            {
                "name": "config_valide",
                "config": {
                    "message_text": "Message de test valide",
                    "message_style": "centered",
                    "options": {"auto_widen": True}
                },
                "should_be_valid": True
            },
            {
                "name": "message_manquant",
                "config": {
                    "message_style": "centered"
                },
                "should_be_valid": False
            },
            {
                "name": "style_manquant",
                "config": {
                    "message_text": "Test"
                },
                "should_be_valid": False
            },
            {
                "name": "style_invalide",
                "config": {
                    "message_text": "Test",
                    "message_style": "invalid_style"
                },
                "should_be_valid": False
            },
            {
                "name": "message_trop_long",
                "config": {
                    "message_text": "M" * 600,  # Message de 600 caractères
                    "message_style": "centered"
                },
                "should_be_valid": True  # Valide mais avec warning
            },
            {
                "name": "keyword_sans_mots_cles",
                "config": {
                    "message_text": "Test avec style keyword",
                    "message_style": "keyword_simple"
                },
                "should_be_valid": True  # Valide mais avec warning
            }
        ]

        for test_case in test_configs:
            try:
                validation = builder._validate_simple_message_config(test_case["config"])

                if test_case["should_be_valid"]:
                    self.assertTrue(validation['valid'],
                                  f"Config '{test_case['name']}' devrait être valide: {validation['errors']}")
                else:
                    self.assertFalse(validation['valid'],
                                   f"Config '{test_case['name']}' devrait être invalide")

                print(f"[SUCCESS] Validation '{test_case['name']}': "
                      f"valide={validation['valid']}, erreurs={len(validation['errors'])}, "
                      f"warnings={len(validation['warnings'])}")

            except Exception as e:
                self.fail(f"Erreur validation config '{test_case['name']}': {e}")

    def test_template_validation(self):
        """Test de validation du template Premier Tech"""
        print(f"\n=== TEST: Validation Template PT ===")

        try:
            builder = SimpleMessageBuilder(self.template_path)
            is_valid = builder.validate_template()

            self.assertTrue(is_valid, "Template Premier Tech devrait être valide")

            # Vérifier les informations du template
            self.assertEqual(len(builder.message_slides), 3, "Nombre de slides de message incorrect")
            self.assertIn(16, builder.message_info, "Slide 17 (index 16) manquante")
            self.assertIn(17, builder.message_info, "Slide 18 (index 17) manquante")
            self.assertIn(18, builder.message_info, "Slide 19 (index 18) manquante")

            print(f"[SUCCESS] Template validé: {len(builder.message_slides)} styles de message disponibles")
            for idx, info in builder.message_info.items():
                print(f"  - Slide {info['slide_number']}: {info['layout_name']} ({info['style']})")

        except Exception as e:
            self.fail(f"Erreur validation template: {e}")

    def test_style_mapping(self):
        """Test du mapping des styles vers les slides"""
        print(f"\n=== TEST: Mapping Styles ===")

        builder = SimpleMessageBuilder(self.template_path)

        # Tester le mapping des styles
        style_tests = [
            ("centered", 16),      # Slide 17
            ("illustrated", 17),   # Slide 18
            ("keyword_simple", 18) # Slide 19
        ]

        for style, expected_index in style_tests:
            try:
                actual_index = builder._get_slide_index_for_style(style)
                self.assertEqual(actual_index, expected_index,
                               f"Style '{style}' devrait mapper à l'index {expected_index}, obtenu {actual_index}")

                print(f"[SUCCESS] Style '{style}' -> Slide {expected_index + 1} (index {expected_index})")

            except Exception as e:
                self.fail(f"Erreur mapping style '{style}': {e}")

        # Tester un style invalide
        invalid_index = builder._get_slide_index_for_style("invalid_style")
        self.assertIsNone(invalid_index, "Style invalide devrait retourner None")
        print(f"[SUCCESS] Style invalide géré correctement")

    def _create_base_presentation(self, test_name: str) -> str:
        """Crée une présentation de base pour les tests directs"""
        try:
            from presentation_builder import PresentationBuilder

            # Configuration minimale pour présentation de base
            config = {
                "presentation_name": f"Test {test_name}",
                "subject": self.test_subject,
                "audience": self.test_audience,
                "title_slide": {
                    "title": f"Test Simple Message Builder - {test_name}",
                    "subtitle": "Tests unitaires architecture JSON 2025",
                    "metadata": "2025-01-15 - Tests Unitaires"
                },
                "slides": [],  # Pas de slides pour test direct
                "build_options": {
                    "auto_widen_text": True,
                    "generate_reports": True
                }
            }

            # Créer un fichier JSON temporaire
            base_filename = f"{self.timestamp}_{test_name}_base"
            temp_json = os.path.join(self.output_dir, f"{base_filename}_config.json")

            with open(temp_json, 'w', encoding='utf-8') as f:
                json.dump(config, f, ensure_ascii=False, indent=2)

            builder = PresentationBuilder()
            # Changer vers le répertoire de sortie pour que le chemin relatif fonctionne
            original_cwd = os.getcwd()
            os.chdir(self.output_dir)

            try:
                output_path = builder.build_presentation(f"{base_filename}_config.json")
                self.assertTrue(output_path, f"Échec création présentation de base pour {test_name}")
                self.assertTrue(os.path.exists(output_path), f"Fichier de base non créé: {output_path}")

                # Retourner le chemin absolu
                return os.path.join(self.output_dir, output_path) if not os.path.isabs(output_path) else output_path
            finally:
                os.chdir(original_cwd)

        except Exception as e:
            self.fail(f"Erreur création présentation de base {test_name}: {e}")

    def tearDown(self):
        """Nettoyage après chaque test"""
        pass  # Conserver les fichiers pour inspection


def run_simple_message_tests():
    """Lance tous les tests de simple message"""
    print("=== TESTS UNITAIRES SIMPLE MESSAGE BUILDER - ARCHITECTURE JSON 2025 ===")
    print(f"Timestamp: {datetime.now().isoformat()}")
    print(f"Sujet: unit-tests | Audience: simple_message_builder_tester")

    # Configuration du test runner
    loader = unittest.TestLoader()
    suite = loader.loadTestsFromTestCase(TestSimpleMessageBuilder)

    # Exécuter les tests avec reporting détaillé
    runner = unittest.TextTestRunner(verbosity=2, buffer=False)
    result = runner.run(suite)

    # Rapport final
    print(f"\n=== RAPPORT FINAL ===")
    print(f"Tests exécutés: {result.testsRun}")
    print(f"Échecs: {len(result.failures)}")
    print(f"Erreurs: {len(result.errors)}")

    if result.failures:
        print("\nÉCHECS:")
        for test, traceback in result.failures:
            print(f"- {test}: {traceback}")

    if result.errors:
        print("\nERREURS:")
        for test, traceback in result.errors:
            print(f"- {test}: {traceback}")

    # Générer rapport JSON
    report_path = os.path.join('output', f'test_simple_message_builder_{datetime.now().strftime("%Y%m%d_%H%M%S")}_report.json')
    test_report = {
        "timestamp": datetime.now().isoformat(),
        "test_suite": "SimpleMessageBuilder JSON Architecture 2025",
        "subject": "unit-tests",
        "audience": "simple_message_builder_tester",
        "results": {
            "tests_run": result.testsRun,
            "failures": len(result.failures),
            "errors": len(result.errors),
            "success_rate": (result.testsRun - len(result.failures) - len(result.errors)) / result.testsRun * 100 if result.testsRun > 0 else 0
        },
        "details": {
            "failures": [{"test": str(test), "error": traceback} for test, traceback in result.failures],
            "errors": [{"test": str(test), "error": traceback} for test, traceback in result.errors]
        }
    }

    with open(report_path, 'w', encoding='utf-8') as f:
        json.dump(test_report, f, ensure_ascii=False, indent=2)

    print(f"\nRapport JSON généré: {report_path}")

    return result.wasSuccessful()


if __name__ == "__main__":
    success = run_simple_message_tests()
    sys.exit(0 if success else 1)