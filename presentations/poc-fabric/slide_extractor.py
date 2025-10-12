#!/usr/bin/env python3
"""
Slide Extractor for POC_Fabric
Extracteur détaillé de contenu PowerPoint pour génération de narration avec Sam AI

Ce module extrait de façon granulaire le contenu de chaque slide d'une présentation
PowerPoint et génère des fichiers JSON structurés par scènes pour la narration.
"""

import argparse
import json
import os
import sys
from pathlib import Path
from typing import Dict, List, Any, Optional
from xml.etree import ElementTree as ET

from pptx import Presentation
from pptx.shapes.base import BaseShape
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches, Pt


class SlideExtractor:
    """
    Extracteur détaillé de contenu PowerPoint avec support des animations
    et création de fichiers JSON structurés par scènes pour Sam AI.
    """

    def __init__(self, presentation_path: str, output_dir: str, slide_range: Optional[tuple] = None):
        """
        Initialise l'extracteur avec le chemin de la présentation et le dossier de sortie.

        Args:
            presentation_path: Chemin vers le fichier .pptx
            output_dir: Dossier de sortie pour les fichiers JSON générés
            slide_range: Tuple (début, fin) pour sélectionner des slides spécifiques (1-indexé)
        """
        self.presentation_path = Path(presentation_path)
        self.output_dir = Path(output_dir)
        self.presentation = None
        self.slide_range = slide_range
        self.metadata = {
            "presentation_name": self.presentation_path.stem,
            "total_slides": 0,
            "extraction_version": "1.0.0",
            "supported_features": [
                "text_extraction",
                "shape_positions",
                "speaker_notes",
                "basic_formatting",
                "shape_properties"
            ],
            "limitations": [
                "complex_animations_via_xml_only",
                "transitions_limited",
                "embedded_media_basic"
            ]
        }

    def load_presentation(self) -> bool:
        """
        Charge la présentation PowerPoint et initialise les métadonnées.

        Returns:
            bool: True si le chargement est réussi
        """
        try:
            self.presentation = Presentation(self.presentation_path)
            self.metadata["total_slides"] = len(self.presentation.slides)
            print(f"[OK] Présentation chargée : {self.metadata['total_slides']} slides")
            return True
        except Exception as e:
            print(f"[ERROR] Erreur lors du chargement de la présentation : {e}")
            return False

    def extract_all_slides(self) -> bool:
        """
        Extrait le contenu de toutes les slides (ou d'une sélection) et génère les fichiers JSON.

        Returns:
            bool: True si l'extraction est réussie
        """
        if not self.presentation:
            print("[ERROR] Présentation non chargée")
            return False

        success_count = 0
        slides_to_process = self._get_slides_to_process()
        total_to_process = len(slides_to_process)

        for slide_num, slide in enumerate(self.presentation.slides, 1):
            # Vérifier si cette slide doit être traitée
            if slide_num not in slides_to_process:
                continue

            try:
                slide_data = self.extract_slide_content(slide, slide_num)
                self.save_slide_json(slide_data, slide_num)
                success_count += 1
                print(f"[OK] Slide {slide_num} extraite et sauvegardée")
            except Exception as e:
                print(f"[ERROR] Erreur lors de l'extraction de la slide {slide_num} : {e}")

        # Sauvegarder les métadonnées globales
        self.save_presentation_metadata()

        print(f"[OK] Extraction terminée : {success_count}/{total_to_process} slides traitées")
        if self.slide_range:
            print(f"[INFO] Plage sélectionnée : slides {self.slide_range[0]} à {self.slide_range[1]}")
        return success_count > 0

    def _get_slides_to_process(self) -> List[int]:
        """
        Détermine quelles slides doivent être traitées selon la plage spécifiée.

        Returns:
            List des numéros de slides à traiter (1-indexé)
        """
        if not self.slide_range:
            # Traiter toutes les slides
            return list(range(1, self.metadata['total_slides'] + 1))

        start, end = self.slide_range
        total_slides = self.metadata['total_slides']

        # Valider et ajuster la plage
        start = max(1, min(start, total_slides))
        end = max(start, min(end, total_slides))

        return list(range(start, end + 1))

    def extract_slide_content(self, slide, slide_number: int) -> Dict[str, Any]:
        """
        Extrait le contenu détaillé d'une slide individuelle.

        Args:
            slide: Objet slide python-pptx
            slide_number: Numéro de la slide (1-indexé)

        Returns:
            Dict contenant toutes les données structurées de la slide
        """
        slide_data = {
            "slide_number": slide_number,
            "slide_title": self.extract_slide_title(slide),
            "layout_name": slide.slide_layout.name if hasattr(slide.slide_layout, 'name') else "Unknown",
            "total_scenes": 1,  # Minimum une scène, sera ajusté si animations
            "scenes": []
        }

        # Scène principale avec contenu statique
        main_scene = {
            "scene_id": 1,
            "scene_type": "static_content",
            "context": "",  # À remplir manuellement
            "technical_description": {
                "visual_elements": self.extract_visual_elements(slide),
                "layout": slide_data["layout_name"],
                "background": self.extract_background_info(slide),
                "automatic_animations": [],  # À compléter avec XML si nécessaire
                "slide_dimensions": self.get_slide_dimensions()
            },
            "speaker_notes": self.extract_speaker_notes(slide)
        }

        slide_data["scenes"].append(main_scene)

        # TODO: Analyser les animations manuelles et créer des scènes supplémentaires
        # Cette fonctionnalité nécessitera l'analyse XML pour les animations complexes

        return slide_data

    def extract_slide_title(self, slide) -> str:
        """
        Extrait le titre de la slide si disponible.

        Args:
            slide: Objet slide python-pptx

        Returns:
            str: Titre de la slide ou titre généré
        """
        try:
            for shape in slide.shapes:
                if hasattr(shape, "placeholder_format") and shape.placeholder_format:
                    try:
                        if shape.placeholder_format.type.name in ['TITLE', 'SUBTITLE']:
                            if hasattr(shape, 'text') and shape.text.strip():
                                return shape.text.strip()
                    except:
                        continue
                elif hasattr(shape, 'text') and shape.text.strip():
                    # Premier text box non-vide comme titre de fallback
                    text = shape.text.strip()
                    if len(text) < 100:  # Probablement un titre
                        return text
        except Exception:
            pass

        return f"Slide {getattr(slide, 'slide_id', 'Unknown')}"

    def extract_visual_elements(self, slide) -> List[Dict[str, Any]]:
        """
        Extrait tous les éléments visuels de la slide avec leurs propriétés détaillées.

        Args:
            slide: Objet slide python-pptx

        Returns:
            List des éléments visuels avec leurs propriétés
        """
        elements = []

        for shape_idx, shape in enumerate(slide.shapes):
            try:
                element = {
                    "element_id": f"shape_{shape_idx + 1}",
                    "shape_id": getattr(shape, 'shape_id', shape_idx + 1),
                    "type": self.get_shape_type(shape),
                    "position": self.get_shape_position(shape),
                    "content": self.extract_shape_content(shape),
                    "formatting": self.extract_shape_formatting(shape),
                    "properties": self.extract_shape_properties(shape)
                }
                elements.append(element)
            except Exception as e:
                # Ajouter un élément avec erreur plutôt que d'échouer complètement
                elements.append({
                    "element_id": f"shape_{shape_idx + 1}",
                    "shape_id": shape_idx + 1,
                    "type": "error",
                    "error": str(e),
                    "position": {"left": 0, "top": 0, "width": 0, "height": 0},
                    "content": {"type": "none", "data": None},
                    "formatting": {},
                    "properties": {}
                })

        return elements

    def get_shape_type(self, shape: BaseShape) -> str:
        """
        Détermine le type de forme PowerPoint.

        Args:
            shape: Forme PowerPoint

        Returns:
            str: Type de forme lisible
        """
        try:
            if hasattr(shape, 'shape_type'):
                shape_type = shape.shape_type
                type_mapping = {
                    MSO_SHAPE_TYPE.AUTO_SHAPE: "auto_shape",
                    MSO_SHAPE_TYPE.CALLOUT: "callout",
                    MSO_SHAPE_TYPE.CHART: "chart",
                    MSO_SHAPE_TYPE.COMMENT: "comment",
                    MSO_SHAPE_TYPE.CONNECTOR: "connector",
                    MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT: "embedded_ole",
                    MSO_SHAPE_TYPE.FORM_CONTROL: "form_control",
                    MSO_SHAPE_TYPE.FREEFORM: "freeform",
                    MSO_SHAPE_TYPE.GROUP: "group",
                    MSO_SHAPE_TYPE.LINE: "line",
                    MSO_SHAPE_TYPE.LINKED_OLE_OBJECT: "linked_ole",
                    MSO_SHAPE_TYPE.LINKED_PICTURE: "linked_picture",
                    MSO_SHAPE_TYPE.MEDIA: "media",
                    MSO_SHAPE_TYPE.OLE_CONTROL_OBJECT: "ole_control",
                    MSO_SHAPE_TYPE.PICTURE: "picture",
                    MSO_SHAPE_TYPE.PLACEHOLDER: "placeholder",
                    MSO_SHAPE_TYPE.SCRIPT_ANCHOR: "script_anchor",
                    MSO_SHAPE_TYPE.TABLE: "table",
                    MSO_SHAPE_TYPE.TEXT_EFFECT: "text_effect",
                    MSO_SHAPE_TYPE.TEXT_BOX: "text_box"
                }
                return type_mapping.get(shape_type, f"unknown_{shape_type}")
        except Exception:
            # Fallback : identifier le type par les propriétés disponibles
            return self.identify_shape_by_properties(shape)

        return "unknown"

    def identify_shape_by_properties(self, shape: BaseShape) -> str:
        """
        Identifie le type de forme par ses propriétés quand shape_type échoue.

        Args:
            shape: Forme PowerPoint

        Returns:
            str: Type de forme identifié
        """
        try:
            # Identifier par les propriétés spécifiques
            if hasattr(shape, 'text') and shape.text:
                return "text_box"
            elif hasattr(shape, 'table'):
                return "table"
            elif hasattr(shape, 'chart'):
                return "chart"
            elif hasattr(shape, 'image'):
                return "picture"
            elif hasattr(shape, 'placeholder_format') and shape.placeholder_format:
                return "placeholder"
            elif hasattr(shape, 'connection_site_count'):
                return "connector"  # Propriété spécifique aux connecteurs
            elif hasattr(shape, 'auto_shape_type'):
                return "auto_shape"
            else:
                # Analyser le nom de classe pour plus d'informations
                class_name = shape.__class__.__name__.lower()
                if 'connector' in class_name:
                    return "connector"
                elif 'group' in class_name:
                    return "group"
                elif 'picture' in class_name:
                    return "picture"
                elif 'textbox' in class_name or 'textframe' in class_name:
                    return "text_box"
                else:
                    return f"identified_as_{class_name}"
        except Exception:
            return "unknown_shape"

    def get_shape_position(self, shape: BaseShape) -> Dict[str, float]:
        """
        Extrait la position et les dimensions d'une forme.

        Args:
            shape: Forme PowerPoint

        Returns:
            Dict avec position et dimensions en pixels
        """
        try:
            # Conversion EMU vers pixels (approximation)
            emu_to_px = 1 / 9525  # 1 pixel ≈ 9525 EMU

            return {
                "left": round(shape.left * emu_to_px, 2),
                "top": round(shape.top * emu_to_px, 2),
                "width": round(shape.width * emu_to_px, 2),
                "height": round(shape.height * emu_to_px, 2),
                "rotation": getattr(shape, 'rotation', 0)
            }
        except Exception:
            return {"left": 0, "top": 0, "width": 0, "height": 0, "rotation": 0}

    def extract_shape_content(self, shape: BaseShape) -> Dict[str, Any]:
        """
        Extrait le contenu textuel ou autres données d'une forme.

        Args:
            shape: Forme PowerPoint

        Returns:
            Dict avec le contenu de la forme
        """
        content = {"type": "none", "data": None}
        shape_type = self.get_shape_type(shape)

        try:
            # Texte - vérifier en premier
            if hasattr(shape, 'text'):
                try:
                    text_content = shape.text
                    if text_content and text_content.strip():
                        content = {
                            "type": "text",
                            "data": text_content.strip(),
                            "has_text_frame": hasattr(shape, 'text_frame'),
                            "paragraphs_count": len(shape.text_frame.paragraphs) if hasattr(shape, 'text_frame') else 0
                        }
                        return content
                except:
                    pass

            # Table
            if shape_type == "table" and hasattr(shape, 'table'):
                try:
                    content = {
                        "type": "table",
                        "data": self.extract_table_content(shape.table),
                        "rows": len(shape.table.rows),
                        "columns": len(shape.table.columns)
                    }
                    return content
                except:
                    pass

            # Image/Picture
            if shape_type == "picture" or hasattr(shape, 'image'):
                try:
                    content = {
                        "type": "image",
                        "data": {
                            "filename": getattr(shape.image, 'filename', 'unknown') if hasattr(shape, 'image') else 'unknown',
                            "content_type": getattr(shape.image, 'content_type', 'unknown') if hasattr(shape, 'image') else 'unknown'
                        }
                    }
                    return content
                except:
                    pass

            # Chart
            if shape_type == "chart" or hasattr(shape, 'chart'):
                try:
                    content = {
                        "type": "chart",
                        "data": {
                            "chart_type": str(shape.chart.chart_type) if hasattr(shape, 'chart') and shape.chart else "unknown",
                            "has_title": hasattr(shape.chart, 'chart_title') if hasattr(shape, 'chart') and shape.chart else False
                        }
                    }
                    return content
                except:
                    pass

            # Connector ou autres formes géométriques
            if shape_type == "connector":
                content = {
                    "type": "connector",
                    "data": {
                        "description": "Ligne de connexion ou forme géométrique",
                        "shape_category": "connector"
                    }
                }
                return content

            # Auto Shape
            if shape_type == "auto_shape":
                try:
                    content = {
                        "type": "auto_shape",
                        "data": {
                            "auto_shape_type": str(shape.auto_shape_type) if hasattr(shape, 'auto_shape_type') else "unknown",
                            "description": "Forme automatique PowerPoint"
                        }
                    }
                    return content
                except:
                    pass

            # Placeholder
            if shape_type == "placeholder":
                try:
                    content = {
                        "type": "placeholder",
                        "data": {
                            "placeholder_type": str(shape.placeholder_format.type) if hasattr(shape, 'placeholder_format') and shape.placeholder_format else "unknown",
                            "description": "Zone de texte de mise en page"
                        }
                    }
                    return content
                except:
                    pass

            # Type générique basé sur le shape_type détecté
            content = {
                "type": shape_type,
                "data": {
                    "description": f"Élément de type {shape_type}",
                    "shape_category": shape_type
                }
            }

        except Exception as e:
            content = {
                "type": "error",
                "data": None,
                "error": str(e)
            }

        return content

    def extract_table_content(self, table) -> List[List[str]]:
        """
        Extrait le contenu d'une table PowerPoint.

        Args:
            table: Objet table python-pptx

        Returns:
            List 2D représentant le contenu de la table
        """
        table_data = []
        try:
            for row in table.rows:
                row_data = []
                for cell in row.cells:
                    row_data.append(cell.text if hasattr(cell, 'text') else "")
                table_data.append(row_data)
        except Exception:
            pass
        return table_data

    def extract_shape_formatting(self, shape: BaseShape) -> Dict[str, Any]:
        """
        Extrait les informations de formatage d'une forme.

        Args:
            shape: Forme PowerPoint

        Returns:
            Dict avec les propriétés de formatage
        """
        formatting = {}

        try:
            # Formatage du texte si disponible
            if hasattr(shape, 'text_frame') and shape.text_frame:
                formatting["text"] = self.extract_text_formatting(shape.text_frame)

            # Couleur de remplissage
            if hasattr(shape, 'fill'):
                formatting["fill"] = self.extract_fill_properties(shape.fill)

            # Contour
            if hasattr(shape, 'line'):
                formatting["line"] = self.extract_line_properties(shape.line)

        except Exception as e:
            formatting["error"] = str(e)

        return formatting

    def extract_text_formatting(self, text_frame) -> Dict[str, Any]:
        """
        Extrait le formatage du texte d'un text_frame.

        Args:
            text_frame: Objet text_frame python-pptx

        Returns:
            Dict avec propriétés de formatage du texte
        """
        try:
            # Prendre le formatage du premier paragraphe comme référence
            if text_frame.paragraphs:
                para = text_frame.paragraphs[0]
                if para.runs:
                    run = para.runs[0]
                    font = run.font

                    return {
                        "font_name": font.name,
                        "font_size": font.size.pt if font.size else None,
                        "bold": font.bold,
                        "italic": font.italic,
                        "underline": font.underline,
                        "color": self.get_color_hex(font.color) if font.color else None,
                        "alignment": str(para.alignment) if para.alignment else None
                    }
        except Exception:
            pass

        return {}

    def get_color_hex(self, color_obj) -> Optional[str]:
        """
        Convertit un objet couleur PowerPoint en hex.

        Args:
            color_obj: Objet couleur python-pptx

        Returns:
            String hex de la couleur ou None
        """
        try:
            if hasattr(color_obj, 'rgb'):
                rgb = color_obj.rgb
                return f"#{rgb.r:02x}{rgb.g:02x}{rgb.b:02x}"
        except Exception:
            pass
        return None

    def extract_fill_properties(self, fill) -> Dict[str, Any]:
        """Extrait les propriétés de remplissage d'une forme."""
        try:
            return {
                "type": str(fill.type) if hasattr(fill, 'type') else "unknown",
                "color": self.get_color_hex(fill.fore_color) if hasattr(fill, 'fore_color') else None
            }
        except Exception:
            return {}

    def extract_line_properties(self, line) -> Dict[str, Any]:
        """Extrait les propriétés de contour d'une forme."""
        try:
            return {
                "color": self.get_color_hex(line.color) if hasattr(line, 'color') else None,
                "width": line.width.pt if hasattr(line, 'width') and line.width else None
            }
        except Exception:
            return {}

    def extract_shape_properties(self, shape: BaseShape) -> Dict[str, Any]:
        """
        Extrait les propriétés additionnelles d'une forme.

        Args:
            shape: Forme PowerPoint

        Returns:
            Dict avec propriétés additionnelles
        """
        properties = {}

        try:
            properties.update({
                "name": shape.name if hasattr(shape, 'name') else None,
                "auto_shape_type": str(shape.auto_shape_type) if hasattr(shape, 'auto_shape_type') else None,
                "is_placeholder": hasattr(shape, 'placeholder_format') and shape.placeholder_format is not None,
                "visible": True  # PowerPoint n'expose pas directement la visibilité
            })

            # Propriétés spécifiques selon le type
            if hasattr(shape, 'placeholder_format') and shape.placeholder_format:
                try:
                    properties["placeholder_type"] = str(shape.placeholder_format.type)
                except Exception:
                    properties["placeholder_type"] = "unknown_placeholder"

        except Exception as e:
            properties["error"] = str(e)

        return properties

    def extract_background_info(self, slide) -> Dict[str, Any]:
        """
        Extrait les informations de background de la slide.

        Args:
            slide: Objet slide python-pptx

        Returns:
            Dict avec informations du background
        """
        # python-pptx a des limitations pour l'accès au background
        # Cette fonction pourrait être étendue avec l'analyse XML
        return {
            "type": "unknown",
            "note": "Background extraction requires XML analysis"
        }

    def get_slide_dimensions(self) -> Dict[str, float]:
        """
        Obtient les dimensions de la slide.

        Returns:
            Dict avec largeur et hauteur de la slide
        """
        try:
            if self.presentation:
                # Conversion EMU vers pixels
                emu_to_px = 1 / 9525
                return {
                    "width": round(self.presentation.slide_width * emu_to_px, 2),
                    "height": round(self.presentation.slide_height * emu_to_px, 2)
                }
        except Exception:
            pass

        # Dimensions standard PowerPoint 16:9
        return {"width": 1280, "height": 720}

    def extract_speaker_notes(self, slide) -> str:
        """
        Extrait les notes de présentation de la slide.

        Args:
            slide: Objet slide python-pptx

        Returns:
            String avec les notes ou chaîne vide
        """
        try:
            if hasattr(slide, 'notes_slide') and slide.notes_slide:
                notes_text_frame = slide.notes_slide.notes_text_frame
                if notes_text_frame and hasattr(notes_text_frame, 'text'):
                    return notes_text_frame.text.strip()
        except Exception:
            pass

        return ""

    def save_slide_json(self, slide_data: Dict[str, Any], slide_number: int) -> bool:
        """
        Sauvegarde les données d'une slide en fichier JSON.

        Args:
            slide_data: Données structurées de la slide
            slide_number: Numéro de la slide

        Returns:
            bool: True si la sauvegarde est réussie
        """
        try:
            output_file = self.output_dir / f"slide_{slide_number:02d}.json"
            with open(output_file, 'w', encoding='utf-8') as f:
                json.dump(slide_data, f, indent=2, ensure_ascii=False)
            return True
        except Exception as e:
            print(f"[ERROR] Erreur lors de la sauvegarde de la slide {slide_number} : {e}")
            return False

    def save_presentation_metadata(self) -> bool:
        """
        Sauvegarde les métadonnées globales de la présentation.

        Returns:
            bool: True si la sauvegarde est réussie
        """
        try:
            metadata_file = self.output_dir.parent / "metadata" / "presentation_info.json"
            metadata_file.parent.mkdir(exist_ok=True)

            with open(metadata_file, 'w', encoding='utf-8') as f:
                json.dump(self.metadata, f, indent=2, ensure_ascii=False)
            return True
        except Exception as e:
            print(f"[ERROR] Erreur lors de la sauvegarde des métadonnées : {e}")
            return False


def parse_slide_range(range_str: str) -> tuple:
    """
    Parse une chaîne de plage de slides (ex: "2-11", "5", "1-3").

    Args:
        range_str: Chaîne représentant la plage (ex: "2-11")

    Returns:
        Tuple (début, fin) ou None si invalide
    """
    try:
        if '-' in range_str:
            start, end = map(int, range_str.split('-', 1))
            return (start, end)
        else:
            # Une seule slide
            slide_num = int(range_str)
            return (slide_num, slide_num)
    except ValueError:
        raise argparse.ArgumentTypeError(f"Format de plage invalide: '{range_str}'. Utilisez '2-11' ou '5'")


def main():
    """Fonction principale pour l'extraction POC_Fabric avec support de sélection de slides."""

    parser = argparse.ArgumentParser(
        description="Extracteur de contenu PowerPoint pour POC_Fabric avec Sam AI",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""Exemples d'utilisation:
  python slide_extractor.py                    # Extraire toutes les slides
  python slide_extractor.py --slides 2-11      # Extraire les slides 2 à 11
  python slide_extractor.py --slides 5         # Extraire seulement la slide 5
  python slide_extractor.py -s 1-3 -o custom  # Slides 1-3 vers dossier 'custom'
        """
    )

    parser.add_argument(
        'presentation',
        nargs='?',
        default='POC_Fabric.pptx',
        help='Chemin vers le fichier PowerPoint (défaut: POC_Fabric.pptx)'
    )

    parser.add_argument(
        '-o', '--output',
        default='scripts',
        help='Dossier de sortie pour les fichiers JSON (défaut: scripts)'
    )

    parser.add_argument(
        '-s', '--slides',
        type=parse_slide_range,
        help='Plage de slides à extraire (ex: "2-11", "5", "1-3")'
    )

    parser.add_argument(
        '--list-slides',
        action='store_true',
        help='Afficher le nombre total de slides et quitter'
    )

    args = parser.parse_args()

    presentation_path = args.presentation
    output_dir = args.output
    slide_range = args.slides

    # Vérifier que le fichier existe
    if not os.path.exists(presentation_path):
        print(f"[ERROR] Fichier non trouvé : {presentation_path}")
        sys.exit(1)

    # Si l'utilisateur veut juste lister les slides
    if args.list_slides:
        try:
            from pptx import Presentation
            pres = Presentation(presentation_path)
            total_slides = len(pres.slides)
            print(f"[INFO] Présentation '{presentation_path}' contient {total_slides} slides")
            print(f"[INFO] Utilisez --slides 1-{total_slides} pour extraire toutes les slides")
            return
        except Exception as e:
            print(f"[ERROR] Erreur lors de l'ouverture de la présentation : {e}")
            sys.exit(1)

    # Créer le dossier de sortie
    os.makedirs(output_dir, exist_ok=True)

    # Initialiser et exécuter l'extraction
    extractor = SlideExtractor(presentation_path, output_dir, slide_range)

    print(f"[INFO] Chargement de la présentation : {presentation_path}")
    if not extractor.load_presentation():
        sys.exit(1)

    # Afficher les informations sur la plage de slides
    if slide_range:
        total_slides = extractor.metadata['total_slides']
        start, end = slide_range
        print(f"[INFO] Plage sélectionnée : slides {start}-{end} (sur {total_slides} total)")
        if start > total_slides or end > total_slides:
            print(f"[WARNING] Certaines slides de la plage dépassent le nombre total ({total_slides})")
    else:
        print(f"[INFO] Extraction de toutes les slides ({extractor.metadata['total_slides']} total)")

    print(f"[INFO] Extraction vers : {output_dir}")
    if extractor.extract_all_slides():
        print("[SUCCESS] Extraction POC_Fabric terminée avec succès !")
    else:
        print("[ERROR] Extraction échouée")
        sys.exit(1)


if __name__ == "__main__":
    main()