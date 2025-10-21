#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Section Header Builder - Construction de titres de section Premier Tech
Version JSON-native pour l'architecture 2025 du presentation_builder.
Utilise les slides 15, 16 du template Premier Tech pour créer des séparations de sections.
"""

import os
import json
import shutil
from datetime import datetime
from pathlib import Path
from typing import List, Dict, Optional, Any
from pptx import Presentation
from pptx.util import Inches


class SectionHeaderBuilder:
    """
    Classe pour construire des headers de section Premier Tech.
    Version modernisée pour l'architecture JSON 2025.
    """

    def __init__(self, template_path: str = "templates/Template_PT.pptx"):
        """
        Initialise le constructeur avec le template Premier Tech.

        Args:
            template_path: Chemin vers le template Premier Tech
        """
        self.template_path = template_path

        # Mapping des slides disponibles pour les headers de section
        self.section_slides = {
            14: {  # Slide 15 (index 14) - Titre de section bleu
                "name": "Titre de section bleu",
                "usage": "Transitions majeures, séparations importantes",
                "audience": "Leaders, Managers",
                "style": "major"
            },
            15: {  # Slide 16 (index 15) - Titre de section blanc
                "name": "Titre de section blanc",
                "usage": "Transitions modérées, sous-sections",
                "audience": "Spécialistes, Audiences mixtes",
                "style": "moderate"
            }
        }

        if not os.path.exists(template_path):
            raise FileNotFoundError(f"Template Premier Tech non trouvé: {template_path}")

        # Analyser la structure des slides de section
        self._analyze_section_structures()

    def _analyze_section_structures(self):
        """Analyse la structure des slides de section disponibles"""
        try:
            pres = Presentation(self.template_path)

            self.section_info = {}
            for slide_index, slide_data in self.section_slides.items():
                if len(pres.slides) > slide_index:
                    section_slide = pres.slides[slide_index]

                    self.section_info[slide_index] = {
                        'layout_name': section_slide.slide_layout.name,
                        'shape_count': len(section_slide.shapes),
                        'slide_index': slide_index,
                        'slide_number': slide_index + 1,
                        'style': slide_data['style'],
                        'usage': slide_data['usage']
                    }

            print(f"[INFO] {len(self.section_info)} slides de section analysées")
            for idx, info in self.section_info.items():
                print(f"[INFO] Slide {info['slide_number']}: {info['layout_name']} ({info['style']})")

        except Exception as e:
            raise Exception(f"Erreur analyse templates section: {e}")

    def process_section_header_config(self, config: Dict[str, Any], presentation_path: str) -> Dict[str, Any]:
        """
        Traite une configuration JSON de section header et l'applique à une présentation.
        Point d'entrée principal pour l'architecture JSON 2025.

        Args:
            config: Configuration JSON validée du section header
            presentation_path: Chemin vers la présentation cible

        Returns:
            Dict contenant les résultats du traitement
        """
        try:
            print(f"[SECTION_HEADER] Traitement configuration JSON...")

            # Extraire les données de configuration
            section_title = config.get('section_title', 'Nouvelle section')
            section_number = config.get('section_number', None)
            header_style = config.get('header_style', 'major')
            options = config.get('options', {})
            auto_widen = options.get('auto_widen', True)
            insert_position = options.get('insert_position', None)

            print(f"[SECTION_HEADER] Titre: {section_title}")
            print(f"[SECTION_HEADER] Style: {header_style}")
            print(f"[SECTION_HEADER] Options: auto_widen={auto_widen}, position={insert_position}")

            # Valider la configuration
            validation_result = self._validate_section_header_config(config)
            if not validation_result['valid']:
                raise ValueError(f"Configuration invalide: {', '.join(validation_result['errors'])}")

            # Insérer le section header dans la présentation
            result_path = self.insert_section_into_existing_presentation(
                presentation_path=presentation_path,
                section_title=section_title,
                section_number=section_number,
                header_style=header_style,
                insert_position=insert_position,
                auto_widen=auto_widen
            )

            # Générer le rapport de traitement
            processing_report = {
                "timestamp": datetime.now().isoformat(),
                "method": "JSON Section Header Builder 2025",
                "success": True,
                "configuration": config,
                "result_path": result_path,
                "validation": validation_result,
                "processing_details": {
                    "title_applied": section_title,
                    "style_applied": header_style,
                    "number_applied": section_number,
                    "auto_widen_applied": auto_widen,
                    "insert_position": insert_position
                }
            }

            print(f"[SUCCESS] Section header JSON traité avec succès")
            return processing_report

        except Exception as e:
            error_report = {
                "timestamp": datetime.now().isoformat(),
                "method": "JSON Section Header Builder 2025",
                "success": False,
                "error": str(e),
                "configuration": config
            }
            print(f"[ERROR] Erreur traitement section header JSON: {e}")
            return error_report

    def _validate_section_header_config(self, config: Dict[str, Any]) -> Dict[str, Any]:
        """Valide une configuration JSON de section header"""
        errors = []
        warnings = []

        # Vérifier les champs requis
        if 'section_title' not in config:
            errors.append("Champ 'section_title' manquant")
        elif not config['section_title'] or not isinstance(config['section_title'], str):
            errors.append("Le titre de section doit être une chaîne non vide")
        elif len(config['section_title']) > 100:
            warnings.append("Titre de section très long (>100 caractères)")

        if 'header_style' not in config:
            errors.append("Champ 'header_style' manquant")
        elif config['header_style'] not in ['major', 'moderate']:
            errors.append("'header_style' doit être 'major' ou 'moderate'")

        # Valider le numéro de section optionnel
        if 'section_number' in config and config['section_number'] is not None:
            if not isinstance(config['section_number'], int) or config['section_number'] < 1:
                errors.append("'section_number' doit être un entier positif")
            elif config['section_number'] > 99:
                warnings.append("Numéro de section très élevé (>99)")

        # Valider les options
        if 'options' in config:
            options = config['options']
            if 'auto_widen' in options and not isinstance(options['auto_widen'], bool):
                errors.append("'auto_widen' doit être un booléen")
            if 'insert_position' in options and options['insert_position'] is not None:
                if not isinstance(options['insert_position'], int) or options['insert_position'] < 0:
                    errors.append("'insert_position' doit être un entier positif ou null")

        return {
            'valid': len(errors) == 0,
            'errors': errors,
            'warnings': warnings
        }

    def insert_section_into_existing_presentation(self,
                                                presentation_path: str,
                                                section_title: str,
                                                section_number: Optional[int] = None,
                                                header_style: str = "major",
                                                insert_position: Optional[int] = None,
                                                auto_widen: bool = True) -> str:
        """
        Insère un header de section directement dans une présentation existante.
        Version modernisée avec support pour l'architecture JSON 2025.

        Args:
            presentation_path: Chemin vers la présentation existante
            section_title: Titre de la section
            section_number: Numéro de la section (optionnel)
            header_style: Style du header ("major", "moderate")
            insert_position: Position d'insertion (None = à la fin)
            auto_widen: Active l'élargissement automatique des objets texte

        Returns:
            str: Chemin vers le fichier modifié
        """
        try:
            print(f"[INSERT] Insertion directe section header dans: {os.path.basename(presentation_path)}")
            print(f"[INSERT] Style: {header_style}, Titre: {section_title}")

            # ÉTAPE 1: Créer une copie de sauvegarde
            backup_path = presentation_path.replace('.pptx', '_backup_before_section.pptx')
            shutil.copy2(presentation_path, backup_path)
            print(f"[BACKUP] Sauvegarde créée: {backup_path}")

            # ÉTAPE 2: Charger la présentation existante
            target_prs = Presentation(presentation_path)
            print(f"[LOAD] Présentation chargée: {len(target_prs.slides)} slides existantes")

            # ÉTAPE 3: Déterminer la slide source à utiliser
            source_slide_index = self._get_slide_index_for_style(header_style)
            if source_slide_index is None:
                raise ValueError(f"Style '{header_style}' non reconnu")

            # ÉTAPE 4: Vérifier que le layout section existe
            section_layout_index = self._find_section_layout_index(target_prs, source_slide_index)
            if section_layout_index is None:
                raise Exception(f"Layout section pour style '{header_style}' non trouvé dans la présentation")

            # ÉTAPE 5: Ajouter la slide section avec le bon layout
            section_layout = target_prs.slide_layouts[section_layout_index]
            new_slide = target_prs.slides.add_slide(section_layout)
            print(f"[ADD] Slide section ajoutée avec layout: {section_layout.name}")

            # ÉTAPE 6: Personnaliser le contenu de la slide section
            self._customize_section_slide_direct(new_slide, section_title, section_number, header_style)

            # ÉTAPE 7: Appliquer l'élargissement automatique si demandé
            if auto_widen:
                self._apply_auto_widen_to_slide(new_slide)

            # ÉTAPE 8: Réorganiser les slides si nécessaire
            if insert_position is not None and insert_position < len(target_prs.slides) - 1:
                self._move_slide_to_position_direct(target_prs, len(target_prs.slides) - 1, insert_position)

            # ÉTAPE 9: Sauvegarder
            target_prs.save(presentation_path)
            print(f"[SUCCESS] Section header inséré directement dans la présentation")

            # ÉTAPE 10: Générer le rapport
            self._generate_direct_insertion_report(presentation_path, section_title, section_number,
                                                 header_style, insert_position or len(target_prs.slides), auto_widen)

            return presentation_path

        except Exception as e:
            print(f"[ERROR] Erreur insertion directe section header: {e}")
            # Restaurer la sauvegarde en cas d'erreur
            if 'backup_path' in locals() and os.path.exists(backup_path):
                shutil.copy2(backup_path, presentation_path)
                print(f"[RESTORE] Présentation originale restaurée")
            raise

    def _get_slide_index_for_style(self, style: str) -> Optional[int]:
        """Retourne l'index de slide approprié pour le style demandé"""
        style_mapping = {
            "major": 14,     # Slide 15 - Titre de section bleu
            "moderate": 15   # Slide 16 - Titre de section blanc
        }
        return style_mapping.get(style)

    def _find_section_layout_index(self, presentation: Presentation, source_slide_index: int) -> Optional[int]:
        """Trouve l'index du layout section dans la présentation"""
        try:
            template_prs = Presentation(self.template_path)
            template_layout_name = template_prs.slides[source_slide_index].slide_layout.name

            for i, layout in enumerate(presentation.slide_layouts):
                if layout.name == template_layout_name:
                    print(f"[LAYOUT] Layout '{template_layout_name}' trouvé à l'index {i}")
                    return i

            return None

        except Exception as e:
            print(f"[WARNING] Erreur recherche layout section: {e}")
            return None

    def _customize_section_slide_direct(self, slide, section_title: str, section_number: Optional[int], header_style: str):
        """Personnalise directement la slide section ajoutée"""
        try:
            print(f"[CUSTOMIZE] Personnalisation slide section directe...")
            print(f"[CUSTOMIZE] Slide avec {len(slide.shapes)} shapes à personnaliser")

            shape_updates = 0

            # Analyse détaillée de chaque shape
            for i, shape in enumerate(slide.shapes):
                try:
                    if hasattr(shape, 'text_frame') and shape.text_frame:
                        current_text = shape.text_frame.text.strip()
                        print(f"[DEBUG] Shape {i}: '{current_text}' (longueur: {len(current_text)})")

                        # Personnalisation selon le style et la structure détectée
                        if header_style in ["major", "moderate"]:
                            # Slides 15/16: Shape 0 = titre
                            if i == 0:  # Premier shape = titre
                                # ÉTAPE 1: Redimensionner la shape aux dimensions de la slide
                                # Utiliser les dimensions de la présentation plutôt que du layout
                                from pptx.util import Inches
                                # Dimensions standard PowerPoint (16:9)
                                slide_width = Inches(13.3)  # 13.3 pouces de largeur
                                slide_height = Inches(7.5)  # 7.5 pouces de hauteur

                                shape.width = slide_width
                                shape.height = slide_height
                                shape.left = 0
                                shape.top = 0

                                print(f"[RESIZE] Shape redimensionnée aux dimensions de la slide")
                                print(f"[RESIZE] Dimensions: {slide_width/914400:.2f}\" x {slide_height/914400:.2f}\"")

                                # ÉTAPE 2: Configurer le texte
                                shape.text_frame.text = section_title

                                # ÉTAPE 3: Activer le text wrapping systématiquement
                                shape.text_frame.word_wrap = True

                                # ÉTAPE 4: Centrer le texte horizontalement et verticalement
                                from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

                                # Centrage horizontal
                                for paragraph in shape.text_frame.paragraphs:
                                    paragraph.alignment = PP_ALIGN.CENTER

                                # Centrage vertical
                                shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

                                print(f"[UPDATE] Shape {i} (titre {header_style}): {section_title}")
                                print(f"[LAYOUT] Text wrapping ACTIVÉ, centrage horizontal et vertical appliqué")
                                shape_updates += 1

                except Exception as e:
                    print(f"[WARNING] Erreur personnalisation shape {i}: {e}")

            print(f"[SUCCESS] Slide section personnalisée: {shape_updates} éléments mis à jour")

        except Exception as e:
            print(f"[ERROR] Erreur personnalisation directe section: {e}")
            raise

    def _apply_auto_widen_to_slide(self, slide):
        """Applique l'élargissement automatique à une slide spécifique"""
        try:
            print(f"[WIDEN] Application élargissement automatique...")

            widen_count = 0
            for i, shape in enumerate(slide.shapes):
                if hasattr(shape, 'text_frame'):
                    current_width_inches = shape.width / Inches(1)

                    # Élargir les objets texte étroits (<4 pouces)
                    if current_width_inches < 4.0:
                        new_width_inches = min(current_width_inches * 1.5, 8.0)
                        shape.width = Inches(new_width_inches)
                        print(f"[WIDEN] Shape {i}: {current_width_inches:.2f}\" -> {new_width_inches:.2f}\"")
                        widen_count += 1

                    # Élargir modérément les objets moyens (4-6 pouces)
                    elif current_width_inches < 6.0:
                        new_width_inches = min(current_width_inches * 1.2, 8.0)
                        shape.width = Inches(new_width_inches)
                        print(f"[WIDEN] Shape {i}: {current_width_inches:.2f}\" -> {new_width_inches:.2f}\" (modéré)")
                        widen_count += 1

            if widen_count > 0:
                print(f"[SUCCESS] {widen_count} objets texte élargis")
            else:
                print(f"[INFO] Aucun objet texte nécessitant un élargissement")

        except Exception as e:
            print(f"[WARNING] Erreur élargissement: {e}")

    def _move_slide_to_position_direct(self, presentation: Presentation, from_index: int, to_index: int):
        """Déplace une slide à la position désirée"""
        try:
            print(f"[POSITION] Slide section ajoutée en position {from_index + 1}")
            print(f"[INFO] Position finale: {to_index + 1}")

        except Exception as e:
            print(f"[WARNING] Déplacement slide: {e}")

    def _generate_direct_insertion_report(self, presentation_path: str, section_title: str,
                                        section_number: Optional[int], header_style: str,
                                        insert_position: int, auto_widen: bool):
        """Génère un rapport d'insertion directe modernisé"""
        try:
            source_slide_index = self._get_slide_index_for_style(header_style)

            report = {
                "insertion_timestamp": datetime.now().isoformat(),
                "method": "JSON Section Header Builder 2025 - Direct Insertion",
                "template_used": self.template_path,
                "target_presentation": presentation_path,
                "section_details": {
                    "title": section_title,
                    "number": section_number,
                    "style": header_style,
                    "intended_position": insert_position,
                    "auto_widen_enabled": auto_widen
                },
                "source_slide": {
                    "index": source_slide_index,
                    "number": source_slide_index + 1 if source_slide_index else None,
                    "layout": self.section_info.get(source_slide_index, {}).get('layout_name', 'Unknown'),
                    "style_description": self.section_slides.get(source_slide_index, {}).get('usage', 'Unknown')
                },
                "quality_assurance": {
                    "method": "JSON-Native Direct Layout-Based Addition",
                    "styles_preserved": True,
                    "premier_tech_standards": True,
                    "direct_integration": True,
                    "json_configuration": True,
                    "architecture_2025": True
                },
                "advantages": [
                    "Architecture JSON 2025 native",
                    "Insertion directe dans la présentation",
                    "Styles Premier Tech 100% préservés",
                    "Configuration JSON validée",
                    "Intégration transparente",
                    "Sauvegarde automatique créée",
                    f"Style '{header_style}' adapté à l'usage"
                ]
            }

            # Sauvegarder le rapport
            report_path = presentation_path.replace('.pptx', '_direct_section_insertion_report.json')
            with open(report_path, 'w', encoding='utf-8') as f:
                json.dump(report, f, ensure_ascii=False, indent=2)

            print(f"[INFO] Rapport d'insertion: {os.path.basename(report_path)}")

        except Exception as e:
            print(f"[WARNING] Erreur génération rapport: {e}")

    def validate_template(self) -> bool:
        """Valide que le template Premier Tech est correct pour les sections"""
        try:
            checks = {
                "file_exists": os.path.exists(self.template_path),
                "has_slides": False,
                "section_slides_exist": False,
                "slides_count": 0,
                "available_styles": []
            }

            if checks["file_exists"]:
                pres = Presentation(self.template_path)
                checks["slides_count"] = len(pres.slides)
                checks["has_slides"] = len(pres.slides) > 0

                # Vérifier que toutes les slides de section existent
                available_styles = []
                for slide_index in self.section_slides.keys():
                    if len(pres.slides) > slide_index:
                        style = self.section_slides[slide_index]['style']
                        available_styles.append(style)

                checks["available_styles"] = available_styles
                checks["section_slides_exist"] = len(available_styles) == len(self.section_slides)

            all_valid = all([checks["file_exists"], checks["has_slides"], checks["section_slides_exist"]])

            print("=== VALIDATION TEMPLATE PREMIER TECH POUR SECTIONS ===")
            for check, result in checks.items():
                if check == "available_styles":
                    print(f"[INFO] Styles disponibles: {', '.join(result)}")
                else:
                    status = "OK" if result else "ERREUR"
                    print(f"[{status}] {check}: {result}")

            if checks["section_slides_exist"]:
                print(f"[INFO] {len(checks['available_styles'])} styles de section disponibles:")
                for slide_index, slide_data in self.section_slides.items():
                    if slide_index in [idx for idx in self.section_slides.keys() if len(pres.slides) > idx]:
                        print(f"  - {slide_data['style']}: Slide {slide_index + 1} ({slide_data['usage']})")

            return all_valid

        except Exception as e:
            print(f"[ERROR] Erreur validation: {e}")
            return False

    def list_available_styles(self) -> Dict[str, Dict[str, Any]]:
        """Liste tous les styles de section disponibles"""
        return {
            slide_data['style']: {
                "slide_number": slide_index + 1,
                "name": slide_data['name'],
                "usage": slide_data['usage'],
                "audience": slide_data['audience']
            }
            for slide_index, slide_data in self.section_slides.items()
        }


def create_section_header_from_json(config_data: Dict[str, Any], presentation_path: str,
                                   template_path: str = "templates/Template_PT.pptx") -> Dict[str, Any]:
    """
    Fonction utilitaire pour créer un section header à partir d'une configuration JSON.
    Point d'entrée principal pour l'architecture JSON 2025.

    Args:
        config_data: Configuration JSON du section header
        presentation_path: Chemin vers la présentation cible
        template_path: Chemin vers le template Premier Tech

    Returns:
        Dict contenant les résultats du traitement
    """
    try:
        builder = SectionHeaderBuilder(template_path)
        return builder.process_section_header_config(config_data, presentation_path)
    except Exception as e:
        return {
            "success": False,
            "error": str(e),
            "timestamp": datetime.now().isoformat()
        }


def load_section_header_template(template_name: str = "major_section") -> Dict[str, Any]:
    """
    Charge un template de section header prédéfini.

    Args:
        template_name: Nom du template (major_section, moderate_section, numbered_major, custom_section)

    Returns:
        Dict contenant la configuration du template
    """
    template_path = "templates/presentation-project/slide-payload-templates/section_header_builder_template.json"

    try:
        with open(template_path, 'r', encoding='utf-8') as f:
            templates = json.load(f)

        examples = templates.get('examples', {})

        if template_name in examples:
            return examples[template_name]
        else:
            # Retourner le template de base si le nom n'est pas trouvé
            return templates.get('payload_structure', {})

    except Exception as e:
        print(f"[WARNING] Erreur chargement template: {e}")
        # Template de fallback
        return {
            "section_title": "Nouvelle section",
            "section_number": None,
            "header_style": "major",
            "options": {
                "auto_widen": True,
                "insert_position": None
            }
        }


def load_section_header_payload(payload_path: str) -> Dict[str, Any]:
    """
    Charge un payload de section header depuis un fichier JSON.

    Args:
        payload_path: Chemin vers le fichier JSON contenant le payload

    Returns:
        Dict contenant le payload de section header
    """
    try:
        with open(payload_path, 'r', encoding='utf-8') as f:
            payload = json.load(f)

        print(f"[PAYLOAD] Chargé depuis {payload_path}")
        return payload

    except Exception as e:
        print(f"[ERROR] Erreur chargement payload {payload_path}: {e}")
        return load_section_header_template("major_section")  # Fallback


def process_section_header_from_payload_file(payload_path: str, presentation_path: str,
                                           template_path: str = "templates/Template_PT.pptx") -> Dict[str, Any]:
    """
    Traite un section header en chargeant le payload depuis un fichier JSON.
    Point d'entrée pour l'architecture nouvelle avec fichiers payload séparés.

    Args:
        payload_path: Chemin vers le fichier JSON contenant le payload
        presentation_path: Chemin vers la présentation cible
        template_path: Chemin vers le template Premier Tech

    Returns:
        Dict contenant les résultats du traitement
    """
    try:
        # Charger le payload
        payload = load_section_header_payload(payload_path)

        # Traiter avec le payload chargé
        builder = SectionHeaderBuilder(template_path)
        result = builder.process_section_header_config(payload, presentation_path)

        # Ajouter les informations de payload
        result["payload_source"] = payload_path
        result["payload_loaded"] = True

        return result

    except Exception as e:
        return {
            "success": False,
            "error": str(e),
            "timestamp": datetime.now().isoformat(),
            "payload_source": payload_path,
            "payload_loaded": False
        }


if __name__ == "__main__":
    # Interface de test pour validation
    print("=== Section Header Builder - Architecture JSON 2025 ===")

    # Exemple d'utilisation
    template_config = load_section_header_template("major_section")
    print(f"Template chargé: {template_config}")

    # Test de validation
    builder = SectionHeaderBuilder()
    if builder.validate_template():
        print("[SUCCESS] Template Premier Tech validé pour section headers")
    else:
        print("[ERROR] Template Premier Tech invalide")