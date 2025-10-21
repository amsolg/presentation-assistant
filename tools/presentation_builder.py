#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Layout-Based Presentation Builder
=================================

Nouvelle version du presentation builder utilisant layout_name au lieu de slide_number.
Cette approche améliore la lisibilité, la maintenabilité et la flexibilité.

Architecture:
1. Parse le schéma JSON avec layout_name dans chaque slide
2. Construit un mapping layout_name → slide_number depuis les templates
3. Copie les slides du template selon le layout demandé
4. Applique toutes les configurations spécifiées dans le schéma

Avantages:
- Configuration plus lisible ("Page titre" vs slide_number: 11)
- Flexibilité totale dans l'ordre des slides
- Réutilisation libre des mêmes layouts
- Validation automatique de l'existence des layouts
"""

import os
import sys
import json
import argparse
import glob
import shutil
from datetime import datetime
from pathlib import Path
from typing import Dict, List, Any, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor


class LayoutBasedPresentationBuilder:
    """
    Constructeur de présentations basé sur les layout_name.

    Workflow:
    1. Charge la configuration JSON avec layout_name
    2. Construit le mapping layout_name → slide_number
    3. Pour chaque slide demandée, copie depuis le template
    4. Applique toutes les configurations spécifiées
    """

    def __init__(self):
        """Initialise le constructeur avec les chemins et mappings."""
        self.script_dir = Path(__file__).parent
        self.template_path = self.script_dir.parent / "templates" / "Template_PT.pptx"
        self.slide_structures_path = self.script_dir.parent / "templates" / "presentation-project" / "slide-structure"
        self.premier_tech_enums_path = self.script_dir.parent / "templates" / "presentation-project" / "premier_tech_schema_enums.json"

        # Vérifications d'existence
        if not self.template_path.exists():
            raise FileNotFoundError(f"Template Premier Tech non trouvé: {self.template_path}")

        if not self.slide_structures_path.exists():
            raise FileNotFoundError(f"Dossier slide-structure non trouvé: {self.slide_structures_path}")

        # Construire le mapping layout_name → slide_number
        self.layout_mapping = self._build_layout_mapping()

        # Charger les enums Premier Tech pour validation
        self.premier_tech_enums = self._load_premier_tech_enums()

        print(f"[INIT] Template Premier Tech: {self.template_path}")
        print(f"[INIT] Structures slides: {self.slide_structures_path}")
        print(f"[INIT] Layouts disponibles: {len(self.layout_mapping)}")

    def _build_layout_mapping(self) -> Dict[str, int]:
        """
        Construit le mapping layout_name → slide_number.

        Parse tous les fichiers JSON dans slide-structure pour créer
        un dictionnaire permettant de trouver le slide_number correspondant
        à un layout_name donné.

        Returns:
            Dict[str, int]: Mapping layout_name → slide_number
        """
        mapping = {}

        # Chercher tous les fichiers slide_*.json (anciens et nouveaux)
        pattern = str(self.slide_structures_path / "slide_*.json")
        structure_files = glob.glob(pattern)

        for file_path in structure_files:
            try:
                with open(file_path, 'r', encoding='utf-8') as f:
                    data = json.load(f)

                layout_name = data.get("layout_name")
                slide_number = data.get("slide_number")

                if layout_name and slide_number:
                    # Si le layout existe déjà, garder le premier trouvé
                    if layout_name not in mapping:
                        mapping[layout_name] = slide_number

            except Exception as e:
                print(f"[WARNING] Erreur lecture structure {file_path}: {e}")

        print(f"[MAPPING] Layouts mappés: {len(mapping)}")
        if len(mapping) < 50:  # On s'attend à ~57 layouts uniques
            print(f"[WARNING] Nombre de layouts faible. Vérifier les structures.")

        return mapping

    def _load_premier_tech_enums(self) -> Dict[str, Any]:
        """Charge les enums Premier Tech pour validation."""
        try:
            if not self.premier_tech_enums_path.exists():
                print(f"[WARNING] Fichier enums Premier Tech non trouvé: {self.premier_tech_enums_path}")
                return {}

            with open(self.premier_tech_enums_path, 'r', encoding='utf-8') as f:
                enums_data = json.load(f)

            print(f"[INIT] Premier Tech enums chargés: {enums_data.get('total_slides_analyzed', 0)} slides analysées")
            return enums_data

        except Exception as e:
            print(f"[WARNING] Erreur chargement enums Premier Tech: {e}")
            return {}

    def _validate_property_value(self, property_name: str, value: Any) -> bool:
        """Valide une valeur contre les enums Premier Tech."""
        if not self.premier_tech_enums or 'enums' not in self.premier_tech_enums:
            return True

        enums = self.premier_tech_enums['enums']
        if property_name not in enums:
            return True

        allowed_values = enums[property_name].get('enum', [])
        if value not in allowed_values:
            print(f"[WARNING] Valeur '{value}' non valide pour '{property_name}'. Valeurs autorisées: {allowed_values}")
            return False

        return True

    def load_presentation_config(self, json_path: str) -> Dict[str, Any]:
        """
        Charge et valide le fichier JSON de configuration.

        Support des deux formats:
        1. Nouveau format avec layout_name (recommandé)
        2. Format legacy avec slide_number (compatibilité)
        """
        try:
            with open(json_path, 'r', encoding='utf-8') as f:
                config = json.load(f)

            # Validation des champs requis
            required_fields = ["presentation_name", "subject", "audience", "slides", "output_path"]
            for field in required_fields:
                if field not in config:
                    raise ValueError(f"Champ requis manquant: {field}")

            # Validation des slides
            if not isinstance(config["slides"], list) or len(config["slides"]) == 0:
                raise ValueError("Le tableau 'slides' ne peut pas être vide")

            # Traitement et normalisation du chemin de sortie
            is_test = config.get("is_test", False)
            subject = config["subject"]
            audience = config["audience"]
            output_path = config["output_path"]

            # Forcer la structure d'output dans le dossier de l'audience
            base_dir = "tests" if is_test else "presentations"

            # Si l'output_path est un nom de fichier simple ou dans le root, le rediriger
            if not output_path.startswith(("presentations/", "tests/")) or "/" not in output_path:
                # Extraire le nom du fichier
                filename = os.path.basename(output_path)
                if not filename.endswith(".pptx"):
                    filename = f"{subject}_{audience}.pptx"

                # Construire le chemin normalisé
                normalized_path = f"{base_dir}/{subject}/{audience}/output/{filename}"
                config["output_path"] = normalized_path
                print(f"[CONFIG] Output path normalisé: {normalized_path}")
            elif is_test:
                # Mode test: rediriger presentations/ vers tests/
                if output_path.startswith("presentations/"):
                    config["output_path"] = output_path.replace("presentations/", "tests/", 1)
                    print(f"[CONFIG] Mode test activé - Redirection vers dossier 'tests'")

            # Validation finale: s'assurer que le path contient /output/
            final_output = config["output_path"]
            if "/output/" not in final_output:
                # Insérer /output/ avant le nom de fichier
                parts = final_output.split("/")
                if len(parts) >= 3:  # base_dir/subject/audience/file.pptx
                    filename = parts[-1]
                    path_prefix = "/".join(parts[:-1])
                    config["output_path"] = f"{path_prefix}/output/{filename}"
                    print(f"[CONFIG] Ajout du dossier output: {config['output_path']}")

            # Validation et conversion des slides
            for i, slide in enumerate(config["slides"]):
                if "layout_name" in slide:
                    # Format moderne avec layout_name
                    layout_name = slide["layout_name"]
                    if layout_name not in self.layout_mapping:
                        raise ValueError(f"Slide {i+1}: Layout '{layout_name}' non trouvé dans les templates")
                elif "slide_number" in slide:
                    # Format legacy - convertir en layout_name
                    slide_number = slide["slide_number"]
                    if not isinstance(slide_number, int) or slide_number < 1 or slide_number > 57:
                        raise ValueError(f"Slide {i+1}: 'slide_number' doit être entre 1 et 57")

                    # Trouver le layout_name correspondant au slide_number
                    layout_name = None
                    for layout, num in self.layout_mapping.items():
                        if num == slide_number:
                            layout_name = layout
                            break

                    if layout_name:
                        slide["layout_name"] = layout_name
                        print(f"[CONVERSION] Slide {i+1}: slide_number {slide_number} -> layout_name '{layout_name}'")
                    else:
                        raise ValueError(f"Slide {i+1}: Aucun layout trouvé pour slide_number {slide_number}")
                else:
                    raise ValueError(f"Slide {i+1}: 'layout_name' ou 'slide_number' requis")

            print(f"[CONFIG] Présentation: {config['presentation_name']}")
            print(f"[CONFIG] Sujet: {config['subject']}")
            print(f"[CONFIG] Audience: {config['audience']}")
            print(f"[CONFIG] Mode test: {is_test}")
            print(f"[CONFIG] Output: {config['output_path']}")
            print(f"[CONFIG] Slides à créer: {len(config['slides'])}")

            return config

        except json.JSONDecodeError as e:
            raise ValueError(f"JSON invalide: {e}")
        except FileNotFoundError:
            raise FileNotFoundError(f"Fichier de configuration non trouvé: {json_path}")

    def _copy_slide_from_template(self, layout_name: str, target_presentation: Presentation) -> Any:
        """
        Copie une slide spécifique du template vers la présentation cible.

        Args:
            layout_name: Nom du layout à copier
            target_presentation: Présentation de destination

        Returns:
            La nouvelle slide copiée
        """
        slide_number = self.layout_mapping[layout_name]
        slide_index = slide_number - 1  # Convertir en index 0-based

        print(f"[COPY] Copie layout '{layout_name}' (slide {slide_number})")

        # Charger le template
        template_prs = Presentation(self.template_path)

        if slide_index >= len(template_prs.slides):
            raise ValueError(f"Slide {slide_number} n'existe pas dans le template")

        # Récupérer la slide source et son layout
        source_slide = template_prs.slides[slide_index]
        source_layout = source_slide.slide_layout

        # Trouver le layout correspondant dans la présentation cible
        target_layout = None
        for layout in target_presentation.slide_layouts:
            if layout.name == source_layout.name:
                target_layout = layout
                break

        if not target_layout:
            print(f"[ERROR] Layout '{source_layout.name}' non trouvé dans la présentation cible")
            # Utiliser le premier layout disponible comme fallback
            target_layout = target_presentation.slide_layouts[0]

        # Créer la nouvelle slide avec le bon layout
        new_slide = target_presentation.slides.add_slide(target_layout)

        print(f"[SUCCESS] Slide '{layout_name}' ajoutée avec {len(new_slide.shapes)} shapes")

        return new_slide

    def _apply_slide_configuration(self, slide: Any, slide_config: Dict[str, Any]) -> bool:
        """
        Applique la configuration complète à une slide.

        Args:
            slide: Slide PowerPoint
            slide_config: Configuration à appliquer

        Returns:
            bool: True si succès
        """
        try:
            layout_name = slide_config.get("layout_name", "Unknown")
            shapes_config = slide_config.get("shapes", [])

            print(f"[CONFIGURE] Application configuration layout '{layout_name}' ({len(shapes_config)} shapes)")

            if not shapes_config:
                print(f"[INFO] Aucune configuration de shapes pour layout '{layout_name}'")
                return True

            configured_count = 0
            for shape_config in shapes_config:
                shape_id = shape_config.get("shape_id")
                if shape_id is None:
                    print(f"[WARNING] Shape sans shape_id, ignoré")
                    continue

                # Trouver le shape correspondant
                target_shape = self._find_shape_by_id(slide, shape_id)
                if target_shape is None:
                    print(f"[WARNING] Shape ID {shape_id} non trouvé dans la slide")
                    continue

                # Appliquer la configuration
                if self._apply_shape_configuration(target_shape, shape_config):
                    configured_count += 1

            print(f"[SUCCESS] {configured_count}/{len(shapes_config)} shapes configurés pour '{layout_name}'")
            return configured_count > 0

        except Exception as e:
            print(f"[ERROR] Erreur configuration slide: {e}")
            return False

    def _find_shape_by_id(self, slide: Any, shape_id: int) -> Any:
        """
        Trouve un shape par son ID dans une slide.

        Les shape_id correspondent aux index dans l'ordre des shapes.
        """
        try:
            shape_index = shape_id - 1  # Convertir en index 0-based

            if 0 <= shape_index < len(slide.shapes):
                return slide.shapes[shape_index]
            else:
                print(f"[WARNING] Shape ID {shape_id} hors limites (disponible: 1-{len(slide.shapes)})")
                return None

        except Exception as e:
            print(f"[ERROR] Erreur recherche shape ID {shape_id}: {e}")
            return None

    def _apply_shape_configuration(self, shape: Any, shape_config: Dict[str, Any]) -> bool:
        """
        Applique la configuration complète à un shape.

        Support de toutes les propriétés Premier Tech:
        - Position et dimensions
        - Texte et formatage
        - Marges et alignement
        - Propriétés PowerPoint spécifiques
        """
        try:
            shape_id = shape_config.get('shape_id', 'unknown')
            print(f"[SHAPE] Configuration shape {shape_id}")

            success = True

            # 1. Propriétés géométriques
            if not self._apply_geometry_properties(shape, shape_config):
                success = False

            # 2. Propriétés de texte
            if not self._apply_text_properties(shape, shape_config):
                success = False

            # 3. Propriétés de formatage avancées
            if not self._apply_advanced_formatting(shape, shape_config):
                success = False

            # 4. Propriétés PowerPoint spécifiques
            if not self._apply_powerpoint_properties(shape, shape_config):
                success = False

            if success:
                print(f"[SUCCESS] Shape {shape_id} configuré avec succès")
            else:
                print(f"[WARNING] Shape {shape_id} partiellement configuré")

            return success

        except Exception as e:
            print(f"[ERROR] Erreur configuration shape {shape_config.get('shape_id')}: {e}")
            return False

    def _apply_geometry_properties(self, shape: Any, shape_config: Dict[str, Any]) -> bool:
        """Applique les propriétés géométriques (position et dimensions)."""
        try:
            # Position
            position = shape_config.get("position", {})

            left = position.get("left")
            if left is not None:
                shape.left = Pt(left)

            top = position.get("top")
            if top is not None:
                shape.top = Pt(top)

            width = position.get("width")
            if width is not None:
                shape.width = Pt(width)

            height = position.get("height")
            if height is not None:
                shape.height = Pt(height)

            return True

        except Exception as e:
            print(f"[ERROR] Erreur propriétés géométriques: {e}")
            return False

    def _apply_text_properties(self, shape: Any, shape_config: Dict[str, Any]) -> bool:
        """Applique les propriétés de texte et formatage."""
        try:
            if not hasattr(shape, 'text_frame') or not shape.text_frame:
                return True

            # Appliquer le texte
            text = shape_config.get("text")
            if text is not None:
                shape.text_frame.text = text
                print(f"[TEXT] Texte appliqué: '{text[:50]}{'...' if len(text) > 50 else ''}'")

            # Formatage du texte
            if shape.text_frame.paragraphs:
                paragraph = shape.text_frame.paragraphs[0]

                # Créer ou récupérer le run
                if paragraph.runs:
                    run = paragraph.runs[0]
                else:
                    run = paragraph.add_run()

                # Police
                font_name = shape_config.get("font_name")
                if font_name:
                    self._validate_property_value("font_name", font_name)
                    run.font.name = font_name

                # Taille
                font_size = shape_config.get("font_size")
                if font_size:
                    self._validate_property_value("font_size", font_size)
                    run.font.size = Pt(font_size)

                # Gras
                bold = shape_config.get("bold")
                if bold is not None:
                    run.font.bold = bold

                # Couleur
                color = shape_config.get("color")
                if color and color.startswith("#"):
                    self._validate_property_value("color", color)
                    hex_color = color.lstrip("#")
                    rgb = tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
                    run.font.color.rgb = RGBColor(*rgb)

                # Alignement horizontal
                alignment = shape_config.get("alignment")
                if alignment:
                    self._validate_property_value("alignment", alignment)
                    if alignment == "LEFT":
                        paragraph.alignment = PP_ALIGN.LEFT
                    elif alignment == "CENTER":
                        paragraph.alignment = PP_ALIGN.CENTER
                    elif alignment == "RIGHT":
                        paragraph.alignment = PP_ALIGN.RIGHT

            return True

        except Exception as e:
            print(f"[ERROR] Erreur propriétés texte: {e}")
            return False

    def _apply_advanced_formatting(self, shape: Any, shape_config: Dict[str, Any]) -> bool:
        """Applique les propriétés de formatage avancées."""
        try:
            if not hasattr(shape, 'text_frame') or not shape.text_frame:
                return True

            text_frame = shape.text_frame

            # Marges
            margin_left = shape_config.get("margin_left")
            if margin_left is not None:
                self._validate_property_value("margin_values", margin_left)
                text_frame.margin_left = Pt(margin_left)

            margin_right = shape_config.get("margin_right")
            if margin_right is not None:
                self._validate_property_value("margin_values", margin_right)
                text_frame.margin_right = Pt(margin_right)

            margin_top = shape_config.get("margin_top")
            if margin_top is not None:
                self._validate_property_value("margin_values", margin_top)
                text_frame.margin_top = Pt(margin_top)

            margin_bottom = shape_config.get("margin_bottom")
            if margin_bottom is not None:
                self._validate_property_value("margin_values", margin_bottom)
                text_frame.margin_bottom = Pt(margin_bottom)

            # Alignement vertical
            vertical_alignment = shape_config.get("vertical_alignment")
            if vertical_alignment:
                self._validate_property_value("vertical_alignment", vertical_alignment)
                if vertical_alignment == "TOP":
                    text_frame.vertical_anchor = MSO_ANCHOR.TOP
                elif vertical_alignment == "MIDDLE":
                    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                elif vertical_alignment == "BOTTOM":
                    text_frame.vertical_anchor = MSO_ANCHOR.BOTTOM

            return True

        except Exception as e:
            print(f"[ERROR] Erreur formatage avancé: {e}")
            return False

    def _apply_powerpoint_properties(self, shape: Any, shape_config: Dict[str, Any]) -> bool:
        """Applique les propriétés PowerPoint spécifiques."""
        try:
            if not hasattr(shape, 'text_frame') or not shape.text_frame:
                return True

            text_frame = shape.text_frame

            # Autofit
            autofit = shape_config.get("autofit", {})
            if isinstance(autofit, dict):
                autofit_type = autofit.get("type")
                if autofit_type:
                    self._validate_property_value("autofit_type", autofit_type)
                    if autofit_type == "none":
                        text_frame.auto_size = MSO_AUTO_SIZE.NONE
                    elif autofit_type == "normal":
                        text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

            return True

        except Exception as e:
            print(f"[ERROR] Erreur propriétés PowerPoint: {e}")
            return False

    def build_presentation(self, json_path: str) -> str:
        """
        Construit une présentation complète à partir du JSON avec layout_name.

        Args:
            json_path: Chemin vers le fichier JSON de configuration

        Returns:
            str: Chemin vers la présentation créée
        """
        try:
            print(f"=== LAYOUT-BASED PRESENTATION BUILDER ===")
            print(f"Configuration: {json_path}")

            # 1. Charger la configuration
            config = self.load_presentation_config(json_path)

            # 2. Créer la présentation en copiant le template
            output_path = config["output_path"]
            os.makedirs(os.path.dirname(output_path), exist_ok=True)

            # Copier le template complet comme base
            shutil.copy2(str(self.template_path), output_path)
            presentation = Presentation(output_path)

            # Supprimer toutes les slides existantes
            slides_to_remove = list(range(len(presentation.slides)))
            for i in reversed(slides_to_remove):
                try:
                    rId = presentation.slides._sldIdLst[i].rId
                    presentation.part.drop_rel(rId)
                    del presentation.slides._sldIdLst[i]
                except Exception as e:
                    print(f"[WARNING] Erreur suppression slide {i}: {e}")

            print(f"[INIT] Présentation vide créée à partir du template")

            # 3. Ajouter chaque slide selon sa configuration
            for i, slide_config in enumerate(config["slides"]):
                layout_name = slide_config["layout_name"]
                print(f"\n[SLIDE {i+1}] Traitement layout '{layout_name}'")

                # Copier la slide du template
                new_slide = self._copy_slide_from_template(layout_name, presentation)

                # Appliquer la configuration
                self._apply_slide_configuration(new_slide, slide_config)

            # 4. Sauvegarder la présentation
            presentation.save(output_path)

            # 5. Vérifier le succès
            if os.path.exists(output_path):
                print(f"\n=== SUCCESS: Présentation créée ===")
                print(f"Fichier: {output_path}")
                print(f"Slides: {len(config['slides'])}")
                return output_path
            else:
                raise Exception("Construction échouée")

        except Exception as e:
            print(f"\n=== ERROR: Construction échouée ===")
            print(f"Erreur: {e}")
            raise


def main():
    """Interface en ligne de commande."""
    parser = argparse.ArgumentParser(
        description='Construction de présentations Premier Tech basée sur layout_name'
    )

    parser.add_argument('json_file', nargs='?', help='Fichier JSON de configuration de la présentation')
    parser.add_argument('--validate', action='store_true', help='Valider seulement le JSON')
    parser.add_argument('--list-layouts', action='store_true', help='Lister tous les layouts disponibles')

    args = parser.parse_args()

    try:
        builder = LayoutBasedPresentationBuilder()

        if args.list_layouts:
            print(f"\n=== LAYOUTS DISPONIBLES ({len(builder.layout_mapping)}) ===")
            for layout_name, slide_number in sorted(builder.layout_mapping.items()):
                print(f"  {slide_number:2d}: {layout_name}")
            sys.exit(0)

        if args.validate:
            if not args.json_file:
                print("Erreur: json_file requis pour --validate")
                sys.exit(1)
            config = builder.load_presentation_config(args.json_file)
            print(f"JSON valide: {args.json_file}")
            sys.exit(0)

        if not args.json_file:
            print("Erreur: json_file requis pour la génération")
            sys.exit(1)

        output_path = builder.build_presentation(args.json_file)
        print(f"\nSUCCES: {output_path}")

    except Exception as e:
        print(f"\nERREUR: {e}")
        sys.exit(1)


if __name__ == "__main__":
    main()